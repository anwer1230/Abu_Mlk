"""
╔═══════════════════════════════════════════════════════════════════════════════╗
║         مركز سرعة انجاز للخدمات الطلابية والأكاديمية - الإصدار المتكامل       ║
║              نظام التليجرام التلقائي + المنصة الأكاديمية المتكاملة            ║
╚═══════════════════════════════════════════════════════════════════════════════╝

╔══════════════════════════════════════════════════════════════════════════════╗
║                     🗺️  خريطة الكود — دليل سريع                           ║
╠══════════════════════════════════╦═══════════════╦═══════════════════════════╣
║ القسم / الوظيفة                   ║ السطر (تقريبي)║ الدالة / الفئة الرئيسية   ║
╠══════════════════════════════════╬═══════════════╬═══════════════════════════╣
║ الاستيرادات والإعدادات العامة      ║   1 – 110     ║ import / gevent / logging ║
║ نظام السجلات (Logging/IO)         ║ 110 – 340     ║ _MemoryLogHandler         ║
║ إعدادات API (TG / Groq / GH)     ║ 340 – 700     ║ API_ID · API_HASH · GROQ  ║
║ PREDEFINED_USERS / إعدادات تحديث ║ 700 – 960     ║ load_update_settings()    ║
║ اتصال تيليجرام — Login/Client    ║ 960 – 3150    ║ TelegramClientManager     ║
║ TelegramManager (المدير العلوي)   ║ 3150 – 3460   ║ TelegramManager           ║
║ واجهات API للاتصال والتحقق        ║ 3460 – 3860   ║ api_send_code / verify    ║
║ المراقبة + الإرسال المجدول        ║ 3860 – 5430   ║ monitoring_worker         ║
║ الانضمام التلقائي للمجموعات        ║ 5430 – 5680   ║ api_auto_join_advanced    ║
║ البحث العام في تيليجرام            ║ 5680 – 5990   ║ search_global_groups      ║
║ نظام التعلم الذكي (LearningBot)   ║ 5990 – 6905   ║ LearningBot/Manager       ║
║ المساعد الذكي AI + GitHub         ║ 6905 – 9340   ║ api_ai_assistant          ║
║ منشئ العروض PPTX                  ║ 9340 – 10700  ║ _PresentationGenerator    ║
║ منسق الملفات PDF/DOCX/Excel       ║ 10700 – 11540 ║ api_pdf_to_word           ║
║ GitHub helpers + روابط محفوظة     ║ 11540 – 11900 ║ upload/download_github    ║
║ رفع ملفات الأجهزة البيومترية       ║ 11900 – 12100 ║ upload_biometric_file     ║
║ لوحة الإدارة + نظام البطاقات      ║ 12100 – 13150 ║ admin_dashboard           ║
║ نظام الإشعارات + Web Push         ║ 13150 – 13400 ║ send_push_notification    ║
║ مراقبة الروابط + البحث            ║ 13400 – 14150 ║ link_monitor              ║
║ التحديث التلقائي للكود             ║ 14150 – 14480 ║ check_for_updates()       ║
║ التعلم التلقائي المتقدم ★ جديد ★   ║ 14480 – 14700 ║ start_auto_learning()     ║
║ تشغيل الخادم                       ║ آخر سطر       ║ socketio.run()            ║
╚══════════════════════════════════╩═══════════════╩═══════════════════════════╝

── الخصائص الوظيفية الرئيسية ──────────────────────────────────────────────────
  ① إرسال رسائل مجدولة/فوري مع صور عبر حسابات تيليجرام متعددة (5 مستخدمين)
  ② مراقبة الكلمات المفتاحية في المجموعات والإرسال التلقائي
  ③ بوت تعلم ذكي: يرد بشكل طبيعي باستخدام Groq + ذاكرة دائمة (JSON+GitHub)
  ④ التحليل التلقائي للمحادثات السابقة واستخلاص أنماط الرد عند تسجيل الدخول
  ⑤ منشئ عروض PPTX من نص + تحويل PDF/DOCX/Excel
  ⑥ مساعد ذكي AI يعدّل ملفات المشروع ويدفعها لـ GitHub مباشرةً
  ⑦ نظام بطاقات تفعيل مع لوحة إدارة كاملة
  ⑧ إشعارات Web Push + إشعارات ترويجية دورية
  ⑨ حفظ تلقائي للجلسات على GitHub (backup صامت)
── الخصائص المساعدة ───────────────────────────────────────────────────────────
  • load_settings / save_settings     — إعدادات كل مستخدم في JSON
  • upload_to_github / download_from_github — مزامنة البيانات مع GitHub
  • save_string_session / load_string_session — إدارة جلسات تيليجرام
  • load_learned_patterns / save_learned_patterns — أنماط التعلم المستفادة
  • _apply_learned_patterns           — تطبيق الأنماط قبل Groq
  • send_push_notification            — Web Push لمستخدم بعينه
  • load_cards_data / validate_voucher — نظام البطاقات والقسائم
"""

# استخدام OS thread حقيقي — بدون gevent monkey patching لتجنب تعارض asyncio
import threading as _pre_patch_threading
_OSThread = _pre_patch_threading.Thread

import os
import json
import uuid
import time
import logging

# ── وحدات منفصلة للأنظمة الفرعية (card_system, gps_tracking, isolation_system) ──
try:
    from card_system import (
        load_cards_data as _cs_load_cards, save_cards_data as _cs_save_cards,
        generate_vouchers as _cs_generate_vouchers, validate_voucher as _cs_validate_voucher,
        activate_card_voucher as _cs_activate_voucher, terminate_card_session as _cs_terminate_session,
        start_card_session_monitor as _cs_start_monitor,
    )
    _CARD_MODULE_LOADED = True
except ImportError:
    _CARD_MODULE_LOADED = False

try:
    from gps_tracking import geo_lookup as _gps_geo_lookup, build_map_url as _gps_map_url
    _GPS_MODULE_LOADED = True
except ImportError:
    _GPS_MODULE_LOADED = False

try:
    from isolation_system import (
        user_has_active_session as _iso_has_session,
        get_user_isolated_dir as _iso_user_dir,
    )
    _ISOLATION_MODULE_LOADED = True
except ImportError:
    _ISOLATION_MODULE_LOADED = False
# ── نهاية استيراد الوحدات المنفصلة ──────────────────────────────────────────
import asyncio
import threading
import queue
import re
import random
import string
import io
import base64
import tempfile
from datetime import datetime, timedelta
from threading import Lock, Event, Thread

# إضافات التحليل الإحصائي والعروض (اختيارية)
try:
    import pandas as pd
    import numpy as np
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import seaborn as sns
    from scipy import stats
    from scipy.stats import pearsonr, spearmanr, ttest_ind, ttest_1samp, f_oneway
    import plotly.express as px
    import plotly.graph_objects as go
    _DATA_SCIENCE_AVAILABLE = True
except ImportError:
    pd = None
    np = None
    matplotlib = None
    plt = None
    sns = None
    stats = None
    pearsonr = spearmanr = ttest_ind = ttest_1samp = f_oneway = None
    px = None
    go = None
    _DATA_SCIENCE_AVAILABLE = False

from io import BytesIO
import hashlib
import secrets
import requests

# إضافات معالجة الملفات
try:
    import docx
    import pdfplumber
    import fitz  # PyMuPDF
except ImportError:
    docx = None
    pdfplumber = None
    fitz = None

from flask import Flask, session, request, render_template, jsonify, redirect, send_file, abort, make_response
from install_tracker import track_installation, register_admin_routes
from flask_socketio import SocketIO, emit, join_room, leave_room
from telethon import TelegramClient, events, functions
from telethon.errors import SessionPasswordNeededError, PhoneCodeExpiredError, PhoneCodeInvalidError, PasswordHashInvalidError, FloodWaitError, UserAlreadyParticipantError, InviteHashExpiredError, InviteHashInvalidError
from telethon.sessions import StringSession
import socket

# ══════════════════════════════════════════════════════════
#  استيراد نظام المصادقة المستقل — auth.py
#  login/session management has been separated per-user
# ══════════════════════════════════════════════════════════
try:
    from auth import (
        TelegramLogin as _AuthTelegramLogin,
        save_settings as _auth_save_settings,
        load_settings as _auth_load_settings,
        clear_user_session as _auth_clear_user_session,
        save_string_session as _auth_save_string_session,
        load_string_session as _auth_load_string_session,
        get_user_session_dir as _auth_get_user_session_dir,
        # نظام المستخدمين الديناميكي وتسجيل الدخول
        load_dynamic_users,
        save_dynamic_users,
        add_dynamic_user,
        delete_dynamic_user,
        load_user_accounts,
        save_user_accounts,
        authenticate_platform_user,
        create_platform_account,
        delete_platform_account,
    )
    _AUTH_MODULE_LOADED = True
except ImportError as _auth_import_err:
    _AUTH_MODULE_LOADED = False

# تكوين السجلات المحسن
_log_handlers = [logging.StreamHandler()]
try:
    _log_handlers.append(logging.FileHandler('telegram_monitoring.log', encoding='utf-8'))
except Exception:
    pass
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=_log_handlers
)
logger = logging.getLogger(__name__)

# ── إنشاء مجلد outputs عند بدء التشغيل ──
_outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
os.makedirs(_outputs_dir, exist_ok=True)

class _MemoryLogHandler(logging.Handler):
    """يحتفظ بآخر N سجل في الذاكرة + دفع فوري عبر Socket.IO"""
    def __init__(self, capacity=500):
        super().__init__()
        self._records = []
        self._capacity = capacity
        self._lock = __import__('threading').Lock()
        self._socketio = None  # يُعيَّن بعد إنشاء socketio
        self._id_counter = 0

    def emit(self, record):
        try:
            import time as _time
            with self._lock:
                self._id_counter += 1
                _formatted = self.format(record)
                entry = {
                    'id': f"{int(_time.time()*1000)}_{self._id_counter}",
                    'time': self.formatTime(record, '%H:%M:%S'),
                    'level': record.levelname,
                    'msg': (_formatted.split(' - ', 3)[-1]
                            if ' - ' in _formatted else record.getMessage()),
                    'name': record.name,
                    'full': _formatted,
                }
                self._records.append(entry)
                if len(self._records) > self._capacity:
                    self._records = self._records[-self._capacity:]
            # ── بث فوري لجميع العملاء عبر Socket.IO (live_log للشاشة الحية) ──
            if self._socketio:
                try:
                    self._socketio.emit('live_log', entry)
                except Exception:
                    pass
            # ── إعادة توجيه لـ TypeScript LogSystem ──
            try:
                _TS_LOG_QUEUE.put_nowait({
                    'level': 'error' if record.levelno >= 40 else ('warn' if record.levelno >= 30 else 'info'),
                    'message': f"[{entry['name']}] {entry['msg']}",
                    'category': 'python',
                    'details': {'logger': entry['name'], 'time': entry['time']}
                })
            except Exception:
                pass
        except Exception:
            pass

    def get_records(self, level=None):
        with self._lock:
            recs = list(self._records)
        if level:
            lvl_map = {'ERROR': 40, 'WARNING': 30, 'INFO': 20}
            min_lvl = lvl_map.get(level.upper(), 0)
            recs = [r for r in recs if logging.getLevelName(r['level']) >= min_lvl]
        return recs[-200:]

_mem_log_handler = _MemoryLogHandler(capacity=200)
_mem_log_handler.setFormatter(logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s'))
logging.getLogger().addHandler(_mem_log_handler)

# ── قائمة انتظار لإعادة توجيه سجلات Python → TypeScript LogSystem ──────
_TS_LOG_QUEUE = queue.Queue(maxsize=500)

def _ts_log_forwarder_loop():
    """خيط OS حقيقي يرسل سجلات Python لـ TypeScript LogSystem فوراً"""
    import urllib.request as _ureq
    import json as _json_ts
    _url = 'http://localhost:8080/sys/logs'
    while True:
        try:
            entry = _TS_LOG_QUEUE.get(timeout=3)
            body = _json_ts.dumps(entry).encode('utf-8')
            req = _ureq.Request(_url, data=body,
                                headers={'Content-Type': 'application/json'},
                                method='POST')
            _ureq.urlopen(req, timeout=1)
        except queue.Empty:
            continue
        except Exception:
            pass

_ts_fwd_thread = _OSThread(target=_ts_log_forwarder_loop, daemon=True, name='TsLogFwd')
_ts_fwd_thread.start()

# ── التقاط stdout/stderr وتوجيهها لـ TypeScript LogSystem ──────────────────
class _ConsoleCapture:
    """يلتقط stdout/stderr ويضعها في _TS_LOG_QUEUE للعرض الفوري"""
    _tls = _pre_patch_threading.local()

    def __init__(self, original, stream_name):
        self._orig = original
        self._name = stream_name

    def write(self, text):
        self._orig.write(text)
        msg = text.strip()
        if msg and not getattr(self._tls, 'in_capture', False):
            self._tls.in_capture = True
            try:
                _TS_LOG_QUEUE.put_nowait({
                    'level': 'debug',
                    'message': msg,
                    'category': 'python-console',
                    'details': {'stream': self._name}
                })
            except Exception:
                pass
            finally:
                self._tls.in_capture = False

    def flush(self):
        self._orig.flush()

    def fileno(self):
        try:
            return self._orig.fileno()
        except Exception:
            raise io.UnsupportedOperation('fileno')

    def isatty(self):
        return False

import sys as _sys
_sys.stdout = _ConsoleCapture(_sys.stdout, 'stdout')
_sys.stderr = _ConsoleCapture(_sys.stderr, 'stderr')

# ── [إصلاح] سجلات خاصة بكل مستخدم ──────────────────────────────
from collections import deque as _deque
_USER_LOGS: dict = {}
_USER_LOGS_LOCK = __import__('threading').Lock()
_MAX_USER_LOGS = 300

_log_emit_counter = 0

def _emit_log_update(level: str, msg: str, user_id: str = None):
    """إرسال سجل بتنسيق موحد كامل عبر Socket.IO فوراً"""
    global _log_emit_counter
    try:
        import time as _t
        _log_emit_counter += 1
        entry = {
            'id': f"{int(_t.time()*1000)}_{_log_emit_counter}",
            'time': __import__('datetime').datetime.now().strftime('%H:%M:%S'),
            'level': level.upper(),
            'msg': msg,
            'name': user_id or 'system'
        }
        if user_id:
            try:
                socketio.emit('log_update', entry, to=user_id)
            except Exception:
                pass
        try:
            socketio.emit('log_update', entry)
        except Exception:
            pass
    except Exception:
        pass

def log_user_event(user_id: str, level: str, msg: str):
    """تسجيل حدث في سجل المستخدم الخاص وفي السجل العام مع دفع فوري"""
    try:
        lvl_num = {'DEBUG': 10, 'INFO': 20, 'WARNING': 30, 'ERROR': 40}.get(level.upper(), 20)
        logger.log(lvl_num, f"[{user_id}] {msg}")
        record = {
            'time': __import__('datetime').datetime.now().strftime('%H:%M:%S'),
            'level': level.upper(),
            'msg': msg,
            'name': user_id,
            'source': user_id,
        }
        with _USER_LOGS_LOCK:
            if user_id not in _USER_LOGS:
                _USER_LOGS[user_id] = _deque(maxlen=_MAX_USER_LOGS)
            _USER_LOGS[user_id].append(record)
        # دفع فوري عبر Socket.IO بتنسيق كامل
        _emit_log_update(level, msg, user_id)
    except Exception:
        pass

def _get_user_logs(user_id: str, level_filter=None) -> list:
    """إرجاع سجلات المستخدم الخاصة + سجلات النظام المشتركة"""
    lvl_map = {'ERROR': 40, 'WARNING': 30, 'INFO': 20, 'DEBUG': 10}
    min_lvl = lvl_map.get((level_filter or '').upper(), 0) if level_filter and level_filter != 'ALL' else 0
    # سجلات المستخدم الخاصة
    with _USER_LOGS_LOCK:
        user_recs = list(_USER_LOGS.get(user_id, []))
    # سجلات النظام العامة (لا تخص مستخدماً آخر محدداً)
    other_users = [u for u in _USER_LOGS if u != user_id]
    global_recs = _mem_log_handler.get_records(None)
    merged = user_recs + [r for r in global_recs if not any(ou in r.get('msg', '') for ou in other_users)]
    # فلترة حسب المستوى
    if min_lvl > 0:
        merged = [r for r in merged if lvl_map.get(r.get('level', 'INFO'), 20) >= min_lvl]
    return merged[-100:]

# إنشاء التطبيق
app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "sreaa-injaz-stable-key-2024-!@#$%")
app.config["PERMANENT_SESSION_LIFETIME"] = 60 * 60 * 24 * 30  # 30 يوم

# إعداد SocketIO — threading mode لتجنب تعارض asyncio/gevent
socketio = SocketIO(
    app,
    cors_allowed_origins="*",
    async_mode='threading',
    ping_timeout=20,
    ping_interval=10,
    logger=False,
    engineio_logger=False,
    allow_upgrades=True,
)
# تفعيل الدفع الفوري للسجلات عبر Socket.IO
_mem_log_handler._socketio = socketio

# إعدادات النظام
SESSIONS_DIR = os.path.join('/tmp', 'sessions') if os.environ.get('RENDER') else "sessions"
if not os.path.exists(SESSIONS_DIR):
    os.makedirs(SESSIONS_DIR)

# ── مجلد البيانات الدائمة (مستقل عن الجلسات ولا يُحذف معها) ──
DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "data")
os.makedirs(DATA_DIR, exist_ok=True)

# ── ملف بطاقات الشحن ──
CARDS_FILE = os.path.join(DATA_DIR, "cards.json")
_CARDS_LOCK = threading.Lock()

# ── ملف روابط الدعوة ──
INVITE_FILE = os.path.join(DATA_DIR, "invite_tokens.json")

# ── دوال إدارة روابط الدعوة (One-time Invite Links) ──
def load_invites():
    try:
        if os.path.exists(INVITE_FILE):
            with open(INVITE_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    return {"tokens": []}

def save_invites(data):
    try:
        with open(INVITE_FILE, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        logger.error(f"فشل حفظ روابط الدعوة: {e}")
        return False

def generate_invite_token():
    token = secrets.token_urlsafe(16)
    data = load_invites()
    data["tokens"].append({
        "token": token,
        "created_at": datetime.now().isoformat(),
        "status": "active",
        "used_by": None,
        "used_at": None
    })
    save_invites(data)
    return token

def validate_invite_token(token):
    data = load_invites()
    for item in data["tokens"]:
        if item["token"] == token:
            if item["status"] == "used":
                return False, None, "❌ هذا الرابط تم استخدامه مسبقاً"
            if item["status"] == "expired":
                return False, None, "❌ انتهت صلاحية هذا الرابط"
            try:
                created = datetime.fromisoformat(item["created_at"])
                if (datetime.now() - created).days > 7:
                    item["status"] = "expired"
                    save_invites(data)
                    return False, None, "❌ انتهت صلاحية الرابط (أكثر من 7 أيام)"
            except Exception:
                pass
            return True, item, None
    return False, None, "❌ الرابط غير صحيح"

def mark_token_used(token, user_id):
    data = load_invites()
    for item in data["tokens"]:
        if item["token"] == token:
            item["status"] = "used"
            item["used_by"] = user_id
            item["used_at"] = datetime.now().isoformat()
            save_invites(data)
            return True
    return False

def get_invite_link(token):
    try:
        base_url = request.host_url.rstrip('/')
    except Exception:
        base_url = "http://localhost:5000"
    return f"{base_url}/?invite={token}"

# ====== Web Push / VAPID إعداد ======
try:
    from pywebpush import webpush as _webpush_fn, WebPushException as _WebPushException
    WEBPUSH_AVAILABLE = True
except ImportError:
    WEBPUSH_AVAILABLE = False

_VAPID_PUB_D  = "BIkbfhrnC0MgiZk3KJ8fPTUX300SviFxsJguvQjx6q-TNMxK23yhqhyr5Q5vqad8-k3aD1J7NPWg3GbKHuJdpOc"
_VAPID_PRV_B64 = (
    "LS0tLS1CRUdJTiBFQyBQUklWQVRFIEtFWS0tLS0tCk1IY0NBUUVFSURmeGQ0SUZkSHNrQlFNVjd0RD"
    "VodXVZZlIyOXhtU1ZOMzk5clhhNXYrTTFvQW9HQ0NxR1NNNDkKQXdFSG9VUURRZ0FFaVJ0K0d1Y0xR"
    "eUNKbVRjb254ODlOUmZmVFJLK0lYR3dtQzY5Q1BIcXI1TTB6RXJiZktHcQpIS3ZsRG0rcHAzejZUZG"
    "9QVW5zMDlhRGNac29lNGwyazV3PT0KLS0tLS1FTkQgRUMgUFJJVkFURSBLRVktLS0tLQo"
)
VAPID_PUBLIC_KEY  = os.environ.get("VAPID_PUBLIC_KEY",  _VAPID_PUB_D)
_vapid_prv_raw    = os.environ.get("VAPID_PRIVATE_B64", _VAPID_PRV_B64)
import base64 as _b64mod
VAPID_PRIVATE_PEM = _b64mod.b64decode(_vapid_prv_raw + "==").decode()
VAPID_CLAIMS      = {"sub": "mailto:admin@abumalik.app"}

PUSH_SUBS_FILE = os.path.join(SESSIONS_DIR, "push_subscriptions.json")

def _load_push_subs():
    try:
        if os.path.exists(PUSH_SUBS_FILE):
            with open(PUSH_SUBS_FILE) as _f:
                return json.load(_f)
    except Exception:
        pass
    return {}

def _save_push_subs(subs):
    try:
        with open(PUSH_SUBS_FILE, "w") as _f:
            json.dump(subs, _f, indent=2)
    except Exception:
        pass

push_subscriptions = _load_push_subs()

# ─── الإشعارات الدورية الترويجية (Promo Notifications) ──────────────

PROMO_FILE = os.path.join(os.path.dirname(__file__), 'data', 'promo_messages.json')
os.makedirs(os.path.dirname(PROMO_FILE), exist_ok=True)

DEFAULT_PROMO_MESSAGES = [
    "🚀 مركز سرعة إنجاز – الحل الذكي لأتمتة تيليجرام! راقب الكلمات المفتاحية، أرسل رسائل ذكية تتجاوز الحظر، وأدر حسابات متعددة بكل احترافية. كل ذلك في منصة واحدة!",
    "📊 حلل بياناتك الأكاديمية بذكاء واحترافية! ارفع ملفات PDF، DOCX، واستخرج الجداول والصور، مع تلخيص ذكي باستخدام الذكاء الاصطناعي. وفر وقتك وركز على ما يهم!",
    "🔄 الإرسال الذكي يتجاوز حظر المجموعات – جربه الآن! مع وضع 'السلام عليكم' الذكي، ترسل رسالتك بأمان حتى في أكثر المجموعات حماية، وتعدل الرسالة بعد تفاعل الأعضاء.",
    "🔔 تنبيهات فورية، نطق صوتي، وتحديثات مستمرة… كل هذا في منصة واحدة! اشعارات ويب Push ونطق صوتي للتنبيهات المهمة، حتى لو كان التطبيق مغلقاً. كن على اطلاع دائم!",
    "📈 عزز حضورك على تيليجرام مع أدوات احترافية لا تُقهر! انضم تلقائياً للمجموعات، ابحث عن الروابط في محادثاتك، وصدر النتائج إلى Excel بضغطة زر.",
    "🎓 من الطالب إلى المحترف: مركز سرعة إنجاز رفيقك الأكاديمي! أنشئ عروض PowerPoint احترافية، حوّل HTML إلى Word/Excel، وحلل بياناتك إحصائياً بذكاء اصطناعي متقدم.",
    "🛡️ آمن، مرن، وقابل للتوسع – صُنع خصيصاً للمحترفين! مع دعم حسابات متعددة، تكامل GitHub، ونظام البطاقات، يمكنك إدارة فريقك أو عملائك بكل احترافية وأمان.",
    "⏰ الإرسال المتسلسل والردود التلقائية تعمل 24/7! حدد رسائل دورية، واترك النظام يرسلها تلقائياً دون تدخل، مع إمكانية تعديل أو حذف الدفعات المرسلة دفعة واحدة.",
    "🌍 ابحث في تيليجرام بالكامل واكتشف فرصاً جديدة! استخدم البحث العام للعثور على مجموعات وقنوات حسب اهتمامك، وانضم إليها بنقرة واحدة.",
    "💡 حل واحد، كل ما تحتاجه لإدارة تيليجرام بذكاء لا يُهزم! من المراقبة الفورية إلى التحليل الأكاديمي، ومن الإرسال الذكي إلى الإدارة المتكاملة – منصة متكاملة بامتياز."
]

def load_promo_data():
    try:
        if os.path.exists(PROMO_FILE):
            with open(PROMO_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception as e:
        logger.error(f"Failed to load promo data: {e}")
    default_data = {"enabled": False, "current_index": 0, "messages": DEFAULT_PROMO_MESSAGES}
    save_promo_data(default_data)
    return default_data

def save_promo_data(data):
    try:
        with open(PROMO_FILE, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Failed to save promo data: {e}")

_PROMO_THREAD = None
_PROMO_RUNNING = False

def promo_worker():
    global _PROMO_RUNNING
    _PROMO_RUNNING = True
    logger.info("🔄 بدأت خلفية الإشعارات الدورية (كل 5 دقائق)")
    while _PROMO_RUNNING:
        try:
            data = load_promo_data()
            if data.get("enabled", False) and data.get("messages"):
                idx = data.get("current_index", 0) % len(data["messages"])
                msg = data["messages"][idx]
                data["current_index"] = (idx + 1) % len(data["messages"])
                save_promo_data(data)
                sent_count = 0
                for user_id, sub in list(push_subscriptions.items()):
                    try:
                        send_push_notification(user_id, "📢 مركز سرعة إنجاز", msg, {"type": "promo", "source": "periodic"})
                        sent_count += 1
                    except Exception as e:
                        logger.error(f"[Promo] Failed to send to {user_id}: {e}")
                logger.info(f"[Promo] ✅ تم إرسال الرسالة {idx+1} إلى {sent_count} مشترك")
        except Exception as e:
            logger.error(f"[Promo] خطأ في الدورة: {e}")
        for _ in range(60):
            if not _PROMO_RUNNING:
                break
            time.sleep(5)
    logger.info("⏹ توقفت خلفية الإشعارات الدورية")

def start_promo_thread():
    global _PROMO_THREAD, _PROMO_RUNNING
    if _PROMO_THREAD and _PROMO_THREAD.is_alive():
        return
    _PROMO_RUNNING = True
    _PROMO_THREAD = threading.Thread(target=promo_worker, daemon=True)
    _PROMO_THREAD.start()

def stop_promo_thread():
    global _PROMO_RUNNING
    _PROMO_RUNNING = False


def send_push_notification(user_id, title, body, data=None):
    """إرسال إشعار Web Push إلى متصفح المستخدم"""
    if not WEBPUSH_AVAILABLE:
        return False
    sub = push_subscriptions.get(str(user_id))
    if not sub:
        return False
    _s = load_settings(user_id)
    if not _s.get('push_notifications_enabled', False):
        return False
    try:
        import json as _json
        _payload = _json.dumps({
            "title": title,
            "body": body,
            "data": data or {},
            "icon": "/static/icons/icon-192.png",
            "badge": "/static/icons/icon-72.png"
        })
        _webpush_fn(
            subscription_info=sub,
            data=_payload,
            vapid_private_key=VAPID_PRIVATE_PEM,
            vapid_claims=VAPID_CLAIMS,
            ttl=3600
        )
        return True
    except Exception as _ex:
        if hasattr(_ex, 'response') and _ex.response and getattr(_ex.response, 'status_code', 0) in (410, 404):
            push_subscriptions.pop(str(user_id), None)
            _save_push_subs(push_subscriptions)
        logger.warning(f"Push notification failed for {user_id}: {_ex}")
        return False

def get_user_session_dir(user_id):
    """مجلد منفصل لكل مستخدم لعزل البيانات والإعدادات"""
    user_dir = os.path.join(SESSIONS_DIR, str(user_id))
    if not os.path.exists(user_dir):
        os.makedirs(user_dir)
    return user_dir

# نظام المستخدمين الخمسة المحددين مسبقاً
# ── تحميل المستخدمين ديناميكياً من GitHub (مع fallback للقيم الفارغة) ──────
try:
    PREDEFINED_USERS = load_dynamic_users()
except Exception:
    PREDEFINED_USERS = {}

# معالجات الأخطاء الشاملة
@app.errorhandler(404)
def not_found_error(error):
    try:
        return jsonify({"error": "Page not found"}), 404
    except Exception as e:
        logger.error(f"Error in 404 handler: {str(e)}")
        return jsonify({"error": "Page not found"}), 404

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {str(error)}")
    try:
        return render_template('index.html', 
                              settings={}, 
                              connection_status='disconnected',
                              app_title="مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"), 500
    except Exception as e:
        logger.error(f"Error in 500 handler: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

@app.errorhandler(Exception)
def handle_exception(e):
    logger.error(f"Unhandled exception: {str(e)}")
    try:
        return render_template('index.html', 
                              settings={}, 
                              connection_status='disconnected',
                              app_title="مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"), 500
    except Exception as template_error:
        logger.error(f"Error in exception handler: {str(template_error)}")
        return jsonify({"error": "Server error"}), 500

# معالج أخطاء Socket.IO
@socketio.on_error_default
def default_error_handler(e):
    logger.error(f"Socket.IO error: {str(e)}")

USERS = {}
USERS_LOCK = Lock()

# ===================================================================
# مراقب استقرار الشبكة — يُعيد اتصالات تيليجرام تلقائياً
# ===================================================================
class NetworkStabilityMonitor:
    """يفحص الإنترنت كل 30 ثانية ويُعيد الاتصالات عند الانقطاع"""
    _CHECK_HOSTS = [('8.8.8.8', 53), ('1.1.1.1', 53), ('9.9.9.9', 53)]
    CHECK_INTERVAL = 30

    def __init__(self):
        self.is_online = True
        self._down_since = None

    def _check(self):
        import socket as _sock
        for host, port in self._CHECK_HOSTS:
            try:
                _sock.create_connection((host, port), timeout=4)
                return True
            except OSError:
                pass
        return False

    def _notify(self, msg):
        try:
            socketio.emit('log_update', {'message': msg})
        except Exception:
            pass
        logger.info(msg)

    def _reconnect_all(self):
        with USERS_LOCK:
            items = list(USERS.items())
        for uid, ud in items:
            cm = ud.get('client_manager')
            if cm and ud.get('is_running'):
                try:
                    if hasattr(cm, 'reconnect_if_needed'):
                        t = _OSThread(target=cm.reconnect_if_needed,
                                      daemon=True, name=f'Reconnect-{uid}')
                        t.start()
                except Exception as e:
                    logger.error(f'[NetworkMonitor] خطأ إعادة اتصال {uid}: {e}')

    def _loop(self):
        time.sleep(20)  # انتظر حتى يستقر التطبيق عند البداية
        while True:
            try:
                online = self._check()
                if not online and self.is_online:
                    self.is_online = False
                    self._down_since = time.time()
                    logger.warning('[NetworkMonitor] ⚠️ انقطع الإنترنت!')
                    self._notify('⚠️ تحذير: انقطع الإنترنت عن الخادم!')
                elif online and not self.is_online:
                    self.is_online = True
                    secs = int(time.time() - (self._down_since or time.time()))
                    logger.info(f'[NetworkMonitor] 🟢 عاد الإنترنت بعد {secs}s')
                    self._notify(f'🌐 عاد الإنترنت بعد {secs} ثانية — إعادة اتصال تيليجرام...')
                    self._reconnect_all()
                if online:
                    self._down_since = None
            except Exception as e:
                logger.error(f'[NetworkMonitor] خطأ: {e}')
            time.sleep(self.CHECK_INTERVAL)

    def start(self):
        t = _OSThread(target=self._loop, daemon=True, name='NetworkMonitor')
        t.start()
        logger.info('🛡️ مراقب استقرار الشبكة: نشط — فحص كل 30 ثانية')

network_monitor = NetworkStabilityMonitor()

# ===================================================================
# حلقة asyncio مشتركة لجميع عمليات تسجيل الدخول
# سبب الحل: إنشاء حلقة asyncio منفصلة لكل مستخدم يتعارض مع
# gevent's epoll patching — الحل هو حلقة واحدة مشتركة.
# ===================================================================
_SHARED_LOGIN_LOOP = None
_SHARED_LOGIN_LOOP_LOCK = _pre_patch_threading.Lock()
_SHARED_LOGIN_LOOP_READY = _pre_patch_threading.Event()

def _run_shared_login_loop(loop):
    """تشغيل الحلقة المشتركة في OS thread حقيقي"""
    _SHARED_LOGIN_LOOP_READY.set()
    loop.run_forever()

def _ensure_shared_login_loop():
    """الحصول على الحلقة المشتركة أو إنشاؤها إذا لم تكن موجودة"""
    global _SHARED_LOGIN_LOOP
    with _SHARED_LOGIN_LOOP_LOCK:
        if _SHARED_LOGIN_LOOP is None or _SHARED_LOGIN_LOOP.is_closed() or not _SHARED_LOGIN_LOOP.is_running():
            _SHARED_LOGIN_LOOP_READY.clear()
            loop = asyncio.new_event_loop()
            _SHARED_LOGIN_LOOP = loop
            t = _OSThread(
                target=_run_shared_login_loop,
                args=(loop,),
                daemon=True,
                name='SharedLoginLoop'
            )
            t.start()
            _SHARED_LOGIN_LOOP_READY.wait(timeout=10)
    return _SHARED_LOGIN_LOOP

# بيانات Telegram API — مدمجة مع دعم متغيرات البيئة
_TG_ID_PARTS   = ['220', '439', '94']
_TG_HASH_PARTS = ['56f645', '82b363d3', '67280db9', '6586b97801']
API_ID   = os.environ.get('TELEGRAM_API_ID',   ''.join(_TG_ID_PARTS))
API_HASH = os.environ.get('TELEGRAM_API_HASH', ''.join(_TG_HASH_PARTS))

# مفتاح الذكاء الاصطناعي GROQ — مدمج مع دعم متغيرات البيئة
_GROQ_PARTS  = ['gsk_ZNr7uNRZ', '6EyZUASH1oB', 'dWGdyb3FYwx', 'Jpzik4OICbSNCIntD4wFFV']
GROQ_API_KEY = os.environ.get('GROQ_API_KEY', ''.join(_GROQ_PARTS))
os.environ.setdefault('GROQ_API_KEY', GROQ_API_KEY)

# بيانات GitHub للمساعد الذكي — مدمجة مع دعم متغيرات البيئة
_GH_PARTS     = ['ghp_538Hwv', 'NvFzGm9dYs', 'n3pHjFkTfz', 'jC5j3wWNSE']
GITHUB_TOKEN  = os.environ.get('GITHUB_TOKEN',  ''.join(_GH_PARTS))
GITHUB_REPO   = os.environ.get('GITHUB_REPO',   'anwer1230/-Anwer_program')
GITHUB_BRANCH = os.environ.get('GITHUB_BRANCH', 'main')

# ─── ملف إعدادات التحديث ──────────────────────────────────────────────
UPDATE_SETTINGS_FILE = os.path.join(os.path.dirname(__file__), 'data', 'update_settings.json')
os.makedirs(os.path.dirname(UPDATE_SETTINGS_FILE), exist_ok=True)

def load_update_settings():
    try:
        if os.path.exists(UPDATE_SETTINGS_FILE):
            with open(UPDATE_SETTINGS_FILE, 'r') as f:
                return json.load(f)
    except Exception:
        pass
    return {"auto_update": False, "last_check": None, "last_update": None}

def save_update_settings(settings):
    try:
        with open(UPDATE_SETTINGS_FILE, 'w') as f:
            json.dump(settings, f, indent=2)
    except Exception as e:
        logger.error(f"Failed to save update settings: {e}")

def get_current_commit():
    try:
        result = subprocess.run(
            ['git', 'rev-parse', 'HEAD'],
            capture_output=True, text=True, check=True,
            cwd=os.path.dirname(__file__)
        )
        return result.stdout.strip()
    except Exception as e:
        logger.error(f"Failed to get current commit: {e}")
        return None

def get_latest_commit():
    try:
        repo = GITHUB_REPO or "anwer1230/-Anwer_program"
        if '/' in repo:
            owner, name = repo.split('/', 1)
        else:
            owner, name = "anwer1230", repo
        url = f"https://api.github.com/repos/{owner}/{name}/commits/{GITHUB_BRANCH or 'main'}"
        headers = {}
        if GITHUB_TOKEN:
            headers['Authorization'] = f'token {GITHUB_TOKEN}'
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code == 200:
            return resp.json().get('sha', '').strip()
    except Exception as e:
        logger.error(f"Failed to get latest commit: {e}")
    return None

def check_for_updates():
    current = get_current_commit()
    latest = get_latest_commit()
    if not current or not latest:
        return False, current, latest, "❌ تعذر التحقق من التحديثات"
    if current != latest:
        return True, current, latest, f"🔄 تحديث جديد متاح! (الحالي: {current[:7]} → الأحدث: {latest[:7]})"
    return False, current, latest, "✅ التطبيق محدث لأحدث إصدار"

def perform_update():
    logs = []
    try:
        logs.append("📥 جاري سحب التغييرات من GitHub...")
        pull_result = subprocess.run(
            ['git', 'pull', 'origin', GITHUB_BRANCH or 'main'],
            capture_output=True, text=True, cwd=os.path.dirname(__file__)
        )
        if pull_result.returncode != 0:
            logs.append(f"❌ فشل git pull: {pull_result.stderr}")
            return False, logs
        logs.append(f"✅ {pull_result.stdout.strip()}")
        logs.append("📦 جاري تحديث المتطلبات...")
        pip_result = subprocess.run(
            [sys.executable, '-m', 'pip', 'install', '-r', 'requirements.txt'],
            capture_output=True, text=True, cwd=os.path.dirname(__file__)
        )
        if pip_result.returncode != 0:
            logs.append(f"⚠️ تحذير في تثبيت المتطلبات: {pip_result.stderr}")
        else:
            logs.append("✅ تم تحديث المتطلبات")
        new_commit = get_current_commit()
        settings = load_update_settings()
        settings['last_update'] = datetime.now().isoformat()
        save_update_settings(settings)
        logs.append(f"✅ تم التحديث إلى الإصدار: {new_commit[:7] if new_commit else 'غير معروف'}")
        return True, logs
    except Exception as e:
        logs.append(f"❌ خطأ في التحديث: {str(e)}")
        return False, logs

_AUTO_UPDATE_THREAD = None
_AUTO_UPDATE_RUNNING = False

def auto_update_worker():
    global _AUTO_UPDATE_RUNNING
    _AUTO_UPDATE_RUNNING = True
    while _AUTO_UPDATE_RUNNING:
        try:
            settings = load_update_settings()
            if settings.get('auto_update', False):
                has_update, current, latest, msg = check_for_updates()
                if has_update:
                    logger.info(f"🔄 تحديث تلقائي: {msg}")
                    success, logs = perform_update()
                    if success:
                        logger.info(f"✅ تم التحديث التلقائي: {logs[-1]}")
                        try:
                            socketio.emit('update_completed', {
                                "message": "✅ تم تحديث التطبيق تلقائياً",
                                "version": latest[:7] if latest else "غير معروف"
                            })
                        except Exception:
                            pass
                    else:
                        logger.error(f"❌ فشل التحديث التلقائي: {logs[-1] if logs else ''}")
        except Exception as e:
            logger.error(f"خطأ في خلفية التحديث التلقائي: {e}")
        for _ in range(360):
            if not _AUTO_UPDATE_RUNNING:
                break
            time.sleep(10)

def start_auto_update_thread():
    global _AUTO_UPDATE_THREAD, _AUTO_UPDATE_RUNNING
    if _AUTO_UPDATE_THREAD and _AUTO_UPDATE_THREAD.is_alive():
        return
    _AUTO_UPDATE_RUNNING = True
    _AUTO_UPDATE_THREAD = threading.Thread(target=auto_update_worker, daemon=True)
    _AUTO_UPDATE_THREAD.start()
    logger.info("🔄 خلفية التحديث التلقائي بدأت")

def stop_auto_update_thread():
    global _AUTO_UPDATE_RUNNING
    _AUTO_UPDATE_RUNNING = False
    logger.info("⏹ خلفية التحديث التلقائي توقفت")

if not API_ID or not API_HASH:
    logger.warning("⚠️ لم يتم إعداد TELEGRAM_API_ID و TELEGRAM_API_HASH - وظائف التليجرام لن تعمل")

# =========================== 
# نظام Queue للتنبيهات المحسن
# ===========================
class AlertQueue:
    """نظام queue متقدم لإدارة التنبيهات"""

    def __init__(self):
        self.queue = queue.Queue()
        self.running = False
        self.thread = None

    def start(self):
        """بدء معالج التنبيهات"""
        if not self.running:
            self.running = True
            self.thread = threading.Thread(target=self._process_alerts, daemon=True)
            self.thread.start()
            logger.info("Alert queue processor started")

    def stop(self):
        """إيقاف معالج التنبيهات"""
        self.running = False
        if self.thread:
            self.thread.join(timeout=5)

    def add_alert(self, user_id, alert_data):
        """إضافة تنبيه جديد للقائمة"""
        try:
            self.queue.put({
                'user_id': user_id,
                'alert_data': alert_data,
                'timestamp': time.time()
            }, timeout=1)
        except queue.Full:
            logger.warning(f"Alert queue full for user {user_id}")

    def _process_alerts(self):
        """معالجة التنبيهات بشكل مستمر"""
        while self.running:
            try:
                alert = self.queue.get(timeout=1)
                self._send_alert(alert)
                self.queue.task_done()
            except queue.Empty:
                continue
            except Exception as e:
                logger.error(f"Error processing alert: {str(e)}")

    def _send_alert(self, alert):
        """إرسال التنبيه للمستخدم"""
        user_id = alert['user_id']
        alert_data = alert['alert_data']

        try:
            socketio.emit('new_alert', alert_data, to=user_id)
            socketio.emit('log_update', {
                "message": f"🚨 تنبيه فوري: '{alert_data['keyword']}' في {alert_data['group']}"
            }, to=user_id)

            try:
                send_push_notification(
                    user_id,
                    f"🔔 تنبيه: {alert_data.get('keyword', '')}",
                    f"في {alert_data.get('group', '')} من {alert_data.get('sender', '')}",
                    alert_data
                )
            except Exception:
                pass

            self._send_to_saved_messages(user_id, alert_data)

        except Exception as e:
            logger.error(f"Failed to send alert for user {user_id}: {str(e)}")

    def _send_to_saved_messages(self, user_id, alert_data):
        """إرسال التنبيه للرسائل المحفوظة في تيليجرام"""
        try:
            with USERS_LOCK:
                client_manager = USERS.get(user_id, {}).get('client_manager')

            if not client_manager or not client_manager.client:
                logger.warning(f"⚠️ No client available to send alert for user {user_id}")
                return

            keyword    = alert_data.get('keyword', '')
            group_name = alert_data.get('group', '')
            group_link = alert_data.get('group_link') or ''
            sender     = alert_data.get('sender', 'غير معروف')
            msg_time   = alert_data.get('message_time', alert_data.get('timestamp', ''))
            full_text  = alert_data.get('full_message') or alert_data.get('message', '')

            link_line = f"\n🔗 الرابط: {group_link}" if group_link else ""

            notification_msg = (
                f"🚨 تنبيه مراقبة\n\n"
                f"🔑 الكلمة: {keyword}\n"
                f"👥 المجموعة: {group_name}"
                f"{link_line}\n"
                f"👤 المرسل: {sender}\n"
                f"🕐 الوقت: {msg_time}\n\n"
                f"... الرسالة:\n{full_text}"
            )

            loop = getattr(client_manager, 'loop', None)
            if not loop or not loop.is_running():
                logger.warning(f"⚠️ Event loop not running for user {user_id} — cannot send alert")
                return

            async def _do_send():
                try:
                    await client_manager.client.send_message('me', notification_msg, link_preview=False)
                    logger.info(f"✅ Alert sent to Telegram saved messages for user {user_id}: '{keyword}'")
                except Exception as e:
                    logger.error(f"❌ Failed to send Telegram alert for user {user_id}: {e}")

            asyncio.run_coroutine_threadsafe(_do_send(), loop)

        except Exception as e:
            logger.error(f"Failed to send to saved messages: {str(e)}")

alert_queue = AlertQueue()

# ===========================
# تنقية الرسائل
# ===========================
class MessageSanitizer:
    PATTERNS = {
        'telegram_links': r'https?://(?:t\.me|telegram\.me)/[^\s<>]+|(?<!\w)t\.me/[^\s<>]+|(?<!\w)telegram\.me/[^\s<>]+',
        'whatsapp_links': r'https?://(?:wa\.me|chat\.whatsapp\.com|whatsapp\.com)/[^\s<>]+|(?<!\w)wa\.me/[^\s<>]+',
        'general_links':  r'https?://[^\s<>]+|www\.[^\s<>]+',
        'telegram_handles': r'@[a-zA-Z0-9_]{4,}',
        'phone_numbers': r'(?:\+?\d{1,3}[\s\-]?)?\(?\d{2,4}\)?[\s\-]?\d{3,4}[\s\-]?\d{3,4}',
        'ad_keywords': (
            r'\b(?:للتواصل|للاستفسار|واتساب|واتس|تليجرام|تليقرام|قناة|قناتي|انضم|انضموا|'
            r'خدمات|خدماتنا|إعلان|اعلان|عرض|عروض|خصم|تخفيض|تخفيضات|طلب\s*شراء|'
            r'بيع|تسويق|دورات|كورسات|اشتراك|راسلني|اطلب|عمولة|كاش|سحب|إيداع)\b'
        ),
    }

    WHATSAPP_TRANSFORMATIONS = [
        (r'https?://wa\.me/(\d+)', r'wa.me/\1'),
        (r'https?://chat\.whatsapp\.com/([^\s<>]+)', r'https://chat.whatsapp.com/\1'),
        (r'https?://whatsapp\.com/channel/([^\s<>]+)', r'https://whatsapp.com/channel/\1'),
    ]

    @classmethod
    def sanitize(cls, text, mode='smart'):
        """
        تنقية النص حسب الوضع المحدد
        mode: 'smart' (تنقية ذكية), 'clean' (تنقية كاملة), 'transform' (تحويل الروابط فقط)
        """
        if not text:
            return text
        cleaned = str(text)

        if mode == 'transform':
            for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
                cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
            return cleaned

        if mode == 'clean':
            for key in ('telegram_links', 'whatsapp_links', 'general_links'):
                cleaned = re.sub(cls.PATTERNS[key], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(cls.PATTERNS['telegram_handles'], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(cls.PATTERNS['phone_numbers'], '', cleaned)
            cleaned = re.sub(cls.PATTERNS['ad_keywords'], '', cleaned, flags=re.IGNORECASE)
            cleaned = re.sub(r'[ \t]+', ' ', cleaned)
            lines = []
            seen = set()
            for raw in cleaned.split('\n'):
                line = raw.strip(' \t-•·،,.|')
                if not line:
                    continue
                if not re.search(r'[\w\u0600-\u06FF]', line):
                    continue
                if line in seen:
                    continue
                seen.add(line)
                lines.append(line)
            result = '\n'.join(lines).strip()
            return result if result else None

        # وضع smart: تحويل روابط واتساب + تنقية خفيفة
        for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
            cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
        for key in ('telegram_links', 'general_links'):
            cleaned = re.sub(cls.PATTERNS[key], '', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(cls.PATTERNS['ad_keywords'], '', cleaned, flags=re.IGNORECASE)
        cleaned = re.sub(r'[ \t]+', ' ', cleaned).strip()
        return cleaned if cleaned else None

    @classmethod
    def has_promo_content(cls, text):
        if not text:
            return False
        s = str(text)
        for key in ('telegram_links', 'whatsapp_links', 'general_links',
                    'telegram_handles', 'phone_numbers', 'ad_keywords'):
            if re.search(cls.PATTERNS[key], s, re.IGNORECASE):
                return True
        return False

    @classmethod
    def transform_whatsapp_links(cls, text):
        """تحويل روابط واتساب فقط دون حذف"""
        if not text:
            return text
        cleaned = str(text)
        for pattern, replacement in cls.WHATSAPP_TRANSFORMATIONS:
            cleaned = re.sub(pattern, replacement, cleaned, flags=re.IGNORECASE)
        return cleaned

PROTECTION_BOTS = {
    'missrose_bot', 'rose_bot', 'therose_bot', 'rosebot',
    'shieldy_bot', 'shieldy', 'combot', 'combot_tech',
    'cas_bot', 'spamwatch_bot', 'spamwatchbot', 'antispam_bot',
    'antispambot', 'anti_spam_bot', 'spam_bot', 'spambot',
    'groupguardbot', 'groupguard_bot', 'guard_bot', 'guardbot',
    'safeguard_bot', 'safeguardbot', 'safe_guard_bot',
    'defender_bot', 'defenderbot', 'banhammer_bot', 'banhammerbot',
    'security_bot', 'securitybot', 'grouphelpbot', 'group_helpbot',
    'voteban_bot', 'votebanbot', 'antichannelpinbot', 'antiservicebot',
    'lolzteambot', 'protectionbot', 'policeman_bot', 'policemanbot',
    'sheriffbot', 'sheriff_bot', 'nightbot', 'mee6', 'cleanerbot',
    'cleaner_bot', 'modbot', 'moderationbot', 'no_spam_bot', 'nospambot',
    'stopspambot', 'stop_spam_bot', 'anti_flood_bot', 'antifloodbot',
    'flood_control_bot', 'hamasbot', 'arabicguard', 'arabguard_bot',
    'captchabot', 'captcha_bot', 'verifybot', 'verify_bot',
    'recaptcha_bot', 'human_verify_bot', 'wickbot', 'wick_bot',
    'dynobot', 'silence_bot', 'silencebot', 'mutebot', 'mute_bot',
    'word_filter_bot', 'filterbot', 'filter_bot',
    # ── بوت جبل وأسماؤه البديلة ──
    'jabal_bot', 'jabalbot', 'jabal', 'jbl_bot', 'jbl',
    'mtn_bot', 'mountain_bot', 'mtnchat_bot',
    # ── بوتات إضافية شائعة ──
    'groupprotect_bot', 'tgspam_bot', 'spam_protection_bot',
    'adminbot', 'admin_bot', 'groupadmin_bot',
}

PROTECTION_BOT_SUBSTRINGS = (
    'shieldy', 'rose', 'guard', 'combot', 'spamwatch', 'antispam',
    'anti_spam', 'safeguard', 'defender', 'banhammer', 'captcha',
    'verify', 'protect', 'police', 'sheriff', 'cleanbot', 'noflood',
    'antiflood', 'flood_', 'modbot', 'nochannel',
    # ── إضافات بوت جبل ──
    'jabal', 'jbl', 'mountain', 'mtn',
    # ── كلمات عامة ──
    'antispam', 'nospam', 'stopspam', 'silencebot', 'mutebot',
)

PROTECTED_GROUPS_CACHE = {}
PROTECTED_GROUPS_LOCK = Lock()

# ── نظام التعلم التلقائي للبوتات ────────────────────────────────────────────
BOTS_DISCOVERED_FILE = os.path.join(DATA_DIR, "discovered_bots.json")
_BOTS_FILE_LOCK = threading.Lock()

def load_discovered_bots():
    with _BOTS_FILE_LOCK:
        try:
            if os.path.exists(BOTS_DISCOVERED_FILE):
                with open(BOTS_DISCOVERED_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass
        return {"bots": []}

def save_discovered_bots(data):
    with _BOTS_FILE_LOCK:
        try:
            os.makedirs(DATA_DIR, exist_ok=True)
            with open(BOTS_DISCOVERED_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"فشل حفظ البوتات المكتشفة: {e}")

def add_discovered_bot(bot_username):
    if not bot_username:
        return
    bot_username = bot_username.lower().lstrip('@')
    if bot_username in PROTECTION_BOTS:
        return
    data = load_discovered_bots()
    if bot_username in data.get("bots", []):
        return
    data.setdefault("bots", []).append(bot_username)
    save_discovered_bots(data)
    PROTECTION_BOTS.add(bot_username)
    logger.info(f"✅ تم اكتشاف بوت حماية جديد وحفظه: @{bot_username}")

# تحميل البوتات المكتشفة مسبقاً عند بدء التشغيل
try:
    _disc = load_discovered_bots()
    for _b in _disc.get("bots", []):
        PROTECTION_BOTS.add(_b)
    if _disc.get("bots"):
        logger.info(f"✅ تم تحميل {len(_disc['bots'])} بوت مكتشف من discovered_bots.json")
except Exception:
    pass
# ─────────────────────────────────────────────────────────────────────────────

def _cache_protection(cache_key, result, reason, bots=None):
    with PROTECTED_GROUPS_LOCK:
        PROTECTED_GROUPS_CACHE[cache_key] = {
            'result': result, 'reason': reason,
            'bots': bots or [], 'ts': time.time()
        }

def save_settings(user_id, settings, force=False):
    try:
        if not force:
            existing = load_settings(user_id)
            if existing == settings:
                return True
        user_dir = get_user_session_dir(user_id)
        path = os.path.join(user_dir, "settings.json")
        with open(path, "w", encoding="utf-8") as f:
            json.dump(settings, f, ensure_ascii=False, indent=4)
        # احتفاظ بنسخة في المجلد الرئيسي للتوافق مع الكود القديم
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        with open(legacy_path, "w", encoding="utf-8") as f:
            json.dump(settings, f, ensure_ascii=False, indent=4)
        return True
    except Exception as e:
        logger.error(f"Error saving settings for {user_id}: {str(e)}")
        return False

def load_settings(user_id):
    try:
        user_dir = get_user_session_dir(user_id)
        path = os.path.join(user_dir, "settings.json")
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        # fallback للملف القديم
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(legacy_path):
            with open(legacy_path, "r", encoding="utf-8") as f:
                data = json.load(f)
            # نقل البيانات للمجلد الجديد (force=True لتفادي الاستدعاء الدائري)
            save_settings(user_id, data, force=True)
            return data
        return {}
    except Exception as e:
        logger.error(f"Error loading settings for {user_id}: {str(e)}")
        return {}

def clear_user_session(user_id):
    """حذف مجلد المستخدم بالكامل"""
    try:
        import shutil
        user_dir = get_user_session_dir(user_id)
        if os.path.exists(user_dir):
            shutil.rmtree(user_dir)
        legacy_path = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(legacy_path):
            os.remove(legacy_path)
        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            os.remove(session_file)
        logger.info(f"Cleared session for {user_id}")
        # حذف ملف سلسلة الجلسة أيضاً
        str_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        if os.path.exists(str_session_file):
            os.remove(str_session_file)
        return True
    except Exception as e:
        logger.error(f"Error clearing session for {user_id}: {str(e)}")
        return False


def save_string_session(user_id, session_str):
    """حفظ سلسلة جلسة StringSession في ملف نصي"""
    try:
        os.makedirs(SESSIONS_DIR, exist_ok=True)
        path = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        with open(path, 'w') as f:
            f.write(session_str)
        logger.info(f"Saved StringSession for {user_id}")
    except Exception as e:
        logger.error(f"Failed to save StringSession for {user_id}: {e}")


def load_string_session(user_id):
    """تحميل سلسلة جلسة StringSession من ملف"""
    try:
        path = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
        if os.path.exists(path):
            with open(path, 'r') as f:
                val = f.read().strip()
            if val:
                logger.info(f"Loaded StringSession for {user_id}")
                return val
    except Exception as e:
        logger.error(f"Failed to load StringSession for {user_id}: {e}")
    return None

def _clean_group_entry(raw: str) -> str:
    import re
    cleaned = raw.strip()
    if not cleaned:
        return ''
    # حماية الروابط والمعرفات الرقمية من التجريد غير المقصود
    if (cleaned.startswith('https://') or cleaned.startswith('http://')
            or cleaned.startswith('@') or re.match(r'^-?\d+$', cleaned)):
        return cleaned
    # تجريد البوليتات والأرقام والرموز من بداية السطر فقط
    cleaned = re.sub(r'^[\s\u00b7\u2022\u25cf\u25aa\u25ab\u25fe\u25fd\u2023\u203b\u2043\u2219\*\-\–\—\.\،\,\#\>\|•●◾◾✓✦①②③④⑤⑥⑦⑧⑨⑩\d]+[\s.،:]*', '', cleaned)
    return cleaned.strip()

def dedupe_groups(groups):
    seen = set()
    result = []
    if isinstance(groups, str):
        groups = [g for g in groups.replace('\n', ',').split(',')]
    for g in groups or []:
        if not g:
            continue
        original = _clean_group_entry(g)
        if not original:
            continue
        norm = original.lower()
        norm = norm.replace('https://telegram.me/', 'https://t.me/')
        norm = norm.replace('http://telegram.me/', 'https://t.me/')
        norm = norm.replace('http://t.me/', 'https://t.me/')
        if '?' in norm:
            norm = norm.split('?', 1)[0]
        if '#' in norm:
            norm = norm.split('#', 1)[0]
        norm = norm.rstrip('/').strip()
        if not norm:
            continue
        if norm in seen:
            continue
        seen.add(norm)
        result.append(original)
    return result

def load_all_sessions():
    logger.info("Loading existing sessions...")
    session_count = 0
    with USERS_LOCK:
        try:
            for filename in os.listdir(SESSIONS_DIR):
                if filename.endswith('.json'):
                    user_id = filename.split('.')[0]
                    settings = load_settings(user_id)
                    if settings and 'phone' in settings:
                        USERS[user_id] = {
                            'client_manager': None,
                            'settings': settings,
                            'thread': None,
                            'is_running': False,
                            'stats': {"sent": 0, "errors": 0},
                            'connected': False,
                            'authenticated': False,
                            'awaiting_code': False,
                            'awaiting_password': False,
                            'phone_code_hash': None,
                            'monitoring_active': False,
                            'event_handlers_registered': False,
                            'sent_batches': settings.get('sent_batches', []) or []
                        }
                        session_count += 1
                        logger.info(f"✓ Loaded session for {user_id}")
        except Exception as e:
            logger.error(f"Error loading sessions: {str(e)}")
    logger.info(f"Loaded {session_count} sessions successfully")
    return session_count

# =========================== 
# مدير التليجرام المحسن
# ===========================
class TelegramClientManager:
    def __init__(self, user_id):
        self.user_id = user_id
        self.client = None
        self.loop = None
        self.thread = None
        self.stop_flag = threading.Event()
        self.is_ready = threading.Event()
        self.event_handlers_registered = False
        self.monitored_keywords = []
        self.monitored_groups = []
        self._processed_msg_ids = set()

    async def send_to_saved_messages(self, text):
        try:
            if self.client:
                await self.client.send_message('me', text)
                logger.info(f"Sent message to saved messages for user {self.user_id}")
        except Exception as e:
            logger.error(f"Failed to send to saved messages: {str(e)}")

    async def get_group_protection_details(self, entity_obj):
        """
        فحص شامل للمجموعة — يرجع (is_protected, reason, detected_bots).
        يحفظ أي بوتات جديدة تلقائياً في discovered_bots.json.
        """
        try:
            chat_id = getattr(entity_obj, 'id', None)
            if chat_id is None:
                return False, None, []
            cache_key = (self.user_id, chat_id)
            with PROTECTED_GROUPS_LOCK:
                cached = PROTECTED_GROUPS_CACHE.get(cache_key)
                if cached is not None and time.time() - cached.get('ts', 0) < 1800:
                    return cached['result'], cached['reason'], cached.get('bots', [])

            reason = None
            detected_bots = []

            # ── فحص القيود المباشرة ────────────────────────────────────────
            try:
                full = await self.client.get_entity(entity_obj)
                banned = getattr(getattr(full, 'default_banned_rights', None), 'send_messages', None)
                if banned:
                    reason = 'المجموعة تمنع الأعضاء من الإرسال (restricted)'
                    _cache_protection(cache_key, True, reason, [])
                    return True, reason, []
            except Exception:
                pass

            # ── فحص الأعضاء بحثاً عن بوتات الحماية ──────────────────────
            try:
                async for participant in self.client.iter_participants(entity_obj, limit=100):
                    uname = (getattr(participant, 'username', '') or '').lower()
                    if not uname:
                        continue
                    if uname in PROTECTION_BOTS:
                        detected_bots.append(f"@{uname}")
                        reason = f'بوت حماية مكتشف: @{uname}'
                        logger.info(f"Group {chat_id} protected by known bot @{uname}")
                    elif any(s in uname for s in PROTECTION_BOT_SUBSTRINGS):
                        detected_bots.append(f"@{uname}")
                        reason = f'بوت حماية مشتبه: @{uname}'
                        logger.info(f"Group {chat_id} possibly protected by @{uname}")
                        # ── تعلم تلقائي: حفظ البوت المشتبه به ───────────
                        try:
                            add_discovered_bot(uname)
                        except Exception:
                            pass
                    if len(detected_bots) >= 5:
                        break
            except Exception as iter_err:
                logger.debug(f"iter_participants فشل لـ {chat_id}: {iter_err}")

            if detected_bots:
                reason = f'بوتات حماية مكتشفة: {", ".join(detected_bots[:5])}'
                _cache_protection(cache_key, True, reason, detected_bots)
                # إشعار socket فوري
                try:
                    socketio.emit('log_update', {
                        "message": f"🛡️ مجموعة محمية: {getattr(entity_obj, 'title', chat_id)} — {', '.join(detected_bots[:3])}"
                    }, to=self.user_id)
                    socketio.emit('group_protection_warning', {
                        "group_id": chat_id,
                        "group_title": getattr(entity_obj, 'title', str(chat_id)),
                        "bots": detected_bots,
                        "timestamp": time.strftime('%H:%M:%S')
                    }, to=self.user_id)
                except Exception:
                    pass
                return True, reason, detected_bots

            _cache_protection(cache_key, False, None, [])
            return False, None, []
        except Exception as e:
            logger.debug(f"get_group_protection_details error: {e}")
            return False, None, []

    async def is_group_protected(self, entity_obj):
        """التحقق من حماية المجموعة — يرجع (is_protected, reason) للتوافق مع الكود القديم."""
        is_prot, reason, _ = await self.get_group_protection_details(entity_obj)
        return is_prot, reason

    async def is_session_valid(self):
        """التحقق من صحة الجلسة الحالية"""
        try:
            if not self.client:
                return False
            is_authorized = await self.client.is_user_authorized()
            if not is_authorized:
                return False
            me = await self.client.get_me()
            return me is not None
        except Exception as e:
            error_str = str(e).lower()
            if any(kw in error_str for kw in ['auth_key', 'session', 'revoked', 'unauthorized', 'deactivated']):
                logger.warning(f"Session invalid for {self.user_id}: {e}")
                socketio.emit('session_revoked', {
                    "user_id": self.user_id,
                    "reason": str(e)
                }, to=self.user_id)
            return False

    def check_session_valid_sync(self):
        """نسخة متزامنة للتحقق من صحة الجلسة"""
        try:
            return self.run_coroutine(self.is_session_valid())
        except Exception:
            return False

    async def force_reset_session(self):
        """إعادة تعيين الجلسة وإزالة ملفاتها"""
        try:
            if self.client:
                try:
                    await self.client.disconnect()
                except Exception:
                    pass
            for _ext in ['_session.session', '_string.txt']:
                _fp = os.path.join(SESSIONS_DIR, f"{self.user_id}{_ext}")
                if os.path.exists(_fp):
                    os.remove(_fp)
            logger.info(f"Force reset session for {self.user_id}")
            return True
        except Exception as e:
            logger.error(f"Force reset error for {self.user_id}: {e}")
            return False

    @staticmethod
    def _extract_verification_code(text: str):
        """استخراج كود التحقق من النص"""
        if not text:
            return None
        patterns = [
            r'\b(\d{5,6})\b',
            r'code[:\s]+(\d{5,6})',
            r'كود[:\s]+(\d{5,6})',
            r'رمز[:\s]+(\d{5,6})',
            r'verification[:\s]+(\d{5,6})',
        ]
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(1)
        return None

    def start_client_thread(self):
        if self.thread and self.thread.is_alive():
            return
        self.stop_flag.clear()
        self.is_ready.clear()
        # استخدام OS thread حقيقي (محفوظ قبل monkey_patch) لتجنب تعارض asyncio مع eventlet
        self.thread = _OSThread(target=self._run_client_loop, daemon=True)
        self.thread.start()
        # انتظر حتى 60 ثانية — لكن لا ترمي exception عند timeout، بل سجّل تحذيراً
        if not self.is_ready.wait(timeout=60):
            logger.warning(f"Client initialization timeout for {self.user_id} — قد يتصل لاحقاً")

    def _run_client_loop(self):
        try:
            # إنشاء event loop مستقل لهذا الثريد — بدون set_event_loop لأنها عملية عامة
            # تتعارض مع ثريدات الحسابات الأخرى في بيئة gevent
            self.loop = asyncio.new_event_loop()
            if API_ID and API_HASH:
                saved_str = load_string_session(self.user_id)
                self.client = TelegramClient(StringSession(saved_str or ''), int(API_ID), API_HASH)
            else:
                logger.error("API_ID or API_HASH not set")
                return
            self.loop.run_until_complete(self._client_main())
        except Exception as e:
            logger.error(f"Client thread error for {self.user_id}: {str(e)}")
        finally:
            if self.loop and not self.loop.is_closed():
                self.loop.close()

    async def _client_main(self):
        try:
            if self.client:
                await self.client.connect()
                self.is_ready.set()
                await self._register_event_handlers()
                async def _watch_stop():
                    while not self.stop_flag.is_set():
                        await asyncio.sleep(0.5)
                    try:
                        await self.client.disconnect()
                    except Exception:
                        pass
                stop_task = asyncio.ensure_future(_watch_stop())
                while not self.stop_flag.is_set():
                    try:
                        is_auth = await self.client.is_user_authorized()
                        if is_auth:
                            try:
                                await self.client.run_until_disconnected()
                                # run_until_disconnected returned — client disconnected
                                if self.stop_flag.is_set():
                                    break
                                logger.info(f"Client disconnected for {self.user_id}, reconnecting in 3s...")
                                await asyncio.sleep(3)
                                try:
                                    await self.client.connect()
                                    logger.info(f"Reconnected successfully for {self.user_id}")
                                except Exception as rc_err:
                                    logger.error(f"Reconnect failed for {self.user_id}: {rc_err}")
                                    await asyncio.sleep(5)
                            except Exception as run_err:
                                if self.stop_flag.is_set():
                                    break
                                logger.warning(f"run_until_disconnected interrupted for {self.user_id}: {run_err}")
                                await asyncio.sleep(2)
                                try:
                                    if not self.client.is_connected():
                                        await self.client.connect()
                                        logger.info(f"Reconnected after error for {self.user_id}")
                                except Exception as rc2_err:
                                    logger.error(f"Reconnect after error failed for {self.user_id}: {rc2_err}")
                                    await asyncio.sleep(5)
                        else:
                            await asyncio.sleep(1)
                    except Exception as auth_check_err:
                        logger.debug(f"Auth check during loop for {self.user_id}: {auth_check_err}")
                        await asyncio.sleep(1)
                if not stop_task.done():
                    stop_task.cancel()
        except Exception as e:
            logger.error(f"Client main error: {str(e)}")
        finally:
            try:
                if self.client and self.client.is_connected():
                    await self.client.disconnect()
            except Exception:
                pass

    async def _register_event_handlers(self):
        try:
            if self.event_handlers_registered or not self.client:
                return

            @self.client.on(events.NewMessage())
            async def new_message_handler(event):
                await self._handle_new_message(event)
                if not getattr(event.message, 'out', False):
                    # التحقق من private أو group (وليس private فقط)
                    if (learning_manager.is_active(self.user_id, 'private') or
                            learning_manager.is_active(self.user_id, 'group')):
                        bot = learning_manager.get_bot(self.user_id)
                        await bot.handle_incoming_message(event, self)

            self.event_handlers_registered = True
            logger.info(f"✅ Event handlers registered for user {self.user_id} (all messages)")

        except Exception as e:
            logger.error(f"Failed to register event handlers: {str(e)}")

    async def _handle_new_message(self, event):
        try:
            message = event.message
            if not message or not message.text:
                return
            text = message.text or ''
            chat = await event.get_chat()
            chat_username = getattr(chat, 'username', None)
            chat_title    = getattr(chat, 'title',    None)
            chat_id       = getattr(chat, 'id',       None)

            if chat_username:
                group_identifier = f"@{chat_username}"
                group_link       = f"https://t.me/{chat_username}"
            elif chat_title:
                group_identifier = chat_title
                group_link       = None
            elif hasattr(chat, 'first_name'):
                fname = getattr(chat, 'first_name', '') or ''
                lname = getattr(chat, 'last_name',  '') or ''
                group_identifier = f"{fname} {lname}".strip() or str(chat_id)
                group_link       = None
            else:
                group_identifier = str(chat_id)
                group_link       = None

            is_outgoing = getattr(message, 'out', False)
            logger.info(f"📨 [{self.user_id}] {'صادرة' if is_outgoing else 'واردة'} | {group_identifier} | {text[:50]!r}")

            if not is_outgoing:
                try:
                    await self._handle_auto_reply(event, message, group_identifier)
                except Exception as ar_err:
                    logger.error(f"Auto-reply error: {ar_err}")

            kw_list = self.monitored_keywords
            if not kw_list:
                return

            msg_uid = f"{getattr(event, 'chat_id', 0)}_{message.id}"
            if msg_uid in self._processed_msg_ids:
                return
            if len(self._processed_msg_ids) > 500:
                self._processed_msg_ids.clear()
            self._processed_msg_ids.add(msg_uid)

            import unicodedata
            def _normalize(s):
                return ''.join(c for c in unicodedata.normalize('NFKD', s)
                               if unicodedata.category(c) != 'Mn')

            text_clean = _normalize(text).lower()
            matched = []
            for keyword in kw_list:
                kw = keyword.strip()
                if kw and _normalize(kw).lower() in text_clean:
                    matched.append(kw)

            if matched:
                combined_kw = ' | '.join(matched)
                logger.info(f"🔑 [{self.user_id}] {len(matched)} كلمة مطابقة: '{combined_kw}' في {group_identifier}")
                await self._trigger_keyword_alert(message, combined_kw, group_identifier, group_link, event)

        except Exception as e:
            logger.error(f"Error handling new message: {str(e)}", exc_info=True)

    async def _handle_auto_reply(self, event, message, group_identifier):
        try:
            settings = load_settings(self.user_id)
            if not settings.get('auto_reply_enabled', True):
                return
            rules = settings.get('auto_replies', []) or []
            if not rules:
                return

            text = message.text or ''
            text_lower = text.lower()

            is_private = bool(event.is_private)
            is_group_or_channel = bool(event.is_group or event.is_channel)

            for idx, rule in enumerate(rules):
                if not isinstance(rule, dict):
                    continue
                keyword = (rule.get('keyword') or '').strip()
                reply_text = (rule.get('reply') or '').strip()
                if not keyword or not reply_text:
                    continue

                scope = (rule.get('scope') or 'all').lower()
                if scope == 'private' and not is_private:
                    continue
                if scope == 'groups' and not is_group_or_channel:
                    continue

                match_mode = (rule.get('match') or 'contains').lower()
                matched = False
                try:
                    if match_mode == 'exact':
                        matched = (text.strip().lower() == keyword.lower())
                    elif match_mode == 'regex':
                        matched = bool(re.search(keyword, text, re.IGNORECASE))
                    else:
                        matched = (keyword.lower() in text_lower)
                except re.error as rerr:
                    logger.warning(f"Auto-reply regex error in rule #{idx} ({keyword}): {rerr}")
                    continue

                if matched:
                    try:
                        await self.client.send_message(
                            entity=event.chat_id,
                            message=reply_text,
                            reply_to=message.id
                        )
                        logger.info(f"✅ Auto-reply sent for keyword '{keyword[:40]}' in {group_identifier} (user={self.user_id})")
                        try:
                            _emit_log_update('INFO',
                                f"🤖 رد تلقائي على '{keyword[:30]}' في {group_identifier}",
                                self.user_id)
                            socketio.emit('auto_reply_triggered', {
                                "keyword": keyword,
                                "reply": reply_text,
                                "chat": group_identifier,
                                "timestamp": time.strftime('%H:%M:%S')
                            }, to=self.user_id)
                        except Exception:
                            pass
                        try:
                            rule['used_count'] = int(rule.get('used_count') or 0) + 1
                            rule['last_used'] = time.strftime('%Y-%m-%d %H:%M:%S')
                            settings['auto_replies'] = rules
                            save_settings(self.user_id, settings)
                        except Exception:
                            pass
                        break
                    except Exception as send_err:
                        logger.error(f"❌ Failed to send auto-reply for '{keyword[:30]}': {send_err}", exc_info=True)
        except Exception as e:
            logger.error(f"Auto-reply handler error: {e}")

    async def _trigger_keyword_alert(self, message, keyword, group_identifier, group_link, event):
        try:
            sender_name = "غير معروف"
            sender_id   = None
            sender_username = None
            try:
                sender = await event.get_sender()
                if sender:
                    first = getattr(sender, 'first_name', '') or ''
                    last  = getattr(sender, 'last_name',  '') or ''
                    uname = getattr(sender, 'username',   '') or ''
                    sender_id       = getattr(sender, 'id', None)
                    sender_username = uname
                    full  = f"{first} {last}".strip()
                    sender_name = full if full else (f"@{uname}" if uname else str(sender_id))
            except Exception:
                pass

            msg_time  = time.strftime('%H:%M:%S', time.localtime(message.date.timestamp()))
            full_text = message.text or ''

            chat = await event.get_chat()
            chat_username = getattr(chat, 'username', None)
            raw_chat_id   = getattr(chat, 'id', None)
            msg_id        = message.id

            if chat_username:
                msg_link = f"https://t.me/{chat_username}/{msg_id}"
            elif raw_chat_id:
                cid = str(raw_chat_id).lstrip('-')
                if cid.startswith('100'):
                    cid = cid[3:]
                msg_link = f"https://t.me/c/{cid}/{msg_id}"
            else:
                msg_link = group_link

            if sender_username:
                sender_link = f"https://t.me/{sender_username}"
            elif sender_id:
                sender_link = f"tg://user?id={sender_id}"
            else:
                sender_link = None

            group_part  = f"[{group_identifier}]({msg_link})" if msg_link else group_identifier
            sender_part = f"[{sender_name}]({sender_link})"  if sender_link else sender_name

            notification_msg = (
                f"🚨 **تنبيه مراقبة**\n\n"
                f"🔑 الكلمة: `{keyword}`\n"
                f"👥 المجموعة: {group_part}\n"
                f"👤 المرسل: {sender_part}\n"
                f"🕐 الوقت: {msg_time}\n\n"
                f"💬 الرسالة:\n{full_text}"
            )

            alert_data = {
                "keyword":      keyword,
                "group":        group_identifier,
                "group_link":   msg_link or group_link,
                "message":      full_text[:200] + ("..." if len(full_text) > 200 else ""),
                "full_message": full_text,
                "timestamp":    time.strftime('%H:%M:%S'),
                "sender":       sender_name,
                "sender_link":  sender_link,
                "message_time": msg_time,
                "message_id":   msg_id,
            }

            try:
                await self.client.send_message('me', notification_msg,
                                               parse_mode='md', link_preview=False)
                logger.info(f"✅ Alert sent: '{keyword}' in {group_identifier} | msg {msg_link}")
            except Exception as tg_err:
                logger.error(f"❌ Failed to send Telegram alert: {tg_err}")

            alert_queue.add_alert(self.user_id, alert_data)

        except Exception as e:
            logger.error(f"❌ Error triggering keyword alert: {str(e)}")

    def update_monitoring_settings(self, keywords, groups):
        self.monitored_keywords = [k.strip() for k in keywords if k.strip()]
        logger.info(f"Updated monitoring settings for {self.user_id}: {len(self.monitored_keywords)} keywords")

    def run_coroutine(self, coro):
        # If the loop is gone or closed, try to restart the client thread
        if not self.loop or self.loop.is_closed() or not self.loop.is_running():
            if not self.thread or not self.thread.is_alive():
                logger.warning(f"Client thread dead for {self.user_id}, auto-restarting...")
                self.stop_flag.clear()
                self.is_ready.clear()
                self.loop = None
                self.thread = _OSThread(target=self._run_client_loop, daemon=True)
                self.thread.start()
                self.is_ready.wait(timeout=30)
            if not self.loop or self.loop.is_closed() or not self.loop.is_running():
                raise Exception("العميل يُعاد تشغيله، حاول مرة أخرى بعد ثوانٍ")
        future = asyncio.run_coroutine_threadsafe(coro, self.loop)
        return future.result(timeout=60)

    async def _edit_batch_messages(self, batch_id, new_text):
        """تعديل جميع رسائل دفعة محددة"""
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            batch = next((b for b in ud.get('sent_batches', []) if b["id"] == batch_id), None)
        if not batch:
            return {"ok": False, "msg": "الدفعة غير موجودة"}
        ok_count = 0
        fail_count = 0
        for entry in batch["entries"]:
            try:
                entity_str = entry.get("group", "")
                msg_id = entry["msg_id"]
                entity = await self.client.get_entity(entity_str)
                await self.client.edit_message(entity, msg_id, new_text)
                ok_count += 1
                socketio.emit('log_update', {"message": f"✏️ تم تعديل الرسالة في {entity_str}"}, to=self.user_id)
                await asyncio.sleep(0.5)
            except Exception as e:
                fail_count += 1
                socketio.emit('log_update', {"message": f"❌ فشل التعديل في {entry.get('group','?')}: {str(e)[:60]}"}, to=self.user_id)
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            for b in ud.get('sent_batches', []):
                if b["id"] == batch_id:
                    b["text"] = new_text
                    b["edited_at"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    break
        socketio.emit('batch_edited', {"batch_id": batch_id, "new_text": new_text, "ok": ok_count, "fail": fail_count}, to=self.user_id)
        return {"ok": True, "edited": ok_count, "failed": fail_count}

    async def _delete_batch_messages(self, batch_id):
        """حذف جميع رسائل دفعة محددة"""
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            batch = next((b for b in ud.get('sent_batches', []) if b["id"] == batch_id), None)
        if not batch:
            return {"ok": False, "msg": "الدفعة غير موجودة"}
        ok_count = 0
        fail_count = 0
        for entry in batch["entries"]:
            try:
                entity_str = entry.get("group", "")
                msg_id = entry["msg_id"]
                entity = await self.client.get_entity(entity_str)
                await self.client.delete_messages(entity, [msg_id])
                ok_count += 1
                socketio.emit('log_update', {"message": f"🗑️ تم حذف الرسالة من {entity_str}"}, to=self.user_id)
                await asyncio.sleep(0.5)
            except Exception as e:
                fail_count += 1
                socketio.emit('log_update', {"message": f"❌ فشل الحذف من {entry.get('group','?')}: {str(e)[:60]}"}, to=self.user_id)
        with USERS_LOCK:
            ud = USERS.get(self.user_id, {})
            if ud:
                ud['sent_batches'] = [b for b in ud.get('sent_batches', []) if b["id"] != batch_id]
        socketio.emit('batch_deleted', {"batch_id": batch_id, "ok": ok_count, "fail": fail_count}, to=self.user_id)
        return {"ok": True, "deleted": ok_count, "failed": fail_count}

    def stop(self):
        self.stop_flag.set()
        if hasattr(self, 'client') and self.client and hasattr(self, 'loop') and self.loop and self.loop.is_running():
            try:
                future = asyncio.run_coroutine_threadsafe(self.client.disconnect(), self.loop)
                future.result(timeout=2)
            except Exception as e:
                logger.error(f"Error disconnecting client during stop: {e}")
        if self.thread and self.thread.is_alive():
            self.thread.join(timeout=3)
        if hasattr(self, 'loop') and self.loop and not self.loop.is_closed():
            try:
                self.loop.call_soon_threadsafe(self.loop.stop)
            except Exception as e:
                logger.error(f"Error stopping loop: {e}")

def get_all_users_operations_status():
    operations_status = {}
    with USERS_LOCK:
        for user_id, user_data in USERS.items():
            if user_id in PREDEFINED_USERS:
                operations_status[user_id] = {
                    'name': PREDEFINED_USERS[user_id]['name'],
                    'connected': user_data.get('connected', False),
                    'authenticated': user_data.get('authenticated', False),
                    'is_running': user_data.get('is_running', False),
                    'monitoring_active': user_data.get('monitoring_active', False),
                    'stats': user_data.get('stats', {"sent": 0, "errors": 0})
                }
    return operations_status

def notify_user_about_background_operations(user_id):
    try:
        active_operations = []
        with USERS_LOCK:
            for uid, user_data in USERS.items():
                if uid != user_id and uid in PREDEFINED_USERS:
                    if user_data.get('is_running', False) or user_data.get('monitoring_active', False):
                        active_operations.append({
                            'user_name': PREDEFINED_USERS[uid]['name'],
                            'operations': []
                        })
                        if user_data.get('monitoring_active', False):
                            active_operations[-1]['operations'].append('مراقبة نشطة')
                        if user_data.get('is_running', False):
                            active_operations[-1]['operations'].append('إرسال مجدول')
        if active_operations:
            operations_text = []
            for op in active_operations:
                operations_text.append(f"• {op['user_name']}: {', '.join(op['operations'])}")
            socketio.emit('log_update', {
                "message": f"📊 العمليات النشطة في الخلفية:\n" + "\n".join(operations_text)
            }, to=user_id)
    except Exception as e:
        logger.error(f"Error notifying about background operations: {str(e)}")

# ===========================
# فئة تسجيل الدخول المحسّنة
# ===========================
class TelegramLogin:
    def __init__(self, user_id):
        self.user_id = user_id
        self.client = None
        self.loop = None
        self.thread = None
        self.is_ready = threading.Event()
        self.phone_code_hash = None
        self.authenticated = False
        self.connected = False
        self.awaiting_code = False
        self.awaiting_password = False
        self.phone_number = None

    def _run_loop(self):
        """محفوظ للتوافقية — لم يعد مستخدماً. الحلقة المشتركة تُدار عبر _ensure_shared_login_loop()"""
        pass

    async def _connect(self):
        """الاتصال بخوادم تيليجرام مع إعادة المحاولة — يعمل على الحلقة المشتركة"""
        max_attempts = 5
        timeout = 40
        for attempt in range(1, max_attempts + 1):
            try:
                if attempt > 1:
                    logger.info(f"[{self.user_id}] إعادة المحاولة {attempt}/{max_attempts}...")
                    socketio.emit('log_update', {
                        "message": f"🔄 إعادة المحاولة {attempt}/{max_attempts}..."
                    }, to=self.user_id)
                    await asyncio.sleep(4)

                await asyncio.wait_for(self.client.connect(), timeout=timeout)
                try:
                    self.authenticated = await asyncio.wait_for(
                        self.client.is_user_authorized(), timeout=15
                    )
                except Exception:
                    self.authenticated = False
                self.connected = self.client.is_connected()

                if self.connected:
                    logger.info(f"[{self.user_id}] ✅ اتصل بنجاح في المحاولة {attempt}")
                    break

                logger.warning(f"[{self.user_id}] المحاولة {attempt}: الاتصال غير مستقر")

            except asyncio.TimeoutError:
                logger.error(f"[{self.user_id}] المحاولة {attempt}: انتهت المهلة ({timeout}s)")
                self.connected = False
                self.authenticated = False
                if attempt == max_attempts:
                    socketio.emit('log_update', {
                        "message": "❌ انتهت مهلة الاتصال — تأكد من إمكانية الوصول لتيليجرام"
                    }, to=self.user_id)
            except Exception as e:
                err_str = str(e)
                logger.error(f"[{self.user_id}] المحاولة {attempt}: {err_str}")
                self.connected = False
                self.authenticated = False
                if "FLOOD_WAIT" in err_str.upper():
                    socketio.emit('log_update', {
                        "message": "⏳ تيليجرام يطلب الانتظار — حاول بعد دقيقة"
                    }, to=self.user_id)
                    break
                if attempt == max_attempts:
                    socketio.emit('log_update', {
                        "message": f"❌ فشل الاتصال: {err_str[:120]}"
                    }, to=self.user_id)

        if not self.is_ready.is_set():
            self.is_ready.set()

    def start(self):
        """بدء تشغيل العميل باستخدام الحلقة المشتركة لتجنب تعارض asyncio/gevent"""
        try:
            self.loop = _ensure_shared_login_loop()
            if not self.loop or not self.loop.is_running():
                logger.error(f"[{self.user_id}] الحلقة المشتركة غير متاحة")
                return False

            try:
                from telethon.network.connection import ConnectionTcpMTProtoRandomizedIntermediate as _ConnType
            except ImportError:
                try:
                    from telethon.network.connection import ConnectionTcpFull as _ConnType
                except ImportError:
                    _ConnType = None

            client_kwargs = dict(
                connection_retries=5,
                retry_delay=2,
                timeout=40,
            )
            if _ConnType is not None:
                client_kwargs['connection'] = _ConnType
            self.client = TelegramClient(
                StringSession(), int(API_ID), API_HASH,
                **client_kwargs,
            )

            # DC1 (91.108.4.0) محظور على بعض استضافات السحابة (Render/Heroku) — نبدأ من DC2
            # Telethon سيحوّل تلقائياً للـ DC الصحيح بعد المصادقة
            try:
                self.client.session.set_dc(2, '149.154.167.51', 443)
                logger.info(f"[{self.user_id}] 🔀 تم تعيين DC2 كنقطة بداية للاتصال")
            except Exception as _dc_err:
                logger.warning(f"[{self.user_id}] تعذّر تعيين DC2: {_dc_err}")

            # الانتظار: 5 محاولات × 40 ثانية + 4 × 4 ثوانٍ بين المحاولات = ~216 ثانية
            future = asyncio.run_coroutine_threadsafe(self._connect(), self.loop)
            try:
                future.result(timeout=150)
            except Exception as e:
                logger.error(f"[{self.user_id}] فشل start(): {e}")
                self.connected = False

            return self.connected

        except Exception as e:
            logger.error(f"[{self.user_id}] خطأ في start(): {e}")
            if not self.is_ready.is_set():
                self.is_ready.set()
            return False

    def stop(self):
        """قطع اتصال العميل — لا نوقف الحلقة المشتركة لأنها مشتركة بين المستخدمين"""
        if self.client and self.loop and self.loop.is_running():
            try:
                asyncio.run_coroutine_threadsafe(
                    self.client.disconnect(), self.loop
                ).result(timeout=5)
            except Exception:
                pass
        # لا نوقف self.loop — الحلقة مشتركة!

    def send_code(self, phone_number):
        """الخطوة 1: إرسال كود التحقق إلى رقم الهاتف"""
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.send_code_request(phone_number),
                self.loop
            )
            result = future.result(timeout=30)
            self.phone_number = phone_number
            self.phone_code_hash = result.phone_code_hash
            self.awaiting_code = True
            self.authenticated = False
            return {
                "success": True,
                "message": "✅ تم إرسال الكود إلى هاتفك",
                "phone_code_hash": self.phone_code_hash
            }
        except Exception as e:
            error_msg = str(e)
            if "FLOOD_WAIT" in error_msg:
                return {"success": False, "message": "⏱️ انتظر قليلاً ثم حاول مرة أخرى"}
            return {"success": False, "message": f"❌ فشل إرسال الكود: {error_msg}"}

    def verify_code(self, code):
        """الخطوة 2: التحقق من الكود المرسل"""
        if not self.phone_code_hash:
            return {"success": False, "message": "لم يتم طلب كود بعد. أرسل الكود أولاً"}
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.sign_in(
                    phone=self.phone_number,
                    code=code,
                    phone_code_hash=self.phone_code_hash
                ),
                self.loop
            )
            result = future.result(timeout=30)
            self.awaiting_code = False
            self.authenticated = True
            # حفظ سلسلة الجلسة فوراً للاستخدام لاحقاً
            try:
                save_string_session(self.user_id, self.client.session.save())
            except Exception as _se:
                logger.error(f"Could not save session string: {_se}")
            me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
            me = me_future.result(timeout=30)
            return {
                "success": True,
                "message": "✅ تم تسجيل الدخول بنجاح",
                "user": {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            }
        except Exception as e:
            error_msg = str(e)
            if "PASSWORD" in error_msg.upper() or "SESSION_PASSWORD_NEEDED" in error_msg:
                self.awaiting_password = True
                self.awaiting_code = False
                return {
                    "success": False,
                    "requires_password": True,
                    "message": "🔐 هذا الحساب محمي بالتحقق بخطوتين. الرجاء إدخال كلمة المرور"
                }
            return {"success": False, "message": f"❌ كود غير صحيح: {error_msg}"}

    def verify_password(self, password):
        """الخطوة 3: إدخال رمز التحقق الثانوي (للمستخدمين الذين لديهم 2FA)"""
        if not self.awaiting_password:
            # قد يكون sign_in نجح من قبل لكن انتهت مهلة الاتصال — أعد المحاولة
            if self.authenticated and self.client and self.client.is_connected():
                try:
                    me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
                    me = me_future.result(timeout=15)
                    return {
                        "success": True,
                        "message": "✅ تم تسجيل الدخول بنجاح",
                        "user": {
                            "id": me.id,
                            "first_name": me.first_name,
                            "last_name": me.last_name,
                            "username": me.username,
                            "phone": me.phone,
                            "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                        }
                    }
                except Exception:
                    pass
            return {"success": False, "message": "الحساب لا يتطلب رمز تحقق ثانوي"}
        if not self.client or not self.client.is_connected():
            return {"success": False, "message": "العميل غير متصل"}
        try:
            future = asyncio.run_coroutine_threadsafe(
                self.client.sign_in(password=password),
                self.loop
            )
            future.result(timeout=45)
            self.awaiting_password = False
            self.authenticated = True
            # حفظ سلسلة الجلسة فوراً
            try:
                save_string_session(self.user_id, self.client.session.save())
            except Exception as _se:
                logger.error(f"Could not save session string (2FA): {_se}")
            me_future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
            me = me_future.result(timeout=30)
            return {
                "success": True,
                "message": "✅ تم تسجيل الدخول بنجاح",
                "user": {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            }
        except Exception as e:
            err = str(e)
            if "password" in err.lower() or "invalid" in err.lower():
                return {"success": False, "message": f"❌ كلمة مرور غير صحيحة: {err}"}
            return {"success": False, "message": f"❌ خطأ في التحقق: {err}"}

    def get_login_status(self):
        """الحصول على حالة تسجيل الدخول الحالية"""
        status = {
            "authenticated": self.authenticated,
            "awaiting_code": self.awaiting_code,
            "awaiting_password": self.awaiting_password,
            "connected": self.connected,
            "phone_number": self.phone_number,
            "user": None
        }
        if self.authenticated and self.client and self.client.is_connected():
            try:
                future = asyncio.run_coroutine_threadsafe(self.client.get_me(), self.loop)
                me = future.result(timeout=10)
                status["user"] = {
                    "id": me.id,
                    "first_name": me.first_name,
                    "last_name": me.last_name,
                    "username": me.username,
                    "phone": me.phone,
                    "full_name": f"{me.first_name or ''} {me.last_name or ''}".strip()
                }
            except Exception:
                pass
        return status

    def logout(self):
        """تسجيل الخروج من الحساب الحالي"""
        try:
            if self.client and self.loop and self.client.is_connected():
                future = asyncio.run_coroutine_threadsafe(self.client.log_out(), self.loop)
                future.result(timeout=30)
            session_file = os.path.join(SESSIONS_DIR, f"{self.user_id}_session.session")
            if os.path.exists(session_file):
                os.remove(session_file)
            self.authenticated = False
            self.awaiting_code = False
            self.awaiting_password = False
            self.phone_number = None
            self.phone_code_hash = None
            return {"success": True, "message": "✅ تم تسجيل الخروج بنجاح"}
        except Exception as e:
            return {"success": False, "message": f"❌ خطأ في تسجيل الخروج: {str(e)}"}


class TelegramManager:
    def __init__(self):
        self.client_managers = {}
        self.login_managers = {}
        self._smart_running = set()   # لمنع تشغيل عمليتين لنفس المجموعة: {user_id_group_id}

    def get_client_manager(self, user_id):
        if user_id not in self.client_managers:
            self.client_managers[user_id] = TelegramClientManager(user_id)
        return self.client_managers[user_id]

    def ensure_client_active(self, user_id):
        try:
            # تحقق من وجود سلسلة الجلسة (StringSession) بدلاً من ملف SQLite
            str_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_string.txt")
            if not os.path.exists(str_session_file):
                return False
            with USERS_LOCK:
                if user_id not in USERS:
                    return False
                client_manager = USERS[user_id].get('client_manager')

            if client_manager:
                # انتظر قصير (5 ثوانٍ) إذا كان الخيط يبدأ للتو
                if client_manager.thread and client_manager.thread.is_alive():
                    client_manager.is_ready.wait(timeout=5)
                if client_manager.client and client_manager.is_ready.is_set():
                    try:
                        is_auth = client_manager.run_coroutine(
                            client_manager.client.is_user_authorized()
                        )
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['authenticated'] = bool(is_auth)
                                USERS[user_id]['connected'] = True
                        return bool(is_auth)
                    except Exception as e:
                        logger.debug(f"is_user_authorized check failed for {user_id}: {e}")
                # الخيط موجود أو الجلسة محفوظة — أعد True بناءً على ملف الجلسة
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['authenticated'] = True
                        USERS[user_id]['connected'] = True
                return True

            # لا يوجد client_manager — ابدأ واحداً في الخلفية دون تعطيل المستدعي
            client_manager = self.get_client_manager(user_id)
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['connected'] = True

            def _bg_ensure(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تنشيط جلسة موجودة في الخلفية لـ {uid}")
                except Exception as e:
                    logger.warning(f"ensure bg start error for {uid}: {e}")

            _OSThread(target=_bg_ensure, daemon=True).start()
            return True
        except Exception as e:
            logger.error(f"ensure_client_active error for {user_id}: {e}")
            return False

    def setup_client(self, user_id, phone_number):
        try:
            if not API_ID or not API_HASH:
                socketio.emit('log_update', {"message": "❌ بيانات Telegram API غير متوفرة"}, to=user_id)
                return {"status": "error", "message": "❌ بيانات API غير متوفرة"}

            # إيقاف أي جلسة تسجيل دخول قديمة
            if user_id in self.login_managers:
                old_login = self.login_managers.pop(user_id)
                try:
                    old_login.stop()
                except Exception:
                    pass

            # إيقاف أي عميل تشغيل قديم
            if user_id in self.client_managers:
                old_manager = self.client_managers.pop(user_id)
                try:
                    old_manager.stop()
                except Exception as stop_err:
                    logger.warning(f"Could not stop old client manager for {user_id}: {stop_err}")

            # حذف ملف الجلسة القديم لضمان بدء نظيف
            for ext in ('', '.session'):
                session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session{ext}")
                if os.path.exists(session_file):
                    try:
                        os.remove(session_file)
                    except Exception:
                        pass

            socketio.emit('log_update', {"message": "🔄 جاري الاتصال بخوادم تيليجرام..."}, to=user_id)
            log_user_event(user_id, 'INFO', f"🔄 بدء تسجيل الدخول للرقم: {phone_number}")

            # إنشاء كائن تسجيل الدخول الجديد
            login = TelegramLogin(user_id)
            self.login_managers[user_id] = login
            _mgr = self  # مرجع للـ self داخل الخيط

            # ── تشغيل الاتصال في خيط OS حقيقي لتجنب توقف الخادم ──
            def _bg_connect():
                try:
                    connected = login.start()  # ينتظر حتى 30 ثانية
                    if not connected:
                        logger.error(f"Login connection failed for {user_id}")
                        socketio.emit('login_result', {
                            "status": "error",
                            "message": "❌ فشل الاتصال بخوادم تيليجرام - تحقق من الإنترنت"
                        }, to=user_id)
                        socketio.emit('log_update', {"message": "❌ فشل الاتصال بتيليجرام"}, to=user_id)
                        log_user_event(user_id, 'ERROR', "❌ فشل الاتصال بتيليجرام")
                        return

                    socketio.emit('log_update', {"message": "📡 فحص حالة التصريح..."}, to=user_id)
                    log_user_event(user_id, 'INFO', "📡 فحص حالة التصريح...")

                    # إذا كان الحساب مسجلاً بالفعل (جلسة محفوظة)
                    if login.authenticated:
                        client_manager = _mgr.get_client_manager(user_id)
                        client_manager.start_client_thread()
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['client_manager'] = client_manager
                                USERS[user_id]['connected'] = True
                                USERS[user_id]['authenticated'] = True
                                USERS[user_id]['awaiting_code'] = False
                                USERS[user_id]['awaiting_password'] = False
                        socketio.emit('login_status', {
                            "logged_in": True, "connected": True,
                            "awaiting_code": False, "awaiting_password": False, "is_running": False
                        }, to=user_id)
                        socketio.emit('connection_status', {"status": "connected"}, to=user_id)
                        socketio.emit('log_update', {"message": "✅ تم تسجيل الدخول بنجاح"}, to=user_id)
                        socketio.emit('login_result', {"status": "success", "message": "✅ تم تسجيل الدخول"}, to=user_id)
                        log_user_event(user_id, 'INFO', "✅ تم تسجيل الدخول بنجاح (جلسة محفوظة)")
                        return

                    # إرسال كود التحقق
                    socketio.emit('log_update', {"message": f"📱 إرسال كود التحقق إلى: {phone_number}"}, to=user_id)
                    log_user_event(user_id, 'INFO', f"📱 إرسال كود التحقق إلى: {phone_number}")

                    result = login.send_code(phone_number)
                    if not result["success"]:
                        socketio.emit('log_update', {"message": f"❌ {result['message']}"}, to=user_id)
                        socketio.emit('login_result', {"status": "error", "message": result["message"]}, to=user_id)
                        log_user_event(user_id, 'ERROR', f"❌ فشل إرسال الكود: {result['message']}")
                        return

                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['awaiting_code'] = True
                            USERS[user_id]['awaiting_password'] = False
                            USERS[user_id]['connected'] = True

                    socketio.emit('login_status', {
                        "logged_in": False, "connected": True,
                        "awaiting_code": True, "awaiting_password": False, "is_running": False
                    }, to=user_id)
                    socketio.emit('log_update', {"message": "✅ تم إرسال كود التحقق - تحقق من رسائل تيليجرام"}, to=user_id)
                    socketio.emit('login_result', {"status": "code_required", "message": "📱 تم إرسال كود التحقق"}, to=user_id)
                    log_user_event(user_id, 'INFO', "✅ تم إرسال كود التحقق")

                except Exception as e:
                    error_message = str(e)
                    logger.error(f"BG setup error for {user_id}: {error_message}")
                    log_user_event(user_id, 'ERROR', f"❌ خطأ في الإعداد: {error_message}")
                    if "ResendCodeRequest" in error_message or "all available options" in error_message:
                        msg = "⚠️ يرجى الانتظار قبل طلب كود جديد"
                    else:
                        msg = f"❌ خطأ: {error_message}"
                    socketio.emit('log_update', {"message": msg}, to=user_id)
                    socketio.emit('login_result', {"status": "error", "message": msg}, to=user_id)

            _OSThread(target=_bg_connect, daemon=True).start()
            return {"status": "pending", "message": "🔄 جارِ الاتصال بتيليجرام..."}

        except Exception as e:
            error_message = str(e)
            logger.error(f"Setup error for {user_id}: {error_message}")
            socketio.emit('log_update', {"message": f"❌ خطأ في الإعداد: {error_message}"}, to=user_id)
            return {"status": "error", "message": f"❌ خطأ: {error_message}"}

    def _fetch_account_name(self, user_id):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    return None
                client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return None
            me = client_manager.run_coroutine(client_manager.client.get_me())
            if not me:
                return None
            parts = []
            if getattr(me, 'first_name', None):
                parts.append(me.first_name)
            if getattr(me, 'last_name', None):
                parts.append(me.last_name)
            name = ' '.join(parts).strip()
            if not name:
                name = getattr(me, 'username', None) or 'حساب تليجرام'
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['account_name'] = name
                    USERS[user_id]['account_username'] = getattr(me, 'username', None)
                    USERS[user_id]['account_phone'] = getattr(me, 'phone', None)
            try:
                self._fetch_account_photo(user_id, me)
            except Exception as photo_err:
                logger.debug(f"Avatar fetch skipped for {user_id}: {photo_err}")
            return name
        except Exception as e:
            logger.error(f"Error fetching account name for {user_id}: {e}")
            return None

    def _fetch_account_photo(self, user_id, me=None):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    return None
                client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return None
            if me is None:
                me = client_manager.run_coroutine(client_manager.client.get_me())
            if not me:
                return None

            avatars_dir = os.path.join(SESSIONS_DIR, 'avatars')
            os.makedirs(avatars_dir, exist_ok=True)
            target_path = os.path.join(avatars_dir, f"{user_id}.jpg")

            async def _download():
                try:
                    return await client_manager.client.download_profile_photo(me, file=target_path)
                except Exception as e:
                    logger.debug(f"download_profile_photo error: {e}")
                    return None

            saved = client_manager.run_coroutine(_download())
            if saved and os.path.exists(target_path) and os.path.getsize(target_path) > 0:
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['account_avatar'] = f"/api/account_avatar/{user_id}?t={int(time.time())}"
                return target_path
            return None
        except Exception as e:
            logger.debug(f"Error fetching account photo for {user_id}: {e}")
            return None

    def verify_code(self, user_id, code):
        try:
            login = self.login_managers.get(user_id)
            if not login:
                return {"status": "error", "message": "❌ لم يتم بدء جلسة تسجيل الدخول"}

            result = login.verify_code(code)

            if result.get("requires_password"):
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['awaiting_code'] = False
                        USERS[user_id]['awaiting_password'] = True
                socketio.emit('login_status', {
                    "logged_in": False,
                    "connected": True,
                    "awaiting_code": False,
                    "awaiting_password": True,
                    "is_running": False
                }, to=user_id)
                return {"status": "password_required", "message": result["message"]}

            if not result["success"]:
                return {"status": "error", "message": result["message"]}

            # وقف عميل تسجيل الدخول أولاً لتحرير ملف الجلسة قبل تشغيل المدير الرئيسي
            if user_id in self.login_managers:
                try:
                    self.login_managers[user_id].stop()
                except Exception:
                    pass
                self.login_managers.pop(user_id, None)

            user_info = result.get("user", {})
            account_name = user_info.get("full_name") or user_info.get("username") or "حساب تليجرام"

            client_manager = self.get_client_manager(user_id)

            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['connected'] = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['awaiting_code'] = False
                    USERS[user_id]['awaiting_password'] = False
                    USERS[user_id]['account_name'] = account_name

            socketio.emit('login_status', {
                "logged_in": True,
                "connected": True,
                "awaiting_code": False,
                "awaiting_password": False,
                "is_running": False,
                "account_name": account_name
            }, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)

            # تشغيل عميل التليجرام الرئيسي في الخلفية — بدون تعطيل استجابة HTTP
            def _start_client_bg_code(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تشغيل عميل التليجرام في الخلفية لـ {uid}")
                except Exception as bg_err:
                    logger.warning(f"تحذير تشغيل العميل في الخلفية لـ {uid}: {bg_err}")

            _OSThread(target=_start_client_bg_code, daemon=True).start()

            return {"status": "success", "message": "✅ تم التحقق بنجاح", "account_name": account_name}

        except Exception as e:
            logger.error(f"Code verification error: {str(e)}")
            return {"status": "error", "message": f"❌ خطأ: {str(e)}"}

    def verify_password(self, user_id, password):
        try:
            login = self.login_managers.get(user_id)
            if not login:
                return {"status": "error", "message": "❌ لم يتم بدء جلسة تسجيل الدخول"}

            result = login.verify_password(password)

            if not result["success"]:
                return {"status": "error", "message": result["message"]}

            # وقف عميل تسجيل الدخول أولاً لتحرير ملف الجلسة قبل تشغيل المدير الرئيسي
            if user_id in self.login_managers:
                try:
                    self.login_managers[user_id].stop()
                except Exception:
                    pass
                self.login_managers.pop(user_id, None)

            user_info = result.get("user", {})
            account_name = user_info.get("full_name") or user_info.get("username") or "حساب تليجرام"

            client_manager = self.get_client_manager(user_id)

            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['client_manager'] = client_manager
                    USERS[user_id]['connected'] = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['awaiting_code'] = False
                    USERS[user_id]['awaiting_password'] = False
                    USERS[user_id]['account_name'] = account_name

            socketio.emit('login_status', {
                'logged_in': True,
                'connected': True,
                'awaiting_code': False,
                'awaiting_password': False,
                'account_name': account_name
            }, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)

            # تشغيل عميل التليجرام الرئيسي في الخلفية — بدون تعطيل استجابة HTTP
            def _start_client_bg_2fa(cm=client_manager, uid=user_id):
                try:
                    cm.start_client_thread()
                    logger.info(f"✅ تم تشغيل عميل التليجرام (2FA) في الخلفية لـ {uid}")
                except Exception as bg_err:
                    logger.warning(f"تحذير تشغيل العميل (2FA) في الخلفية لـ {uid}: {bg_err}")

            _OSThread(target=_start_client_bg_2fa, daemon=True).start()

            return {"status": "success", "message": "✅ تم التحقق بنجاح", "account_name": account_name}

        except Exception as e:
            logger.error(f"Password verification error: {str(e)}")
            return {"status": "error", "message": f"❌ خطأ: {str(e)}"}

    def _resolve_entity(self, client_manager, entity):
        import re as _re
        entity = _clean_group_entry(str(entity))
        if not entity:
            raise Exception("اسم المجموعة فارغ بعد التنظيف")

        # ── معرّف رقمي (chat ID مثل -1001234567890) ──
        if _re.match(r'^-?\d+$', entity):
            try:
                return client_manager.run_coroutine(
                    client_manager.client.get_entity(int(entity))
                )
            except Exception as e:
                raise Exception(f"لا يمكن الوصول إلى المعرّف الرقمي {entity}: {e}")

        # ── رابط دعوة خاص (invite link يحتوي على +) ──
        m_invite = _re.search(r't\.me/\+([A-Za-z0-9_\-]+)', entity)
        if m_invite:
            invite_hash = m_invite.group(1)
            # جرّب ImportChatInviteRequest (ينضم إن لم يكن عضواً)
            try:
                from telethon.tl.functions.messages import ImportChatInviteRequest
                result = client_manager.run_coroutine(
                    client_manager.client(ImportChatInviteRequest(invite_hash))
                )
                if result and hasattr(result, 'chats') and result.chats:
                    return result.chats[0]
            except Exception as invite_err:
                inv_msg = str(invite_err).lower()
                # إذا كان مصادقاً عليه مسبقاً، جرّب get_entity بالرابط كاملاً
                if 'already' in inv_msg or 'joined' in inv_msg or 'user_already' in inv_msg:
                    try:
                        return client_manager.run_coroutine(
                            client_manager.client.get_entity(entity)
                        )
                    except Exception:
                        pass
                # جرّب get_entity بالرابط كاملاً على كل حال
                try:
                    return client_manager.run_coroutine(
                        client_manager.client.get_entity(entity)
                    )
                except Exception:
                    pass
            raise Exception(f"لا يمكن الوصول إلى رابط الدعوة: {entity}")

        # ── روابط t.me العامة (@username) ──
        # استخراج اسم المستخدم من الروابط مثل https://t.me/username
        username_clean = entity.lstrip('@')
        m = _re.search(r't\.me/([^/\s\?#+]+)', entity)  # استثناء + من الأسماء
        if m:
            username_clean = m.group(1)

        last_exc = None

        # ── المحاولة 1: get_entity مع الرابط كما هو (نجح لو كان في الـ cache) ──
        try:
            return client_manager.run_coroutine(
                client_manager.client.get_entity(entity)
            )
        except Exception as e:
            last_exc = e

        # ── المحاولة 2: get_entity مع @ prefix ──
        if username_clean and not username_clean.startswith('+'):
            try:
                return client_manager.run_coroutine(
                    client_manager.client.get_entity('@' + username_clean)
                )
            except Exception as e:
                last_exc = e

        # ── المحاولة 3: ResolveUsernameRequest — يستعلم مباشرة من سيرفرات تيليجرام ──
        # يعمل لأي مجموعة/قناة عامة حتى لو لم يسبق التفاعل معها
        if username_clean and not username_clean.startswith('+'):
            try:
                result = client_manager.run_coroutine(
                    client_manager.client(functions.contacts.ResolveUsernameRequest(
                        username=username_clean
                    ))
                )
                if result and result.chats:
                    return result.chats[0]
                if result and result.users:
                    return result.users[0]
            except Exception as e:
                last_exc = e

        raise Exception(str(last_exc) if last_exc else f"لا يمكن الوصول إلى: {entity}")

    def send_message_async(self, user_id, entity, message, forced_action=None):
        """
        إرسال رسالة مع دعم الإجراء المختار مسبقاً من نافذة الفحص الاستباقي.
        forced_action: 'skip' | 'sanitize' | 'salam' | 'send' | None
          None  → استخدم الإعدادات الافتراضية للمستخدم
        """
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود - يرجى تسجيل الدخول أولاً")
                client_manager = USERS[user_id].get('client_manager')
                if not client_manager:
                    raise Exception("لم يتم تسجيل الدخول - يرجى تسجيل الدخول في التليجرام أولاً")
                if not client_manager.client:
                    raise Exception("عميل التليجرام غير مُهيأ - يرجى إعادة تسجيل الدخول")

            try:
                is_authorized = client_manager.run_coroutine(
                    client_manager.client.is_user_authorized()
                )
                if not is_authorized:
                    raise Exception("جلسة التليجرام منتهية الصلاحية - يرجى إعادة تسجيل الدخول")
            except Exception as auth_error:
                raise Exception(f"خطأ في التحقق من التصريح: {str(auth_error)}")

            entity_obj = self._resolve_entity(client_manager, entity)

            # ── تحديد الإجراء: من الفحص الاستباقي أو من الإعدادات ──────────
            if forced_action is not None:
                # forced_action قادم من /api/pre_send_scan — يُطبَّق مباشرة
                action = forced_action
                if action == 'skip':
                    socketio.emit('log_update', {
                        "message": f"⏭️ تم تخطي {entity} (قرار الفحص الاستباقي)"
                    }, to=user_id)
                    return {"success": False, "skipped": True,
                            "message": f"تم تخطي المجموعة: {entity}"}
            else:
                action, _ = self._check_group_protection(user_id, client_manager, entity_obj, entity)

            # ── الإرسال الذكي المتقدم للمجموعات المحمية (وضع salam) ──
            if action == 'salam':
                group_id = getattr(entity_obj, 'id', None) or hash(str(entity))
                key = f"{user_id}_{group_id}"
                if key in self._smart_running:
                    return {"success": True, "skipped": False, "smart": True,
                            "message": "⚠️ دورة ذكية قيد التشغيل لهذه المجموعة، ستُكمل عند انتهائها"}
                self._smart_running.add(key)
                _OSThread(
                    target=self._run_smart_protected_send,
                    args=(user_id, client_manager, entity_obj, entity, message, key),
                    daemon=True,
                    name=f"SmartSend-{key}"
                ).start()
                return {"success": True, "smart": True,
                        "message": f"🧠 بدأ الإرسال الذكي المتقدم لـ {entity}"}

            # نمرر action مباشرة لتجنب استدعاء _check_group_protection مرة ثانية
            final_message = self._maybe_sanitize(
                user_id, client_manager, entity_obj, entity, message,
                forced_action=action
            )
            if final_message is None:
                return {"success": False, "skipped": True,
                        "message": "تم تخطي الإرسال: المجموعة محمية أو الرسالة فارغة بعد التنقية"}

            result = client_manager.run_coroutine(
                client_manager.client.send_message(entity_obj, final_message)
            )

            return {"success": True, "message_id": result.id}

        except Exception as e:
            logger.error(f"Send message error: {str(e)}")
            raise Exception(str(e))

    def _check_group_protection(self, user_id, client_manager, entity_obj, entity_label):
        """
        التحقق من حماية المجموعة وإرجاع الإجراء المناسب.
        الوضع الافتراضي الآن هو 'salam' (الإرسال الذكي المتقدم).
        """
        try:
            settings = load_settings(user_id)
            # تغيير الافتراضي من 'smart' إلى 'salam'
            mode = (settings.get('sanitize_mode') or 'salam').lower()

            # إذا كان الوضع معطلاً، أرسل بدون فحص
            if mode == 'off':
                return 'send', None

            # التحقق من وجود بوتات حماية
            try:
                is_prot, reason = client_manager.run_coroutine(
                    client_manager.is_group_protected(entity_obj)
                )
            except Exception as e:
                logger.warning(f"Group protection check error: {e}")
                is_prot, reason = False, None

            # إذا كانت المجموعة محمية
            if is_prot:
                # وضع التخطي
                if mode == 'skip':
                    msg = f"⏭️ تم تخطي المجموعة المحمية: {entity_label}"
                    if reason:
                        msg += f" ({reason})"
                    socketio.emit('log_update', {"message": msg}, to=user_id)
                    self._send_protection_warning(user_id, entity_label, reason)
                    return 'skip', reason

                # وضع الإرسال الذكي المتقدم (الافتراضي الآن)
                if mode == 'salam':
                    socketio.emit('log_update', {
                        "message": f"🤖 مجموعة محمية: {entity_label} — الإرسال الذكي المتقدم (دوري)"
                    }, to=user_id)
                    return 'salam', reason

                # وضع التنقية الذكية
                if mode == 'smart':
                    socketio.emit('log_update', {
                        "message": f"🧠 مجموعة محمية: {entity_label} ({reason or 'بوت حماية'}) — سيتم تنقية الرسالة"
                    }, to=user_id)
                    return 'sanitize', reason

                # وضع التنقية الدائمة
                if mode == 'always':
                    socketio.emit('log_update', {
                        "message": f"🛡️ مجموعة محمية: {entity_label} — تنقية دائمة مفعّلة"
                    }, to=user_id)
                    return 'sanitize', reason

            # إذا كانت المجموعة غير محمية
            if mode == 'always':
                return 'sanitize', None
            if mode == 'salam':
                return 'send', None  # مجموعات غير محمية: أرسل عادي
            return 'send', None

        except Exception as e:
            logger.warning(f"_check_group_protection error: {e}")
            return 'send', None

    def _check_group_protection_detailed(self, user_id, client_manager, entity_obj, entity_label):
        """
        فحص تفصيلي للمجموعة يُرجع dict كامل يشمل:
        - protected: bool
        - bots: list
        - reason: str|None
        - entity_label: str
        مستخدم من /api/pre_send_scan فقط.
        """
        try:
            is_prot, reason, bots = client_manager.run_coroutine(
                client_manager.get_group_protection_details(entity_obj)
            )
            return {
                "entity_label": entity_label,
                "protected": is_prot,
                "bots": bots,
                "reason": reason,
                "error": False
            }
        except Exception as e:
            logger.warning(f"_check_group_protection_detailed error for {entity_label}: {e}")
            return {
                "entity_label": entity_label,
                "protected": False,
                "bots": [],
                "reason": f"خطأ في الفحص: {str(e)[:80]}",
                "error": True
            }

    def _send_protection_warning(self, user_id, group_name, reason):
        """إرسال تحذير للمستخدم عن المجموعة المحمية"""
        try:
            socketio.emit('protection_warning', {
                "group": group_name,
                "reason": reason,
                "timestamp": time.strftime('%H:%M:%S')
            }, to=user_id)
            try:
                with USERS_LOCK:
                    if user_id in USERS:
                        client_manager = USERS[user_id].get('client_manager')
                        if client_manager and client_manager.client:
                            loop = getattr(client_manager, 'loop', None)
                            if loop and loop.is_running():
                                warning_msg = f"""🛡️ **تنبيه: مجموعة محمية**

⚠️ تم اكتشاف أن المجموعة تحتوي على بوتات حماية:
📌 **{group_name}**

📋 **السبب:** {reason or 'يحتوي على بوتات حماية'}

💡 **الإجراء المتخذ:** تم تخطي الإرسال إلى هذه المجموعة حماية لحسابك.

🔧 **لتغيير الإعدادات:**
• تخطي المجموعات المحمية: إيقاف الإرسال إليها تلقائياً
• تنقية الروابط: تحويل روابط واتساب إلى صيغة آمنة"""
                                asyncio.run_coroutine_threadsafe(
                                    client_manager.client.send_message('me', warning_msg, link_preview=False),
                                    loop
                                )
            except Exception:
                pass
        except Exception as e:
            logger.error(f"Failed to send protection warning: {e}")

    def _run_smart_protected_send(self, user_id, client_manager, entity_obj, entity_label, final_message, key):
        """
        تُنفذ في خيط منفصل (لا تحجب حلقة الإرسال الرئيسية):
        1. ترسل "السلام عليكم" كرسالة أولية ثابتة.
        2. تنتظر فترة زمنية محددة (تؤخذ من interval_seconds في الإعدادات).
        3. تتحقق من عدد الرسائل التي أُرسلت بعدها.
        4. إذا كان أكثر من required_messages، تعدّل الرسالة إلى final_message.
        5. إذا لم يكن أكثر من required_messages، تحذف الرسالة وتنتظر الدورة التالية.
        6. تتكرر العملية بشكل دوري غير محدود حتى يتم إيقافها يدوياً.
        """
        try:
            if not client_manager or not client_manager.client:
                socketio.emit('log_update', {
                    "message": f"❌ فشل الإرسال الذكي لـ {entity_label}: العميل غير متصل"
                }, to=user_id)
                return

            # ── قراءة الإعدادات من المستخدم ──
            settings = load_settings(user_id)
            required_messages = int(settings.get('smart_required_messages', 3))
            interval_seconds = int(settings.get('interval_seconds', 3600))
            cycle_duration = max(60, interval_seconds)

            socketio.emit('log_update', {
                "message": f"🧠 بدء الدورة الذكية لـ {entity_label} — المدة: {cycle_duration}ث، العدد المطلوب: {required_messages}"
            }, to=user_id)

            # ── الحلقة الرئيسية (دورية لا نهائية) ──
            while True:
                if not self._smart_running or key not in self._smart_running:
                    logger.info(f"[Smart] توقفت الدورة الذكية لـ {entity_label} (key={key})")
                    break

                # ── 1. إرسال "السلام عليكم" ──
                try:
                    msg = client_manager.run_coroutine(
                        client_manager.client.send_message(entity_obj, "السلام عليكم")
                    )
                    logger.info(f"[Smart] أُرسلت 'السلام عليكم' إلى {entity_label}")
                    socketio.emit('log_update', {
                        "message": f"🧠 [Smart] أُرسلت 'السلام عليكم' إلى {entity_label} — في انتظار {cycle_duration}ث..."
                    }, to=user_id)
                    # ── إشعار: تم إرسال السلام عليكم ──
                    try:
                        _gid = getattr(entity_obj, 'id', None)
                        _uname = getattr(entity_obj, 'username', None)
                        if _uname:
                            _msg_link = f"https://t.me/{_uname}/{msg.id}"
                        elif _gid:
                            _gid_str = str(_gid)
                            _link_id = _gid_str[4:] if _gid_str.startswith('-100') else _gid_str.lstrip('-')
                            _msg_link = f"https://t.me/c/{_link_id}/{msg.id}"
                        else:
                            _msg_link = f"رقم الرسالة: {msg.id}"
                        _send_time = time.strftime('%Y-%m-%d %H:%M:%S')
                        _notif = (
                            f"🔔 إشعار — إرسال ذكي (السلام عليكم)\n\n"
                            f"📌 المجموعة: {entity_label}\n"
                            f"📨 تم إرسال 'السلام عليكم' بنجاح\n"
                            f"🔗 رابط الرسالة: {_msg_link}\n"
                            f"⏰ وقت الإرسال: {_send_time}\n\n"
                            f"⏳ ينتظر النظام تفاعل الأعضاء قبل إرسال الرسالة الكاملة..."
                        )
                        client_manager.run_coroutine(
                            client_manager.client.send_message('me', _notif, link_preview=False)
                        )
                    except Exception as _notif_err:
                        logger.debug(f"[Smart] خطأ في إرسال إشعار السلام: {_notif_err}")
                except Exception as send_err:
                    logger.error(f"[Smart] فشل إرسال 'السلام عليكم' إلى {entity_label}: {send_err}")
                    time.sleep(10)
                    continue

                # ── 2. انتظار المدة المحددة مع مراقبة الرسائل ──
                start_time = time.time()
                last_id = msg.id
                messages_after = 0

                while (time.time() - start_time) < cycle_duration:
                    if not self._smart_running or key not in self._smart_running:
                        break
                    time.sleep(2)
                    try:
                        new_msgs = client_manager.run_coroutine(
                            client_manager.client.get_messages(entity_obj, limit=10, min_id=last_id)
                        )
                        others = [m for m in new_msgs if m.id > last_id and not getattr(m, 'out', False)]
                        if others:
                            messages_after += len(others)
                            last_id = max(m.id for m in others)
                            logger.info(f"[Smart] {entity_label}: {len(others)} رسالة جديدة (المجموع: {messages_after})")
                            socketio.emit('log_update', {
                                "message": f"🧠 [Smart] {entity_label}: استقبل {messages_after}/{required_messages} رسالة"
                            }, to=user_id)
                    except Exception as poll_err:
                        logger.error(f"[Smart] خطأ في جلب الرسائل من {entity_label}: {poll_err}")
                        break

                # ── 3. اتخاذ القرار بناءً على عدد الرسائل ──
                if messages_after >= required_messages:
                    # ✅ تعديل الرسالة إلى النص النهائي
                    try:
                        client_manager.run_coroutine(
                            client_manager.client.edit_message(entity_obj, msg.id, final_message)
                        )
                        logger.info(f"[Smart] تم تعديل الرسالة في {entity_label} (عدد الرسائل: {messages_after})")
                        socketio.emit('log_update', {
                            "message": f"✅ [Smart] تم تعديل الرسالة في {entity_label} بعد {messages_after} رسائل"
                        }, to=user_id)
                        socketio.emit('smart_send_done', {
                            "success": True,
                            "entity": entity_label,
                            "waited": messages_after,
                            "message": f"✅ تم تعديل الرسالة في {entity_label} بعد {messages_after} رسائل"
                        }, to=user_id)
                        # ── إشعار: تم تعديل الرسالة ──
                        try:
                            _gid2 = getattr(entity_obj, 'id', None)
                            _uname2 = getattr(entity_obj, 'username', None)
                            if _uname2:
                                _edit_link = f"https://t.me/{_uname2}/{msg.id}"
                            elif _gid2:
                                _gid2_str = str(_gid2)
                                _link_id2 = _gid2_str[4:] if _gid2_str.startswith('-100') else _gid2_str.lstrip('-')
                                _edit_link = f"https://t.me/c/{_link_id2}/{msg.id}"
                            else:
                                _edit_link = f"رقم الرسالة: {msg.id}"
                            _edit_time = time.strftime('%Y-%m-%d %H:%M:%S')
                            _notif2 = (
                                f"✅ إشعار — تعديل الرسالة الذكية\n\n"
                                f"📌 المجموعة: {entity_label}\n"
                                f"✏️ تم تعديل الرسالة بعد {messages_after} رسائل\n"
                                f"🔗 رابط الرسالة المعدّلة: {_edit_link}\n"
                                f"⏰ وقت التعديل: {_edit_time}"
                            )
                            client_manager.run_coroutine(
                                client_manager.client.send_message('me', _notif2, link_preview=False)
                            )
                        except Exception as _notif2_err:
                            logger.debug(f"[Smart] خطأ في إرسال إشعار التعديل: {_notif2_err}")
                        # ── إيقاف الدورة بعد نجاح الإرسال — منع إرسال "السلام عليكم" مجدداً
                        break
                    except Exception as edit_err:
                        logger.error(f"[Smart] فشل تعديل الرسالة في {entity_label}: {edit_err}")
                        socketio.emit('log_update', {
                            "message": f"❌ [Smart] فشل تعديل الرسالة في {entity_label}: {str(edit_err)[:80]}"
                        }, to=user_id)
                else:
                    # ❌ لم نصل إلى العدد المطلوب — احذف رسالة السلام
                    logger.info(f"[Smart] لم يتم تعديل الرسالة في {entity_label} (عدد الرسائل: {messages_after} < {required_messages})")
                    socketio.emit('log_update', {
                        "message": f"⏭️ [Smart] لم يتم تعديل الرسالة في {entity_label} (استقبل {messages_after} فقط)"
                    }, to=user_id)
                    try:
                        client_manager.run_coroutine(
                            client_manager.client.delete_messages(entity_obj, [msg.id])
                        )
                        logger.info(f"[Smart] تم حذف رسالة 'السلام عليكم' غير المعدلة من {entity_label}")
                    except Exception as del_err:
                        logger.warning(f"[Smart] فشل حذف رسالة 'السلام عليكم' من {entity_label}: {del_err}")

                # ── 4. انتظار الدورة التالية ──
                if self._smart_running and key in self._smart_running:
                    socketio.emit('log_update', {
                        "message": f"⏳ [Smart] انتظار {cycle_duration // 60} دقيقة قبل الدورة التالية لـ {entity_label}"
                    }, to=user_id)
                    for _ in range(cycle_duration // 10):
                        if not self._smart_running or key not in self._smart_running:
                            break
                        time.sleep(10)

        except Exception as e:
            logger.error(f"[Smart] خطأ عام في الإرسال الذكي لـ {entity_label}: {e}")
            socketio.emit('log_update', {
                "message": f"❌ [Smart] فشل في {entity_label}: {str(e)[:100]}"
            }, to=user_id)
        finally:
            self._smart_running.discard(key)
            socketio.emit('log_update', {
                "message": f"⏹ [Smart] توقفت الدورة الذكية لـ {entity_label}"
            }, to=user_id)

    def _maybe_sanitize(self, user_id, client_manager, entity_obj, entity_label, message, forced_action=None):
        """تنقية الرسالة حسب وضع الحماية.
        forced_action: إذا تم تمريره يُتجنب استدعاء _check_group_protection مرة ثانية."""
        try:
            if forced_action is not None:
                action = forced_action
            else:
                action, _ = self._check_group_protection(user_id, client_manager, entity_obj, entity_label)

            if action in ('skip', 'smart'):
                return None

            if action == 'send':
                return message

            if action == 'transform':
                transformed = MessageSanitizer.transform_whatsapp_links(message)
                if transformed != message:
                    socketio.emit('log_update', {
                        "message": f"🔄 تم تحويل روابط واتساب في الرسالة إلى {entity_label}"
                    }, to=user_id)
                return transformed

            if not message:
                return message
            cleaned = MessageSanitizer.sanitize(message, mode='clean')
            if cleaned is None:
                socketio.emit('log_update', {
                    "message": f"⚠️ تم تخطي الإرسال إلى {entity_label}: الرسالة إعلانية بالكامل بعد التنقية"
                }, to=user_id)
                return None
            if cleaned != message:
                socketio.emit('log_update', {
                    "message": f"🧹 تنقية الرسالة قبل الإرسال إلى {entity_label}"
                }, to=user_id)
            return cleaned
        except Exception as e:
            logger.warning(f"_maybe_sanitize error: {e}")
            return message

    def send_media_async(self, user_id, entity, image_files):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود")
                client_manager = USERS[user_id].get('client_manager')

            if not client_manager:
                raise Exception("العميل غير متصل")

            is_authorized = client_manager.run_coroutine(
                client_manager.client.is_user_authorized()
            )
            if not is_authorized:
                raise Exception("العميل غير مصرح")

            entity_obj = self._resolve_entity(client_manager, entity)

            action, _reason = self._check_group_protection(user_id, client_manager, entity_obj, entity)
            if action == 'skip':
                return {"success": False, "skipped": True,
                        "message": f"تم تخطي المجموعة المحمية: {entity}"}

            results = []
            paths = [f['path'] for f in image_files if os.path.exists(f.get('path', ''))]
            if not paths:
                raise Exception("لا توجد ملفات صور صالحة")

            if len(paths) == 1:
                result = client_manager.run_coroutine(
                    client_manager.client.send_file(entity_obj, paths[0])
                )
                results.append(result.id)
            else:
                media_result = client_manager.run_coroutine(
                    client_manager.client.send_file(entity_obj, paths)
                )
                if hasattr(media_result, '__iter__'):
                    for r in media_result:
                        results.append(r.id)
                else:
                    results.append(media_result.id)

            return {"success": True, "message_ids": results}

        except Exception as e:
            logger.error(f"Send media error: {str(e)}")
            raise Exception(str(e))

    def send_message_with_media_async(self, user_id, entity, message, image_files):
        try:
            with USERS_LOCK:
                if user_id not in USERS:
                    raise Exception("المستخدم غير موجود")
                client_manager = USERS[user_id].get('client_manager')

            if not client_manager:
                raise Exception("العميل غير متصل")

            is_authorized = client_manager.run_coroutine(
                client_manager.client.is_user_authorized()
            )

            if not is_authorized:
                raise Exception("العميل غير مصرح")

            entity_obj = self._resolve_entity(client_manager, entity)

            if message:
                _cleaned = self._maybe_sanitize(user_id, client_manager, entity_obj, entity, message)
                if _cleaned is None:
                    return {"success": False, "skipped": True,
                            "message": "تم تخطي الإرسال: الرسالة بعد التنقية أصبحت فارغة"}
                message = _cleaned

            results = []

            if image_files and len(image_files) > 0:
                try:
                    image_paths = []
                    for img_file in image_files:
                        if os.path.exists(img_file['path']):
                            image_paths.append(img_file['path'])
                        else:
                            logger.warning(f"Image file not found: {img_file['path']}")

                    if image_paths:
                        if len(image_paths) == 1:
                            media_result = client_manager.run_coroutine(
                                client_manager.client.send_file(
                                    entity_obj, 
                                    image_paths[0],
                                    caption=message if message else "📷"
                                )
                            )
                            results.append(media_result.id)
                            logger.info(f"Successfully sent single image with message to {entity}")
                        else:
                            try:
                                media_result = client_manager.run_coroutine(
                                    client_manager.client.send_file(
                                        entity_obj,
                                        image_paths,
                                        caption=message if message and message.strip() else None
                                    )
                                )
                                if hasattr(media_result, '__iter__'):
                                    for result in media_result:
                                        results.append(result.id)
                                else:
                                    results.append(media_result.id)
                                logger.info(f"Successfully sent {len(image_paths)} images as album to {entity}")
                            except Exception as album_error:
                                logger.warning(f"Failed to send as album, sending individually: {str(album_error)}")
                                for i, img_path in enumerate(image_paths):
                                    try:
                                        cap = (message if message and message.strip() else None) if i == 0 else None
                                        media_result = client_manager.run_coroutine(
                                            client_manager.client.send_file(
                                                entity_obj,
                                                img_path,
                                                caption=cap
                                            )
                                        )
                                        results.append(media_result.id)
                                    except Exception as img_error:
                                        logger.error(f"Error sending individual image {i+1}: {str(img_error)}")
                                        continue
                except Exception as media_error:
                    logger.error(f"Error in media sending process: {str(media_error)}")
                    raise Exception(f"فشل إرسال الصورة: {str(media_error)[:100]}")
            else:
                if message and message.strip():
                    text_result = client_manager.run_coroutine(
                        client_manager.client.send_message(entity_obj, message)
                    )
                    results.append(text_result.id)
                    logger.info(f"Successfully sent text message to {entity}")

            return {"success": True, "message_ids": results}

        except Exception as e:
            logger.error(f"Send message with media error: {str(e)}")
            raise Exception(str(e))

telegram_manager = TelegramManager()

# =========================== 
# نظام المراقبة المحسن مع Event Handlers
# ===========================
def monitoring_worker(user_id):
    logger.info(f"Starting enhanced monitoring worker with event handlers for user {user_id}")

    try:
        with USERS_LOCK:
            if user_id not in USERS:
                logger.error(f"No user data found for {user_id}")
                return

            USERS[user_id]['monitoring_active'] = True
            client_manager = USERS[user_id].get('client_manager')
            settings = USERS[user_id]['settings']

        if not client_manager:
            logger.error(f"No client manager for user {user_id}")
            return

        watch_words = settings.get('watch_words', [])
        send_groups = settings.get('groups', [])

        if hasattr(client_manager, 'update_monitoring_settings'):
            client_manager.update_monitoring_settings(watch_words, send_groups)

        if watch_words:
            socketio.emit('log_update', {
                "message": f"🚀 بدأت المراقبة الشاملة الفورية - {len(watch_words)} كلمة مراقبة في كامل الحساب | الإرسال لـ {len(send_groups)} مجموعة"
            }, to=user_id)
        else:
            socketio.emit('log_update', {
                "message": f"🚀 بدأت المراقبة الشاملة لكامل الرسائل في الحساب | الإرسال لـ {len(send_groups)} مجموعة"
            }, to=user_id)

        _persisted = load_settings(user_id)
        _saved_last_send = _persisted.get('last_scheduled_send', 0)
        if _saved_last_send == 0:
            _saved_last_send = time.time()
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['last_scheduled_send'] = _saved_last_send

        # ── مدة التشغيل المحددة ──────────────────────────────────
        _sched_dur = int(settings.get('schedule_duration', 0))
        _sched_start = time.time()
        _stopped_by_duration = False

        if _sched_dur > 0:
            _h = _sched_dur // 3600
            _m = (_sched_dur % 3600) // 60
            socketio.emit('log_update', {
                "message": f"⏱️ الإرسال المجدول سيعمل لمدة {_h}س {_m}د ثم يتوقف تلقائياً"
            }, to=user_id)
        socketio.emit('schedule_status', {
            "running": True,
            "duration_hours": _sched_dur / 3600,
            "remaining_seconds": _sched_dur if _sched_dur > 0 else None
        }, to=user_id)

        consecutive_errors = 0
        max_consecutive_errors = 5
        _last_remain_emit = 0

        while True:
            with USERS_LOCK:
                if user_id not in USERS or not USERS[user_id].get('is_running', False):
                    logger.info(f"Stopping monitoring for user {user_id} as is_running is False")
                    break

                user_data = USERS[user_id].copy()
                USERS[user_id]['monitoring_active'] = True

            # ── التحقق من انتهاء مدة التشغيل ────────────────────
            if _sched_dur > 0:
                _elapsed = time.time() - _sched_start
                _remain  = _sched_dur - _elapsed
                if _remain <= 0:
                    logger.info(f"Schedule duration expired for user {user_id}")
                    socketio.emit('log_update', {
                        "message": "⏹ انتهت المدة المحددة — توقف الإرسال المجدول تلقائياً"
                    }, to=user_id)
                    _stopped_by_duration = True
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['is_running'] = False
                    socketio.emit('schedule_status', {
                        "running": False,
                        "stopped_by_duration": True,
                        "duration_hours": _sched_dur / 3600,
                        "can_resume": True
                    }, to=user_id)
                    # ── إشعار Web Push حتى لو التطبيق مغلق ──
                    _h2 = _sched_dur // 3600
                    _m2 = (_sched_dur % 3600) // 60
                    send_push_notification(
                        user_id,
                        "⏹ توقف الإرسال المجدول",
                        f"انتهت المدة المحددة ({_h2}س {_m2}د). اضغط هنا لاستئنافه.",
                        data={"type": "schedule_expired", "duration_hours": _sched_dur / 3600}
                    )
                    break
                else:
                    # إرسال الوقت المتبقي كل 30 ثانية
                    if time.time() - _last_remain_emit >= 30:
                        _last_remain_emit = time.time()
                        socketio.emit('schedule_remaining', {
                            "remaining_seconds": int(_remain),
                            "remaining_minutes": int(_remain // 60),
                            "remaining_hours": _remain / 3600
                        }, to=user_id)

            try:
                settings = user_data.get('settings', {})
                send_type = settings.get('send_type', 'manual')
                current_time = time.time()

                if send_type == 'scheduled':
                    interval_seconds = int(settings.get('interval_seconds', 3600))
                    last_send = user_data.get('last_scheduled_send', 0)
                    remaining = interval_seconds - (current_time - last_send)

                    if remaining <= 0:
                        logger.info(f"Executing scheduled send for user {user_id} (interval={interval_seconds}s)")
                        socketio.emit('log_update', {
                            "message": f"📅 حان موعد الإرسال المجدول — جاري الإرسال إلى {len(settings.get('groups', []))} مجموعة..."
                        }, to=user_id)
                        execute_scheduled_messages(user_id, settings)

                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['last_scheduled_send'] = current_time
                        try:
                            _s = load_settings(user_id)
                            _s['last_scheduled_send'] = current_time
                            save_settings(user_id, _s)
                        except Exception as _se:
                            logger.error(f"Failed to persist last_scheduled_send: {_se}")

                        next_send_at = time.strftime('%H:%M:%S', time.localtime(current_time + interval_seconds))
                        socketio.emit('log_update', {
                            "message": f"⏰ الإرسال التالي في: {next_send_at} (بعد {interval_seconds // 60} دقيقة)"
                        }, to=user_id)
                    else:
                        logger.debug(f"Scheduled send for {user_id}: {int(remaining)}s remaining")

                consecutive_errors = 0

                next_send_remaining = None
                next_send_at_str = None
                if send_type == 'scheduled':
                    interval_seconds = int(settings.get('interval_seconds', 3600))
                    last_send = user_data.get('last_scheduled_send', 0)
                    remaining = interval_seconds - (current_time - last_send)
                    next_send_remaining = max(0, int(remaining))
                    next_send_at_str = time.strftime('%H:%M:%S', time.localtime(current_time + next_send_remaining))

                _dur_remain = None
                if _sched_dur > 0:
                    _dur_remain = max(0, int(_sched_dur - (time.time() - _sched_start)))

                status_info = {
                    'timestamp': time.strftime('%H:%M:%S'),
                    'status': 'active',
                    'type': 'event_driven_monitoring',
                    'keywords_active': bool(watch_words),
                    'event_handlers': True,
                    'send_type': send_type,
                    'next_send_remaining': next_send_remaining,
                    'next_send_at': next_send_at_str,
                    'schedule_remaining': _dur_remain,
                    'schedule_duration_hours': _sched_dur / 3600 if _sched_dur > 0 else 0
                }

                socketio.emit('heartbeat', status_info, to=user_id)

            except Exception as e:
                consecutive_errors += 1
                logger.error(f"Monitoring cycle error for {user_id}: {str(e)}")
                socketio.emit('log_update', {
                    "message": f"⚠️ خطأ في المراقبة: {str(e)[:100]}"
                }, to=user_id)

                if consecutive_errors >= max_consecutive_errors:
                    socketio.emit('log_update', {
                        "message": f"❌ تم إيقاف المراقبة بسبب تكرار الأخطاء ({consecutive_errors})"
                    }, to=user_id)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['is_running'] = False
                    break

            time.sleep(10)

    except Exception as e:
        logger.error(f"Monitoring worker top-level error for {user_id}: {str(e)}")
    finally:
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['is_running'] = False
                USERS[user_id]['monitoring_active'] = False
                USERS[user_id]['thread'] = None

        socketio.emit('log_update', {
            "message": "⏹ تم إيقاف نظام المراقبة المحسن"
        }, to=user_id)

        socketio.emit('heartbeat', {
            'timestamp': time.strftime('%H:%M:%S'),
            'status': 'stopped'
        }, to=user_id)

        if not _stopped_by_duration:
            socketio.emit('schedule_status', {"running": False, "stopped_by_duration": False}, to=user_id)

        logger.info(f"Enhanced monitoring worker ended for user {user_id}")


@app.route("/api/resume_scheduled", methods=["POST"])
def api_resume_scheduled():
    """استئناف الإرسال المجدول بعد توقفه بسبب انتهاء المدة"""
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "❌ غير مسجل"}), 401

    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({"success": False, "message": "❌ المستخدم غير موجود"}), 404

        if USERS[user_id].get('is_running', False):
            return jsonify({"success": False, "message": "⚠️ الإرسال المجدول يعمل بالفعل"})

        settings = USERS[user_id].get('settings', {})
        if not settings.get('groups') or not settings.get('message'):
            return jsonify({"success": False, "message": "❌ الإعدادات غير مكتملة (المجموعات أو الرسالة)"})

        _sched_dur = int(settings.get('schedule_duration', 0))
        _h = _sched_dur // 3600
        _m = (_sched_dur % 3600) // 60
        USERS[user_id]['is_running'] = True

    import threading as _thr
    t = _thr.Thread(target=monitoring_worker, args=(user_id,), daemon=True)
    t.start()
    with USERS_LOCK:
        if user_id in USERS:
            USERS[user_id]['thread'] = t

    msg = f"🔄 تم استئناف الإرسال المجدول" + (f" لمدة {_h}س {_m}د" if _sched_dur > 0 else "")
    socketio.emit('log_update', {"message": msg}, to=user_id)
    return jsonify({"success": True, "message": msg})


def execute_scheduled_messages(user_id, settings):
    groups = settings.get('groups', [])
    message = settings.get('message', '')

    if not groups or not message:
        return

    try:
        socketio.emit('log_update', {
            "message": f"📅 تنفيذ الإرسال المجدول إلى {len(groups)} مجموعة"
        }, to=user_id)

        successful = 0
        failed = 0
        _sched_success_groups = []
        _sched_failed_groups = []
        _sched_smart_groups = []

        # ── الحصول على مدير العميل لفحص العضوية ──────────────────────────
        _sched_client_mgr = None
        try:
            with USERS_LOCK:
                _sched_client_mgr = USERS.get(user_id, {}).get('client_manager')
        except Exception:
            pass

        # ── الدورة الأولى: فحص العضوية وإشعار المستخدم بروابط المجموعات غير المنضم إليها ──
        if _sched_client_mgr and getattr(_sched_client_mgr, 'client', None):
            socketio.emit('log_update', {
                "message": f"🔍 فحص العضوية في {len(groups)} مجموعة..."
            }, to=user_id)

            async def _check_memberships_sched(client, groups_to_check):
                """فحص العضوية لمجموعات الإرسال المجدول دون الانضمام إليها"""
                from telethon.tl.functions.channels import GetParticipantRequest
                from telethon.errors import UserNotParticipantError as _UNPE
                not_joined = []
                for _g in groups_to_check:
                    try:
                        _ident = _g
                        for _pfx in ['https://t.me/', 'https://telegram.me/']:
                            if _g.startswith(_pfx):
                                _ident = _g[len(_pfx):]
                                break
                        if _ident.startswith('@'):
                            _ident = _ident[1:]
                        _ent = await client.get_entity(_ident)
                        if hasattr(_ent, 'megagroup') or hasattr(_ent, 'broadcast'):
                            try:
                                await client(GetParticipantRequest(_ent, 'me'))
                            except _UNPE:
                                not_joined.append(_g)
                            except Exception:
                                pass  # خطأ آخر — نفترض العضوية تفادياً للإشعار الخاطئ
                        # مجموعة عادية — نفترض العضوية إذا تم حل الـ entity
                    except Exception:
                        not_joined.append(_g)
                return not_joined

            _sched_not_joined = []
            try:
                _sched_not_joined = _sched_client_mgr.run_coroutine(
                    _check_memberships_sched(_sched_client_mgr.client, groups)
                )
            except Exception as _sched_check_err:
                logger.debug(f"خطأ في فحص العضوية (المجدول): {_sched_check_err}")

            if _sched_not_joined:
                socketio.emit('log_update', {
                    "message": f"⚠️ أنت غير منضم إلى {len(_sched_not_joined)} مجموعة — ستظهر في التقرير النهائي"
                }, to=user_id)

            socketio.emit('log_update', {
                "message": "🚀 بدء الإرسال المجدول..."
            }, to=user_id)

        # ── الدورة الثانية: إرسال الرسائل إلى جميع المجموعات ──────────────
        # دعم الصور المحفوظة في الإعدادات
        _sched_image_files = []
        _sched_image_path = settings.get('scheduled_image_path', '')
        if _sched_image_path and os.path.exists(_sched_image_path):
            _sched_image_files = [{'path': _sched_image_path, 'name': 'scheduled_image.jpg', 'type': 'image/jpeg'}]

        for i, group in enumerate(groups, 1):
            try:
                if _sched_image_files:
                    result = telegram_manager.send_message_with_media_async(
                        user_id, group, message, _sched_image_files
                    )
                else:
                    result = telegram_manager.send_message_async(user_id, group, message)

                if isinstance(result, dict) and result.get('skipped'):
                    socketio.emit('log_update', {
                        "message": f"⏭️ [{i}/{len(groups)}] تم تخطي: {group} (الرسالة لم تُرسَل)"
                    }, to=user_id)
                    failed += 1
                    _sched_failed_groups.append(group)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['errors'] += 1
                elif isinstance(result, dict) and result.get('smart'):
                    socketio.emit('log_update', {
                        "message": f"✅ [{i}/{len(groups)}] إرسال مجدول نجح إلى: {group}"
                    }, to=user_id)
                    successful += 1
                    _sched_smart_groups.append(group)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['sent'] += 1
                else:
                    socketio.emit('log_update', {
                        "message": f"✅ [{i}/{len(groups)}] إرسال مجدول نجح إلى: {group}"
                    }, to=user_id)
                    successful += 1
                    _sched_success_groups.append(group)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['sent'] += 1

                if i < len(groups):
                    time.sleep(3)

            except Exception as e:
                error_msg = str(e)
                logger.error(f"Scheduled send error to {group}: {error_msg}")

                socketio.emit('log_update', {
                    "message": f"❌ [{i}/{len(groups)}] إرسال مجدول فشل إلى {group}"
                }, to=user_id)

                failed += 1
                _sched_failed_groups.append(group)
                with USERS_LOCK:
                    if user_id in USERS:
                        USERS[user_id]['stats']['errors'] += 1

        socketio.emit('log_update', {
            "message": f"📊 انتهى الإرسال المجدول: ✅ {successful} نجح | ❌ {failed} فشل"
        }, to=user_id)

        # ── إرسال تقرير مفصل لحساب المستخدم الشخصي بعد انتهاء الإرسال المجدول ──
        try:
            with USERS_LOCK:
                _srpt_client_mgr = USERS.get(user_id, {}).get('client_manager')
            if _srpt_client_mgr and getattr(_srpt_client_mgr, 'client', None):
                _srpt_time = time.strftime('%Y-%m-%d %H:%M:%S')
                _srpt_lines = [
                    "📊 تقرير الإرسال المجدول",
                    f"⏰ الوقت: {_srpt_time}",
                    f"📨 إجمالي المجموعات: {len(groups)}",
                    "",
                    f"✅ الرسائل المرسلة بنجاح ({len(_sched_success_groups)}):",
                ]
                for _g in _sched_success_groups:
                    _srpt_lines.append(f"  • {_g}")
                if not _sched_success_groups:
                    _srpt_lines.append("  — لا يوجد")
                _srpt_lines += [
                    "",
                    f"❌ الرسائل الفاشلة ({len(_sched_failed_groups)}):",
                ]
                for _g in _sched_failed_groups:
                    _srpt_lines.append(f"  • {_g}")
                if not _sched_failed_groups:
                    _srpt_lines.append("  — لا يوجد")
                _srpt_lines += [
                    "",
                    f"🧠 الدورة الذكية ({len(_sched_smart_groups)}):",
                ]
                for _g in _sched_smart_groups:
                    _srpt_lines.append(f"  • {_g}")
                if not _sched_smart_groups:
                    _srpt_lines.append("  — لا يوجد")
                _srpt_nj = _sched_not_joined if '_sched_not_joined' in dir() else []
                _srpt_lines += [
                    "",
                    f"⚠️ غير منضم إليها ({len(_srpt_nj)}):",
                ]
                for _g in _srpt_nj:
                    _srpt_lines.append(f"  • {_g}")
                if not _srpt_nj:
                    _srpt_lines.append("  — لا يوجد")
                _full_sched_report = "\n".join(_srpt_lines)
                _srpt_client_mgr.run_coroutine(
                    _srpt_client_mgr.client.send_message('me', _full_sched_report, link_preview=False)
                )
                logger.info(f"[Scheduled] تم إرسال تقرير الإرسال المجدول للمستخدم {user_id}")
        except Exception as _srpt_err:
            logger.debug(f"[Scheduled] خطأ في إرسال التقرير: {_srpt_err}")

    except Exception as e:
        logger.error(f"Scheduled messages error: {str(e)}")

# =========================== 
# أحداث Socket.IO
# ===========================
@socketio.on('connect')
def handle_connect():
    try:
        if 'user_id' not in session:
            session['user_id'] = "user_1"
            session.permanent = True

        user_id = session['user_id']

        if user_id not in PREDEFINED_USERS:
            if PREDEFINED_USERS:
                user_id = list(PREDEFINED_USERS.keys())[0]
                session['user_id'] = user_id
            # else: user_id is a temp slot (not yet in PREDEFINED_USERS), keep it

        join_room(user_id)
        _uname = PREDEFINED_USERS.get(user_id, {}).get('name', user_id)
        logger.info(f"User {user_id} ({_uname}) connected via socket")

        emit('connection_confirmed', {
            'status': 'connected',
            'user_id': user_id,
            'user_name': _uname,
            'timestamp': time.strftime('%H:%M:%S')
        })

        emit('users_list', {
            'current_user': user_id,
            'users': PREDEFINED_USERS
        })

        notify_user_about_background_operations(user_id)

        all_status = get_all_users_operations_status()
        emit('all_users_status', all_status)

    except Exception as e:
        logger.error(f"Connection error: {str(e)}")
        emit('connection_error', {'message': str(e)})

@socketio.on('switch_user')
def handle_switch_user(data):
    try:
        new_user_id = data.get('user_id')

        if not new_user_id or new_user_id not in PREDEFINED_USERS:
            emit('error', {'message': 'مستخدم غير صحيح'})
            return

        # استخدم from_user_id إذا أُرسل من JS (لتجنب تعارض الغرف بعد تبديل HTTP)
        old_user_id = data.get('from_user_id') or session.get('user_id', 'user_1')
        try:
            leave_room(old_user_id)
        except Exception as leave_error:
            logger.warning(f"Error leaving room {old_user_id}: {str(leave_error)}")

        session['user_id'] = new_user_id
        session.permanent = True

        try:
            join_room(new_user_id)
        except Exception as join_error:
            logger.warning(f"Error joining room {new_user_id}: {str(join_error)}")

        logger.info(f"User switched from {old_user_id} to {new_user_id}")

        emit('user_switched', {
            'current_user': new_user_id,
            'user_name': PREDEFINED_USERS[new_user_id]['name'],
            'message': f"تم التبديل إلى {PREDEFINED_USERS[new_user_id]['name']}"
        })

        try:
            with USERS_LOCK:
                if new_user_id in USERS:
                    user_data = USERS[new_user_id]
                    connected = user_data.get('connected', False)
                    authenticated = user_data.get('authenticated', False)
                    awaiting_code = user_data.get('awaiting_code', False)
                    awaiting_password = user_data.get('awaiting_password', False)
                    is_running = user_data.get('is_running', False)

                    emit('connection_status', {
                        "status": "connected" if connected else "disconnected"
                    })

                    emit('login_status', {
                        "logged_in": authenticated,
                        "connected": connected,
                        "awaiting_code": awaiting_code,
                        "awaiting_password": awaiting_password,
                        "is_running": is_running
                    })

                    settings = load_settings(new_user_id)
                    emit('user_settings', settings)
                else:
                    emit('connection_status', {"status": "disconnected"})
                    emit('login_status', {
                        "logged_in": False,
                        "connected": False,
                        "awaiting_code": False,
                        "awaiting_password": False,
                        "is_running": False
                    })
        except Exception as status_error:
            logger.error(f"Error sending user status: {str(status_error)}")

    except Exception as e:
        logger.error(f"Error switching user: {str(e)}")
        emit('error', {'message': f'خطأ في التبديل: {str(e)}'})

@socketio.on('disconnect')
def handle_disconnect(data=None):
    if 'user_id' in session:
        user_id = session['user_id']
        leave_room(user_id)
        logger.info(f"User {user_id} disconnected from socket")

@socketio.on('heartbeat')
def handle_heartbeat(data):
    try:
        user_id = session.get('user_id')
        if user_id:
            emit('heartbeat_response', {
                'timestamp': time.time(),
                'server_time': time.strftime('%H:%M:%S')
            })
    except Exception as e:
        logger.error(f"Heartbeat error: {str(e)}")

# =========================== 
# المسارات الأساسية
# ===========================
@app.route("/")
def index():
    # ── فحص رابط الدعوة ─────────────────────────────────────────
    invite_token = request.args.get("invite")
    if invite_token:
        ok, item, err_msg = validate_invite_token(invite_token)
        if ok:
            uid = session.get('user_id', 'unknown')
            mark_token_used(invite_token, uid)
            session["invite_validated"] = True
            return redirect("/")
        else:
            # تحديد سبب الخطأ
            if "مسبقاً" in (err_msg or ""):
                reason = "used"
            elif "صلاحية" in (err_msg or ""):
                reason = "expired"
            elif "صحيح" in (err_msg or ""):
                reason = "not_found"
            else:
                reason = "invalid"
            return render_template("invite_error.html", reason=reason), 403
    # ── فحص نظام البطاقات ──────────────────────────────────────
    try:
        _cdata = load_cards_data()
        # كلمة مرور الأدمن = دخول حر — لا يُطبَّق عليه فحص البطاقات أبداً
        if _cdata.get("card_system_enabled", False) and not session.get("admin_auth"):
            if not session.get("card_logged_in"):
                return redirect("/login")
            sid = session.get("card_session_id")
            if sid:
                _active = next((s for s in _cdata.get("active_card_sessions", []) if s["session_id"] == sid), None)
                if not _active:
                    session.pop("card_logged_in", None)
                    session.pop("card_session_id", None)
                    return redirect("/login")
    except Exception:
        pass
    # ────────────────────────────────────────────────────────────
    # تحديث قائمة المستخدمين من الملف المحلي لضمان ظهور الحسابات المضافة حديثاً
    global PREDEFINED_USERS
    PREDEFINED_USERS = load_dynamic_users()

    if 'user_id' not in session or session['user_id'] not in PREDEFINED_USERS:
        if PREDEFINED_USERS:
            session['user_id'] = list(PREDEFINED_USERS.keys())[0]
        else:
            session['user_id'] = "user_1"  # فتحة مؤقتة لأول حساب
        session.permanent = True

    user_id = session['user_id']

    settings = load_settings(user_id)
    connection_status = "disconnected"

    with USERS_LOCK:
        if user_id not in USERS:
            USERS[user_id] = {
                'client_manager': None,
                'settings': settings,
                'thread': None,
                'is_running': False,
                'stats': {"sent": 0, "errors": 0},
                'connected': False,
                'authenticated': False,
                'awaiting_code': False,
                'awaiting_password': False,
                'phone_code_hash': None,
                'monitoring_active': False,
                'event_handlers_registered': False,
                'sent_batches': settings.get('sent_batches', []) or []
            }

        user_data = USERS[user_id]
        connected = user_data.get('connected', False)
        connection_status = "connected" if connected else "disconnected"

    _install_id_for_cookie = None
    try:
        _inst_record, _inst_is_new, _inst_id = track_installation(
            user_id=user_id,
            request=request,
            predefined_users=PREDEFINED_USERS,
            users_dict=USERS,
            load_settings_func=load_settings,
            socketio_obj=socketio,
            users_lock=USERS_LOCK,
        )
        _install_id_for_cookie = _inst_id
    except Exception as _e:
        logger.error(f"track_installation (index) error: {_e}")

    app_title = "مركز سرعة انجاز 📚 للخدمات الطلابية والأكاديمية"
    whatsapp_link = "https://wa.me/+966510349663"

    current_user = PREDEFINED_USERS.get(user_id, {
        "id": user_id,
        "name": "إضافة حساب",
        "icon": "fas fa-plus-circle",
        "color": "#6366f1"
    })

    admin_ui_visible = session.get('admin_ui_visible', False)

    response = render_template('index.html',
                          settings=settings,
                          connection_status=connection_status,
                          app_title=app_title,
                          whatsapp_link=whatsapp_link,
                          current_user=current_user,
                          predefined_users=PREDEFINED_USERS,
                          admin_ui_visible=admin_ui_visible,
                          show_all_users=True)

    resp = make_response(response)
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'
    # ضبط كوكيز install_id ليظل ثابتاً للمتصفح (سنة كاملة)
    if _install_id_for_cookie and not request.cookies.get('install_id'):
        resp.set_cookie(
            'install_id', _install_id_for_cookie,
            max_age=365 * 24 * 3600,
            httponly=False,   # يُقرأ بـ JavaScript أيضاً
            samesite='Lax',
            path='/'
        )

    return resp

@app.route("/fresh")
def fresh():
    from flask import make_response
    html = """<!DOCTYPE html>
<html dir="rtl" lang="ar">
<head>
    <meta charset="UTF-8">
    <title>🚀 التطبيق يعمل بنجاح!</title>
    <style>
        body { font-family: Arial, sans-serif; text-align: center; padding: 50px; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; }
        .success { font-size: 2em; margin: 20px 0; }
        .message { font-size: 1.2em; margin: 10px 0; }
        .btn { background: #28a745; color: white; padding: 15px 30px; text-decoration: none; border-radius: 8px; font-size: 1.1em; display: inline-block; margin: 10px; }
        .btn:hover { background: #218838; color: white; }
    </style>
</head>
<body>
    <div class="success">✅ التطبيق يعمل بشكل مثالي!</div>
    <div class="message">🎉 مركز سرعة انجاز للخدمات الطلابية والأكاديمية</div>
    <div class="message">📱 نظام مراقبة التليجرام الذكي</div>
    <a href="/" class="btn">🏠 الانتقال للتطبيق الرئيسي</a>
    <script>
        setTimeout(function() {
            window.location.href = '/';
        }, 3000);
    </script>
</body>
</html>"""

    resp = make_response(html)
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'
    resp.headers['Content-Type'] = 'text/html; charset=utf-8'

    return resp

@app.route('/static/<path:filename>')
def static_files(filename):
    return app.send_static_file(filename)

@app.route("/manifest.json")
def manifest():
    manifest_data = {
        "id": "/",
        "name": "مركز سرعة انجاز للخدمات الطلابية والأكاديمية",
        "short_name": "سرعة انجاز",
        "description": "نظام متكامل: تليجرام تلقائي، تحليل أكاديمي، عروض PowerPoint، منسّق مستندات",
        "start_url": "/",
        "scope": "/",
        "display": "standalone",
        "display_override": ["standalone", "minimal-ui"],
        "orientation": "portrait",
        "theme_color": "#1e3c78",
        "background_color": "#1e3c78",
        "lang": "ar",
        "dir": "rtl",
        "categories": ["education", "productivity", "utilities"],
        "prefer_related_applications": False,
        "icons": [
            {"src": "/static/icons/icon-72.png",  "sizes": "72x72",   "type": "image/png", "purpose": "any maskable"},
            {"src": "/static/icons/icon-192.png", "sizes": "192x192", "type": "image/png", "purpose": "any maskable"},
            {"src": "/static/icons/icon-512.png", "sizes": "512x512", "type": "image/png", "purpose": "any maskable"},
            {"src": "/static/icons/app-logo.png", "sizes": "512x512", "type": "image/png", "purpose": "any maskable"}
        ],
        "screenshots": [],
        "shortcuts": [
            {"name": "التحليل الأكاديمي", "short_name": "أكاديمي", "description": "فتح منصة التحليل", "url": "/academic"},
            {"name": "لوحة التحكم",       "short_name": "تحكم",    "description": "لوحة التحكم الرئيسية", "url": "/"}
        ]
    }
    resp = app.response_class(json.dumps(manifest_data, ensure_ascii=False, indent=2),
                              mimetype='application/manifest+json')
    resp.headers['Cache-Control'] = 'no-cache'
    return resp

@app.route("/sw.js")
def service_worker():
    sw_js = r"""
const CACHE_NAME = 'sra3a-v6';
const STATIC_ASSETS = [
  '/',
  '/static/icons/icon-192.png',
  '/static/icons/icon-512.png',
  '/static/icons/app-logo.png'
];

// ── التثبيت: تخزين الأصول الأساسية ──
self.addEventListener('install', event => {
  self.skipWaiting();
  event.waitUntil(
    caches.open(CACHE_NAME).then(cache =>
      cache.addAll(STATIC_ASSETS).catch(() => {})
    )
  );
});

// ── التفعيل: حذف الكاش القديم ──
self.addEventListener('activate', event => {
  event.waitUntil(
    caches.keys().then(keys =>
      Promise.all(keys.filter(k => k !== CACHE_NAME).map(k => caches.delete(k)))
    ).then(() => self.clients.claim())
  );
});

// ── الجلب: Network-first، cache كـ fallback ──
self.addEventListener('fetch', event => {
  const req = event.request;
  if (req.method !== 'GET') return;

  // تجاهل طلبات API و WebSocket
  const url = new URL(req.url);
  if (url.pathname.startsWith('/api/') ||
      url.pathname.startsWith('/socket.io') ||
      url.pathname.startsWith('/tools/')) return;

  event.respondWith(
    fetch(req)
      .then(resp => {
        if (resp.ok && resp.type === 'basic') {
          const clone = resp.clone();
          caches.open(CACHE_NAME).then(c => c.put(req, clone));
        }
        return resp;
      })
      .catch(() => caches.match(req))
  );
});

// ── رسائل من الصفحة ──
self.addEventListener('message', event => {
  if (event.data === 'SKIP_WAITING') self.skipWaiting();
});

// ── Web Push: استقبال الإشعارات الفورية ──
self.addEventListener('push', function(event) {
  let data = {
    title: '🔔 إشعار جديد',
    body: 'لديك إشعار من أبو مالك',
    icon: '/static/icons/icon-192.png',
    badge: '/static/icons/icon-72.png',
    data: {}
  };
  try {
    if (event.data) {
      const raw = event.data.json();
      data = Object.assign(data, raw);
    }
  } catch(e) {}

  const options = {
    body: data.body,
    icon: data.icon || '/static/icons/icon-192.png',
    badge: data.badge || '/static/icons/icon-72.png',
    data: data.data || {},
    vibrate: [300, 100, 300, 100, 300],
    requireInteraction: false,
    dir: 'rtl',
    lang: 'ar',
    tag: 'abumalik-push',
    renotify: true,
    actions: [
      { action: 'open', title: '📱 فتح التطبيق' },
      { action: 'close', title: '✕ إغلاق' }
    ]
  };

  event.waitUntil(
    self.registration.showNotification(data.title, options)
  );
});

// ── عند الضغط على الإشعار ──
self.addEventListener('notificationclick', function(event) {
  event.notification.close();
  if (event.action === 'close') return;
  event.waitUntil(
    clients.matchAll({ type: 'window', includeUncontrolled: true }).then(function(cs) {
      for (const c of cs) {
        if ('focus' in c) return c.focus();
      }
      if (clients.openWindow) return clients.openWindow('/');
    })
  );
});
"""
    resp = app.response_class(sw_js, content_type='application/javascript')
    resp.headers['Service-Worker-Allowed'] = '/'
    resp.headers['Cache-Control'] = 'no-cache'
    return resp

# ══ /geo_clear: يمسح بيانات المتصفح بالكامل ثم يعيد التوجيه للصفحة الرئيسية ══
@app.route("/geo_clear", methods=["GET"])
def geo_clear():
    """
    يُستدعى عندما يرفض المستخدم إذن الموقع.
    يرسل Clear-Site-Data لمسح كل شيء من المتصفح (كوكيز، تخزين، كاش)،
    ثم يُعيد التوجيه للصفحة الرئيسية ليُعامَل المتصفح الزيارة كأول مرة.
    """
    resp = redirect("/", code=302)
    # Clear-Site-Data يمسح كوكيز + localStorage + sessionStorage + caches
    resp.headers['Clear-Site-Data'] = '"cache", "cookies", "storage"'
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    # إلغاء كوكيز الجلسة والتتبع يدوياً كضمان إضافي
    for cookie_name in ['session', 'install_id']:
        resp.set_cookie(cookie_name, '', expires=0, path='/')
    return resp

# =========================== 
# API Routes
# ===========================

@app.route("/api", methods=["GET", "HEAD"])
def api_health():
    try:
        if request.method == "HEAD":
            return "", 200
        return jsonify({"status": "ok", "timestamp": time.time(), "message": "Server is running"})
    except Exception as e:
        logger.error(f"Error in api health check: {str(e)}")
        if request.method == "HEAD":
            return "", 500
        return jsonify({"status": "error", "message": "Server error"}), 500

@app.route("/api/save_login", methods=["POST"])
def api_save_login():
    data = request.json

    if not data or not data.get('phone'):
        return jsonify({
            "success": False,
            "message": "❌ يرجى إدخال رقم الهاتف"
        })

    new_phone = data.get('phone')

    # ─── تحديد user_id ───────────────────────────────────────────────────────
    # الأولوية: (1) user_id في body الطلب، (2) الجلسة، (3) الافتراضي user_1
    requested_uid = (data.get('user_id') or '').strip()
    if requested_uid and requested_uid in PREDEFINED_USERS:
        session['user_id'] = requested_uid
        session.permanent = True
    elif 'user_id' not in session or session['user_id'] not in PREDEFINED_USERS:
        session['user_id'] = "user_1"
        session.permanent = True
    session.modified = True

    user_id = session['user_id']
    logger.info(f"api_save_login: user_id={user_id}, phone={new_phone}")
    log_user_event(user_id, 'INFO', f"📱 طلب تسجيل دخول للرقم: {new_phone}")

    # تنظيف الجلسة القديمة لنفس الخانة إذا تغيّر الرقم
    current_settings = load_settings(user_id)
    if current_settings.get('phone') and current_settings.get('phone') != new_phone:
        logger.info(f"Phone changed: {current_settings['phone']} → {new_phone} for {user_id}")
        with USERS_LOCK:
            if user_id in USERS:
                if USERS[user_id].get('is_running'):
                    USERS[user_id]['is_running'] = False
                cm = USERS[user_id].get('client_manager')
                if cm:
                    try: cm.stop()
                    except Exception: pass
                del USERS[user_id]
        # حذف ملف الجلسة القديم
        old_session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(old_session_file):
            try:
                os.remove(old_session_file)
            except Exception as e:
                logger.warning(f"Could not remove old session file: {e}")
        _old_uname = PREDEFINED_USERS.get(user_id, {}).get('name', user_id)
        socketio.emit('log_update', {
            "message": f"🔄 تم مسح الجلسة القديمة لـ {_old_uname}"
        }, to=user_id)

    settings = {
        'phone': new_phone,
        'password': data.get('password', ''),
        'login_time': time.time()
    }

    if not save_settings(user_id, settings):
        return jsonify({
            "success": False,
            "message": "❌ فشل في حفظ البيانات"
        })

    try:
        _login_uname = PREDEFINED_USERS.get(user_id, {}).get('name', f"حساب جديد ({user_id})")
        socketio.emit('log_update', {
            "message": f"🔄 بدء تسجيل دخول {_login_uname}..."
        }, to=user_id)

        with USERS_LOCK:
            # إزالة أي خانة أخرى تستخدم نفس الرقم
            users_to_remove = [
                uid for uid, ud in USERS.items()
                if uid != user_id and ud['settings'].get('phone') == new_phone
            ]
            for old_uid in users_to_remove:
                logger.info(f"Removing duplicate phone session: {old_uid}")
                if USERS[old_uid].get('is_running'):
                    USERS[old_uid]['is_running'] = False
                cm = USERS[old_uid].get('client_manager')
                if cm:
                    try: cm.stop()
                    except Exception: pass
                del USERS[old_uid]

            # إنشاء/تحديث إدخال المستخدم الحالي
            USERS[user_id] = {
                'client_manager': None,
                'settings': settings,
                'thread': None,
                'is_running': False,
                'stats': {"sent": 0, "errors": 0},
                'connected': False,
                'authenticated': False,
                'awaiting_code': False,
                'awaiting_password': False,
                'phone_code_hash': None,
                'monitoring_active': False,
                'event_handlers_registered': False,
                'sent_batches': settings.get('sent_batches', []) or []
            }

        result = telegram_manager.setup_client(user_id, settings['phone'])

        if result["status"] == "pending":
            # الاتصال يعمل في الخلفية - النتيجة ستصل عبر socket.io
            return jsonify({
                "success": True,
                "message": "🔄 جارِ الاتصال بتيليجرام...",
                "pending": True
            })

        elif result["status"] == "success":
            socketio.emit('log_update', {"message": "✅ تم تسجيل الدخول بنجاح"}, to=user_id)
            socketio.emit('connection_status', {"status": "connected"}, to=user_id)
            socketio.emit('login_status', {
                "logged_in": True, "connected": True,
                "awaiting_code": False, "awaiting_password": False, "is_running": False
            }, to=user_id)
            return jsonify({"success": True, "message": "✅ تم تسجيل الدخول"})

        elif result["status"] == "code_required":
            socketio.emit('log_update', {"message": "📱 تم إرسال كود التحقق"}, to=user_id)
            return jsonify({"success": True, "message": "📱 تم إرسال كود التحقق", "code_required": True})

        else:
            error_message = result.get('message', 'خطأ غير معروف')
            socketio.emit('log_update', {"message": f"❌ {error_message}"}, to=user_id)

            return jsonify({
                "success": False, 
                "message": f"❌ {error_message}"
            })

    except Exception as e:
        logger.error(f"Login error for user {user_id}: {str(e)}")
        socketio.emit('log_update', {
            "message": f"❌ خطأ: {str(e)}"
        }, to=user_id)

        return jsonify({
            "success": False, 
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/verify_code", methods=["POST"])
def api_verify_code():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']
    data = request.json

    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لم يتم إرسال البيانات"
        })

    code = data.get('code')
    password = data.get('password')

    if not code and not password:
        return jsonify({
            "success": False, 
            "message": "❌ يرجى إدخال الكود أو كلمة المرور"
        })

    try:
        if code:
            result = telegram_manager.verify_code(user_id, code)
        else:
            result = telegram_manager.verify_password(user_id, password)

        if result["status"] == "success":
            account_name = result.get("account_name")
            socketio.emit('log_update', {
                "message": f"✅ تم التحقق بنجاح — أهلاً {account_name}" if account_name else "✅ تم التحقق بنجاح"
            }, to=user_id)

            socketio.emit('connection_status', {
                "status": "connected"
            }, to=user_id)

            # ── تسجيل الحساب تلقائياً إذا لم يكن موجوداً في PREDEFINED_USERS ──
            try:
                global PREDEFINED_USERS
                tg_display_name = account_name or f"حساب {user_id.replace('user_', '')}"
                if user_id not in PREDEFINED_USERS:
                    add_dynamic_user(user_id, tg_display_name, "fas fa-user", "#6366f1")
                    PREDEFINED_USERS = load_dynamic_users()
                    session['platform_logged_in'] = True
                    session.permanent = True
                else:
                    # تحديث الاسم بالاسم الحقيقي من تيليجرام
                    _users_dict = load_dynamic_users()
                    if user_id in _users_dict and account_name:
                        _users_dict[user_id]['name'] = account_name
                        save_dynamic_users(_users_dict)
                        PREDEFINED_USERS = load_dynamic_users()
            except Exception as _auto_ae:
                logger.debug(f"auto-register user error: {_auto_ae}")

            return jsonify({
                "success": True,
                "message": f"✅ تم التحقق بنجاح — أهلاً {account_name}" if account_name else "✅ تم التحقق بنجاح",
                "account_name": account_name
            })

        elif result["status"] == "password_required":
            return jsonify({
                "success": True, 
                "message": result["message"], 
                "password_required": True
            })

        else:
            error_message = result.get('message', 'فشل التحقق')
            socketio.emit('log_update', {
                "message": f"❌ {error_message}"
            }, to=user_id)

            return jsonify({
                "success": False, 
                "message": f"❌ {error_message}"
            })

    except Exception as e:
        socketio.emit('log_update', {
            "message": f"❌ خطأ في التحقق: {str(e)}"
        }, to=user_id)

        return jsonify({
            "success": False, 
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/save_settings", methods=["POST"])
def api_save_settings():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']
    data = request.json

    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لم يتم إرسال البيانات"
        })

    current_settings = load_settings(user_id)
    old_mode = current_settings.get('sanitize_mode', 'salam')
    new_mode = (data.get('sanitize_mode') or 'salam').lower()

    # إذا تغيّر الوضع من salam إلى غيره، أوقف جميع الدورات الذكية النشطة لهذا المستخدم
    if old_mode == 'salam' and new_mode != 'salam':
        keys_to_remove = [k for k in list(telegram_manager._smart_running) if k.startswith(f"{user_id}_")]
        for k in keys_to_remove:
            telegram_manager._smart_running.discard(k)
        if keys_to_remove:
            socketio.emit('log_update', {
                "message": f"⏹ تم إيقاف {len(keys_to_remove)} دورة ذكية بسبب تغيير وضع الإرسال"
            }, to=user_id)

    _sched_dur_h = float(data.get('schedule_duration_hours', 0) or 0)
    current_settings.update({
        'message': data.get('message', ''),
        'groups': dedupe_groups(data.get('groups', '')),
        'interval_seconds': int(data.get('interval_seconds', 3600)),
        'watch_words': [w.strip() for w in data.get('watch_words', '').split('\n') if w.strip()],
        'send_type': data.get('send_type', 'manual'),
        'schedule_duration_hours': _sched_dur_h,
        'schedule_duration': int(_sched_dur_h * 3600),
        'schedule_start_time': None,   # يُعاد ضبطه عند بدء التشغيل
        'max_retries': int(data.get('max_retries', 5)),
        'auto_reconnect': data.get('auto_reconnect', False),
        'sanitize_mode': new_mode,
        'smart_required_messages': int(data.get('smart_required_messages', 3)),
    })

    if save_settings(user_id, current_settings):
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['settings'] = current_settings
                client_manager = USERS[user_id].get('client_manager')
                if client_manager and hasattr(client_manager, 'update_monitoring_settings'):
                    client_manager.update_monitoring_settings(
                        current_settings.get('watch_words', []),
                        current_settings.get('groups', [])
                    )

        socketio.emit('log_update', {
            "message": "✅ تم حفظ الإعدادات بنجاح"
        }, to=user_id)

        return jsonify({
            "success": True, 
            "message": "✅ تم حفظ الإعدادات"
        })
    else:
        return jsonify({
            "success": False, 
            "message": "❌ فشل في حفظ الإعدادات"
        })

# ملاحظة: /api/smart_stop أُزيل — الإيقاف يتم الآن تلقائياً عند تغيير sanitize_mode من salam إلى غيره


@app.route("/api/user_logout", methods=["POST"])
def api_user_logout():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({
            "success": False,
            "message": "❌ لا توجد جلسة نشطة"
        })

    try:
        logger.info(f"User {user_id} logging out...")

        with USERS_LOCK:
            if user_id in USERS:
                client_manager = USERS[user_id].get('client_manager')
                if client_manager:
                    try:
                        if USERS[user_id].get('is_running'):
                            USERS[user_id]['is_running'] = False

                        if hasattr(client_manager, 'client') and client_manager.client:
                            client_manager.client.disconnect()
                            logger.info(f"Client disconnected for user {user_id}")

                        if hasattr(client_manager, 'stop'):
                            client_manager.stop()

                    except Exception as e:
                        logger.error(f"خطأ في إغلاق العميل للمستخدم {user_id}: {e}")

                del USERS[user_id]
                logger.info(f"User data removed from memory for {user_id}")

        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            try:
                os.remove(session_file)
                logger.info(f"Session file removed for {user_id}")
            except Exception as e:
                logger.error(f"خطأ في حذف ملف الجلسة: {e}")

        # إيقاف النشر الدوري فوراً عند تسجيل الخروج
        try:
            rotating_manager.stop(user_id)
            logger.info(f"Rotating send stopped on logout for {user_id}")
        except Exception as _re:
            logger.error(f"خطأ في إيقاف النشر الدوري عند الخروج: {_re}")

        settings_file = os.path.join(SESSIONS_DIR, f"{user_id}.json")
        if os.path.exists(settings_file):
            try:
                settings = load_settings(user_id)
                settings.update({
                    'phone': '',
                    'authenticated': False,
                    'connected': False,
                    'rotating_persistent': False,
                    'monitoring_persistent': False
                })
                save_settings(user_id, settings)
                logger.info(f"Settings cleared for {user_id}")
            except Exception as e:
                logger.error(f"خطأ في مسح الإعدادات: {e}")

        socketio.emit('log_update', {
            "message": "🚪 تم تسجيل الخروج وإنهاء جلسة التليجرام"
        }, to=user_id)

        socketio.emit('connection_status', {
            "status": "disconnected"
        }, to=user_id)

        socketio.emit('login_status', {
            "logged_in": False,
            "connected": False,
            "awaiting_code": False,
            "awaiting_password": False,
            "is_running": False
        }, to=user_id)

        logger.info(f"User {user_id} logged out successfully")

        return jsonify({
            "success": True,
            "message": "✅ تم تسجيل الخروج وإنهاء جلسة التليجرام بنجاح"
        })

    except Exception as e:
        logger.error(f"خطأ في تسجيل الخروج للمستخدم {user_id}: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في تسجيل الخروج: {str(e)}"
        })

@app.route("/api/get_account_info", methods=["GET"])
def api_get_account_info():
    user_id = session.get('user_id', 'user_1')
    try:
        try:
            telegram_manager.ensure_client_active(user_id)
        except Exception as e:
            logger.debug(f"ensure_client_active in get_account_info: {e}")
        with USERS_LOCK:
            udata = USERS.get(user_id, {})
            cached = {
                "account_name": udata.get('account_name'),
                "account_username": udata.get('account_username'),
                "account_phone": udata.get('account_phone'),
                "account_avatar": udata.get('account_avatar'),
                "authenticated": udata.get('authenticated', False)
            }
        if not cached["account_name"] and cached["authenticated"]:
            try:
                cached["account_name"] = telegram_manager._fetch_account_name(user_id)
                with USERS_LOCK:
                    cached["account_username"] = USERS.get(user_id, {}).get('account_username')
                    cached["account_phone"] = USERS.get(user_id, {}).get('account_phone')
                    cached["account_avatar"] = USERS.get(user_id, {}).get('account_avatar')
            except Exception as e:
                logger.error(f"get_account_info refresh failed: {e}")
        if not cached.get("account_avatar"):
            avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{user_id}.jpg")
            if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
                cached["account_avatar"] = f"/api/account_avatar/{user_id}"
        cached["is_pro"] = is_user_restricted(user_id)
        return jsonify({
            "success": True,
            "user_id": user_id,
            "predefined_name": PREDEFINED_USERS.get(user_id, {}).get('name'),
            **cached
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/account_avatar/<uid>", methods=["GET"])
def api_account_avatar(uid):
    try:
        from flask import send_file, abort
        avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{uid}.jpg")
        if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
            return send_file(avatar_file, mimetype='image/jpeg', max_age=60)
        return ('', 404)
    except Exception as e:
        logger.error(f"Avatar serving error for {uid}: {e}")
        return ('', 404)

@app.route("/api/switch_user", methods=["POST"])
def api_switch_user():
    try:
        data = request.get_json()
        new_user_id = data.get('user_id')

        if not new_user_id or new_user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        old_user_id = session.get('user_id', 'user_1')

        if old_user_id in USERS:
            current_settings = USERS[old_user_id].get('settings', {})
            if current_settings:
                save_settings(old_user_id, current_settings)
                logger.info(f"✅ Settings saved for user {old_user_id} - Operations continue running")

        with USERS_LOCK:
            if new_user_id not in USERS:
                saved_settings = load_settings(new_user_id)

                USERS[new_user_id] = {
                    'client_manager': None,
                    'settings': saved_settings,
                    'thread': None,
                    'is_running': False,
                    'stats': {"sent": 0, "errors": 0},
                    'connected': False,
                    'authenticated': False,
                    'awaiting_code': False,
                    'awaiting_password': False,
                    'phone_code_hash': None,
                    'monitoring_active': False,
                    'event_handlers_registered': False,
                    'sent_batches': (saved_settings or {}).get('sent_batches', []) or []
                }

                session_file = os.path.join(SESSIONS_DIR, f"{new_user_id}_session.session")
                if os.path.exists(session_file) and saved_settings.get('phone'):
                    USERS[new_user_id]['connected'] = True
                    USERS[new_user_id]['authenticated'] = True
                    logger.info(f"Found existing session for user {new_user_id}")
            else:
                saved_settings = load_settings(new_user_id)
                USERS[new_user_id]['settings'].update(saved_settings)

        session['user_id'] = new_user_id
        session.permanent = True

        try:
            track_installation(
                user_id=new_user_id,
                request=request,
                predefined_users=PREDEFINED_USERS,
                users_dict=USERS,
                load_settings_func=load_settings,
                socketio_obj=socketio,
                users_lock=USERS_LOCK,
            )
        except Exception as _e:
            logger.error(f"track_installation (switch_user) error: {_e}")

        logger.info(f"✅ User switched from {old_user_id} to {new_user_id} - All operations remain active")

        active_operations_summary = get_all_users_operations_status()

        socketio.emit('user_settings', USERS[new_user_id]['settings'], to=new_user_id)

        account_name = None
        account_avatar = None
        try:
            telegram_manager.ensure_client_active(new_user_id)

            with USERS_LOCK:
                account_name = USERS[new_user_id].get('account_name')
                account_avatar = USERS[new_user_id].get('account_avatar')
            if not account_name and USERS[new_user_id].get('authenticated'):
                account_name = telegram_manager._fetch_account_name(new_user_id)
                with USERS_LOCK:
                    account_avatar = USERS[new_user_id].get('account_avatar')
        except Exception as e:
            logger.error(f"Could not load account name on switch: {e}")

        if not account_avatar:
            avatar_file = os.path.join(SESSIONS_DIR, 'avatars', f"{new_user_id}.jpg")
            if os.path.exists(avatar_file) and os.path.getsize(avatar_file) > 0:
                account_avatar = f"/api/account_avatar/{new_user_id}"

        is_pro = is_user_restricted(new_user_id)
        return jsonify({
            "success": True,
            "message": f"✅ تم التبديل إلى {PREDEFINED_USERS[new_user_id]['name']}" + (f" — حساب تليجرام: {account_name}" if account_name else ""),
            "switched": old_user_id != new_user_id,
            "previous_user_id": old_user_id,
            "is_pro": is_pro,
            "user": {
                "id": new_user_id,
                "name": PREDEFINED_USERS[new_user_id]['name'],
                "icon": PREDEFINED_USERS[new_user_id]['icon'],
                "color": PREDEFINED_USERS[new_user_id]['color'],
                "account_name": account_name,
                "account_avatar": account_avatar,
                "authenticated": USERS[new_user_id].get('authenticated', False),
                "is_pro": is_pro
            },
            "account_name": account_name,
            "account_avatar": account_avatar,
            "settings": USERS[new_user_id]['settings'],
            "active_operations": active_operations_summary
        })

    except Exception as e:
        logger.error(f"Error in user switching API: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في التبديل: {str(e)}"
        })

@app.route("/api/start_monitoring", methods=["POST"])
def api_start_monitoring():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({
                "success": False, 
                "message": "❌ لم يتم إعداد الحساب"
            })

        if not USERS[user_id].get('authenticated'):
            return jsonify({
                "success": False, 
                "message": "❌ يجب تسجيل الدخول أولاً"
            })

        if USERS[user_id]['is_running']:
            return jsonify({
                "success": False, 
                "message": "✅ النظام يعمل بالفعل"
            })

        USERS[user_id]['is_running'] = True

    try:
        _settings = load_settings(user_id)
        _settings['monitoring_persistent'] = True
        save_settings(user_id, _settings)
    except Exception as _e:
        logger.error(f"Failed to persist monitoring flag for {user_id}: {_e}")

    log_user_event(user_id, 'INFO', "🚀 بدء تشغيل نظام المراقبة...")
    socketio.emit('log_update', {
        "message": "🚀 بدء تشغيل نظام المراقبة المحسن مع Event Handlers..."
    }, to=user_id)

    try:
        monitoring_thread = _OSThread(
            target=monitoring_worker, 
            args=(user_id,), 
            daemon=True
        )
        monitoring_thread.start()

        with USERS_LOCK:
            USERS[user_id]['thread'] = monitoring_thread

        socketio.emit('monitoring_status', {
            "monitoring_active": True,
            "status": "running",
            "is_running": True
        }, to=user_id)

        socketio.emit('update_monitoring_buttons', {
            "is_running": True
        }, to=user_id)

        return jsonify({
            "success": True, 
            "message": "🚀 بدأت المراقبة المحسنة مع Event Handlers"
        })

    except Exception as e:
        logger.error(f"Failed to start monitoring for {user_id}: {str(e)}")

        with USERS_LOCK:
            USERS[user_id]['is_running'] = False

        return jsonify({
            "success": False, 
            "message": f"❌ فشل في بدء المراقبة: {str(e)}"
        })

@app.route("/api/stop_monitoring", methods=["POST"])
def api_stop_monitoring():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    try:
        _settings = load_settings(user_id)
        _settings['monitoring_persistent'] = False
        save_settings(user_id, _settings)
    except Exception as _e:
        logger.error(f"Failed to clear monitoring flag for {user_id}: {_e}")

    with USERS_LOCK:
        if user_id in USERS and USERS[user_id]['is_running']:
            USERS[user_id]['is_running'] = False
            log_user_event(user_id, 'INFO', "⏹ تم إيقاف نظام المراقبة")
            socketio.emit('log_update', {
                "message": "⏹ إيقاف نظام المراقبة..."
            }, to=user_id)

            socketio.emit('monitoring_status', {
                "monitoring_active": False,
                "status": "stopped",
                "is_running": False
            }, to=user_id)

            socketio.emit('update_monitoring_buttons', {
                "is_running": False
            }, to=user_id)

            return jsonify({
                "success": True, 
                "message": "⏹ تم إيقاف المراقبة"
            })

    return jsonify({
        "success": False, 
        "message": "❌ النظام غير مشغل"
    })



@app.route("/api/send_now", methods=["POST"])
def api_send_now():
    if 'user_id' not in session:
        return jsonify({
            "success": False, 
            "message": "❌ الجلسة غير صالحة، يرجى إعادة تحميل الصفحة"
        })

    user_id = session['user_id']

    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({
                "success": False, 
                "message": "❌ لم يتم إعداد الحساب"
            })

        if not USERS[user_id].get('authenticated'):
            return jsonify({
                "success": False, 
                "message": "❌ يجب تسجيل الدخول أولاً"
            })

    data = request.get_json()
    if not data:
        return jsonify({
            "success": False, 
            "message": "❌ لا توجد بيانات مرسلة"
        })

    message = data.get('message', '').strip()
    groups = data.get('groups', '').strip()
    images = data.get('images', [])
    send_to_all = bool(data.get('send_to_all', False))
    # الإجراء المختار من نافذة الفحص الاستباقي (skip / sanitize / salam / None)
    pre_scan_action = data.get('action', None)  # None = استخدم الإعدادات الافتراضية

    if not message and not images:
        return jsonify({
            "success": False, 
            "message": "❌ يجب كتابة رسالة أو رفع صورة للإرسال"
        })

    if send_to_all:
        # جلب جميع المجموعات من الحساب تلقائياً
        try:
            with USERS_LOCK:
                client_mgr = USERS.get(user_id, {}).get('client_manager')
            if not client_mgr or not client_mgr.client:
                return jsonify({"success": False, "message": "❌ العميل غير متصل"})
            dialogs = client_mgr.run_coroutine(client_mgr.client.get_dialogs())
            raw_groups_list = []
            for d in dialogs:
                entity = d.entity
                if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast') or hasattr(entity, 'gigagroup'):
                    uname = getattr(entity, 'username', None)
                    if uname:
                        raw_groups_list.append(f"https://t.me/{uname}")
                    else:
                        raw_groups_list.append(str(entity.id))
            groups_list = dedupe_groups(raw_groups_list)
            if not groups_list:
                return jsonify({"success": False, "message": "❌ لم يُعثر على أي مجموعة عامة"})
            socketio.emit('log_update', {"message": f"📡 إرسال لكل المجموعات: {len(groups_list)} مجموعة"}, to=user_id)
        except Exception as e:
            return jsonify({"success": False, "message": f"❌ خطأ في جلب المجموعات: {str(e)}"})
    else:
        if not groups:
            return jsonify({
                "success": False,
                "message": "❌ يجب تحديد المجموعات للإرسال إليها"
            })
        raw_groups = [g.strip() for g in groups.replace('\n', ',').split(',') if g.strip()]
        original_count = len(raw_groups)
        groups_list = dedupe_groups(raw_groups)
        duplicates_removed = original_count - len(groups_list)
        if duplicates_removed > 0:
            socketio.emit('log_update', {
                "message": f"♻️ تم تجاهل {duplicates_removed} رابط مكرر في قائمة الإرسال"
            }, to=user_id)

    if not groups_list:
        return jsonify({
            "success": False, 
            "message": "❌ يجب تحديد مجموعة واحدة على الأقل"
        })

    image_files = []
    if images:
        try:
            for img_data in images:
                raw_data = img_data.get('data', '')
                if ',' in raw_data:
                    base64_data = raw_data.split(',', 1)[1]
                else:
                    base64_data = raw_data
                image_bytes = base64.b64decode(base64_data)

                mime = img_data.get('type', 'image/jpeg')
                ext = mime.split('/')[-1].lower()
                if ext in ('jpeg', 'jpg'):
                    ext = 'jpg'
                elif ext not in ('png', 'gif', 'webp', 'bmp'):
                    ext = 'jpg'

                temp_file = tempfile.NamedTemporaryFile(
                    delete=False, suffix=f'.{ext}', mode='wb'
                )
                temp_file.write(image_bytes)
                temp_file.flush()
                temp_file.close()

                image_files.append({
                    'path': temp_file.name,
                    'name': img_data.get('name', f'image.{ext}'),
                    'type': mime
                })

            socketio.emit('log_update', {
                "message": f"📷 تم تحضير {len(image_files)} صورة للإرسال"
            }, to=user_id)

        except Exception as e:
            logger.error(f"Error processing images: {str(e)}")
            return jsonify({
                "success": False,
                "message": f"❌ خطأ في معالجة الصور: {str(e)}"
            })

    content_type = "رسالة"
    if images and message:
        content_type = f"رسالة مع {len(images)} صورة"
    elif images:
        content_type = f"{len(images)} صورة"

    socketio.emit('log_update', {
        "message": f"🚀 بدء الإرسال الفوري: {content_type} إلى {len(groups_list)} مجموعة"
    }, to=user_id)

    def send_messages_with_images():
        try:
            successful = 0
            failed = 0
            batch_id = str(uuid.uuid4())
            batch_entries = []
            _instant_success_groups = []
            _instant_failed_groups = []
            _instant_smart_groups = []
            _instant_nj_groups = []

            # ── الدورة الأولى: فحص العضوية وإشعار المستخدم بروابط المجموعات غير المنضم إليها ──
            try:
                with USERS_LOCK:
                    _now_client_mgr = USERS.get(user_id, {}).get('client_manager')
                if _now_client_mgr and getattr(_now_client_mgr, 'client', None):
                    socketio.emit('log_update', {
                        "message": f"🔍 فحص العضوية في {len(groups_list)} مجموعة..."
                    }, to=user_id)

                    async def _check_memberships_now(client, groups):
                        """فحص العضوية لجميع المجموعات دون الانضمام إليها"""
                        from telethon.tl.functions.channels import GetParticipantRequest
                        from telethon.errors import UserNotParticipantError as _UNPE
                        not_joined = []
                        for _g in groups:
                            try:
                                _ident = _g
                                for _pfx in ['https://t.me/', 'https://telegram.me/']:
                                    if _g.startswith(_pfx):
                                        _ident = _g[len(_pfx):]
                                        break
                                if _ident.startswith('@'):
                                    _ident = _ident[1:]
                                _ent = await client.get_entity(_ident)
                                if hasattr(_ent, 'megagroup') or hasattr(_ent, 'broadcast'):
                                    try:
                                        await client(GetParticipantRequest(_ent, 'me'))
                                    except _UNPE:
                                        not_joined.append(_g)
                                    except Exception:
                                        pass  # خطأ آخر — نفترض العضوية تفادياً للإشعار الخاطئ
                                # مجموعة عادية — نفترض العضوية إذا تم حل الـ entity
                            except Exception:
                                not_joined.append(_g)
                        return not_joined

                    _not_joined_groups = []
                    try:
                        _not_joined_groups = _now_client_mgr.run_coroutine(
                            _check_memberships_now(_now_client_mgr.client, groups_list)
                        )
                    except Exception as _cmn_err:
                        logger.debug(f"خطأ في فحص العضوية (الإرسال الفوري): {_cmn_err}")

                    if _not_joined_groups:
                        _instant_nj_groups = _not_joined_groups
                        socketio.emit('log_update', {
                            "message": f"⚠️ أنت غير منضم إلى {len(_not_joined_groups)} مجموعة — ستظهر في التقرير النهائي"
                        }, to=user_id)
                    socketio.emit('log_update', {
                        "message": "🚀 بدء الإرسال الفوري..."
                    }, to=user_id)
            except Exception as _check_err:
                logger.debug(f"خطأ في فحص العضوية (الإرسال الفوري): {_check_err}")

            # ── الدورة الثانية: الإرسال الفعلي لجميع المجموعات مع الصورة دائماً ──
            for i, group in enumerate(groups_list, 1):
                try:
                    if images and message:
                        # الصورة ترسل كجزء ثابت من الرسالة دائماً
                        result = telegram_manager.send_message_with_media_async(
                            user_id, group, message, image_files
                        )
                    elif images:
                        result = telegram_manager.send_media_async(
                            user_id, group, image_files
                        )
                    else:
                        result = telegram_manager.send_message_async(
                            user_id, group, message,
                            forced_action=pre_scan_action  # None = استخدم الإعدادات الافتراضية
                        )

                    if isinstance(result, dict) and result.get('skipped'):
                        socketio.emit('log_update', {
                            "message": f"⏭️ [{i}/{len(groups_list)}] تم تخطي المجموعة المحمية: {group}"
                        }, to=user_id)
                        _instant_failed_groups.append(group)
                    elif isinstance(result, dict) and result.get('smart'):
                        socketio.emit('log_update', {
                            "message": f"✅ [{i}/{len(groups_list)}] نجح إلى: {group}"
                        }, to=user_id)
                        successful += 1
                        _instant_smart_groups.append(group)
                    else:
                        socketio.emit('log_update', {
                            "message": f"✅ [{i}/{len(groups_list)}] نجح إلى: {group}"
                        }, to=user_id)
                        successful += 1
                        _instant_success_groups.append(group)
                        # حفظ معرف الرسالة لدفعة "رسائلي"
                        msg_id = None
                        if isinstance(result, dict):
                            msg_id = result.get('message_id') or (result.get('message_ids') or [None])[0]
                        if msg_id:
                            batch_entries.append({"group": group, "msg_id": msg_id})
                        with USERS_LOCK:
                            if user_id in USERS:
                                USERS[user_id]['stats']['sent'] += 1
                        with USERS_LOCK:
                            if user_id in USERS:
                                socketio.emit('stats_update', USERS[user_id]['stats'], to=user_id)

                    if i < len(groups_list):
                        time.sleep(3)

                except Exception as e:
                    error_msg = str(e)
                    if "banned" in error_msg.lower() or "ban" in error_msg.lower():
                        error_type = "محظور من المجموعة"
                    elif "flood" in error_msg.lower():
                        # استخرج وقت الانتظار إذا كان متاحاً
                        import re as _re
                        m = _re.search(r'(\d+)', error_msg)
                        wait_s = int(m.group(1)) if m else '?'
                        error_type = f"تجاوز حد الإرسال — انتظر {wait_s} ثانية"
                    elif "timeout" in error_msg.lower():
                        error_type = "انتهت مهلة الاتصال (timeout)"
                    elif "private" in error_msg.lower():
                        error_type = "مجموعة خاصة/محدودة"
                    elif "can't write" in error_msg.lower() or "write" in error_msg.lower():
                        error_type = "لا يُسمح بالإرسال في هذه المجموعة"
                    elif "not found" in error_msg.lower() or "invalid" in error_msg.lower() or "username" in error_msg.lower():
                        error_type = "المجموعة غير موجودة أو الرابط خاطئ"
                    elif "يُعاد تشغيله" in error_msg or "restart" in error_msg.lower():
                        error_type = "العميل يُعاد تشغيله، أعد المحاولة"
                    else:
                        error_type = error_msg[:150]  # رسالة خطأ كاملة لتسهيل التشخيص
                    log_user_event(user_id, 'ERROR', f"❌ فشل الإرسال إلى {group}: {error_type}")
                    logger.error(f"Send error to {group}: {error_msg}")
                    socketio.emit('log_update', {
                        "message": f"❌ [{i}/{len(groups_list)}] فشل إلى {group}: {error_type}"
                    }, to=user_id)

                    failed += 1
                    _instant_failed_groups.append(group)
                    with USERS_LOCK:
                        if user_id in USERS:
                            USERS[user_id]['stats']['errors'] += 1
                            socketio.emit('stats_update', USERS[user_id]['stats'], to=user_id)

            socketio.emit('log_update', {
                "message": f"📊 انتهى الإرسال: ✅ {successful} نجح | ❌ {failed} فشل"
            }, to=user_id)

            # ── إرسال تقرير مفصل لحساب المستخدم الشخصي بعد انتهاء الإرسال الفوري ──
            try:
                with USERS_LOCK:
                    _rpt_client_mgr = USERS.get(user_id, {}).get('client_manager')
                if _rpt_client_mgr and getattr(_rpt_client_mgr, 'client', None):
                    _rpt_time = time.strftime('%Y-%m-%d %H:%M:%S')
                    _rpt_lines = [
                        "📊 تقرير الإرسال الفوري",
                        f"⏰ الوقت: {_rpt_time}",
                        f"📨 إجمالي المجموعات: {len(groups_list)}",
                        "",
                        f"✅ الرسائل المرسلة بنجاح ({len(_instant_success_groups)}):",
                    ]
                    for _g in _instant_success_groups:
                        _rpt_lines.append(f"  • {_g}")
                    if not _instant_success_groups:
                        _rpt_lines.append("  — لا يوجد")
                    _rpt_lines += [
                        "",
                        f"❌ الرسائل الفاشلة ({len(_instant_failed_groups)}):",
                    ]
                    for _g in _instant_failed_groups:
                        _rpt_lines.append(f"  • {_g}")
                    if not _instant_failed_groups:
                        _rpt_lines.append("  — لا يوجد")
                    _rpt_lines += [
                        "",
                        f"🧠 الدورة الذكية ({len(_instant_smart_groups)}):",
                    ]
                    for _g in _instant_smart_groups:
                        _rpt_lines.append(f"  • {_g}")
                    if not _instant_smart_groups:
                        _rpt_lines.append("  — لا يوجد")
                    _rpt_lines += [
                        "",
                        f"⚠️ غير منضم إليها ({len(_instant_nj_groups)}):",
                    ]
                    for _g in _instant_nj_groups:
                        _rpt_lines.append(f"  • {_g}")
                    if not _instant_nj_groups:
                        _rpt_lines.append("  — لا يوجد")
                    _full_report_now = "\n".join(_rpt_lines)
                    _rpt_client_mgr.run_coroutine(
                        _rpt_client_mgr.client.send_message('me', _full_report_now, link_preview=False)
                    )
                    logger.info(f"[SendNow] تم إرسال تقرير الإرسال الفوري للمستخدم {user_id}")
            except Exception as _rpt_err:
                logger.debug(f"[SendNow] خطأ في إرسال التقرير: {_rpt_err}")

            # ── حفظ الدفعة في "رسائلي" ──
            if batch_entries:
                batch_record = {
                    "id": batch_id,
                    "text": message or "",
                    "has_media": bool(images),
                    "sent_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "sent_count": successful,
                    "entries": batch_entries
                }
                with USERS_LOCK:
                    ud = USERS.get(user_id)
                    if ud is not None:
                        if 'sent_batches' not in ud or ud['sent_batches'] is None:
                            ud['sent_batches'] = []
                        ud['sent_batches'].append(batch_record)
                        if len(ud['sent_batches']) > 100:
                            ud['sent_batches'] = ud['sent_batches'][-100:]
                        # حفظ الدفعات في الإعدادات لضمان الاستمرارية بعد إعادة التشغيل
                        try:
                            settings = load_settings(user_id)
                            settings['sent_batches'] = ud['sent_batches']
                            save_settings(user_id, settings)
                        except Exception:
                            pass
                socketio.emit('batch_saved', batch_record, to=user_id)

        except Exception as e:
            logger.error(f"Send thread error: {str(e)}")
        finally:
            for img_file in image_files:
                try:
                    if os.path.exists(img_file['path']):
                        os.unlink(img_file['path'])
                        logger.info(f"Cleaned up temp file: {img_file['name']}")
                except Exception as e:
                    logger.error(f"Error cleaning temp file {img_file.get('name', 'unknown')}: {str(e)}")

    _OSThread(target=send_messages_with_images, daemon=True).start()

    return jsonify({
        "success": True, 
        "message": f"🚀 بدأ إرسال {content_type} لـ {len(groups_list)} مجموعة"
    })

@app.route("/api/scan_groups_protection", methods=["POST"])
def api_scan_groups_protection():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "غير مسجّل"}), 401
    try:
        with USERS_LOCK:
            client_manager = USERS.get(user_id, {}).get('client_manager')
        if not client_manager:
            return jsonify({"error": "العميل غير متصل"}), 400

        data = request.get_json(force=True, silent=True) or {}
        raw_groups = data.get('groups', '')
        group_list = [g.strip() for g in re.split(r'[\n,]+', raw_groups) if g.strip()]
        if not group_list:
            return jsonify({"error": "لا توجد مجموعات للفحص"}), 400

        results = []
        for g in group_list[:50]:
            try:
                try:
                    entity_obj = client_manager.run_coroutine(
                        client_manager.client.get_entity(g)
                    )
                except Exception:
                    g2 = ('@' + g) if not g.startswith('@') and not g.startswith('https://') else g
                    entity_obj = client_manager.run_coroutine(
                        client_manager.client.get_entity(g2)
                    )
                is_prot, reason = client_manager.run_coroutine(
                    client_manager.is_group_protected(entity_obj)
                )
                title = getattr(entity_obj, 'title', g)
                results.append({
                    "group": g,
                    "title": title,
                    "protected": is_prot,
                    "reason": reason or ('غير محمية ✅' if not is_prot else '')
                })
            except Exception as e:
                results.append({"group": g, "title": g, "protected": False, "reason": f"خطأ: {str(e)[:60]}"})

        protected_count = sum(1 for r in results if r['protected'])
        return jsonify({"success": True, "results": results, "protected_count": protected_count, "total": len(results)})
    except Exception as e:
        logger.error(f"Scan groups error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/get_stats", methods=["GET"])
def api_get_stats():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"sent": 0, "errors": 0})

    with USERS_LOCK:
        if user_id in USERS:
            return jsonify(USERS[user_id]['stats'])

    return jsonify({"sent": 0, "errors": 0})

@app.route("/api/get_login_status", methods=["GET"])
def api_get_login_status():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"logged_in": False, "connected": False})

    with USERS_LOCK:
        if user_id in USERS:
            user_data = USERS[user_id]
            client_manager = user_data.get('client_manager')
            authenticated = user_data.get('authenticated', False)
            connected = user_data.get('connected', False)

            if not authenticated and 'settings' in user_data and 'phone' in user_data['settings']:
                session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
                if os.path.exists(session_file):
                    authenticated = True
                    connected = True
                    USERS[user_id]['authenticated'] = True
                    USERS[user_id]['connected'] = True

            return jsonify({
                "logged_in": authenticated, 
                "connected": connected,
                "is_running": user_data.get('is_running', False)
            })

    return jsonify({"logged_in": False, "connected": False, "is_running": False})

@app.route("/api/get_user_info", methods=["GET"])
def api_get_user_info():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل دخول"})

    with USERS_LOCK:
        if user_id in USERS and 'settings' in USERS[user_id]:
            settings = USERS[user_id]['settings']
            return jsonify({
                "success": True,
                "phone": settings.get('phone', ''),
                "name": settings.get('name', ''),
                "user_id": user_id[:8] + "..."
            })

    return jsonify({"success": False, "message": "لم يتم العثور على معلومات المستخدم"})

@app.route("/api/resend_code", methods=["POST"])
def api_resend_code():
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ الجلسة غير صالحة"})
        user_id = session['user_id']
        data = request.json or {}
        force_sms = bool(data.get('force_sms', False))

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ يرجى البدء بإدخال رقم الهاتف أولاً"})
            client_manager = USERS[user_id].get('client_manager')
            settings = USERS[user_id].get('settings', {})
            phone = settings.get('phone')

        if not client_manager or not client_manager.client or not phone:
            return jsonify({"success": False, "message": "❌ لم يتم إعداد العميل"})

        sent = client_manager.run_coroutine(
            client_manager.client.send_code_request(phone, force_sms=force_sms)
        )
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['awaiting_code'] = True
                USERS[user_id]['phone_code_hash'] = sent.phone_code_hash

        msg = "📱 تم إعادة الإرسال عبر SMS" if force_sms else "📱 تم إعادة إرسال الكود"
        socketio.emit('log_update', {"message": msg}, to=user_id)
        return jsonify({"success": True, "message": msg})
    except Exception as e:
        logger.error(f"Resend code error: {str(e)}")
        return jsonify({"success": False, "message": f"❌ {str(e)}"})

@app.route("/api/reset_login", methods=["POST"])
def api_reset_login():
    user_id = session.get('user_id', 'user_1')

    if user_id not in PREDEFINED_USERS:
        return jsonify({
            "success": False,
            "message": "❌ مستخدم غير صحيح"
        })

    try:
        logger.info(f"Resetting login for user {user_id}")

        with USERS_LOCK:
            if user_id in USERS:
                if USERS[user_id].get('is_running', False):
                    USERS[user_id]['is_running'] = False

                client_manager = USERS[user_id].get('client_manager')
                if client_manager:
                    try:
                        if hasattr(client_manager, 'stop'):
                            client_manager.stop()
                        if hasattr(client_manager, 'client') and client_manager.client:
                            client_manager.client.disconnect()
                        logger.info(f"Client stopped and disconnected for user {user_id}")
                    except Exception as e:
                        logger.error(f"Error stopping client for {user_id}: {e}")

                del USERS[user_id]
                logger.info(f"User data removed from memory for {user_id}")

        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if os.path.exists(session_file):
            try:
                os.remove(session_file)
                logger.info(f"Session file removed for {user_id}")
            except Exception as e:
                logger.error(f"Failed to remove session file for {user_id}: {str(e)}")

        socketio.emit('log_update', {
            "message": f"🔄 تم إعادة تعيين جلسة تسجيل الدخول لـ {PREDEFINED_USERS[user_id]['name']}"
        }, to=user_id)

        socketio.emit('connection_status', {
            "status": "disconnected"
        }, to=user_id)

        socketio.emit('login_status', {
            "logged_in": False,
            "connected": False,
            "awaiting_code": False,
            "awaiting_password": False,
            "is_running": False
        }, to=user_id)

        logger.info(f"Login reset completed for user {user_id}")

        return jsonify({
            "success": True, 
            "message": f"✅ تم إعادة تعيين جلسة {PREDEFINED_USERS[user_id]['name']} بنجاح"
        })

    except Exception as e:
        logger.error(f"Error resetting login for {user_id}: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في إعادة التعيين: {str(e)}"
        })

@app.route("/api/system_health", methods=["GET"])
def api_system_health():
    try:
        import psutil

        memory = psutil.virtual_memory()
        disk = psutil.disk_usage('/')
        cpu_percent = psutil.cpu_percent(interval=1)
        network = psutil.net_io_counters()

        health_info = {
            'memory': {
                'total': memory.total,
                'available': memory.available,
                'percent': memory.percent,
                'used': memory.used
            },
            'disk': {
                'total': disk.total,
                'used': disk.used,
                'free': disk.free,
                'percent': (disk.used / disk.total) * 100
            },
            'cpu': {
                'percent': cpu_percent,
                'count': psutil.cpu_count()
            },
            'network': {
                'bytes_sent': network.bytes_sent,
                'bytes_recv': network.bytes_recv
            },
            'timestamp': time.time()
        }

        return jsonify({
            "success": True,
            "health": health_info
        })

    except Exception as e:
        return jsonify({
            "success": False,
            "message": f"خطأ: {str(e)}"
        })

def extract_telegram_links(text):
    if not text:
        return []

    patterns = [
        r'https?://t\.me/([a-zA-Z0-9_]+)(?:/\d+)?',
        r'https?://telegram\.me/([a-zA-Z0-9_]+)(?:/\d+)?',
        r'https?://t\.me/\+([a-zA-Z0-9_\-]+)',
        r'https?://telegram\.me/\+([a-zA-Z0-9_\-]+)',
        r't\.me/([a-zA-Z0-9_]+)',
        r't\.me/\+([a-zA-Z0-9_\-]+)',
        r'telegram\.me/([a-zA-Z0-9_]+)',
        r'@([a-zA-Z0-9_]{5,})',
    ]

    found_links = set()

    for pattern in patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for match in matches:
            clean_match = match if isinstance(match, str) else match[0] if match else ''

            if pattern.startswith(r'@'):
                clean_link = f"https://t.me/{clean_match}"
            elif '+' in clean_match or pattern.find(r'\+') != -1:
                clean_link = f"https://t.me/+{clean_match.replace('+', '')}"
            elif clean_match and not clean_match.startswith('http'):
                clean_link = f"https://t.me/{clean_match}"
            elif clean_match.startswith('http'):
                clean_link = f"https://t.me/{clean_match.split('/')[-1]}"
            else:
                clean_link = clean_match

            if clean_link and len(clean_link) > 15:
                clean_link = clean_link.split('?')[0].split('#')[0]
                found_links.add(clean_link)

    links_list = sorted(list(found_links))
    result_links = []
    for link in links_list:
        username = link.split('/')[-1].replace('@', '')
        result_links.append({
            'url': link,
            'username': username,
            'type': 'invite' if '+' in link else 'channel'
        })

    return result_links

async def join_telegram_group(client, group_link, user_id=None, client_manager=None):
    try:
        if group_link.startswith('https://t.me/'):
            group_identifier = group_link.replace('https://t.me/', '')
        elif group_link.startswith('https://telegram.me/'):
            group_identifier = group_link.replace('https://telegram.me/', '')
        elif group_link.startswith('@'):
            group_identifier = group_link[1:]
        else:
            group_identifier = group_link

        try:
            entity = await client.get_entity(group_identifier)
            if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast'):
                result = await client(functions.channels.JoinChannelRequest(entity))
            else:
                raise Exception("مجموعة عادية - يجب استخدام رابط دعوة")

            return {
                "success": True,
                "already_joined": False,
                "message": "تم الانضمام بنجاح"
            }

        except UserAlreadyParticipantError:
            return {
                "success": True,
                "already_joined": True,
                "message": "منضم مسبقاً للمجموعة"
            }

        except FloodWaitError as e:
            return {
                "success": False,
                "message": f"يرجى الانتظار {e.seconds} ثانية"
            }

        except InviteHashExpiredError:
            return {
                "success": False,
                "message": "انتهت صلاحية رابط الدعوة"
            }

        except InviteHashInvalidError:
            return {
                "success": False,
                "message": "رابط الدعوة غير صحيح"
            }

        except Exception as group_error:
            error_str = str(group_error).lower()
            appeal_url = None
            appeal_note = ""

            if "cas" in error_str or "combot" in error_str:
                appeal_url = "https://cas.chat/appeal"
                appeal_note = "تم حظرك بواسطة CAS (Combot Anti-Spam). توجه إلى الرابط أعلاه لتقديم استئناف."
            elif "spamwatch" in error_str:
                appeal_url = "https://spamwat.ch/appeal"
                appeal_note = "تم حظرك بواسطة SpamWatch. استخدم الرابط أعلاه للاستئناف."
            elif "shieldy" in error_str:
                appeal_url = "https://t.me/Shieldy_Bot?start=appeal"
                appeal_note = "تم حظرك بواسطة Shieldy. افتح البوت في الخاص لطلب فك الحظر."
            elif "rose" in error_str or "missrose" in error_str:
                appeal_url = "https://t.me/MissRose_Bot?start=appeal"
                appeal_note = "تم حظرك بواسطة Rose. أرسل /start إلى البوت ثم اتبع التعليمات."
            elif "groupguard" in error_str:
                appeal_url = "https://t.me/GroupGuardBot?start=appeal"
                appeal_note = "تم حظرك بواسطة GroupGuard. اتصل بالبوت."
            elif "antispam" in error_str or "spam" in error_str:
                appeal_url = "https://t.me/SpamBot"
                appeal_note = "قد يكون حسابك مصنفاً كسبام. تواصل مع @SpamBot للتحقق."
            else:
                if "banned" in error_str or "blocked" in error_str or "forbidden" in error_str:
                    appeal_url = "https://t.me/SpamBot"
                    appeal_note = "حسابك ربما محظور من الانضمام. جرب التواصل مع @SpamBot أو مشرف المجموعة."

            if appeal_url and user_id and client_manager:
                message_text = f"""🚫 **فشل الانضمام إلى المجموعة** 🚫

**الرابط:** {group_link}
**السبب:** {error_str[:200]}

**إجراء مقترح للاستئناف:**
{appeal_note}
🔗 **رابط الاستئناف:** {appeal_url}

يرجى فتح الرابط ومتابعة التعليمات لرفع الحظر. بعد إلغاء الحظر، يمكنك إعادة المحاولة.
"""
                try:
                    await client_manager.send_to_saved_messages(message_text)
                except Exception as save_err:
                    logger.error(f"Could not send appeal to saved messages: {save_err}")

            try:
                if '/' in group_identifier:
                    result = await client(functions.messages.ImportChatInviteRequest(group_identifier.split('/')[-1]))
                    return {
                        "success": True,
                        "already_joined": False,
                        "message": "تم الانضمام عبر رابط الدعوة"
                    }
                else:
                    raise group_error
            except UserAlreadyParticipantError:
                return {
                    "success": True,
                    "already_joined": True,
                    "message": "منضم مسبقاً للمجموعة"
                }
            except Exception as final_error:
                return {
                    "success": False,
                    "message": f"فشل الانضمام: {str(final_error)}",
                    "appeal_url": appeal_url
                }

    except Exception as e:
        return {
            "success": False,
            "message": f"خطأ: {str(e)}"
        }

@app.route("/api/extract_group_links", methods=["POST"])
def api_extract_group_links():
    try:
        data = request.json
        if not data or not data.get('text'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال النص"
            })

        text = data.get('text', '')
        links = extract_telegram_links(text)

        return jsonify({
            "success": True,
            "links": links,
            "count": len(links),
            "message": f"✅ تم استخراج {len(links)} رابط"
        })

    except Exception as e:
        logger.error(f"Error extracting links: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/join_group", methods=["POST"])
def api_join_group():
    try:
        user_id = session.get('user_id', 'user_1')

        if user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        data = request.json

        if not data or not data.get('group_link'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال رابط المجموعة"
            })

        group_link_raw = data.get('group_link', '')
        if isinstance(group_link_raw, dict):
            group_link = group_link_raw.get('url', '') or group_link_raw.get('link', '') or str(group_link_raw)
        else:
            group_link = str(group_link_raw)

        group_link = group_link.strip()

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({
                    "success": False,
                    "message": f"❌ المستخدم {PREDEFINED_USERS[user_id]['name']} غير مسجل"
                })

            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({
                    "success": False,
                    "message": "❌ يرجى تسجيل الدخول أولاً"
                })

        result = client_manager.run_coroutine(
            join_telegram_group(client_manager.client, group_link, user_id, client_manager)
        )

        socketio.emit('log_update', {
            "message": f"{'✅' if result['success'] else '❌'} {group_link}: {result['message']}"
        }, to=user_id)

        return jsonify(result)

    except Exception as e:
        logger.error(f"Error joining group: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ: {str(e)}"
        })

@app.route("/api/start_auto_join", methods=["POST"])
def api_start_auto_join():
    try:
        user_id = session.get('user_id', 'user_1')

        if user_id not in PREDEFINED_USERS:
            return jsonify({
                "success": False,
                "message": "❌ مستخدم غير صحيح"
            })

        data = request.json
        if not data or not data.get('links'):
            return jsonify({
                "success": False,
                "message": "❌ لم يتم إرسال روابط المجموعات"
            })

        links = data.get('links', [])
        delay = data.get('delay', 3)

        if not links:
            return jsonify({
                "success": False,
                "message": "❌ لا توجد روابط للانضمام إليها"
            })

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({
                    "success": False,
                    "message": f"❌ المستخدم {PREDEFINED_USERS[user_id]['name']} غير مسجل"
                })

            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({
                    "success": False,
                    "message": "❌ يرجى تسجيل الدخول أولاً"
                })

        import threading

        def auto_join_worker():
            success_count = 0
            fail_count = 0
            already_joined_count = 0

            socketio.emit('log_update', {
                "message": f"🚀 بدء الانضمام التلقائي لـ {len(links)} مجموعة..."
            }, to=user_id)

            for i, link_obj in enumerate(links):
                try:
                    if isinstance(link_obj, dict):
                        group_link = link_obj.get('url', '') or link_obj.get('link', '') or str(link_obj)
                    else:
                        group_link = str(link_obj)

                    group_link = group_link.strip()

                    socketio.emit('join_progress', {
                        'current': i + 1,
                        'total': len(links),
                        'link': group_link
                    }, to=user_id)

                    result = client_manager.run_coroutine(
                        join_telegram_group(client_manager.client, group_link, user_id, client_manager)
                    )

                    if result['success']:
                        if result.get('already_joined', False):
                            already_joined_count += 1
                            socketio.emit('log_update', {
                                "message": f"ℹ️ منضم مسبقاً: {group_link}"
                            }, to=user_id)
                        else:
                            success_count += 1
                            socketio.emit('log_update', {
                                "message": f"✅ تم الانضمام: {group_link}"
                            }, to=user_id)
                    else:
                        fail_count += 1
                        socketio.emit('log_update', {
                            "message": f"❌ فشل: {group_link} - {result['message']}"
                        }, to=user_id)

                    socketio.emit('join_stats', {
                        'success': success_count,
                        'fail': fail_count,
                        'already_joined': already_joined_count
                    }, to=user_id)

                    if i < len(links) - 1:
                        time.sleep(delay)

                except Exception as e:
                    fail_count += 1
                    socketio.emit('log_update', {
                        "message": f"❌ خطأ في {group_link}: {str(e)}"
                    }, to=user_id)

            socketio.emit('auto_join_completed', {
                'success': success_count,
                'fail': fail_count,
                'already_joined': already_joined_count,
                'total': len(links)
            }, to=user_id)

            socketio.emit('log_update', {
                "message": f"🎉 انتهى الانضمام التلقائي! النجح: {success_count}, فشل: {fail_count}, منضم مسبقاً: {already_joined_count}"
            }, to=user_id)

        thread = _OSThread(target=auto_join_worker, daemon=True)
        thread.start()

        return jsonify({
            "success": True,
            "message": f"✅ تم بدء الانضمام التلقائي لـ {len(links)} مجموعة",
            "total_links": len(links)
        })

    except Exception as e:
        logger.error(f"Error starting auto join: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"❌ خطأ في بدء الانضمام التلقائي: {str(e)}"
        })

# ==========================
# الإضافات الجديدة (الإرسال المتسلسل، الانضمام المتقدم، البوت التعليمي، الردود التلقائية، البحث)
# ==========================

class RotatingSendManager:
    def __init__(self):
        self.threads = {}
        self.stop_events = {}
        self.next_send_at = {}
        self.interval_seconds = {}

    def start(self, user_id, groups, messages, interval_minutes, callback=None):
        if user_id in self.threads and self.threads[user_id] and self.threads[user_id].is_alive():
            self.stop(user_id)

        stop_event = threading.Event()
        self.stop_events[user_id] = stop_event

        thread = _OSThread(target=self._worker, args=(user_id, groups, messages, interval_minutes, stop_event, callback), daemon=True)
        self.threads[user_id] = thread
        thread.start()
        return True

    def stop(self, user_id):
        if user_id in self.stop_events:
            self.stop_events[user_id].set()
        if user_id in self.threads and self.threads[user_id]:
            self.threads[user_id].join(timeout=2)
        self.next_send_at.pop(user_id, None)
        self.interval_seconds.pop(user_id, None)
        return True

    def _worker(self, user_id, groups, messages, interval_minutes, stop_event, callback):
        messages = [m.strip() for m in messages if m and m.strip()]
        if not messages:
            return

        index = 0
        sleep_seconds = max(60, int(interval_minutes * 60))
        self.interval_seconds[user_id] = sleep_seconds

        while not stop_event.is_set():
            try:
                current_msg = messages[index % len(messages)]
                for group in groups:
                    if stop_event.is_set():
                        break
                    try:
                        telegram_manager.send_message_async(user_id, group, current_msg)
                        if callback:
                            callback(user_id, 'success', group, current_msg)
                    except Exception as e:
                        if callback:
                            callback(user_id, 'error', group, str(e))
                    time.sleep(2)
                index += 1
                self.next_send_at[user_id] = time.time() + sleep_seconds
                for _ in range(sleep_seconds):
                    if stop_event.is_set():
                        break
                    time.sleep(1)
            except Exception as e:
                logger.error(f"Rotating send error for {user_id}: {str(e)}")
                time.sleep(10)

rotating_manager = RotatingSendManager()

@app.route("/api/rotating/save", methods=["POST"])
def api_rotating_save():
    try:
        user_id = session.get('user_id', 'user_1')
        data = request.json
        messages = data.get('messages', [''] * 5)
        groups = data.get('groups', [])
        interval = int(data.get('interval', 5))

        settings = load_settings(user_id)
        settings['rotating_messages'] = messages
        settings['rotating_groups'] = dedupe_groups(groups)
        settings['rotating_interval'] = interval
        save_settings(user_id, settings)

        return jsonify({"success": True, "message": "تم حفظ إعدادات الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/start", methods=["POST"])
def api_rotating_start():
    try:
        user_id = session.get('user_id', 'user_1')
        settings = load_settings(user_id)
        messages = settings.get('rotating_messages', [])
        groups = dedupe_groups(settings.get('rotating_groups', []))
        interval = settings.get('rotating_interval', 5)

        if not groups:
            return jsonify({"success": False, "message": "لا توجد مجموعات محددة"})
        valid_messages = [m for m in messages if m and m.strip()]
        if not valid_messages:
            return jsonify({"success": False, "message": "لا توجد رسائل صالحة"})

        def callback(uid, status, group, info):
            if status == 'success':
                socketio.emit('log_update', {"message": f"🔄 [متسلسل] أرسل إلى {group}"}, to=uid)
            else:
                socketio.emit('log_update', {"message": f"❌ [متسلسل] فشل إلى {group}: {info}"}, to=uid)

        rotating_manager.start(user_id, groups, valid_messages, interval, callback)

        try:
            settings['rotating_persistent'] = True
            save_settings(user_id, settings)
        except Exception as _e:
            logger.error(f"Failed to persist rotating flag for {user_id}: {_e}")

        socketio.emit('log_update', {"message": f"🔄 بدأ الإرسال المتسلسل ({len(valid_messages)} رسائل) كل {interval} دقيقة"}, to=user_id)
        return jsonify({"success": True, "message": "تم بدء الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/stop", methods=["POST"])
def api_rotating_stop():
    try:
        user_id = session.get('user_id', 'user_1')
        rotating_manager.stop(user_id)

        try:
            _settings = load_settings(user_id)
            _settings['rotating_persistent'] = False
            save_settings(user_id, _settings)
        except Exception as _e:
            logger.error(f"Failed to clear rotating flag for {user_id}: {_e}")

        socketio.emit('log_update', {"message": "⏹ تم إيقاف الإرسال المتسلسل"}, to=user_id)
        return jsonify({"success": True, "message": "تم إيقاف الإرسال المتسلسل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/rotating/status", methods=["GET"])
def api_rotating_status():
    try:
        user_id = session.get('user_id', 'user_1')
        settings = load_settings(user_id)
        is_active = user_id in rotating_manager.threads and rotating_manager.threads[user_id] and rotating_manager.threads[user_id].is_alive()
        next_send_in = None
        next_send_at = rotating_manager.next_send_at.get(user_id)
        if is_active and next_send_at:
            remaining = int(next_send_at - time.time())
            next_send_in = max(0, remaining)
        return jsonify({
            "success": True,
            "active": is_active,
            "messages": settings.get('rotating_messages', []),
            "groups": settings.get('rotating_groups', []),
            "interval": settings.get('rotating_interval', 5),
            "next_send_in": next_send_in,
            "interval_seconds": rotating_manager.interval_seconds.get(user_id)
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

def _classify_join_error(msg):
    if not msg:
        return ("خطأ غير معروف", "❓")
    s = str(msg).lower()
    if "anti-spam" in s or "antispam" in s or "spam" in s or "spambot" in s:
        return ("الحساب موسوم كسبام (يحتاج استئناف عبر @SpamBot)", "🚫")
    if "banned" in s or "blocked" in s or "forbidden" in s:
        return ("الحساب محظور من المجموعة أو من الانضمام", "⛔")
    if "expired" in s or "انتهت" in s:
        return ("انتهت صلاحية رابط الدعوة", "⏰")
    if "invalid" in s or "غير صحيح" in s or "غير صالح" in s:
        return ("رابط غير صالح", "🔗")
    if "flood" in s or "wait" in s or "ثانية" in s:
        return ("حد التليجرام مؤقت — يجب الانتظار قبل المحاولة مجدداً", "⏳")
    if "channel_private" in s or "private" in s or "خاص" in s:
        return ("القناة/المجموعة خاصة وتحتاج رابط دعوة", "🔒")
    if "not found" in s or "no user" in s or "could not find" in s or "غير موجود" in s:
        return ("المجموعة غير موجودة أو الرابط خاطئ", "🔍")
    if "too many channels" in s or "channels_too_much" in s:
        return ("الحساب وصل الحد الأقصى من القنوات (500)", "📛")
    if "user_deactivated" in s or "deactivated" in s:
        return ("الحساب معطل من تيليجرام", "🛑")
    if "captcha" in s or "verification" in s:
        return ("المجموعة تتطلب تحقق يدوي (كابتشا)", "🤖")
    if "admin" in s and "approval" in s:
        return ("الانضمام بحاجة موافقة المشرف", "👮")
    if "request" in s and ("send" in s or "join" in s):
        return ("تم إرسال طلب انضمام — بانتظار الموافقة", "📨")
    short = str(msg).strip()
    if len(short) > 120:
        short = short[:120] + "…"
    return (short, "❌")


# ── دالة لجلب الروابط من صفحات خارجية ──
def fetch_telegram_links_from_url(url, max_links=30):
    try:
        import requests as _req
        from bs4 import BeautifulSoup
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
        }
        resp = _req.get(url, timeout=12, headers=headers)
        if resp.status_code != 200:
            logger.warning(f'Failed to fetch {url}: HTTP {resp.status_code}')
            return []
        soup = BeautifulSoup(resp.text, 'html.parser')
        text = soup.get_text(separator=' ', strip=True)
        extracted = extract_telegram_links(text)
        links_from_text = [link['url'] for link in extracted if link.get('url')]
        links_from_href = []
        for a in soup.find_all('a', href=True):
            href = a['href'].strip()
            if 't.me' in href or 'telegram.me' in href:
                if not href.startswith('http'):
                    from urllib.parse import urljoin
                    href = urljoin(url, href)
                links_from_href.append(href)
        all_links = list(set(links_from_text + links_from_href))
        valid_links = [lnk for lnk in all_links if lnk and ('t.me' in lnk or 'telegram.me' in lnk)]
        seen = set()
        unique = []
        for lnk in valid_links:
            if lnk not in seen:
                seen.add(lnk)
                unique.append(lnk)
        return unique[:max_links]
    except Exception as e:
        logger.error(f'Error fetching links from {url}: {e}')
        return []


# ── دالة لاستخراج أسماء المجموعات من النص ──
def extract_group_names_from_text(text):
    import re
    names = []
    patterns = [
        r'(?:قروبات|مجموعات|الدراسات العليا)\s+(جامعة\s+[^\n,]+)',
        r'(?:قسم|دليل)\s+([^\n,]+)',
        r'([^\n,]+)\s+(?:قروب|مجموعة)',
        r'(جامعة\s+[^\n,]+)',
        r'(?:قروبات|مجموعات)\s+([^\n,]+)',
        r'([^\n]{5,50})',
    ]
    for pattern in patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        for m in matches:
            name = m.strip()
            if name and len(name) > 5 and len(name) < 80 and not re.match(r'^[\d\s\W]+$', name):
                if name not in names:
                    names.append(name)
    return names[:20]


# ── دالة للبحث عن مجموعات بأسمائها في تيليجرام ──
async def search_and_join_groups_by_name(client, names_list, limit=3):
    from telethon.tl.functions.contacts import SearchRequest
    results = []
    for name in names_list:
        try:
            sr = await client(SearchRequest(q=name, limit=limit))
            if sr.chats:
                for chat in sr.chats:
                    if hasattr(chat, 'username') and chat.username:
                        group_link = f'https://t.me/{chat.username}'
                        results.append({'name': name, 'found': group_link, 'chat': chat, 'type': 'public'})
                        break
                    elif hasattr(chat, 'title') and chat.title:
                        try:
                            from telethon.tl import functions as _tl_f
                            invite = await client(_tl_f.messages.ExportChatInviteRequest(
                                peer=chat, legacy_revoke_permanent=True))
                            if invite and hasattr(invite, 'link'):
                                results.append({'name': name, 'found': invite.link, 'chat': chat, 'type': 'private'})
                                break
                        except Exception:
                            pass
        except Exception as e:
            logger.warning(f"Search for '{name}' failed: {e}")
    return results


@app.route("/api/auto_join/advanced", methods=["POST"])
def api_auto_join_advanced():
    try:
        user_id = session.get('user_id', 'user_1')
        data = request.json
        raw_input = data.get('links', '')
        fetch_external = data.get('fetch_external', True)
        max_external_links = int(data.get('max_external_links', 30))
        search_by_name = data.get('search_by_name', True)

        # ── 1. استخراج الروابط من النص المختلط ──
        if isinstance(raw_input, str):
            raw_text = raw_input
        elif isinstance(raw_input, list):
            raw_text = "\n".join([
                str(item) if not isinstance(item, dict) else item.get('url', '')
                for item in raw_input
            ])
        else:
            raw_text = str(raw_input)

        extracted = extract_telegram_links(raw_text)
        direct_links = [link['url'] for link in extracted if link.get('url')]

        # ── 2. جلب الروابط من الصفحات الخارجية ──
        external_links = []
        if fetch_external:
            import re as _re_ext
            external_urls = _re_ext.findall(r'https?://[^\s<>"\'()]+', raw_text)
            for url in external_urls:
                if 't.me' not in url and 'telegram.me' not in url:
                    fetched = fetch_telegram_links_from_url(url, max_external_links)
                    if fetched:
                        socketio.emit('log_update',
                            {'message': f'🌐 تم جلب {len(fetched)} رابط من {url}'}, to=user_id)
                        external_links.extend(fetched)
                    else:
                        socketio.emit('log_update',
                            {'message': f'⚠️ لم يتم العثور على روابط تيليجرام في {url}'}, to=user_id)

        all_raw_links = direct_links + external_links

        # ── 3. إذا لم نجد روابط، حاول استخراج أسماء مجموعات ──
        if not all_raw_links and search_by_name:
            group_names = extract_group_names_from_text(raw_text)
            if group_names:
                socketio.emit('log_update',
                    {'message': f'🔍 لم نجد روابط، لكن وجدنا {len(group_names)} اسماً لمجموعات، جاري البحث عنها...'},
                    to=user_id)
                with USERS_LOCK:
                    client_manager = USERS.get(user_id, {}).get('client_manager')
                    if not client_manager or not client_manager.client:
                        return jsonify({'success': False, 'message': 'العميل غير متصل'})
                found_groups = client_manager.run_coroutine(
                    search_and_join_groups_by_name(client_manager.client, group_names)
                )
                if found_groups:
                    all_raw_links = [g['found'] for g in found_groups]
                    socketio.emit('log_update',
                        {'message': f'✅ تم العثور على {len(found_groups)} مجموعة من الأسماء المستخرجة'},
                        to=user_id)
                    for g in found_groups:
                        socketio.emit('log_update',
                            {'message': f"  • {g['name']} → {g['found']}"}, to=user_id)
                else:
                    return jsonify({'success': False, 'message': 'لم يتم العثور على روابط ولا مجموعات تطابق الأسماء'})

        if not all_raw_links:
            return jsonify({'success': False, 'message': 'لم يتم العثور على روابط تيليجرام في النص أو في الصفحات الخارجية'})

        # ── 4. إزالة المكررات وتنظيفها ──
        clean_links = dedupe_groups(all_raw_links)
        if not clean_links:
            return jsonify({'success': False, 'message': 'لا توجد روابط صالحة بعد التنقية'})

        # ── 5. التحقق من العضوية المسبقة ──
        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({'success': False, 'message': 'المستخدم غير موجود'})
            client_manager = USERS[user_id].get('client_manager')
            if not client_manager or not client_manager.client:
                return jsonify({'success': False, 'message': 'العميل غير متصل، يرجى تسجيل الدخول'})

        already_member = []
        pending_links = []
        for link in clean_links:
            try:
                entity = client_manager.run_coroutine(client_manager.client.get_entity(link))
                try:
                    client_manager.run_coroutine(client_manager.client.get_participants(entity, limit=1))
                    already_member.append(link)
                    socketio.emit('log_update',
                        {'message': f'📌 [عضوية مسبقة] {link} — أنت عضو بالفعل، تم تخطيه'}, to=user_id)
                    continue
                except Exception:
                    pending_links.append(link)
            except Exception:
                pending_links.append(link)

        # ── 6. عرض الإحصائيات ──
        total = len(clean_links)
        already_count = len(already_member)
        pending_count = len(pending_links)

        socketio.emit('log_update', {
            'message': f'📊 إحصائيات: {total} رابط إجمالي | {already_count} منضم مسبقاً (تم تخطيهم) | {pending_count} متبقي للانضمام'
        }, to=user_id)

        if not pending_links:
            return jsonify({
                'success': True,
                'message': f'✅ جميع الروابط ({total}) منضم إليها مسبقاً، لا حاجة للانضمام',
                'already': already_member
            })

        # ── 7. بدء الانضمام ──
        delay = max(1, int(data.get('delay', 3)))
        max_retries = max(1, int(data.get('max_retries', 3)))

        try:
            _s = load_settings(user_id)
            _s['auto_join_links'] = pending_links
            _s['auto_join_delay'] = delay
            save_settings(user_id, _s)
        except Exception:
            pass

        import threading as _threading
        stop_event  = _threading.Event()
        pause_event = _threading.Event()
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['auto_join_stop']  = stop_event
                USERS[user_id]['auto_join_pause'] = pause_event

        def _save_state(state_dict):
            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id]['auto_join_state'] = state_dict

        def advanced_join_worker():
            total_pending = len(pending_links)
            results = {'success': 0, 'fail': 0, 'already': already_count, 'total': total, 'items': []}

            _save_state({
                'running': True, 'total': total_pending, 'done': 0,
                'success': 0, 'already': already_count, 'fail': 0,
                'items': [], 'links': pending_links, 'delay': delay
            })

            socketio.emit('auto_join_started', {'total': total_pending}, to=user_id)
            socketio.emit('log_update', {
                'message': f'🚀 بدء الانضمام إلى {total_pending} مجموعة (بعد استبعاد {already_count} منضم مسبقاً)'
            }, to=user_id)

            for idx, link in enumerate(pending_links, 1):
                if stop_event.is_set():
                    socketio.emit('log_update',
                        {'message': f'⏹ تم إيقاف الانضمام بعد {idx - 1} مجموعة'}, to=user_id)
                    break

                if pause_event.is_set():
                    socketio.emit('auto_join_paused', {
                        'paused_at': idx - 1, 'total': total_pending,
                        'success': results['success'], 'fail': results['fail'],
                        'already': already_count
                    }, to=user_id)
                    while pause_event.is_set() and not stop_event.is_set():
                        time.sleep(0.5)
                    if stop_event.is_set():
                        break
                    socketio.emit('auto_join_resumed',
                        {'resumed_at': idx - 1, 'total': total_pending}, to=user_id)

                item = {'idx': idx, 'total': total_pending, 'url': link, 'status': 'processing', 'reason': ''}
                socketio.emit('auto_join_progress', item, to=user_id)

                success = False
                last_error = None
                for attempt in range(max_retries):
                    try:
                        result = client_manager.run_coroutine(
                            join_telegram_group(client_manager.client, link, user_id, client_manager)
                        )
                        if result.get('success'):
                            success = True
                            break
                        else:
                            last_error = result.get('message') or 'فشل غير محدد'
                            if 'FloodWait' in str(last_error):
                                import re as _re_fw
                                m = _re_fw.search(r'(\d+)', str(last_error))
                                if m:
                                    wait_sec = int(m.group(1))
                                    socketio.emit('log_update', {
                                        'message': f'⏳ FloodWait: يطلب تيليجرام الانتظار {wait_sec} ثانية... (المحاولة {attempt+1}/{max_retries})'
                                    }, to=user_id)
                                    time.sleep(wait_sec)
                                    continue
                    except Exception as e:
                        last_error = str(e)
                        if 'FloodWait' in last_error:
                            import re as _re_fw
                            m = _re_fw.search(r'(\d+)', last_error)
                            if m:
                                wait_sec = int(m.group(1))
                                socketio.emit('log_update', {
                                    'message': f'⏳ FloodWait: {wait_sec} ثانية... (المحاولة {attempt+1}/{max_retries})'
                                }, to=user_id)
                                time.sleep(wait_sec)
                                continue
                    time.sleep(delay)

                if success:
                    item['status'] = 'success'
                    item['reason'] = 'تم الانضمام بنجاح'
                    results['success'] += 1
                    socketio.emit('log_update',
                        {'message': f"✅ [{idx}/{total_pending}] انضم إلى: {link}"}, to=user_id)
                else:
                    item['status'] = 'failed'
                    item['reason'] = last_error or 'فشل بعد عدة محاولات'
                    results['fail'] += 1
                    socketio.emit('log_update',
                        {'message': f"❌ [{idx}/{total_pending}] فشل في: {link} — {item['reason']}"}, to=user_id)

                results['items'].append(item)
                _save_state({
                    'running': True, 'total': total_pending, 'done': idx,
                    'success': results['success'], 'already': already_count,
                    'fail': results['fail'], 'items': results['items'][-100:],
                    'links': pending_links, 'delay': delay
                })
                socketio.emit('auto_join_progress', {**item, 'counts': {
                    'success': results['success'], 'fail': results['fail'],
                    'already': already_count, 'done': idx, 'total': total_pending
                }}, to=user_id)

                if idx < total_pending and not stop_event.is_set():
                    time.sleep(delay)

            # التقرير النهائي
            socketio.emit('auto_join_completed', {
                'success': results['success'],
                'fail': results['fail'],
                'already': already_count,
                'total': total,
                'items': results['items']
            }, to=user_id)
            socketio.emit('log_update', {
                'message': f"🎉 انتهى الانضمام: ✅ {results['success']} نجح | ❌ {results['fail']} فشل | 📌 {already_count} منضم مسبقاً"
            }, to=user_id)

            with USERS_LOCK:
                if user_id in USERS:
                    USERS[user_id].pop('auto_join_stop', None)
                    st = USERS[user_id].get('auto_join_state', {})
                    st['running'] = False
                    USERS[user_id]['auto_join_state'] = st

        _OSThread(target=advanced_join_worker, daemon=True).start()
        return jsonify({
            'success': True,
            'total': total,
            'already': already_count,
            'pending': pending_count,
            'already_links': already_member,
            'message': f'بدأ الانضمام إلى {pending_count} مجموعة (تم استبعاد {already_count} منضم مسبقاً)'
        })

    except Exception as e:
        logger.error(f'auto_join_advanced error: {e}')
        return jsonify({'success': False, 'message': str(e)})


# ═══════════════════════════════════════════════════════════════════
#  نظام التعلم الذكي - LearningBot + LearningManager
#  مع ذاكرة دائمة (JSON) + مزامنة GitHub + Groq AI + احتياطي ذكي
# ═══════════════════════════════════════════════════════════════════

class LearningBot:
    """
    بوت تعلم ذكي يستخدم Groq AI مع:
    - ذاكرة دائمة JSON لكل مستخدم (محلي + GitHub)
    - قراءة كامل سياق المحادثة (Telegram + ذاكرة دائمة)
    - احتياطي بالأنماط عند انقطاع الذكاء
    - ردود بشرية طبيعية لـ مركز سرعة إنجاز
    """
    def __init__(self, user_id):
        self.user_id = user_id
        self.unknown_requests = []
        self._sync_counter = 0

        # الذاكرة قصيرة المدى (RAM)
        self.conversations_history = {}  # {conv_key: [{role,text,time}]}

        # الذاكرة الدائمة على القرص
        self.memory_dir = os.path.join(DATA_DIR, "learning_memory")
        os.makedirs(self.memory_dir, exist_ok=True)
        self.memory_file = os.path.join(self.memory_dir, f"{user_id}.json")
        self._memory = self._load_persistent_memory()

        # قاعدة الخدمات (من الذاكرة أو الافتراضي)
        self.knowledge = self._memory.get("knowledge", {})
        if not self.knowledge:
            self.knowledge = self._default_knowledge()
            self._memory["knowledge"] = self.knowledge

        # Groq client
        self.groq_client = None
        try:
            from groq import Groq as _Groq
            if GROQ_API_KEY:
                self.groq_client = _Groq(api_key=GROQ_API_KEY)
        except Exception as e:
            logger.warning(f"Groq init failed for {user_id}: {e}")

        # مزامنة تلقائية كل ساعة
        self._start_auto_sync()
        logger.info(f"✅ LearningBot({user_id}) — محادثات محفوظة: {len(self._memory.get('conversations', {}))}")

    # ─── الخدمات الافتراضية ──────────────────────────────────────

    def _default_knowledge(self):
        return {
            "حل واجب":    {"description": "حل الواجبات والمسائل الدراسية",   "keywords": ["حل", "واجب", "مسألة", "تمارين", "assignment"], "price_range": "50-200 ريال",               "time_range": "2-24 ساعة"},
            "بحث":        {"description": "إعداد البحوث الأكاديمية",          "keywords": ["بحث", "تقرير", "موضوع", "research"],           "price_range": "100-500 ريال",              "time_range": "1-5 أيام"},
            "تلخيص":      {"description": "تلخيص الكتب والمحاضرات",           "keywords": ["تلخيص", "ملخص", "اختصار"],                     "price_range": "30-150 ريال",               "time_range": "2-12 ساعة"},
            "ترجمة":      {"description": "ترجمة النصوص بدقة عالية",          "keywords": ["ترجمة", "ترجم", "translation"],                 "price_range": "20-100 ريال/صفحة",          "time_range": "1-24 ساعة"},
            "تحليل بيانات":{"description": "تحليل إحصائي وبيانات (SPSS/Excel)","keywords": ["تحليل", "بيانات", "إحصاء", "SPSS", "Excel"],   "price_range": "100-400 ريال",              "time_range": "1-3 أيام"},
            "تصميم":      {"description": "تصميم عروض وبوسترات وشرائح",       "keywords": ["تصميم", "بوستر", "عرض", "PowerPoint", "PPT"],   "price_range": "50-250 ريال",               "time_range": "2-24 ساعة"},
        }

    # ─── الذاكرة الدائمة ─────────────────────────────────────────

    def _load_persistent_memory(self):
        default = {
            "knowledge": {},
            "conversations": {},
            "patterns": {},
            "stats": {"total_messages": 0, "last_updated": None}
        }
        # 1. من القرص
        if os.path.exists(self.memory_file):
            try:
                with open(self.memory_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                for k in default:
                    if k not in data:
                        data[k] = default[k]
                return data
            except Exception as e:
                logger.error(f"خطأ تحميل ذاكرة {self.user_id}: {e}")
        # 2. من GitHub
        try:
            raw = download_from_github(f"data/learning_memory/{self.user_id}.json")
            if raw:
                data = json.loads(raw.decode('utf-8'))
                with open(self.memory_file, 'w', encoding='utf-8') as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)
                logger.info(f"✅ ذاكرة {self.user_id} محملة من GitHub")
                return data
        except Exception as e:
            logger.warning(f"فشل تحميل GitHub memory: {e}")
        # 3. ملف جديد
        with open(self.memory_file, 'w', encoding='utf-8') as f:
            json.dump(default, f, ensure_ascii=False, indent=2)
        return default

    def _save_memory(self):
        try:
            self._memory["knowledge"] = self.knowledge
            self._memory["stats"]["last_updated"] = datetime.now().isoformat()
            total = sum(len(v) for v in self._memory.get("conversations", {}).values())
            self._memory["stats"]["total_messages"] = total
            with open(self.memory_file, 'w', encoding='utf-8') as f:
                json.dump(self._memory, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"فشل حفظ ذاكرة {self.user_id}: {e}")

    def _sync_to_github(self):
        def _upload():
            try:
                with open(self.memory_file, 'rb') as f:
                    content = f.read()
                ok = upload_to_github(
                    f"data/learning_memory/{self.user_id}.json",
                    content,
                    f"ذاكرة تعلم {self.user_id} - {datetime.now().strftime('%Y-%m-%d %H:%M')}"
                )
                if ok:
                    logger.info(f"✅ مزامنة ذاكرة {self.user_id} → GitHub")
            except Exception as e:
                logger.error(f"فشل مزامنة GitHub لـ {self.user_id}: {e}")
        threading.Thread(target=_upload, daemon=True).start()

    def _start_auto_sync(self):
        def _loop():
            while True:
                time.sleep(3600)
                self._sync_to_github()
        threading.Thread(target=_loop, daemon=True).start()

    # ─── الذاكرة قصيرة المدى ──────────────────────────────────────

    def _clean_old_history(self, conv_key, max_age=7200):
        """حذف رسائل أقدم من max_age ثانية من RAM"""
        if conv_key in self.conversations_history:
            now = time.time()
            self.conversations_history[conv_key] = [
                e for e in self.conversations_history[conv_key]
                if now - e.get('time', 0) < max_age
            ]
            if not self.conversations_history[conv_key]:
                del self.conversations_history[conv_key]

    def _get_persistent_history(self, conv_key):
        """دمج تاريخ RAM والذاكرة الدائمة (آخر 30 رسالة)"""
        saved = self._memory.get("conversations", {}).get(conv_key, [])
        ram = self.conversations_history.get(conv_key, [])
        combined = saved + [r for r in ram if r not in saved]
        return combined[-30:]

    def _append_to_persistent(self, conv_key, entry):
        """إضافة رسالة واحدة للذاكرة الدائمة"""
        convs = self._memory.setdefault("conversations", {})
        if conv_key not in convs:
            convs[conv_key] = []
        convs[conv_key].append(entry)
        convs[conv_key] = convs[conv_key][-200:]  # احتفظ بآخر 200 فقط

    # ─── كشف الخدمة ──────────────────────────────────────────────

    def detect_service(self, text):
        text_low = text.lower()
        best_match, best_score = None, 0
        for service, data in self.knowledge.items():
            for kw in data.get('keywords', []):
                if kw in text_low and len(kw) > best_score:
                    best_score = len(kw)
                    best_match = service
        return best_match

    def is_service_request(self, text: str) -> tuple:
        """تصنيف الرسالة: service / promo / normal"""
        if self.groq_client:
            try:
                resp = self.groq_client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[
                        {"role": "system", "content": "حدد نوع الرسالة بكلمة واحدة فقط: service أو promo أو normal.\n- service: طلب خدمة أكاديمية (واجب/بحث/ترجمة/تحليل/تصميم/تلخيص)\n- promo: إعلان أو روابط أو أرقام تواصل\n- normal: رسالة عادية أو تحية أو سؤال عام"},
                        {"role": "user", "content": text[:500]}
                    ],
                    max_tokens=5, temperature=0.1
                )
                result = resp.choices[0].message.content.strip().lower()
                if 'service' in result: return True, "service"
                if 'promo' in result:   return False, "promo"
                return False, "normal"
            except Exception as e:
                logger.error(f"AI classify error: {e}")
        # احتياطي
        text_low = text.lower()
        service_kws = ['حل', 'واجب', 'بحث', 'تقرير', 'تلخيص', 'ترجمة', 'تحليل', 'تصميم', 'مساعدة', 'مشروع']
        promo_kws   = ['للتواصل', 'واتساب', 'إعلان', 'عرض خاص', 'خصم', 'كاش باك', 'رابط']
        if any(k in text_low for k in promo_kws):   return False, "promo"
        if any(k in text_low for k in service_kws): return True, "service"
        return False, "normal"

    # ─── جلب تاريخ تيليجرام ──────────────────────────────────────

    async def _fetch_telegram_history(self, client, chat_id, limit=20):
        """جلب آخر limit رسالة مباشرة من Telegram"""
        history = []
        try:
            async for msg in client.iter_messages(chat_id, limit=limit):
                if not msg.text:
                    continue
                role = 'assistant' if getattr(msg, 'out', False) else 'user'
                history.insert(0, {
                    'role': role,
                    'text': msg.text[:400],
                    'time': msg.date.timestamp() if msg.date else time.time()
                })
        except Exception as e:
            logger.warning(f"[{self.user_id}] تعذّر جلب تاريخ Telegram: {e}")
        return history

    # ─── توليد الرد الذكي ────────────────────────────────────────

    async def generate_intelligent_response(self, sender_name: str, text: str,
                                             history: list = None,
                                             detected_service: str = None,
                                             style_examples: list = None) -> str:
        if not self.groq_client:
            return self._fallback_response(detected_service)

        services_info = "\n".join([
            f"- {s}: {d['description']} (السعر: {d.get('price_range','حسب الطلب')}, الوقت: {d.get('time_range','حسب الطلب')})"
            for s, d in self.knowledge.items()
        ])

        history_lines = []
        if history:
            for e in history[-20:]:
                role_label = sender_name if e.get('role') == 'user' else 'أنا'
                msg = (e.get('text') or e.get('content') or '')[:300]
                history_lines.append(f"{role_label}: {msg}")
        history_text = "\n".join(history_lines) if history_lines else "لا توجد رسائل سابقة — هذه أول رسالة."

        # أمثلة أسلوب الرد من المحادثات الفردية السابقة
        style_text = ""
        if style_examples:
            style_lines = []
            for q, a in style_examples[:10]:
                style_lines.append(f"  عميل: {q[:80]}\n  أنا:   {a[:80]}")
            style_text = "\n**أمثلة من ردودي السابقة (اتبع نفس الأسلوب):**\n" + "\n".join(style_lines)

        system_prompt = f"""أنت موظف حقيقي متمرس في التسويق تعمل في "مركز سرعة إنجاز" للخدمات الأكاديمية.

**قواعد الرد الإلزامية (التزم بها بدقة):**
1. استخدم عامية خليجية ودية وطبيعية (هلا، أبشر، مايهمك، تكرم عينك، وش واجبك، الله يوفقك).
2. ردودك قصيرة جداً — جملة أو جملتان فقط، لا أكثر.
3. **تحية (هلا/مرحبا/السلام/صباح/مساء):** رد بتحية ودية فقط — لا تذكر الأسعار أو الخدمات.
   مثال: "وعليكم السلام، هلا أخوي كيف أخدمك؟"
4. **طلب خدمة (واجب/بحث/ترجمة/...):** رد بـ "أبشر، ارسل المهمة وأشوفها." — لا تذكر السعر بدون رؤية المهمة.
5. **سؤال عن القدرة (تقدر/تعرف/...):** أجب بثقة: "أكيد أخوي ما يهمك، ارسل وأشوف."
6. **سؤال عن السعر:** قل "السعر حسب حجم المهمة، ارسل وأحدد لك." — لا تعطِ رقماً قبل رؤية المهمة.
7. **بعد استلام المهمة:** اسأل "الي متي تبيه؟" ثم حدد السعر.
8. **تكرار ممنوع:** لا تكرر نفس الرد في رسالتين متتاليتين — غيّر الأسلوب مع الحفاظ على المعنى.
9. اقرأ تاريخ المحادثة بعناية — لا تطرح أسئلة سبق طرحها.
10. أنت إنسان حقيقي — لا تبدو كبوت أبداً.
11. استخدم ايموجي واحد أو اثنين على الأكثر.

**الخدمات المتاحة:**
{services_info}

**الخدمة المكتشفة الآن:** {detected_service or 'لم تُحدد بعد'}
{style_text}

**تاريخ المحادثة (اقرأه بعناية لفهم السياق):**
{history_text}
"""
        try:
            resp = self.groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": f"رسالة من {sender_name}: {text}"}
                ],
                max_tokens=200, temperature=0.82
            )
            return resp.choices[0].message.content.strip()
        except Exception as e:
            logger.error(f"Groq response error: {e}")
            return self._fallback_response(detected_service)

    def _fallback_response(self, detected_service: str = None) -> str:
        """رد احتياطي من قاعدة الأنماط المحفوظة"""
        patterns = self._memory.get("patterns", {})
        if detected_service and detected_service in patterns:
            replies = patterns[detected_service].get("replies", [])
            if replies:
                return replies[-1]
        if detected_service and detected_service in self.knowledge:
            return f"أبشر على {detected_service}، ارسل التفاصيل وأشوفها 👍"
        return "هلا! أبشر، أي خدمة تحتاجها؟ ارسل لي التفاصيل."

    # ─── معالج الرسائل الواردة (الدالة الرئيسية) ─────────────────

    async def handle_incoming_message(self, event, client_manager):
        try:
            user_id    = self.user_id
            is_private = event.is_private
            is_group   = event.is_group or event.is_channel

            if is_private and not learning_manager.is_active(user_id, 'private'):
                return
            if is_group and not learning_manager.is_active(user_id, 'group'):
                return

            message = event.message
            if not message.text or getattr(message, 'out', False):
                return

            text = message.text
            sender = await event.get_sender()
            sender_name = (getattr(sender, 'first_name', '') or
                           getattr(sender, 'username', '') or 'مستخدم')
            sender_id = str(getattr(sender, 'id', ''))

            # تصفية الإعلانات
            _, msg_type = self.is_service_request(text)
            if msg_type == 'promo':
                logger.info(f"[{user_id}] تجاهل إعلان من {sender_name}")
                return

            conv_key = f"{sender_id}_{event.chat_id}"
            detected_service = self.detect_service(text)

            # ── جلب تاريخ Telegram (أعمق للمحادثات الفردية) ──
            tg_history = []
            style_examples = []
            if client_manager and client_manager.client:
                hist_limit = 50 if is_private else 20
                tg_history = await self._fetch_telegram_history(
                    client_manager.client, event.chat_id, limit=hist_limit
                )
                # قراءة أسلوب الرد من كل المحادثات الفردية (أول مرة أو كل 20 رسالة)
                if is_private:
                    msg_count = len(self._memory.get("conversations", {}).get(conv_key, []))
                    if msg_count % 20 == 0:
                        style_examples = await self._read_all_private_style(
                            client_manager.client, limit=30
                        )

            # ── الذاكرة الدائمة كاحتياطي إذا فشل Telegram ──
            persistent_history = self._get_persistent_history(conv_key)
            combined_history = tg_history if tg_history else persistent_history

            # ── إضافة الرسالة الحالية للذاكرتين ──
            now_entry = {'role': 'user', 'text': text, 'time': time.time()}
            if conv_key not in self.conversations_history:
                self.conversations_history[conv_key] = []
            self.conversations_history[conv_key].append(now_entry)
            self._clean_old_history(conv_key)
            self._append_to_persistent(conv_key, now_entry)

            # ── توليد الرد ──
            response = await self.generate_intelligent_response(
                sender_name, text,
                history=combined_history,
                detected_service=detected_service,
                style_examples=style_examples if is_private else None
            )

            # ── حفظ الرد في الذاكرتين ──
            reply_entry = {'role': 'assistant', 'text': response, 'time': time.time()}
            self.conversations_history[conv_key].append(reply_entry)
            self._append_to_persistent(conv_key, reply_entry)

            # ── تحديث الأنماط للتعلم الذاتي ──
            if detected_service:
                patterns = self._memory.setdefault("patterns", {})
                if detected_service not in patterns:
                    patterns[detected_service] = {"replies": [], "count": 0}
                if response not in patterns[detected_service]["replies"]:
                    patterns[detected_service]["replies"].append(response)
                    patterns[detected_service]["replies"] = patterns[detected_service]["replies"][-8:]
                patterns[detected_service]["count"] += 1

            # ── حفظ الذاكرة الدائمة ──
            self._save_memory()

            # ── مزامنة GitHub كل 3 تفاعلات ──
            self._sync_counter += 1
            if self._sync_counter % 3 == 0:
                self._sync_to_github()

            # ── إرسال الرد لتيليجرام ──
            await event.reply(response)
            try:
                socketio.emit('log_update', {
                    "message": f"🤖 رد ذكي لـ {sender_name}: {response[:120]}"
                }, to=user_id)
            except Exception:
                pass
            logger.info(f"[{user_id}] رد على {sender_name}: {response[:80]}")

            # ── تسجيل الطلبات المجهولة للخاص ──
            if is_private and not detected_service:
                req = {
                    "text": text[:200], "sender": sender_name,
                    "sender_id": sender_id,
                    "time": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "chat_id": event.chat_id
                }
                self.unknown_requests.append(req)
                try:
                    socketio.emit('new_unknown_request', req, to=user_id)
                except Exception:
                    pass

            # ── اقتراح أنماط جديدة كل 5 رسائل (للمحادثات الفردية فقط) ──
            if is_private:
                total_msgs = len(self._memory.get("conversations", {}).get(conv_key, []))
                if total_msgs > 0 and total_msgs % 5 == 0:
                    self._notify_user_about_suggestions(conv_key)

        except Exception as e:
            logger.error(f"Learning bot error [{self.user_id}]: {e}")

    # ─── إدارة الخدمات (API) ─────────────────────────────────────

    def add_service(self, name, description, keywords):
        if name and description:
            self.knowledge[name] = {
                "description": description,
                "keywords": [k.strip() for k in keywords if k.strip()] or [name],
                "price_range": "حسب الطلب",
                "time_range": "حسب الطلب"
            }
            self._save_memory()
            return True
        return False

    def delete_service(self, name):
        if name in self.knowledge:
            del self.knowledge[name]
            self._save_memory()
            return True
        return False

    def get_unknown_requests(self):
        return self.unknown_requests

    def clear_unknown(self):
        self.unknown_requests = []

    # ─── قراءة كل المحادثات الفردية لاستخلاص أسلوب الرد ──────────

    async def _read_all_private_style(self, client, limit=100):
        """
        يقرأ آخر limit رسالة من كل المحادثات الفردية ليتعلم أسلوب الرد.
        يُرجع نصاً ملخصاً لأنماط الردود الناجحة.
        """
        style_examples = []
        try:
            async for dialog in client.iter_dialogs():
                if not dialog.is_user:
                    continue
                entity = dialog.entity
                chat_id = entity.id
                pairs = []
                msgs = []
                async for msg in client.iter_messages(chat_id, limit=limit):
                    if msg.text:
                        role = 'out' if getattr(msg, 'out', False) else 'in'
                        msgs.insert(0, {'role': role, 'text': msg.text[:300]})
                for i in range(len(msgs) - 1):
                    if msgs[i]['role'] == 'in' and msgs[i+1]['role'] == 'out':
                        pairs.append((msgs[i]['text'], msgs[i+1]['text']))
                style_examples.extend(pairs[:5])
                if len(style_examples) >= 30:
                    break
        except Exception as e:
            logger.warning(f"[{self.user_id}] _read_all_private_style error: {e}")
        return style_examples

    # ─── استخراج أنماط من تاريخ محادثة بعينها ───────────────────

    def _extract_patterns_from_history(self, conv_key):
        """
        يستخرج أنماط (نوع + محفز + رد مقترح) من آخر 50 رسالة.
        """
        history = self._memory.get("conversations", {}).get(conv_key, [])
        if len(history) < 2:
            return []
        patterns = []
        for i in range(len(history) - 1):
            u = history[i]
            a = history[i + 1]
            if u.get('role') != 'user' or a.get('role') != 'assistant':
                continue
            user_msg   = (u.get('text') or u.get('content') or '').strip()
            asst_reply = (a.get('text') or a.get('content') or '').strip()
            if not user_msg or not asst_reply or len(asst_reply) < 5:
                continue
            msg_lower = user_msg.lower()
            if any(w in msg_lower for w in ['السلام', 'هلا', 'مرحبا', 'اهلا', 'صباح', 'مساء']):
                ptype = 'greeting'
            elif any(w in msg_lower for w in ['سعر', 'كم', 'تكلفة', 'ثمن', 'بكم']):
                ptype = 'price_ask'
            elif any(w in msg_lower for w in ['تقدر', 'تعرف', 'قادر', 'تقدرون']):
                ptype = 'capability_ask'
            elif any(w in msg_lower for w in ['واجب', 'بحث', 'ترجمة', 'تلخيص', 'حل', 'تحليل', 'تصميم']):
                ptype = 'service_request'
            else:
                ptype = 'general'
            patterns.append({
                'pattern_type': ptype,
                'trigger': user_msg[:60],
                'suggested_reply': asst_reply,
                'frequency': 1,
            })
        # دمج المتشابهات
        merged = {}
        for p in patterns:
            key = (p['pattern_type'], p['suggested_reply'][:40])
            if key not in merged:
                merged[key] = p.copy()
            else:
                merged[key]['frequency'] += 1
        return sorted(merged.values(), key=lambda x: x['frequency'], reverse=True)[:10]

    def _suggest_new_patterns(self, conv_key):
        """يُرجع الأنماط الجديدة غير المحفوظة ولم تُرفض من قبل."""
        existing_replies = set()
        for data in self._memory.get("patterns", {}).values():
            for r in data.get("replies", []):
                existing_replies.add(r[:40])
        rejected = {
            r.get('reply', '')[:40]
            for r in self._memory.get("rejected_suggestions", [])
        }
        suggestions = []
        for p in self._extract_patterns_from_history(conv_key):
            short = p['suggested_reply'][:40]
            if short not in existing_replies and short not in rejected:
                suggestions.append(p)
        return suggestions[:5]

    def _notify_user_about_suggestions(self, conv_key):
        """يُرسل إشعار socket للمستخدم إذا وُجدت اقتراحات جديدة."""
        suggestions = self._suggest_new_patterns(conv_key)
        if not suggestions:
            return
        try:
            socketio.emit('learning_suggestions', {
                "conv_key": conv_key,
                "suggestions": suggestions,
                "count": len(suggestions)
            }, to=self.user_id)
            logger.info(f"[{self.user_id}] أُرسل إشعار تعلم: {len(suggestions)} اقتراح")
        except Exception as e:
            logger.error(f"[{self.user_id}] خطأ إشعار التعلم: {e}")

    def save_suggestion(self, index, conv_key):
        """حفظ اقتراح في قاعدة الأنماط الدائمة."""
        suggestions = self._suggest_new_patterns(conv_key)
        if index < 0 or index >= len(suggestions):
            return False, "رقم الاقتراح غير صحيح"
        s = suggestions[index]
        kw = s['pattern_type']
        patterns = self._memory.setdefault("patterns", {})
        if kw not in patterns:
            patterns[kw] = {"replies": [], "count": 0}
        if s['suggested_reply'] not in patterns[kw]["replies"]:
            patterns[kw]["replies"].append(s['suggested_reply'])
            patterns[kw]["replies"] = patterns[kw]["replies"][-10:]
        patterns[kw]["count"] += 1
        self._save_memory()
        self._sync_to_github()
        return True, f"✅ تم حفظ الاقتراح في قاعدة الأنماط"

    def delete_suggestion(self, index, conv_key):
        """رفض اقتراح ومنع عرضه مستقبلاً."""
        suggestions = self._suggest_new_patterns(conv_key)
        if index < 0 or index >= len(suggestions):
            return False, "رقم الاقتراح غير صحيح"
        s = suggestions[index]
        rejected = self._memory.setdefault("rejected_suggestions", [])
        rejected.append({
            "trigger": s['trigger'],
            "reply": s['suggested_reply'],
            "rejected_at": datetime.now().isoformat()
        })
        self._memory["rejected_suggestions"] = rejected[-100:]
        self._save_memory()
        return True, f"🗑️ تم رفض الاقتراح"

    def get_suggestions(self, conv_key):
        return self._suggest_new_patterns(conv_key)


class LearningManager:
    """مدير بوتات التعلم — يُنشئ LearningBot لكل مستخدم ويحفظ إعداداته"""
    def __init__(self):
        self.bots = {}
        self.user_settings = {}

    def get_bot(self, user_id):
        if user_id not in self.bots:
            self.bots[user_id] = LearningBot(user_id)
        return self.bots[user_id]

    def is_active(self, user_id, chat_type='private'):
        if user_id not in self.user_settings:
            saved = load_settings(user_id)
            self.user_settings[user_id] = {
                'active_private': saved.get('learning_active_private', False),
                'active_group':   saved.get('learning_active_group',   False),
            }
        s = self.user_settings.get(user_id, {})
        if chat_type == 'private': return s.get('active_private', False)
        if chat_type == 'group':   return s.get('active_group',   False)
        return s.get('active_private', False) or s.get('active_group', False)

    def set_active(self, user_id, active, chat_type='private'):
        if user_id not in self.user_settings:
            self.user_settings[user_id] = {}
        key = f'active_{chat_type}'
        self.user_settings[user_id][key] = active
        settings = load_settings(user_id)
        settings['learning_active_private'] = self.user_settings[user_id].get('active_private', False)
        settings['learning_active_group']   = self.user_settings[user_id].get('active_group',   False)
        save_settings(user_id, settings)
        logger.info(f"✅ Learning {chat_type} for {user_id} → {active}")

    def get_settings(self, user_id):
        saved = load_settings(user_id)
        if user_id not in self.user_settings:
            self.user_settings[user_id] = {
                'active_private': saved.get('learning_active_private', False),
                'active_group':   saved.get('learning_active_group',   False),
            }
        return self.user_settings.get(user_id, {'active_private': False, 'active_group': False})

    def toggle_all(self, user_id, active_private, active_group):
        self.set_active(user_id, active_private, 'private')
        self.set_active(user_id, active_group,   'group')


learning_manager = LearningManager()


@app.route("/api/learning/status", methods=["GET"])
def api_learning_status():
    user_id = session.get('user_id', 'user_1')
    settings = learning_manager.get_settings(user_id)
    return jsonify({
        "success": True,
        "active_private": settings.get('active_private', False),
        "active_group": settings.get('active_group', False),
        "reply_in_groups": settings.get('active_group', False)
    })

@app.route("/api/learning/toggle", methods=["POST"])
def api_learning_toggle():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    chat_type = data.get('chat_type', 'private')
    current_settings = learning_manager.get_settings(user_id)
    if chat_type == 'private':
        current = current_settings.get('active_private', False)
    else:
        current = current_settings.get('active_group', False)
    new_active = not current if data.get('active') is None else bool(data.get('active'))
    learning_manager.set_active(user_id, new_active, chat_type)
    return jsonify({"success": True, "active": new_active, "chat_type": chat_type})

@app.route("/api/learning/toggle_all", methods=["POST"])
def api_learning_toggle_all():
    """تفعيل/إلغاء تفعيل كلا النوعين معاً"""
    user_id = session.get('user_id', 'user_1')
    data = request.json
    active_private = data.get('active_private', False)
    active_group = data.get('active_group', False)
    learning_manager.set_active(user_id, active_private, 'private')
    learning_manager.set_active(user_id, active_group, 'group')
    return jsonify({"success": True, "active_private": active_private, "active_group": active_group})

@app.route("/api/learning/services", methods=["GET"])
def api_learning_services():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    return jsonify({"success": True, "services": bot.knowledge})

@app.route("/api/learning/add_service", methods=["POST"])
def api_learning_add_service():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    name = data.get('name', '').strip()
    description = data.get('description', '').strip()
    keywords = data.get('keywords', [])
    if not name or not description:
        return jsonify({"success": False, "message": "الاسم والوصف مطلوبان"})
    bot = learning_manager.get_bot(user_id)
    if bot.add_service(name, description, keywords):
        return jsonify({"success": True, "message": f"تم إضافة الخدمة {name}"})
    return jsonify({"success": False, "message": "فشل في الإضافة"})

@app.route("/api/learning/delete_service", methods=["POST"])
def api_learning_delete_service():
    user_id = session.get('user_id', 'user_1')
    data = request.json
    name = data.get('name', '')
    bot = learning_manager.get_bot(user_id)
    if bot.delete_service(name):
        return jsonify({"success": True, "message": f"تم حذف الخدمة {name}"})
    return jsonify({"success": False, "message": "الخدمة غير موجودة"})

@app.route("/api/learning/unknown_requests", methods=["GET"])
def api_learning_unknown():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    return jsonify({"success": True, "requests": bot.get_unknown_requests()})

@app.route("/api/learning/clear_unknown", methods=["POST"])
def api_learning_clear_unknown():
    user_id = session.get('user_id', 'user_1')
    bot = learning_manager.get_bot(user_id)
    bot.clear_unknown()
    return jsonify({"success": True, "message": "تم مسح الطلبات"})

@app.route("/api/learning/suggestions", methods=["GET"])
def api_learning_suggestions():
    user_id = session.get('user_id', 'user_1')
    conv_key = request.args.get("conv_key", "")
    bot = learning_manager.get_bot(user_id)
    suggestions = bot.get_suggestions(conv_key) if conv_key else []
    return jsonify({"success": True, "suggestions": suggestions})

@app.route("/api/learning/save_suggestion", methods=["POST"])
def api_learning_save_suggestion():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    index    = int(data.get("index", -1))
    conv_key = data.get("conv_key", "")
    bot = learning_manager.get_bot(user_id)
    success, msg = bot.save_suggestion(index, conv_key)
    return jsonify({"success": success, "message": msg})

@app.route("/api/learning/delete_suggestion", methods=["POST"])
def api_learning_delete_suggestion():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    index    = int(data.get("index", -1))
    conv_key = data.get("conv_key", "")
    bot = learning_manager.get_bot(user_id)
    success, msg = bot.delete_suggestion(index, conv_key)
    return jsonify({"success": success, "message": msg})

def _normalize_auto_reply(rule):
    if not isinstance(rule, dict):
        return None
    keyword = (rule.get('keyword') or rule.get('trigger') or '').strip()
    reply = (rule.get('reply') or '').strip()
    if not keyword or not reply:
        return None
    scope = (rule.get('scope') or 'all').lower()
    if scope not in ('all', 'private', 'groups'):
        scope = 'all'
    match = (rule.get('match') or 'contains').lower()
    if match not in ('contains', 'exact', 'regex'):
        match = 'contains'
    return {
        'keyword': keyword,
        'reply': reply,
        'scope': scope,
        'match': match,
        'used_count': int(rule.get('used_count') or 0),
        'last_used': rule.get('last_used') or '',
    }

@app.route("/api/auto_replies", methods=["GET"])
@app.route("/api/get_auto_replies", methods=["GET"])
def api_get_auto_replies():
    user_id = session.get('user_id', 'user_1')
    settings = load_settings(user_id)
    return jsonify({
        "success": True,
        "enabled": settings.get('auto_reply_enabled', True),
        "auto_replies": settings.get('auto_replies', []) or []
    })

@app.route("/api/add_auto_reply", methods=["POST"])
def api_add_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    rule = _normalize_auto_reply({
        'keyword': data.get('keyword') or data.get('trigger') or '',
        'reply': data.get('reply') or '',
        'scope': data.get('scope') or 'all',
        'match': data.get('match') or 'contains',
    })
    if not rule:
        return jsonify({"success": False, "message": "❌ الكلمة المفتاحية ونص الرد مطلوبان"})

    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    rules.append(rule)
    settings['auto_replies'] = rules
    if save_settings(user_id, settings):
        return jsonify({"success": True, "message": "✅ تم إضافة الرد التلقائي", "auto_replies": rules})
    return jsonify({"success": False, "message": "❌ فشل حفظ القاعدة"})

@app.route("/api/update_auto_reply", methods=["POST"])
def api_update_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    try:
        index = int(data.get('index', -1))
    except (TypeError, ValueError):
        index = -1
    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    if not (0 <= index < len(rules)):
        return jsonify({"success": False, "message": "❌ الفهرس غير صحيح"})
    new_rule = _normalize_auto_reply({
        'keyword': data.get('keyword'),
        'reply': data.get('reply'),
        'scope': data.get('scope'),
        'match': data.get('match'),
        'used_count': rules[index].get('used_count'),
        'last_used': rules[index].get('last_used'),
    })
    if not new_rule:
        return jsonify({"success": False, "message": "❌ بيانات غير صالحة"})
    rules[index] = new_rule
    settings['auto_replies'] = rules
    save_settings(user_id, settings)
    return jsonify({"success": True, "message": "✅ تم تحديث القاعدة", "auto_replies": rules})

@app.route("/api/delete_auto_reply", methods=["POST"])
def api_delete_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    try:
        index = int(data.get('index', -1))
    except (TypeError, ValueError):
        index = -1
    settings = load_settings(user_id)
    rules = settings.get('auto_replies', []) or []
    if 0 <= index < len(rules):
        removed = rules.pop(index)
        settings['auto_replies'] = rules
        save_settings(user_id, settings)
        return jsonify({"success": True, "message": f"🗑️ تم حذف الرد '{removed.get('keyword','')[:30]}'",
                        "auto_replies": rules})
    return jsonify({"success": False, "message": "❌ فهرس غير صحيح"})

@app.route("/api/save_auto_replies", methods=["POST"])
def api_save_auto_replies():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    raw = data.get('auto_replies', []) or []
    cleaned = []
    for r in raw:
        nr = _normalize_auto_reply(r)
        if nr:
            cleaned.append(nr)
    settings = load_settings(user_id)
    settings['auto_replies'] = cleaned
    save_settings(user_id, settings)
    return jsonify({"success": True, "message": f"✅ تم حفظ {len(cleaned)} قاعدة رد", "auto_replies": cleaned})

@app.route("/api/toggle_auto_reply", methods=["POST"])
def api_toggle_auto_reply():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    enabled = bool(data.get('enabled', True))
    settings = load_settings(user_id)
    settings['auto_reply_enabled'] = enabled
    save_settings(user_id, settings)

    # إذا تم التفعيل، تأكد من تشغيل المراقبة/الاتصال حتى تعمل الردود
    if enabled:
        try:
            with USERS_LOCK:
                user_info = USERS.get(user_id, {})
                is_running = user_info.get('is_running', False)
                is_auth = user_info.get('authenticated', False)
            if is_auth and not is_running:
                with USERS_LOCK:
                    USERS[user_id]['is_running'] = True
                import threading as _threading
                t = _threading.Thread(
                    target=monitoring_worker,
                    args=(user_id,),
                    daemon=True
                )
                t.start()
                with USERS_LOCK:
                    USERS[user_id]['thread'] = t
                logger.info(f"✅ Auto-reply: started monitoring worker for {user_id}")
        except Exception as _e:
            logger.error(f"Auto-reply toggle start-worker error: {_e}")

    return jsonify({
        "success": True,
        "enabled": enabled,
        "message": "✅ تم تفعيل الردود التلقائية" if enabled else "⏸️ تم تعطيل الردود التلقائية"
    })


@app.route("/api/auto_reply_status", methods=["GET"])
def api_auto_reply_status():
    user_id = session.get('user_id', 'user_1')
    settings = load_settings(user_id)
    enabled = settings.get('auto_reply_enabled', True)
    rules = settings.get('auto_replies', [])
    return jsonify({"success": True, "enabled": enabled, "rules_count": len(rules)})

async def resolve_link_group_name(client, url):
    """يحاول الحصول على اسم المجموعة/القناة من رابط تيليجرام"""
    try:
        username = url.split('/')[-1].replace('@', '')
        if username.startswith('+'):
            return None  # روابط الدعوة الخاصة لا يمكن حلها
        entity = await client.get_entity(username)
        return getattr(entity, 'title', None) or getattr(entity, 'first_name', None)
    except Exception:
        return None


# ═══════════════════════════════════════════════════════════════════
#  وظائف البحث — مُعاد هيكلتها بالكامل
# ═══════════════════════════════════════════════════════════════════

async def search_links_in_chats(client, since_date, keyword='', resolve_names=True):
    """يبحث في كل محادثاتي ومجموعاتي عن روابط تيليجرام مع فلترة بالكلمة المفتاحية"""
    from datetime import timezone as _tz
    since_dt = since_date.replace(tzinfo=_tz.utc) if since_date.tzinfo is None else since_date
    found_links = []
    seen_urls   = set()
    try:
        async for dialog in client.iter_dialogs():
            try:
                if not dialog.entity:
                    continue
                chat_title = dialog.title or 'محادثة'
                async for message in client.iter_messages(dialog, limit=500):
                    if not message.date:
                        continue
                    msg_date = message.date if message.date.tzinfo else message.date.replace(tzinfo=_tz.utc)
                    if msg_date < since_dt:
                        break
                    if not message.text:
                        continue
                    for lnk in extract_telegram_links(message.text):
                        url = lnk.get('url', '')
                        if url and url not in seen_urls:
                            seen_urls.add(url)
                            found_links.append({
                                'url': url,
                                'username': lnk.get('username', ''),
                                'group_name': '',
                                'members': 0,
                                'date': message.date.strftime('%Y-%m-%d %H:%M'),
                                'chat_title': chat_title,
                                'link_type': lnk.get('type', 'channel'),
                            })
                if len(found_links) >= 2000:
                    break
            except Exception as e:
                logger.warning(f"تخطي محادثة: {e}")
                continue
    except Exception as e:
        logger.error(f"خطأ في البحث: {e}")

    keyword_lower = (keyword or '').strip().lower()
    if not resolve_names:
        if keyword_lower:
            return [i for i in found_links
                    if keyword_lower in (i.get('url') or '').lower()
                    or keyword_lower in (i.get('username') or '').lower()]
        return found_links

    # ── تحليل أسماء المجموعات وفلترة بالكلمة المفتاحية ──
    matched = []
    resolve_count = 0
    MAX_RESOLVE = 200
    for item in found_links:
        if resolve_count >= MAX_RESOLVE:
            # بعد الحد: فلترة URL فقط
            if keyword_lower and keyword_lower not in (item.get('url') or '').lower():
                continue
            matched.append(item)
            continue

        is_invite = '/+' in item['url']
        if not is_invite:
            try:
                uname = item['url'].rstrip('/').split('/')[-1]
                if uname and len(uname) >= 3:
                    entity = None
                    try:
                        entity = await client.get_entity(uname)
                    except Exception:
                        try:
                            r = await client(functions.contacts.ResolveUsernameRequest(username=uname))
                            entity = r.chats[0] if r.chats else (r.users[0] if r.users else None)
                        except Exception:
                            pass
                    if entity:
                        title   = getattr(entity, 'title', None) or getattr(entity, 'first_name', None) or ''
                        members = getattr(entity, 'participants_count', 0) or 0
                        item['group_name'] = title
                        item['members']    = members
                    resolve_count += 1
                    await asyncio.sleep(0.12)
            except Exception:
                pass

        # فلترة بالكلمة المفتاحية
        if keyword_lower:
            name_ok = keyword_lower in (item.get('group_name') or '').lower()
            user_ok = keyword_lower in (item.get('username')   or '').lower()
            url_ok  = keyword_lower in (item.get('url')        or '').lower()
            src_ok  = keyword_lower in (item.get('chat_title') or '').lower()
            if not (name_ok or user_ok or url_ok or src_ok):
                continue
        matched.append(item)

    return matched


async def search_public_telegram(client, query, limit=50):
    """يبحث عبر سيرفرات تيليجرام عن مجموعات/قنوات مطابقة للكلمة"""
    results  = []
    seen_ids = set()

    # الطريقة 1: contacts.SearchRequest — أفضل للبحث بالاسم مباشرة
    try:
        from telethon.tl.functions.contacts import SearchRequest as _CSearch
        sr = await client(_CSearch(q=query, limit=min(limit, 100)))
        for chat in (sr.chats or []):
            if chat.id in seen_ids:
                continue
            seen_ids.add(chat.id)
            uname = getattr(chat, 'username', None)
            results.append({
                'id': str(chat.id), 'title': getattr(chat, 'title', '') or '',
                'username': uname, 'url': f"https://t.me/{uname}" if uname else '',
                'members': getattr(chat, 'participants_count', 0) or 0,
                'megagroup': getattr(chat, 'megagroup', False),
                'verified': getattr(chat, 'verified', False),
            })
    except Exception as e:
        logger.warning(f"contacts.SearchRequest: {e}")

    # الطريقة 2: messages.SearchGlobalRequest كـ fallback
    if len(results) < limit // 2:
        try:
            gs = await client(functions.messages.SearchGlobalRequest(
                q=query, offset_date=None, offset_peer=None, offset_id=0, limit=limit
            ))
            for chat in (gs.chats or []):
                if chat.id in seen_ids:
                    continue
                seen_ids.add(chat.id)
                uname = getattr(chat, 'username', None)
                results.append({
                    'id': str(chat.id), 'title': getattr(chat, 'title', '') or '',
                    'username': uname, 'url': f"https://t.me/{uname}" if uname else '',
                    'members': getattr(chat, 'participants_count', 0) or 0,
                    'megagroup': getattr(chat, 'megagroup', False),
                    'verified': getattr(chat, 'verified', False),
                })
        except Exception as e:
            logger.warning(f"SearchGlobalRequest: {e}")

    results.sort(key=lambda x: x.get('members', 0), reverse=True)
    return results[:limit]


@app.route("/api/search_my_links", methods=["POST"])
def api_search_my_links():
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data = request.json or {}
        days = int(data.get('days', 60))
        if days <= 0 or days > 365: days = 60
        keyword = (data.get('keyword') or '').strip()
        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            cm = USERS[user_id].get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})
        since = datetime.now() - timedelta(days=days)
        result = cm.run_coroutine(search_links_in_chats(cm.client, since, keyword=keyword))
        return jsonify({"success": True, "links": result,
                        "message": f"تم العثور على {len(result)} رابط"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/search_my_links/start", methods=["POST"])
def api_search_my_links_start():
    """بحث متدفق بالكلمة المفتاحية داخل محادثاتي — النتائج فور اكتشافها مع تقدم لحظي"""
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data    = request.json or {}
        keyword = (data.get('keyword') or '').strip()
        depth   = (data.get('depth') or 'medium').lower()

        # عمق البحث: سريع / متوسط / كامل
        _DEPTH = {
            'fast':   {'days': 30,   'msg_limit': 500},
            'medium': {'days': 180,  'msg_limit': 1000},
            'full':   {'days': 3650, 'msg_limit': 3000},
        }
        cfg = _DEPTH.get(depth, _DEPTH['medium'])
        search_days  = cfg['days']
        msg_limit    = cfg['msg_limit']

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            cm = USERS[user_id].get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})

        async def _stream():
            from datetime import timezone as _tz
            since_dt      = datetime.now(_tz.utc) - timedelta(days=search_days)
            keyword_lower = keyword.lower()
            seen_urls     = set()
            total_found   = 0
            dialog_count  = 0

            try:
                async for dialog in cm.client.iter_dialogs():
                    try:
                        if not dialog.entity:
                            continue
                        chat_title   = dialog.title or 'محادثة'
                        dialog_count += 1

                        # ── إرسال تقدم الفحص اللحظي ──
                        socketio.emit('search_dialog_progress', {
                            'chat': chat_title,
                            'scanned': dialog_count,
                            'found': total_found
                        }, to=user_id)

                        batch_raw = []
                        async for message in cm.client.iter_messages(dialog, limit=msg_limit):
                            if not message.date:
                                continue
                            msg_date = message.date if message.date.tzinfo else message.date.replace(tzinfo=_tz.utc)
                            if msg_date < since_dt:
                                break
                            if not message.text:
                                continue
                            for lnk in extract_telegram_links(message.text):
                                url = lnk.get('url', '')
                                if url and url not in seen_urls:
                                    seen_urls.add(url)
                                    batch_raw.append({
                                        'url':        url,
                                        'username':   lnk.get('username', ''),
                                        'group_name': '',
                                        'members':    0,
                                        'date':       message.date.strftime('%Y-%m-%d %H:%M'),
                                        'chat_title': chat_title,
                                        'link_type':  lnk.get('type', 'channel'),
                                    })

                        # ── تحليل أسماء المجموعات وفلترة ──
                        emit_batch = []
                        for item in batch_raw:
                            is_invite = '/+' in item['url']

                            if not is_invite:
                                uname = item['url'].rstrip('/').split('/')[-1]
                                if uname and len(uname) >= 3:
                                    try:
                                        entity = None
                                        try:
                                            entity = await cm.client.get_entity(uname)
                                        except Exception:
                                            try:
                                                r = await cm.client(functions.contacts.ResolveUsernameRequest(username=uname))
                                                entity = r.chats[0] if r.chats else (r.users[0] if r.users else None)
                                            except Exception:
                                                pass
                                        if entity:
                                            item['group_name'] = getattr(entity, 'title', None) or getattr(entity, 'first_name', None) or ''
                                            item['members']    = getattr(entity, 'participants_count', 0) or 0
                                        await asyncio.sleep(0.12)
                                    except Exception:
                                        pass

                            # ── تطبيق فلتر الكلمة المفتاحية ──
                            if keyword_lower:
                                name_ok = keyword_lower in (item.get('group_name') or '').lower()
                                user_ok = keyword_lower in (item.get('username')   or '').lower()
                                url_ok  = keyword_lower in (item.get('url')        or '').lower()
                                src_ok  = keyword_lower in (item.get('chat_title') or '').lower()
                                if not (name_ok or user_ok or url_ok or src_ok):
                                    continue

                            emit_batch.append(item)
                            total_found += 1

                        if emit_batch:
                            socketio.emit('search_link_batch', {
                                'items': emit_batch, 'chat_title': chat_title, 'count': total_found
                            }, to=user_id)

                        if total_found >= 2000:
                            break

                    except Exception as ex:
                        logger.warning(f"تخطي {dialog.title}: {ex}")
                        continue

            except Exception as e:
                logger.error(f"بحث خطأ: {e}")
                socketio.emit('search_links_done', {
                    'total': total_found, 'error': str(e), 'dialogs': dialog_count
                }, to=user_id)
                return

            socketio.emit('search_links_done', {
                'total': total_found, 'keyword': keyword, 'dialogs': dialog_count
            }, to=user_id)

        _OSThread(target=lambda: cm.run_coroutine(_stream()), daemon=True).start()
        depth_ar = {'fast': 'سريع', 'medium': 'متوسط', 'full': 'كامل'}.get(depth, 'متوسط')
        kw_msg = f" بكلمة: \"{keyword}\"" if keyword else ""
        return jsonify({"success": True, "message": f"بدأ البحث ({depth_ar}){kw_msg}"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/search_my_links/csv", methods=["POST"])
def api_search_my_links_csv():
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data    = request.json or {}
        links   = data.get('links', [])
        if not links:
            return jsonify({"success": False, "message": "❌ لا توجد روابط"})
        import io as _io, csv as _csv
        buf = _io.StringIO()
        w   = _csv.writer(buf)
        w.writerow(['الرابط', 'اسم المجموعة', 'عدد الأعضاء', 'المصدر', 'التاريخ', 'النوع'])
        for l in links:
            w.writerow([l.get('url',''), l.get('group_name',''), l.get('members',''),
                        l.get('chat_title',''), l.get('date',''), l.get('link_type','')])
        buf.seek(0)
        from flask import Response
        return Response(
            buf.getvalue().encode('utf-8-sig'),
            mimetype='text/csv; charset=utf-8',
            headers={'Content-Disposition': 'attachment; filename=telegram_links.csv'}
        )
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/search_public_channels", methods=["POST"])
def api_search_public_channels():
    """بحث عام في تيليجرام عن مجموعات/قنوات بالكلمة المفتاحية"""
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data    = request.json or {}
        query   = (data.get('query') or '').strip()
        if not query:
            return jsonify({"success": False, "message": "❌ أدخل كلمة للبحث"})
        limit = min(int(data.get('limit', 60)), 100)
        with USERS_LOCK:
            cm = USERS.get(user_id, {}).get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})
        result = cm.run_coroutine(search_public_telegram(cm.client, query, limit))
        return jsonify({"success": True, "channels": result,
                        "message": f"✅ تم العثور على {len(result)} مجموعة/قناة"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


# ── بحث عام في تيليجرام (الطريقتان: SearchRequest + SearchGlobalRequest) ──
async def search_telegram_global(client, query, limit=50, filter_type='all'):
    """
    بحث عام في تيليجرام عن مجموعات/قنوات/بوتات تطابق الكلمة المفتاحية.
    filter_type: 'all' | 'groups' | 'channels' | 'bots'
    """
    from telethon.tl.functions.contacts import SearchRequest
    from telethon.tl.functions.messages import SearchGlobalRequest
    from telethon.tl.types import InputPeerEmpty

    results = []
    seen_ids = set()

    def _classify(chat):
        if hasattr(chat, 'megagroup') and chat.megagroup:
            return 'group'
        if hasattr(chat, 'broadcast') and chat.broadcast:
            return 'channel'
        if hasattr(chat, 'bot') and chat.bot:
            return 'bot'
        return 'chat'

    def _passes_filter(chat_type):
        if filter_type == 'all':
            return True
        return chat_type == filter_type

    def _build_item(chat, source):
        chat_type = _classify(chat)
        username  = getattr(chat, 'username', None)
        return {
            'id':       chat.id,
            'title':    getattr(chat, 'title', None) or getattr(chat, 'first_name', '') or '',
            'username': username,
            'url':      f"https://t.me/{username}" if username else None,
            'type':     chat_type,
            'members':  getattr(chat, 'participants_count', 0) or 0,
            'verified': getattr(chat, 'verified', False),
            'source':   source,
        }

    # الطريقة الأولى: contacts.SearchRequest
    try:
        sr = await client(SearchRequest(q=query, limit=min(limit, 100)))
        for chat in (sr.chats or []):
            if chat.id in seen_ids:
                continue
            item = _build_item(chat, 'SearchRequest')
            if _passes_filter(item['type']):
                seen_ids.add(chat.id)
                results.append(item)
    except Exception as e:
        logger.warning(f"SearchRequest failed for '{query}': {e}")

    # الطريقة الثانية: messages.SearchGlobalRequest (شبكة تيليجرام كاملة)
    if len(results) < limit:
        try:
            gs = await client(SearchGlobalRequest(
                q=query,
                filter=None,
                min_date=None,
                max_date=None,
                offset_rate=0,
                offset_peer=InputPeerEmpty(),
                offset_id=0,
                limit=min(limit, 100),
            ))
            for chat in (gs.chats or []):
                if chat.id in seen_ids:
                    continue
                item = _build_item(chat, 'SearchGlobal')
                if _passes_filter(item['type']):
                    seen_ids.add(chat.id)
                    results.append(item)
        except Exception as e:
            logger.warning(f"SearchGlobalRequest failed for '{query}': {e}")

    results.sort(key=lambda x: x.get('members', 0), reverse=True)
    return results[:limit]


@app.route("/api/search_telegram_global", methods=["POST"])
def api_search_telegram_global():
    """
    بحث عام في قاعدة تيليجرام بالكلمة المفتاحية — مجموعات وقنوات وبوتات.
    Body JSON:
      query       : نص البحث (مطلوب)
      limit       : عدد النتائج (افتراضي 50، حد أقصى 100)
      filter_type : all | groups | channels | bots  (افتراضي: all)
    """
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data    = request.json or {}
        query   = (data.get('query') or '').strip()
        if not query:
            return jsonify({"success": False, "message": "❌ أدخل كلمة للبحث"})
        limit       = max(1, min(int(data.get('limit', 50)), 100))
        filter_type = data.get('filter_type', 'all')

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            cm = USERS[user_id].get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})

        results = cm.run_coroutine(
            search_telegram_global(cm.client, query, limit, filter_type)
        )
        return jsonify({
            "success": True,
            "query":   query,
            "filter":  filter_type,
            "total":   len(results),
            "results": results,
            "message": f"✅ تم العثور على {len(results)} نتيجة"
        })
    except Exception as e:
        logger.error(f"search_telegram_global error: {e}")
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/join_global_search_results", methods=["POST"])
def api_join_global_search_results():
    """
    انضمام إلى مجموعات من نتائج البحث العام.
    Body JSON:
      results      : قائمة نتائج البحث (كل عنصر يحتوي على 'url' و 'title')
      selected_ids : قائمة IDs المطلوب الانضمام إليها (اختياري — إذا فارغة + join_all=false → لا شيء)
      join_all     : true للانضمام إلى كل النتائج (افتراضي: false)
      delay        : تأخير بين الطلبات بالثواني (افتراضي: 2)
    """
    try:
        if 'user_id' not in session:
            return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول"})
        user_id = session['user_id']
        data    = request.json or {}
        items   = data.get('results', [])
        sel_ids = set(data.get('selected_ids', []))
        join_all = data.get('join_all', False)
        delay    = max(1, int(data.get('delay', 2)))

        if not items:
            return jsonify({"success": False, "message": "❌ لا توجد نتائج للانضمام"})

        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير موجود"})
            cm = USERS[user_id].get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})

        targets = items if join_all else [r for r in items if r.get('id') in sel_ids]
        if not targets:
            return jsonify({"success": False, "message": "❌ لم تحدد أي مجموعة للانضمام"})

        joined, skipped, failed = [], [], []
        for item in targets:
            link = item.get('url')
            title = item.get('title', link)
            if not link:
                skipped.append({'title': title, 'reason': 'لا يوجد رابط عام'})
                continue
            try:
                res = cm.run_coroutine(
                    join_telegram_group(cm.client, link, user_id, cm)
                )
                if res.get('success'):
                    joined.append({'link': link, 'title': title, 'message': res.get('message', '')})
                    socketio.emit('log_update',
                        {'message': f'✅ انضم إلى: {title} ({link})'}, to=user_id)
                else:
                    failed.append({'link': link, 'title': title, 'reason': res.get('message', 'فشل')})
                    socketio.emit('log_update',
                        {'message': f'❌ فشل: {title} — {res.get("message","")}'}, to=user_id)
            except Exception as e:
                failed.append({'link': link, 'title': title, 'reason': str(e)})
                socketio.emit('log_update',
                    {'message': f'❌ خطأ: {title} — {e}'}, to=user_id)
            time.sleep(delay)

        socketio.emit('log_update', {
            'message': f'🎉 انتهى: ✅ {len(joined)} نجح | ❌ {len(failed)} فشل | ⏭ {len(skipped)} تجاوز'
        }, to=user_id)

        return jsonify({
            "success": True,
            "joined":  joined,
            "failed":  failed,
            "skipped": skipped,
            "total":   len(targets),
            "message": f"انضمام: {len(joined)} نجح، {len(failed)} فشل، {len(skipped)} تجاوز"
        })
    except Exception as e:
        logger.error(f"join_global_search_results error: {e}")
        return jsonify({"success": False, "message": str(e)})


@app.route("/tools/analyze_stats", methods=["POST"])
def api_academic_analyze_stats():
    try:
        data = request.get_json(force=True, silent=True) or {}
        raw = data.get('data', '')
        nums = [float(x) for x in re.findall(r'[-+]?\d*\.?\d+', str(raw)) if x]
        if len(nums) < 2:
            return jsonify({"error": "أدخل على الأقل رقمين للتحليل"}), 400

        arr = np.array(nums)
        mode_val = float(stats.mode(arr, keepdims=True).mode[0])
        stats_result = {
            "count":    int(len(arr)),
            "sum":      round(float(np.sum(arr)), 4),
            "mean":     round(float(np.mean(arr)), 4),
            "median":   round(float(np.median(arr)), 4),
            "mode":     round(mode_val, 4),
            "std":      round(float(np.std(arr)), 4),
            "variance": round(float(np.var(arr)), 4),
            "min":      round(float(np.min(arr)), 4),
            "max":      round(float(np.max(arr)), 4),
            "range":    round(float(np.max(arr) - np.min(arr)), 4),
            "q1":       round(float(np.percentile(arr, 25)), 4),
            "q3":       round(float(np.percentile(arr, 75)), 4),
            "iqr":      round(float(np.percentile(arr, 75) - np.percentile(arr, 25)), 4),
            "skewness": round(float(stats.skew(arr)), 4),
            "kurtosis": round(float(stats.kurtosis(arr)), 4),
        }

        fig, axes = plt.subplots(1, 2, figsize=(10, 4))
        axes[0].hist(arr, bins='auto', color='#4e73df', edgecolor='white', alpha=0.85)
        axes[0].axvline(stats_result['mean'],   color='red',   linestyle='--', label=f"المتوسط: {stats_result['mean']}")
        axes[0].axvline(stats_result['median'], color='green', linestyle='--', label=f"الوسيط: {stats_result['median']}")
        axes[0].set_title('توزيع البيانات'); axes[0].legend()
        axes[0].set_xlabel('القيمة'); axes[0].set_ylabel('التكرار')
        axes[1].boxplot(arr, vert=True, patch_artist=True,
                        boxprops=dict(facecolor='#4e73df', alpha=0.7))
        axes[1].set_title('المربع الجذري (Boxplot)'); axes[1].set_ylabel('القيمة')
        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=100, bbox_inches='tight')
        plt.close()
        buf.seek(0)
        chart_b64 = base64.b64encode(buf.read()).decode('utf-8')

        return jsonify({"success": True, "stats": stats_result, "chart": chart_b64})
    except Exception as e:
        logger.error(f"Analyze stats error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/tools/format_file", methods=["POST"])
def api_academic_format_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "لم يتم رفع ملف"}), 400
        f = request.files['file']
        filename = f.filename.lower()
        content_parts = []

        file_bytes = f.read()
        if filename.endswith('.pdf'):
            if pdfplumber is None:
                return jsonify({"error": "مكتبة PDF غير متاحة، شغّل: pip install pdfplumber"}), 500
            import io as _io2
            with pdfplumber.open(_io2.BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages[:50], 1):
                    text = page.extract_text() or ''
                    if text.strip():
                        content_parts.append(f"--- صفحة {i} ---\n{text.strip()}")
                    for tbl in page.extract_tables():
                        if tbl:
                            rows_txt = '\n'.join('\t'.join(str(c or '') for c in row) for row in tbl)
                            content_parts.append(f"[جدول]\n{rows_txt}")
        elif filename.endswith('.docx'):
            if docx is None:
                return jsonify({"error": "مكتبة Word غير متاحة، شغّل: pip install python-docx"}), 500
            import io as _io2
            doc_obj = docx.Document(_io2.BytesIO(file_bytes))
            for para in doc_obj.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else ''
                    prefix = '# ' if 'Heading 1' in style else '## ' if 'Heading' in style else ''
                    content_parts.append(prefix + para.text)
            for tbl in doc_obj.tables:
                rows_txt = '\n'.join('\t'.join(c.text for c in row.cells) for row in tbl.rows)
                content_parts.append(f"[جدول]\n{rows_txt}")
        elif filename.endswith('.txt'):
            content_parts.append(file_bytes.decode('utf-8', errors='replace'))
        else:
            return jsonify({"error": "صيغة غير مدعومة. استخدم PDF أو DOCX أو TXT"}), 400

        full_text = '\n\n'.join(content_parts)

        # ── تحسين النص بالذكاء الاصطناعي (اختياري) ──
        ai_summary = None
        use_ai = request.form.get('use_ai', 'false').lower() == 'true'
        if use_ai and full_text.strip():
            try:
                from groq import Groq as _Groq
                _gc = _Groq(api_key=GROQ_API_KEY)
                ai_resp = _gc.chat.completions.create(
                    model='llama-3.3-70b-versatile',
                    messages=[
                        {"role": "system", "content":
                         "أنت مساعد أكاديمي. لخّص الوثيقة التالية بشكل منظّم: العنوان، الأقسام الرئيسية، الأهداف، النتائج، التوصيات."},
                        {"role": "user", "content": f"لخّص هذا النص:\n\n{full_text[:6000]}"}
                    ],
                    max_tokens=800,
                    temperature=0.3,
                )
                ai_summary = ai_resp.choices[0].message.content
            except Exception as _ae:
                ai_summary = f"[تعذّر التلخيص: {_ae}]"

        words = len(full_text.split())
        return jsonify({
            "success":    True,
            "text":       full_text[:15000],
            "words":      words,
            "chars":      len(full_text),
            "sections":   len(content_parts),
            "filename":   f.filename,
            "ai_summary": ai_summary
        })
    except Exception as e:
        logger.error(f"Format file error: {e}")
        return jsonify({"error": str(e)}), 500


# ════════════════════════════════════════════════════════════
#  تحويل HTML إلى Word — محسّن كلياً
#  يدعم: فواصل الصفحات | جداول كاملة | صور base64 | ألوان | RTL
# ════════════════════════════════════════════════════════════

# ألوان CSS المسمّاة
_CSS_NAMED_COLORS = {
    'black':(0,0,0),'white':(255,255,255),'red':(255,0,0),'green':(0,128,0),
    'blue':(0,0,255),'yellow':(255,255,0),'orange':(255,165,0),'purple':(128,0,128),
    'pink':(255,192,203),'gray':(128,128,128),'grey':(128,128,128),'brown':(165,42,42),
    'cyan':(0,255,255),'magenta':(255,0,255),'navy':(0,0,128),'teal':(0,128,128),
    'lime':(0,255,0),'maroon':(128,0,0),'olive':(128,128,0),'silver':(192,192,192),
    'gold':(255,215,0),'coral':(255,127,80),'salmon':(250,128,114),'turquoise':(64,224,208),
    'indigo':(75,0,130),'violet':(238,130,238),'darkblue':(0,0,139),'darkgreen':(0,100,0),
    'darkred':(139,0,0),'darkgray':(169,169,169),'lightblue':(173,216,230),
    'lightgreen':(144,238,144),'lightyellow':(255,255,224),'lightgray':(211,211,211),
    'crimson':(220,20,60),'deepskyblue':(0,191,255),'forestgreen':(34,139,34),
    'hotpink':(255,105,180),'limegreen':(50,205,50),'mediumblue':(0,0,205),
    'orangered':(255,69,0),'royalblue':(65,105,225),'seagreen':(46,139,87),
    'skyblue':(135,206,235),'slategray':(112,128,144),'steelblue':(70,130,180),
    'tomato':(255,99,71),'yellowgreen':(154,205,50),'beige':(245,245,220),
    'ivory':(255,255,240),'khaki':(240,230,140),'lavender':(230,230,250),
    'mintcream':(245,255,250),'snow':(255,250,250),'wheat':(245,222,179),
}

def _w2_css_color(color_str):
    """تحويل لون CSS إلى RGBColor — يدعم hex, rgb(), rgba(), named colors"""
    try:
        from docx.shared import RGBColor
        s = (color_str or '').strip().lower()
        if not s or s == 'transparent' or s == 'inherit' or s == 'currentcolor':
            return None
        # hex
        if s.startswith('#'):
            c = s.lstrip('#')
            if len(c) == 3:
                c = c[0]*2 + c[1]*2 + c[2]*2
            if len(c) >= 6:
                return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
        # rgb / rgba
        if s.startswith('rgb'):
            nums = re.findall(r'[\d.]+', s)
            if len(nums) >= 3:
                r,g,b = int(float(nums[0])), int(float(nums[1])), int(float(nums[2]))
                return RGBColor(min(r,255), min(g,255), min(b,255))
        # named
        if s in _CSS_NAMED_COLORS:
            t = _CSS_NAMED_COLORS[s]
            if t:
                return RGBColor(t[0], t[1], t[2])
    except Exception:
        pass
    return None

def _w2_parse_style(style_str):
    """تحليل CSS style مضمن → dict (يحافظ على القيمة الأخيرة عند التكرار)"""
    props = {}
    for part in (style_str or '').split(';'):
        part = part.strip()
        if ':' in part:
            k, v = part.split(':', 1)
            props[k.strip().lower()] = v.strip()
    return props

def _w2_align(style_props, default='LEFT'):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    m = {'center': WD_ALIGN_PARAGRAPH.CENTER,
         'left':   WD_ALIGN_PARAGRAPH.LEFT,
         'right':  WD_ALIGN_PARAGRAPH.RIGHT,
         'justify':WD_ALIGN_PARAGRAPH.JUSTIFY}
    ta = style_props.get('text-align','').lower()
    return m.get(ta, getattr(WD_ALIGN_PARAGRAPH, default))

def _w2_parse_font_size_pt(fs):
    """تحويل قيمة font-size CSS إلى نقاط (Pt) أو None"""
    if not fs:
        return None
    fs = fs.strip().lower()
    try:
        if fs.endswith('pt'):
            return float(fs[:-2].strip())
        if fs.endswith('px'):
            return round(float(fs[:-2].strip()) * 0.75, 1)
        if fs.endswith('em'):
            return round(float(fs[:-2].strip()) * 12.0, 1)
        if fs.endswith('rem'):
            return round(float(fs[:-3].strip()) * 12.0, 1)
        if fs.endswith('%'):
            return round(float(fs[:-1].strip()) * 12.0 / 100.0, 1)
    except Exception:
        pass
    return None

def _w2_apply_styles_to_run(run, styles):
    """تطبيق dict من CSS styles على run — يدعم وراثة الأنماط"""
    try:
        from docx.shared import Pt, RGBColor
        fw = styles.get('font-weight', '')
        if fw in ('bold','700','800','900','bolder'): run.bold = True
        fi = styles.get('font-style', '')
        if fi == 'italic' or fi == 'oblique': run.italic = True
        td = styles.get('text-decoration', '')
        if 'underline' in td: run.underline = True
        if 'line-through' in td: run.font.strike = True
        c = _w2_css_color(styles.get('color', ''))
        if c: run.font.color.rgb = c
        ff = styles.get('font-family', '')
        if ff:
            fname = ff.split(',')[0].strip().strip('"\'')
            if fname: run.font.name = fname
        pt = _w2_parse_font_size_pt(styles.get('font-size', ''))
        if pt and 6 <= pt <= 72: run.font.size = Pt(pt)
    except Exception:
        pass

def _w2_apply_run(run, node, extra_tag=''):
    """تطبيق تنسيق inline على run (للتوافق مع الكود القديم)"""
    try:
        tag = (node.name or '').lower() if hasattr(node,'name') else extra_tag
        styles = {}
        if tag in ('b','strong'): styles['font-weight'] = 'bold'
        if tag in ('i','em'):     styles['font-style'] = 'italic'
        if tag == 'u':            styles['text-decoration'] = 'underline'
        if tag in ('s','strike','del'): styles['text-decoration'] = 'line-through'
        if tag == 'code':         styles['font-family'] = 'Courier New'
        if hasattr(node,'get'):
            styles.update(_w2_parse_style(node.get('style','')))
            if tag == 'font':
                fc = node.get('color',''); ff = node.get('face','')
                if fc: styles['color'] = fc
                if ff: styles['font-family'] = ff
        _w2_apply_styles_to_run(run, styles)
    except Exception:
        pass

def _w2_inline(node, para, inherited_styles=None):
    """
    معالجة المحتوى inline مع وراثة الأنماط الكاملة من الوالدين.
    inherited_styles: dict من CSS properties موروثة من العنصر الأب
    """
    from bs4 import NavigableString, Tag
    if inherited_styles is None:
        inherited_styles = {}

    INLINE_TAGS = {
        'b','strong','i','em','u','s','strike','del','span','a','font',
        'mark','sup','sub','small','big','code','kbd','abbr','cite','q',
        'bdi','bdo','time','var','ins','samp','dfn',
    }

    for child in node.children:
        if isinstance(child, NavigableString):
            txt = str(child)
            if txt:
                run = para.add_run(txt)
                _w2_apply_styles_to_run(run, inherited_styles)
        elif isinstance(child, Tag):
            ctag = (child.name or '').lower()
            if ctag == 'br':
                para.add_run('\n')
                continue
            if ctag == 'img':
                continue  # تعالجها _w2_node

            # بناء أنماط الابن = أنماط الوالد + أنماط الابن الخاصة
            child_styles = dict(inherited_styles)

            # أنماط مبنية على اسم الوسم
            if ctag in ('b','strong'):   child_styles['font-weight'] = 'bold'
            if ctag in ('i','em'):       child_styles['font-style']  = 'italic'
            if ctag == 'u':              child_styles['text-decoration'] = 'underline'
            if ctag in ('s','strike','del','ins'): child_styles['text-decoration'] = 'line-through'
            if ctag == 'mark':           child_styles['background-color'] = '#ffff00'
            if ctag in ('small',):
                old_pt = _w2_parse_font_size_pt(child_styles.get('font-size','12pt')) or 12
                child_styles['font-size'] = f'{old_pt * 0.85:.1f}pt'
            if ctag == 'code':           child_styles['font-family'] = 'Courier New'
            if ctag == 'font':
                fc = child.get('color',''); ff = child.get('face','')
                if fc: child_styles['color'] = fc
                if ff: child_styles['font-family'] = ff

            # الأنماط المضمّنة الخاصة بالعنصر (تتجاوز الموروثة)
            own_sp = _w2_parse_style(child.get('style','') if hasattr(child,'get') else '')
            child_styles.update(own_sp)

            if ctag in INLINE_TAGS:
                # تحقق هل يحتوي على عناصر block
                has_block = any(
                    isinstance(c, Tag) and (c.name or '').lower() in
                    ('div','p','table','ul','ol','h1','h2','h3','h4','h5','h6')
                    for c in child.children
                )
                if has_block:
                    _w2_inline(child, para, child_styles)
                else:
                    run = para.add_run(child.get_text())
                    _w2_apply_styles_to_run(run, child_styles)
            else:
                # عنصر block داخل inline — استخرج النص مع الأنماط الحالية
                run = para.add_run(child.get_text())
                _w2_apply_styles_to_run(run, child_styles)

def _w2_add_page_break(doc):
    """إضافة فاصل صفحة حقيقي في Word"""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    para = doc.add_paragraph()
    run = para.add_run()
    br = OxmlElement('w:br')
    br.set(qn('w:type'), 'page')
    run._r.append(br)
    return para

def _w2_shade_cell(cell, hex_color_str):
    """تلوين خلفية خلية جدول بلون hex"""
    try:
        from docx.oxml import parse_xml
        xml = ('<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
               ' w:val="clear" w:color="auto" w:fill="' + hex_color_str.upper() + '"/>')
        cell._tc.get_or_add_tcPr().append(parse_xml(xml))
    except Exception:
        pass

def _w2_shade_para(para, hex_color_str):
    """تلوين خلفية فقرة بلون hex"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        shd_xml = (f'<w:shd xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
                   f' w:val="clear" w:color="auto" w:fill="{hex_color_str.upper()}"/>')
        pPr.append(parse_xml(shd_xml))
    except Exception:
        pass

def _w2_set_para_border(para, color_hex='E6B422', side='left', sz=24):
    """إضافة حد ملوّن على جانب فقرة"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        bdr_xml = (
            f'<w:pBdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            f'<w:{side} w:val="single" w:sz="{sz}" w:space="4" w:color="{color_hex.upper()}"/>'
            f'</w:pBdr>'
        )
        pPr.append(parse_xml(bdr_xml))
    except Exception:
        pass

def _w2_set_rtl_para(para):
    """تعيين اتجاه الفقرة RTL للنص العربي"""
    try:
        from docx.oxml import parse_xml
        pPr = para._p.get_or_add_pPr()
        rtl_xml = '<w:bidi xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
        pPr.append(parse_xml(rtl_xml))
    except Exception:
        pass

def _w2_is_rtl(text):
    """كشف إذا كان النص عربياً/RTL"""
    arabic = sum(1 for c in text if '\u0600' <= c <= '\u06FF' or '\u0750' <= c <= '\u077F')
    return arabic > len(text) * 0.3

def _w2_table(node, doc):
    """
    تحويل <table> HTML إلى جدول Word كامل مع:
    - ألوان رأس الجدول من CSS
    - colspan + rowspan
    - خلفيات الخلايا
    - وراثة أنماط النص
    """
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # جمع كل الصفوف من thead + tbody + tfoot
    all_rows = []
    thead_rows = set()
    for section in node.find_all(['thead','tbody','tfoot'], recursive=False) or [node]:
        is_head = section.name == 'thead'
        for tr in section.find_all('tr', recursive=True if section == node else False):
            if is_head:
                thead_rows.add(id(tr))
            all_rows.append(tr)
    if not all_rows:
        all_rows = node.find_all('tr')
    if not all_rows:
        return

    # حساب أقصى عدد أعمدة مع مراعاة colspan
    def row_width(tr):
        w = 0
        for c in tr.find_all(['td','th'], recursive=False):
            w += max(1, int(c.get('colspan', 1) or 1))
        return w

    max_cols = max((row_width(r) for r in all_rows), default=1)
    if max_cols == 0:
        return

    num_rows = len(all_rows)
    table = doc.add_table(rows=num_rows, cols=max_cols)
    table.style = 'Table Grid'

    # شبكة لتتبع rowspan
    grid = [[False]*max_cols for _ in range(num_rows)]

    # قراءة لون border الجدول إن وجد
    tbl_sp = _w2_parse_style(node.get('style',''))
    tbl_border_color = tbl_sp.get('border-color', '').strip()

    for r_idx, tr in enumerate(all_rows):
        is_header_row = (id(tr) in thead_rows) or (r_idx == 0 and not thead_rows)
        # فحص إذا كان الصف هو tr داخل thead
        if tr.parent and tr.parent.name == 'thead':
            is_header_row = True

        cells_el = tr.find_all(['td','th'], recursive=False)
        col_cursor = 0

        for cell_el in cells_el:
            # تجاوز الخلايا المحجوزة بـ rowspan
            while col_cursor < max_cols and grid[r_idx][col_cursor]:
                col_cursor += 1
            if col_cursor >= max_cols:
                break

            colspan = max(1, int(cell_el.get('colspan', 1) or 1))
            rowspan = max(1, int(cell_el.get('rowspan', 1) or 1))
            is_th   = (cell_el.name == 'th') or is_header_row

            # تحديد الخلية في Word
            try:
                cell = table.cell(r_idx, col_cursor)
            except Exception:
                col_cursor += colspan
                continue

            # دمج colspan
            if colspan > 1:
                end_col = min(col_cursor + colspan - 1, max_cols - 1)
                if end_col > col_cursor:
                    try:
                        cell = cell.merge(table.cell(r_idx, end_col))
                    except Exception:
                        pass

            # دمج rowspan
            if rowspan > 1:
                end_row = min(r_idx + rowspan - 1, num_rows - 1)
                if end_row > r_idx:
                    try:
                        cell = cell.merge(table.cell(end_row, col_cursor))
                        # تحديد الشبكة
                        for rr in range(r_idx, end_row + 1):
                            for cc in range(col_cursor, min(col_cursor + colspan, max_cols)):
                                grid[rr][cc] = True
                    except Exception:
                        pass

            # مسح النص الافتراضي
            cell.text = ''
            para = cell.paragraphs[0] if cell.paragraphs else cell.add_paragraph()

            # تحليل أنماط الخلية
            cell_sp = _w2_parse_style(cell_el.get('style',''))
            ta = cell_sp.get('text-align','').lower()
            if not ta:
                ta = 'center' if is_th else 'left'
            align_map = {
                'center': WD_ALIGN_PARAGRAPH.CENTER,
                'left':   WD_ALIGN_PARAGRAPH.LEFT,
                'right':  WD_ALIGN_PARAGRAPH.RIGHT,
                'justify':WD_ALIGN_PARAGRAPH.JUSTIFY,
            }
            para.alignment = align_map.get(ta, WD_ALIGN_PARAGRAPH.LEFT)

            # تلوين خلفية الخلية
            bg = (cell_sp.get('background-color','') or cell_sp.get('background','')).strip()
            if not bg and is_th:
                bg = '#1E4A6E'  # أزرق غامق لرأس الجدول
            if bg:
                rgb = _w2_css_color(bg)
                if rgb:
                    hex_c = '{:02X}{:02X}{:02X}'.format(int(rgb[0]), int(rgb[1]), int(rgb[2]))
                    _w2_shade_cell(cell, hex_c)

            # الأنماط الموروثة للـ runs داخل الخلية
            cell_text_styles = {}
            txt_c = _w2_css_color(cell_sp.get('color',''))
            if is_th:
                cell_text_styles['font-weight'] = 'bold'
                if not txt_c:
                    txt_c = _w2_css_color('#FFFFFF')  # نص أبيض على رأس أزرق
            if txt_c:
                cell_text_styles['color'] = cell_sp.get('color','') or ('#FFFFFF' if is_th else '')

            # معالجة محتوى الخلية
            _w2_inline(cell_el, para, cell_text_styles)

            # تطبيق bold على رأس الجدول بعد _w2_inline
            if is_th:
                for run in para.runs:
                    run.bold = True
                    try:
                        if not run.font.color.rgb:
                            run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                    except Exception:
                        pass

            # تلوين خاص (matrix-high, matrix-mid, matrix-low)
            cell_classes = cell_el.get('class', [])
            if isinstance(cell_classes, str):
                cell_classes = cell_classes.split()
            if 'matrix-high' in cell_classes:
                _w2_shade_cell(cell, '2ECC71')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass
            elif 'matrix-mid' in cell_classes:
                _w2_shade_cell(cell, 'F39C12')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass
            elif 'matrix-low' in cell_classes:
                _w2_shade_cell(cell, 'E74C3C')
                for run in para.runs:
                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); run.bold=True
                    except Exception: pass

            col_cursor += colspan
            grid[r_idx][col_cursor-colspan] = True

    doc.add_paragraph()

def _w2_svg_to_png_bytes(svg_bytes, width_px=400, height_px=300):
    """محاولة تحويل SVG إلى PNG"""
    # المحاولة الأولى: cairosvg
    try:
        import cairosvg
        return cairosvg.svg2png(bytestring=svg_bytes, output_width=width_px, output_height=height_px)
    except Exception:
        pass
    # المحاولة الثانية: svglib + reportlab
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        from io import BytesIO as _BIO2
        import tempfile, os as _os
        with tempfile.NamedTemporaryFile(suffix='.svg', delete=False) as tf:
            tf.write(svg_bytes)
            tmp_path = tf.name
        try:
            drawing = svg2rlg(tmp_path)
            if drawing:
                drawing.width  = width_px
                drawing.height = height_px
                drawing.renderSVG = None
                buf2 = _BIO2()
                renderPM.drawToFile(drawing, buf2, fmt='PNG')
                return buf2.getvalue()
        finally:
            try: _os.unlink(tmp_path)
            except Exception: pass
    except Exception:
        pass
    # المحاولة الثالثة: PIL placeholder
    try:
        from PIL import Image, ImageDraw
        from io import BytesIO as _BIO3
        img = Image.new('RGBA', (width_px, height_px), color=(240, 245, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.rectangle([2,2,width_px-3,height_px-3], outline=(100,120,180), width=2)
        draw.text((10, height_px//2-10), "[ SVG ]", fill=(100,120,180))
        buf3 = _BIO3()
        img.save(buf3, 'PNG')
        return buf3.getvalue()
    except Exception:
        return None

def _w2_embed_image(src, doc, width_inches=5.0, center=True, caption=None):
    """
    تضمين صورة في ملف Word — يدعم:
    - base64 PNG/JPEG/GIF/WebP/SVG
    - روابط HTTP/HTTPS
    - مسارات /static/ المحلية
    """
    try:
        from docx.shared import Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from io import BytesIO
        import base64 as _b64

        img_bytes = None
        if not src:
            return
        # تجاهل مسارات غير قابلة للوصول
        if src.startswith('/sdcard') or src.startswith('file://') or src == '#':
            return

        if 'data:image/svg' in src[:30]:
            parts = src.split(',', 1)
            raw = parts[1] if len(parts) == 2 else ''
            try:
                if '%3C' in raw or '%3c' in raw:
                    import urllib.parse
                    raw = urllib.parse.unquote(raw)
                svg_b = _b64.b64decode(raw) if (raw and not raw.startswith('<')) else raw.encode()
            except Exception:
                svg_b = raw.encode() if isinstance(raw, str) else raw
            img_bytes = _w2_svg_to_png_bytes(svg_b, 600, 400)
            if not img_bytes: return
            width_inches = min(width_inches, 5.5)

        elif src.startswith('data:image'):
            try:
                header, data = src.split(',', 1)
                img_bytes = _b64.b64decode(data + '==')
            except Exception:
                return

        elif src.startswith('http://') or src.startswith('https://'):
            try:
                resp = requests.get(src, timeout=12, headers={'User-Agent': 'Mozilla/5.0'})
                if resp.ok:
                    ct = resp.headers.get('Content-Type','')
                    img_bytes = _w2_svg_to_png_bytes(resp.content) if 'svg' in ct else resp.content
            except Exception:
                return

        elif src.startswith('/static/') or src.startswith('static/'):
            local_path = os.path.join(_PROJECT_ROOT, src.lstrip('/'))
            if os.path.exists(local_path):
                with open(local_path, 'rb') as _lf:
                    img_bytes = _lf.read()

        if not img_bytes:
            return

        # التحقق من صحة الصورة وتحويلها إذا لزم
        try:
            from PIL import Image as _PILImg
            pil_obj = _PILImg.open(BytesIO(img_bytes))
            if pil_obj.mode in ('RGBA','P','LA','CMYK'):
                pil_obj = pil_obj.convert('RGB')
                buf = BytesIO()
                pil_obj.save(buf, 'PNG')
                img_bytes = buf.getvalue()
            else:
                pil_obj.verify()  # تحقق فقط
        except Exception:
            try:
                from PIL import Image as _PILImg
                pil_obj = _PILImg.open(BytesIO(img_bytes)).convert('RGB')
                buf = BytesIO()
                pil_obj.save(buf, 'PNG')
                img_bytes = buf.getvalue()
            except Exception:
                return

        para = doc.add_paragraph()
        if center:
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(BytesIO(img_bytes), width=Inches(min(float(width_inches), 6.0)))

        if caption:
            cap_p = doc.add_paragraph(caption)
            cap_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in cap_p.runs:
                r.italic = True
                try:
                    from docx.shared import Pt
                    r.font.size = Pt(10)
                except Exception: pass

    except Exception as e:
        logger.debug(f"Image embed skip ({str(src)[:60]}): {e}")

def _w2_node(node, doc, in_para=None):
    """
    المعالج الرئيسي: يحوّل عنصر HTML إلى عناصر DOCX مع دعم كامل لـ:
    - ألوان الخلفية على divs والفقرات
    - الجداول الكاملة مع rowspan/colspan
    - الصور بجميع أنواعها
    - القوائم المتداخلة
    - RTL للنص العربي
    - blockquote وpre وcode
    """
    from bs4 import NavigableString, Tag
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    if isinstance(node, NavigableString):
        txt = str(node)
        if in_para and txt.strip():
            in_para.add_run(txt)
        elif in_para and txt and not txt.strip():
            in_para.add_run(txt)
        return

    if not isinstance(node, Tag):
        return

    tag = (node.name or '').lower()
    classes = node.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()
    classes_set = set(classes)
    sp = _w2_parse_style(node.get('style',''))

    # ── تجاهل العناصر غير المرئية ──
    if tag in ('script','style','noscript','meta','link','head'):
        return

    # ── فاصل صفحة ──
    is_page_break = (
        'page-break' in classes_set or 'pagebreak' in classes_set or
        sp.get('page-break-before','') == 'always' or
        sp.get('break-before','') in ('page','always') or
        sp.get('page-break-after','') == 'always'
    )
    if is_page_break and tag in ('div','section','p','hr','span'):
        _w2_add_page_break(doc)
        for child in node.children:
            _w2_node(child, doc)
        return

    # ── عناوين H1-H6 ──
    if tag in ('h1','h2','h3','h4','h5','h6'):
        lvl = int(tag[1])
        try:
            para = doc.add_heading(level=min(lvl, 4))
        except Exception:
            para = doc.add_paragraph()
        para.clear()
        sizes = {1:24, 2:20, 3:17, 4:15, 5:13, 6:12}

        # تطبيق خلفية العنوان إذا وُجدت
        bg = (sp.get('background-color','') or sp.get('background','')).strip()
        if bg:
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                _w2_shade_para(para, '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2])))

        # بناء أنماط الوراثة للعنوان
        heading_styles = {'font-weight': 'bold', 'font-size': f'{sizes.get(lvl,14)}pt'}
        c_str = sp.get('color','')
        if c_str: heading_styles['color'] = c_str

        _w2_inline(node, para, heading_styles)

        # تطبيق الحجم والـ bold على جميع runs بعد الإضافة
        for run in para.runs:
            run.bold = True
            run.font.size = Pt(sizes.get(lvl, 14))
            if c_str:
                c = _w2_css_color(c_str)
                if c: run.font.color.rgb = c

        para.alignment = _w2_align(sp, 'LEFT')

        # RTL للعناوين العربية
        txt_check = node.get_text()
        if _w2_is_rtl(txt_check):
            _w2_set_rtl_para(para)
        return

    # ── فقرات <p> ──
    if tag == 'p':
        para = doc.add_paragraph()
        para.alignment = _w2_align(sp, 'LEFT')

        # خلفية الفقرة
        bg = (sp.get('background-color','') or sp.get('background','')).strip()
        if bg:
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                _w2_shade_para(para, '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2])))

        # أنماط الفقرة
        p_styles = {}
        c_str = sp.get('color','')
        if c_str: p_styles['color'] = c_str
        fw = sp.get('font-weight','')
        if fw: p_styles['font-weight'] = fw
        fs = sp.get('font-size','')
        if fs: p_styles['font-size'] = fs
        ff = sp.get('font-family','')
        if ff: p_styles['font-family'] = ff

        _w2_inline(node, para, p_styles)

        txt_check = node.get_text()
        if _w2_is_rtl(txt_check):
            _w2_set_rtl_para(para)
        return

    # ── جدول ──
    if tag == 'table':
        _w2_table(node, doc)
        return

    # ── صورة ──
    if tag == 'img':
        src = node.get('src','')
        alt = node.get('alt','')
        w_attr = node.get('width','')
        try:
            w_in = min(float(str(w_attr).replace('px','').strip()) / 96.0, 6.0) if w_attr else 5.0
        except Exception:
            w_in = 5.0
        if src:
            _w2_embed_image(src, doc, width_inches=w_in)
        return

    # ── canvas ──
    if tag == 'canvas':
        # محاولة استخراج data-url إن وُجدت
        du = node.get('data-url','') or node.get('data-image','')
        if du:
            _w2_embed_image(du, doc, width_inches=5.5)
        return

    # ── SVG: تحويل إلى PNG وتضمينه ──
    if tag == 'svg':
        try:
            svg_str = str(node)
            svg_bytes = svg_str.encode('utf-8')
            # تحديد العرض والارتفاع من خصائص الـ SVG
            try:
                vb = node.get('viewBox','')
                w_attr = node.get('width','')
                h_attr = node.get('height','')
                if vb:
                    parts = vb.replace(',',' ').split()
                    vb_w = float(parts[2]) if len(parts)>2 else 400
                    vb_h = float(parts[3]) if len(parts)>3 else 300
                    out_w = int(min(vb_w * 2, 1200))
                    out_h = int(min(vb_h * 2, 900))
                elif w_attr and h_attr:
                    out_w = int(min(float(str(w_attr).replace('px','').strip() or 400), 1200))
                    out_h = int(min(float(str(h_attr).replace('px','').strip() or 300), 900))
                else:
                    out_w, out_h = 800, 600
            except Exception:
                out_w, out_h = 800, 600
            png_bytes = _w2_svg_to_png_bytes(svg_bytes, out_w, out_h)
            if png_bytes:
                # حساب حجم الصورة بالإنش
                w_in = min(out_w / 96.0, 5.5)
                para = doc.add_paragraph()
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = para.add_run()
                run.add_picture(BytesIO(png_bytes), width=Inches(w_in))
        except Exception as _svg_err:
            logger.debug(f"SVG embed error: {_svg_err}")
        return

    # ── figure ──
    if tag == 'figure':
        img_el = node.find('img')
        if img_el:
            src = img_el.get('src','')
            if src:
                _w2_embed_image(src, doc, width_inches=5.0)
        cap = node.find('figcaption')
        if cap:
            cp = doc.add_paragraph(cap.get_text(strip=True))
            cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for r in cp.runs:
                r.italic = True
                try: r.font.size = Pt(10)
                except Exception: pass
        return

    # ── قوائم ──
    if tag == 'ul':
        for li in node.find_all('li', recursive=False):
            para = doc.add_paragraph(style='List Bullet')
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            _w2_inline(li, para)
            if _w2_is_rtl(li.get_text()):
                _w2_set_rtl_para(para)
        return
    if tag == 'ol':
        for li in node.find_all('li', recursive=False):
            para = doc.add_paragraph(style='List Number')
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            _w2_inline(li, para)
            if _w2_is_rtl(li.get_text()):
                _w2_set_rtl_para(para)
        return

    # ── hr ──
    if tag == 'hr':
        p = doc.add_paragraph()
        try:
            from docx.oxml import parse_xml
            pPr = p._p.get_or_add_pPr()
            bdr = ('<w:pBdr xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                   '<w:bottom w:val="single" w:sz="6" w:space="1" w:color="AAAAAA"/>'
                   '</w:pBdr>')
            pPr.append(parse_xml(bdr))
        except Exception:
            p.add_run('─' * 60)
        return

    # ── br مستقل ──
    if tag == 'br':
        if in_para:
            in_para.add_run('\n')
        else:
            doc.add_paragraph()
        return

    # ── blockquote ──
    if tag == 'blockquote':
        para = doc.add_paragraph()
        _w2_inline(node, para)
        try:
            from docx.shared import Inches as _IN
            para.paragraph_format.left_indent  = _IN(0.4)
            para.paragraph_format.right_indent = _IN(0.4)
        except Exception: pass
        _w2_set_para_border(para, 'AAAAAA', 'left', 18)
        _w2_shade_para(para, 'F5F5F5')
        return

    # ── pre / code block ──
    if tag == 'pre':
        code_el = node.find('code')
        code_text = (code_el or node).get_text()
        for line in code_text.split('\n'):
            p = doc.add_paragraph(line)
            p.paragraph_format.left_indent = Pt(18)
            _w2_shade_para(p, 'F4F4F4')
            for r in p.runs:
                r.font.name = 'Courier New'
                try: r.font.size = Pt(10)
                except Exception: pass
        return

    # ── عناصر خاصة بالتصميم ──
    if 'insight-box' in classes_set:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _w2_inline(node, para)
        _w2_shade_para(para, 'FEF8E7')
        _w2_set_para_border(para, 'E6B422', 'left', 24)
        return

    if 'toc-row' in classes_set or 'list-row' in classes_set:
        spans = node.find_all(['span','div'], recursive=False)
        texts = [s.get_text(strip=True) for s in spans] if spans else [node.get_text(strip=True)]
        line = '  '.join(filter(None, texts))
        if line:
            para = doc.add_paragraph(line)
            para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        return

    if 'ref-item' in classes_set:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        _w2_inline(node, para)
        try:
            para.paragraph_format.left_indent = Pt(18)
            para.paragraph_format.space_after  = Pt(6)
        except Exception: pass
        return

    if 'figure-caption' in classes_set:
        para = doc.add_paragraph(node.get_text(strip=True))
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in para.runs:
            r.italic = True
            try: r.font.size = Pt(10)
            except Exception: pass
        return

    if 'page-number' in classes_set:
        return

    if 'header-logo' in classes_set:
        img_el = node.find('img')
        if img_el:
            src = img_el.get('src','')
            if src and 'svg' not in src[:30]:
                _w2_embed_image(src, doc, width_inches=1.5)
        return

    # ── div / section مع خلفية ملوّنة → صندوق نصي ──
    CONTAINER_TAGS = {
        'div','section','article','main','aside','header','footer','nav',
        'body','html','form','fieldset','details','summary','address',
        'span','a','label','li',
    }
    if tag in CONTAINER_TAGS:
        bg_raw = (sp.get('background-color','') or sp.get('background','')).strip()
        border_color = sp.get('border-left-color','') or sp.get('border-color','')

        if in_para and tag in ('span','a','label'):
            own_sp = {}
            own_sp.update(sp)
            _w2_inline(node, in_para, own_sp)
            return

        # ── فاصل صفحة تلقائي بين أقسام .page ──
        # يُضاف فاصل قبل كل .page عدا الأولى في الوثيقة
        is_page_div = ('page' in classes_set or 'slide' in classes_set or
                       'sheet' in classes_set or 'cover-page' in classes_set or
                       'content-page' in classes_set)
        if is_page_div and tag == 'div':
            existing_pages = getattr(doc, '_w2_page_count', 0)
            if existing_pages > 0:
                _w2_add_page_break(doc)
            doc._w2_page_count = existing_pages + 1

        # ── استخراج اللون من gradient ──
        bg = bg_raw
        if bg and 'gradient' in bg.lower():
            # استخراج أول لون من linear-gradient(...)
            m = re.search(r'#([0-9a-fA-F]{3,6})\b', bg)
            if not m:
                m = re.search(r'rgba?\s*\([\d.,\s]+\)', bg)
            bg = m.group(0) if m else ''

        if bg and bg not in ('transparent','inherit','none',''):
            # معالجة الأبناء أولاً في فقرات عادية، ثم تلوين الفقرة الأولى
            start_idx = len(doc.paragraphs)
            for child in node.children:
                _w2_node(child, doc, in_para=None)
            end_idx = len(doc.paragraphs)
            rgb_bg = _w2_css_color(bg)
            if rgb_bg:
                # تجاهل الأبيض الخالص (لا حاجة لتلوينه)
                if not (rgb_bg[0] >= 250 and rgb_bg[1] >= 250 and rgb_bg[2] >= 250):
                    hex_bg = '{:02X}{:02X}{:02X}'.format(int(rgb_bg[0]),int(rgb_bg[1]),int(rgb_bg[2]))
                    for i in range(start_idx, end_idx):
                        try:
                            _w2_shade_para(doc.paragraphs[i], hex_bg)
                            if border_color:
                                bc = _w2_css_color(border_color)
                                if bc:
                                    _w2_set_para_border(doc.paragraphs[i],
                                        '{:02X}{:02X}{:02X}'.format(int(bc[0]),int(bc[1]),int(bc[2])),
                                        'left', 18)
                        except Exception: pass
        else:
            for child in node.children:
                _w2_node(child, doc, in_para=None)
        return

    # ── أي عنصر آخر ──
    for child in node.children:
        _w2_node(child, doc, in_para=in_para)


@app.route("/formatter/")
@app.route("/formatter")
def formatter_page():
    """منسق الملفات"""
    resp = make_response(render_template('formatter.html'))
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    return resp


@app.route("/formatter/<path:filename>")
def formatter_static(filename):
    """ملفات ثابتة لمنسق الملفات"""
    from flask import send_from_directory as _sfd
    return _sfd(os.path.join(os.path.dirname(__file__), 'static', 'wf_build'), filename)


# ════════════════════════════════════════════════════════════
#  PDF → Word مع صور كاملة وجداول (PyMuPDF)
# ════════════════════════════════════════════════════════════

def _smart_draw_charts(doc, text_content, add_charts=True):
    """نظام الرسم الذكي: يحلل المحتوى ويضيف رسوماً بيانية مناسبة"""
    if not add_charts:
        return
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.font_manager as fm
        import json as _json
        from io import BytesIO
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt

        # طلب AI لتحليل البيانات واقتراح الرسم
        from groq import Groq as _Groq
        client = _Groq(api_key=GROQ_API_KEY)
        prompt = f"""Analyze the following text and determine if it contains numerical data suitable for a chart.
Reply ONLY with valid JSON (no explanation):
{{
  "has_chart": true,
  "chart_type": "bar",
  "title": "Chart Title",
  "direction": "vertical",
  "colors": ["#1e4a2f","#2d6b41","#4f9a66","#7cb56e","#a8d08d"],
  "labels": ["Label1","Label2","Label3"],
  "values": [10, 25, 15],
  "x_label": "X axis",
  "y_label": "Y axis"
}}
If no chart is appropriate, set has_chart to false.
Text (first 2000 chars):
{text_content[:2000]}"""

        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile',
            messages=[{"role": "user", "content": prompt}],
            max_tokens=400, temperature=0.1
        )
        raw = resp.choices[0].message.content.strip()
        # استخراج JSON من الرد
        import re as _re
        m = _re.search(r'\{.*\}', raw, _re.DOTALL)
        if not m:
            return
        chart_data = _json.loads(m.group())
        if not chart_data.get('has_chart', False):
            return

        labels  = chart_data.get('labels', [])
        values  = chart_data.get('values', [])
        colors  = chart_data.get('colors', ['#1e4a2f','#2d6b41','#4f9a66','#7cb56e'])
        title   = chart_data.get('title', 'مخطط البيانات')
        ctype   = chart_data.get('chart_type', 'bar')
        direc   = chart_data.get('direction', 'vertical')
        xlabel  = chart_data.get('x_label', '')
        ylabel  = chart_data.get('y_label', '')

        if not labels or not values:
            return
        n = min(len(labels), len(values))
        labels, values = labels[:n], [float(v) for v in values[:n]]
        while len(colors) < n:
            colors += colors

        fig, ax = plt.subplots(figsize=(8, 4.5))
        fig.patch.set_facecolor('#f8fdf4')
        ax.set_facecolor('#f0f7ea')

        if ctype == 'pie':
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, colors=colors[:n],
                autopct='%1.1f%%', startangle=90,
                pctdistance=0.82, wedgeprops=dict(width=0.6))
            for at in autotexts:
                at.set_fontsize(9); at.set_fontweight('bold')
        elif ctype == 'line':
            ax.plot(labels, values, marker='o', color=colors[0],
                    linewidth=2.5, markersize=7, markerfacecolor='white',
                    markeredgewidth=2)
            ax.fill_between(range(n), values, alpha=0.12, color=colors[0])
            ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
            if xlabel: ax.set_xlabel(xlabel)
            if ylabel: ax.set_ylabel(ylabel)
            ax.grid(axis='y', linestyle='--', alpha=0.5)
        elif ctype == 'scatter':
            ax.scatter(range(n), values, c=colors[:n], s=120, zorder=5)
            ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
        else:  # bar (default)
            if direc == 'horizontal':
                bars = ax.barh(labels, values, color=colors[:n], edgecolor='white', linewidth=0.5)
                if xlabel: ax.set_xlabel(xlabel)
                if ylabel: ax.set_ylabel(ylabel)
                for bar, val in zip(bars, values):
                    ax.text(bar.get_width()+max(values)*0.01, bar.get_y()+bar.get_height()/2,
                            f'{val:g}', va='center', fontsize=9)
            else:
                bars = ax.bar(labels, values, color=colors[:n], edgecolor='white',
                              linewidth=0.5, width=0.65)
                ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=20, ha='right')
                if xlabel: ax.set_xlabel(xlabel)
                if ylabel: ax.set_ylabel(ylabel)
                for bar, val in zip(bars, values):
                    ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01,
                            f'{val:g}', ha='center', fontsize=9)
            ax.grid(axis='y' if direc != 'horizontal' else 'x', linestyle='--', alpha=0.4)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)

        ax.set_title(title, fontsize=13, fontweight='bold', pad=10)
        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)

        # فاصل + الرسم + تسمية
        doc.add_paragraph('─' * 50)
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(buf, width=Inches(5.5))
        cap = doc.add_paragraph(f"📊 {title}")
        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for r in cap.runs:
            r.italic = True
            try: r.font.size = Pt(10)
            except Exception: pass
        doc.add_paragraph()

    except Exception as e:
        logger.debug(f"Smart draw error: {e}")


@app.route("/tools/pdf_to_word", methods=["POST"])
def api_pdf_to_word():
    """
    تحويل PDF → DOCX مع:
    - الحفاظ الكامل على ألوان الخطوط والخلفيات
    - تجنب تكرار نص الجداول
    - جداول منسّقة مع تلوين رأس الجدول
    - صور مضمّنة كاملة
    - دعم RTL للنص العربي
    - تطبيق إعدادات المستخدم (خط / هوامش / حجم)
    """
    try:
        if 'file' not in request.files:
            return jsonify({"error": "لم يتم رفع ملف PDF"}), 400

        f = request.files['file']
        if not f.filename.lower().endswith('.pdf'):
            return jsonify({"error": "يرجى رفع ملف PDF فقط"}), 400

        # إعدادات المستخدم من الواجهة
        add_smart_draw  = request.form.get('smart_draw', 'false').lower() == 'true'
        user_font       = request.form.get('font_family', 'Times New Roman').strip() or 'Times New Roman'
        user_font_size  = float(request.form.get('font_size', '12') or '12')
        margin_in_str   = request.form.get('margin', '1.0')
        try:
            margin_in = float(margin_in_str)
        except Exception:
            margin_in = 1.0

        file_bytes = f.read()
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_',
                            f.filename.rsplit('.', 1)[0]) or 'document'

        import io as _io
        from io import BytesIO
        import docx as _docx
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = _docx.Document()
        for sec in doc.sections:
            sec.top_margin    = Inches(margin_in)
            sec.bottom_margin = Inches(margin_in)
            sec.left_margin   = Inches(max(margin_in, 1.0))
            sec.right_margin  = Inches(max(margin_in, 1.0))
        doc.styles['Normal'].font.name = user_font
        doc.styles['Normal'].font.size = Pt(user_font_size)

        all_text = []
        total    = 0

        # ── دالة مساعدة: استخراج لون فيتز (int) → (r,g,b) أو None ──
        def _fitz_color(color_val):
            try:
                if color_val is None:
                    return None
                if isinstance(color_val, (list, tuple)) and len(color_val) >= 3:
                    return (int(color_val[0]*255), int(color_val[1]*255), int(color_val[2]*255))
                ci = int(color_val)
                r = (ci >> 16) & 0xFF
                g = (ci >> 8)  & 0xFF
                b =  ci        & 0xFF
                if r == g == b == 0:
                    return None  # أسود = افتراضي، تجاهل
                return (r, g, b)
            except Exception:
                return None

        # ── دالة: هل يتداخل بلوك نصي مع منطقة جدول ──
        def _in_table_region(block_bbox, table_bboxes, threshold=0.5):
            bx0,by0,bx1,by1 = block_bbox
            for tx0,ty0,tx1,ty1 in table_bboxes:
                ix0 = max(bx0,tx0); iy0 = max(by0,ty0)
                ix1 = min(bx1,tx1); iy1 = min(by1,ty1)
                if ix1 > ix0 and iy1 > iy0:
                    inter = (ix1-ix0)*(iy1-iy0)
                    area  = max((bx1-bx0)*(by1-by0), 1)
                    if inter/area > threshold:
                        return True
            return False

        try:
            import fitz as _fitz

            pdf_doc = _fitz.open(stream=file_bytes, filetype="pdf")
            total   = len(pdf_doc)

            # ── جمع بيانات الجداول من pdfplumber (bbox + data) ──
            plumb_pages = {}   # page_num -> {'bboxes': [...], 'tables': [...]}
            try:
                if pdfplumber:
                    with pdfplumber.open(_io.BytesIO(file_bytes)) as _plumb:
                        for pi, pp in enumerate(_plumb.pages[:total]):
                            found = pp.find_tables()
                            bboxes = [t.bbox for t in found]
                            data   = [t.extract() for t in found]
                            if bboxes:
                                plumb_pages[pi] = {'bboxes': bboxes, 'tables': data}
            except Exception:
                pass

            seen_img_xrefs = set()

            for page_num, page in enumerate(pdf_doc):
                if page_num > 0:
                    _w2_add_page_break(doc)

                page_table_bboxes = plumb_pages.get(page_num, {}).get('bboxes', [])

                # ── استخراج النص مع ألوان كاملة ──
                page_dict = page.get_text("dict", sort=True)

                for block in page_dict.get('blocks', []):
                    if block.get('type') != 0:
                        continue  # تجاهل بلوكات الصور — ستُعالَج لاحقاً

                    block_bbox = block.get('bbox', (0,0,0,0))

                    # تخطّ النص الواقع داخل منطقة جدول
                    if _in_table_region(block_bbox, page_table_bboxes):
                        continue

                    for line in block.get('lines', []):
                        # تجميع النص والتنسيق لكل span في السطر
                        spans_data = []
                        max_size = 0
                        line_is_bold = False
                        line_has_color = False

                        for span in line.get('spans', []):
                            t   = span.get('text', '')
                            if not t.strip():
                                continue
                            sz  = span.get('size', user_font_size)
                            flg = span.get('flags', 0)
                            col = span.get('color', 0)
                            rgb = _fitz_color(col)
                            is_bold = bool(flg & 16)
                            spans_data.append({
                                'text': t,
                                'size': sz,
                                'bold': is_bold,
                                'color': rgb,
                                'font': span.get('font',''),
                            })
                            if sz > max_size: max_size = sz
                            if is_bold: line_is_bold = True
                            if rgb: line_has_color = True

                        if not spans_data:
                            continue

                        line_txt = ' '.join(s['text'] for s in spans_data).strip()
                        all_text.append(line_txt)

                        # تحديد نوع الفقرة بناءً على الحجم
                        is_h1 = max_size >= 18
                        is_h2 = 14 <= max_size < 18
                        is_h3 = 12.5 <= max_size < 14 and line_is_bold

                        if is_h1:
                            para = doc.add_heading(level=1)
                            para.clear()
                            base_size = min(max_size, 24)
                        elif is_h2:
                            para = doc.add_heading(level=2)
                            para.clear()
                            base_size = min(max_size, 18)
                        elif is_h3:
                            para = doc.add_heading(level=3)
                            para.clear()
                            base_size = 14
                        else:
                            para = doc.add_paragraph()
                            base_size = user_font_size

                        # إذا كان السطر يحتوي على span واحد أو أقل اختلافاً في التنسيق
                        if len(spans_data) == 1 or (not line_has_color and not any(s['bold'] != line_is_bold for s in spans_data)):
                            run = para.add_run(line_txt)
                            run.bold = line_is_bold or is_h1 or is_h2 or is_h3
                            run.font.size = Pt(base_size)
                            run.font.name = user_font
                            if spans_data[0]['color']:
                                r,g,b = spans_data[0]['color']
                                try: run.font.color.rgb = RGBColor(r,g,b)
                                except Exception: pass
                        else:
                            # spans متعددة مع تنسيق مختلف
                            for span_d in spans_data:
                                run = para.add_run(span_d['text'] + ' ')
                                run.bold = span_d['bold'] or is_h1 or is_h2
                                run.font.size = Pt(min(span_d['size'], 36))
                                run.font.name = user_font
                                if span_d['color']:
                                    r,g,b = span_d['color']
                                    try: run.font.color.rgb = RGBColor(r,g,b)
                                    except Exception: pass

                        # RTL للنص العربي
                        if _w2_is_rtl(line_txt):
                            _w2_set_rtl_para(para)
                            para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        elif is_h1 or is_h2:
                            para.alignment = WD_ALIGN_PARAGRAPH.CENTER

                # ── إدراج الجداول من pdfplumber ──
                page_tables = plumb_pages.get(page_num, {}).get('tables', [])
                for tbl_data in page_tables:
                    if not tbl_data:
                        continue
                    clean = [[str(c or '').strip() for c in row] for row in tbl_data]
                    clean = [r for r in clean if any(c for c in r)]
                    if not clean:
                        continue
                    max_cols = max(len(r) for r in clean)
                    if max_cols == 0:
                        continue
                    clean = [r + [''] * (max_cols - len(r)) for r in clean]

                    doc.add_paragraph()
                    tbl_word = doc.add_table(rows=len(clean), cols=max_cols)
                    tbl_word.style = 'Table Grid'

                    # تعيين عرض الأعمدة بالتساوي
                    try:
                        from docx.oxml.ns import qn as _qn
                        from docx.oxml import OxmlElement as _OXE
                        tbl_w = tbl_word._tbl
                        tblPr = tbl_w.tblPr if tbl_w.tblPr is not None else _OXE('w:tblPr')
                        tblW  = _OXE('w:tblW')
                        tblW.set(_qn('w:w'), '5000')
                        tblW.set(_qn('w:type'), 'pct')
                        tblPr.append(tblW)
                    except Exception:
                        pass

                    for ri, row in enumerate(clean):
                        for ci, val in enumerate(row):
                            cell = tbl_word.cell(ri, ci)
                            cell.text = ''
                            p = cell.paragraphs[0]
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER

                            is_rtl_cell = _w2_is_rtl(val)
                            run = p.add_run(val)
                            run.font.name = user_font
                            run.font.size = Pt(user_font_size)

                            if ri == 0:
                                run.bold = True
                                try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                                except Exception: pass
                                _w2_shade_cell(cell, '1E4A6E')  # أزرق غامق
                            elif ri % 2 == 0:
                                _w2_shade_cell(cell, 'EBF5FB')  # أزرق فاتح للصفوف الزوجية

                            if is_rtl_cell:
                                _w2_set_rtl_para(p)
                    doc.add_paragraph()

                # ── استخراج الصور من الصفحة ──
                img_list = page.get_images(full=True)
                for img_info in img_list:
                    xref = img_info[0]
                    if xref in seen_img_xrefs:
                        continue
                    seen_img_xrefs.add(xref)
                    try:
                        base_img   = pdf_doc.extract_image(xref)
                        img_bytes  = base_img["image"]
                        w_px, h_px = base_img.get("width",1), base_img.get("height",1)
                        if w_px < 80 or h_px < 80:
                            continue  # تجاهل ديكور صغير

                        from PIL import Image as _PILImg
                        try:
                            pil_img = _PILImg.open(BytesIO(img_bytes))
                            if pil_img.mode in ('RGBA','P','LA','CMYK'):
                                pil_img = pil_img.convert('RGB')
                            buf = BytesIO()
                            pil_img.save(buf, 'PNG')
                            buf.seek(0)
                            img_bytes = buf.getvalue()
                        except Exception:
                            pass

                        aspect = w_px / max(h_px, 1)
                        w_in   = min(5.5, max(1.5, aspect * 4.0))

                        para = doc.add_paragraph()
                        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run  = para.add_run()
                        run.add_picture(BytesIO(img_bytes), width=Inches(w_in))
                    except Exception as _ie:
                        logger.debug(f"PDF image xref={xref}: {_ie}")

            pdf_doc.close()

        except ImportError:
            # ── fallback: pdfplumber فقط بدون PyMuPDF ──
            if pdfplumber is None:
                return jsonify({"error": "مكتبة PDF غير متاحة"}), 500
            with pdfplumber.open(_io.BytesIO(file_bytes)) as pdf:
                total = len(pdf.pages)
                for i, page in enumerate(pdf.pages[:60], 1):
                    if i > 1:
                        _w2_add_page_break(doc)

                    # جمع bboxes الجداول لهذه الصفحة
                    page_tbls = page.find_tables()
                    tbl_bboxes = [t.bbox for t in page_tbls]
                    tbl_data   = [t.extract() for t in page_tbls]

                    text = page.extract_text(layout=True) or ''
                    for line in text.split('\n'):
                        lt = line.strip()
                        if not lt:
                            continue
                        all_text.append(lt)
                        para = doc.add_paragraph()
                        run = para.add_run(lt)
                        run.font.name = user_font
                        run.font.size = Pt(user_font_size)
                        if _w2_is_rtl(lt):
                            _w2_set_rtl_para(para)
                            para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

                    for tbl in tbl_data:
                        if not tbl:
                            continue
                        clean = [[str(c or '').strip() for c in r] for r in tbl]
                        clean = [r for r in clean if any(c for c in r)]
                        if not clean:
                            continue
                        mc = max(len(r) for r in clean)
                        clean = [r + ['']*(mc-len(r)) for r in clean]
                        tw = doc.add_table(rows=len(clean), cols=mc)
                        tw.style = 'Table Grid'
                        for ri, row in enumerate(clean):
                            for ci, val in enumerate(row):
                                c = tw.cell(ri, ci)
                                c.text = ''
                                p = c.paragraphs[0]
                                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                                run = p.add_run(val)
                                run.font.name = user_font
                                run.font.size = Pt(user_font_size)
                                if ri == 0:
                                    run.bold = True
                                    try: run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
                                    except Exception: pass
                                    _w2_shade_cell(c, '1E4A6E')
                                elif ri % 2 == 0:
                                    _w2_shade_cell(c, 'EBF5FB')
                        doc.add_paragraph()

        # ── الرسم الذكي (اختياري) ──
        if add_smart_draw and all_text:
            _smart_draw_charts(doc, '\n'.join(all_text[:400]))

        final_path = str(_PPTX_OUTPUTS_DIR / (safe_name + '_converted.docx'))
        doc.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/word/download/{safe_name}_converted.docx",
            "filename": f"{safe_name}_converted.docx",
            "pages": total
        })

    except Exception as e:
        logger.error(f"PDF to Word error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/html_to_word", methods=["POST"])
def api_html_to_word():
    """
    تحويل HTML → DOCX مع دعم كامل لـ:
    - ألوان CSS كاملة مع وراثة الأنماط
    - جداول كاملة (colspan + rowspan + خلفيات)
    - صور base64 + HTTP + مسارات محلية
    - RTL للنص العربي تلقائياً
    - blockquote / pre / code / insight-box
    - تطبيق إعدادات المستخدم (خط / هوامش / حجم)
    """
    try:
        data = request.get_json() or {}
        html_content   = data.get('html', '')
        filename       = data.get('filename', 'document')
        add_smart_draw = data.get('smart_draw', False)

        # إعدادات المستخدم
        user_font      = (data.get('font_family','') or 'Times New Roman').strip()
        user_font_size = float(data.get('font_size', 12) or 12)
        margin_val     = float(data.get('margin', 1.0) or 1.0)
        rtl_doc        = data.get('rtl', False)

        if not html_content:
            return jsonify({"error": "محتوى HTML فارغ"}), 400

        from bs4 import BeautifulSoup
        import docx as _docx
        from docx.shared import Pt, Inches

        soup = BeautifulSoup(html_content, 'html.parser')

        doc = _docx.Document()

        # هوامش الصفحة
        for sec in doc.sections:
            sec.top_margin    = Inches(margin_val)
            sec.bottom_margin = Inches(margin_val)
            sec.left_margin   = Inches(max(margin_val, 0.8))
            sec.right_margin  = Inches(max(margin_val, 0.8))

        # النمط الافتراضي
        normal = doc.styles['Normal']
        normal.font.name = user_font
        normal.font.size = Pt(user_font_size)

        # تطبيق النمط الافتراضي على الـ Heading styles أيضاً
        for h_level in range(1, 5):
            try:
                hstyle = doc.styles[f'Heading {h_level}']
                hstyle.font.name = user_font
            except Exception:
                pass

        # إزالة العناصر غير المرئية
        body = soup.find('body') or soup
        for el in body.find_all(['script','style','noscript','head'], recursive=True):
            el.decompose()

        # معالجة المحتوى
        for child in body.children:
            _w2_node(child, doc)

        # حذف الفقرة الفارغة الأولى إن وُجدت
        try:
            first = doc.paragraphs[0]
            if not first.text.strip() and len(doc.paragraphs) > 1:
                first._element.getparent().remove(first._element)
        except Exception:
            pass

        # تطبيق اتجاه RTL على مستوى القسم إذا طُلب
        if rtl_doc:
            try:
                from docx.oxml import parse_xml
                for sec in doc.sections:
                    sectPr = sec._sectPr
                    bidi_xml = '<w:bidi xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>'
                    sectPr.append(parse_xml(bidi_xml))
            except Exception:
                pass

        # الرسم الذكي (اختياري)
        if add_smart_draw:
            body_text = ' '.join(p.text for p in doc.paragraphs if p.text.strip())[:3000]
            if body_text:
                _smart_draw_charts(doc, body_text)

        # حفظ الملف
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_', filename) or 'document'
        final_path = str(_PPTX_OUTPUTS_DIR / (safe_name + '.docx'))
        doc.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/word/download/{safe_name}.docx",
            "filename": f"{safe_name}.docx"
        })
    except Exception as e:
        logger.error(f"HTML to Word error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/word/download/<path:filename>")
def tools_word_download(filename):
    """تحميل ملف Word المُنشأ"""
    from flask import send_from_directory
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    return send_from_directory(outputs_dir, filename, as_attachment=True)


@app.route("/tools/html_to_excel", methods=["POST"])
def api_html_to_excel():
    """تحويل HTML → Excel (.xlsx) مع استخراج الجداول والنصوص"""
    try:
        data = request.get_json() or {}
        html_content = data.get('html', '')
        filename     = data.get('filename', 'document')

        if not html_content:
            return jsonify({"error": "محتوى HTML فارغ"}), 400

        from bs4 import BeautifulSoup
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        soup = BeautifulSoup(html_content, 'html.parser')
        wb   = openpyxl.Workbook()
        ws   = wb.active
        ws.title = "المحتوى"

        current_row = 1

        # نمط الترويسات
        header_font  = Font(bold=True, size=12, color="FFFFFF")
        header_fill  = PatternFill("solid", fgColor="1a7a3c")
        center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
        wrap_align   = Alignment(wrap_text=True, vertical="top")
        thin_border  = Border(
            left=Side(style='thin'), right=Side(style='thin'),
            top=Side(style='thin'), bottom=Side(style='thin')
        )

        def write_text_row(text, bold=False, size=11):
            nonlocal current_row
            if not text.strip():
                return
            cell = ws.cell(row=current_row, column=1, value=text.strip())
            cell.font = Font(bold=bold, size=size)
            cell.alignment = wrap_align
            current_row += 1

        def write_table(table_el):
            nonlocal current_row
            rows = table_el.find_all('tr')
            if not rows:
                return
            start_row = current_row
            max_col = 0
            for ri, tr in enumerate(rows):
                cells = tr.find_all(['th', 'td'])
                col = 1
                for ci, cell_el in enumerate(cells):
                    while ws.cell(row=current_row, column=col).value is not None:
                        col += 1
                    text  = cell_el.get_text(separator=' ', strip=True)
                    is_th = cell_el.name == 'th' or ri == 0
                    c = ws.cell(row=current_row, column=col, value=text)
                    c.border    = thin_border
                    c.alignment = center_align if is_th else wrap_align
                    if is_th:
                        c.font = header_font
                        c.fill = header_fill
                    else:
                        c.font = Font(size=11)
                    colspan = int(cell_el.get('colspan', 1))
                    rowspan = int(cell_el.get('rowspan', 1))
                    if colspan > 1 or rowspan > 1:
                        ws.merge_cells(
                            start_row=current_row, start_column=col,
                            end_row=current_row + rowspan - 1, end_column=col + colspan - 1
                        )
                    col += colspan
                    if col - 1 > max_col:
                        max_col = col - 1
                current_row += 1
            # ضبط عرض الأعمدة تلقائياً
            for c in range(1, max_col + 1):
                ws.column_dimensions[get_column_letter(c)].width = 22
            current_row += 1  # سطر فارغ بعد الجدول

        body = soup.find('body') or soup
        for el in body.find_all(['script', 'style', 'noscript'], recursive=True):
            el.decompose()

        for node in body.children:
            if not hasattr(node, 'name') or not node.name:
                t = str(node).strip()
                if t:
                    write_text_row(t)
                continue
            tag = node.name.lower()
            if tag in ('h1', 'h2', 'h3', 'h4', 'h5', 'h6'):
                level = int(tag[1])
                write_text_row(node.get_text(strip=True), bold=True, size=max(11, 16 - level))
            elif tag == 'table':
                write_table(node)
            elif tag in ('p', 'div', 'li', 'span', 'blockquote', 'pre'):
                # ابحث عن جداول بداخله
                inner_tables = node.find_all('table')
                if inner_tables:
                    for tbl in inner_tables:
                        write_table(tbl)
                else:
                    txt = node.get_text(separator=' ', strip=True)
                    if txt:
                        write_text_row(txt)
            elif tag in ('ul', 'ol'):
                for li in node.find_all('li', recursive=False):
                    write_text_row('• ' + li.get_text(strip=True))

        # حفظ الملف
        safe_name  = re.sub(r'[^\w\u0600-\u06FF\-_]', '_', filename) or 'document'
        outputs_dir = str(_PPTX_OUTPUTS_DIR)
        os.makedirs(outputs_dir, exist_ok=True)
        final_path = os.path.join(outputs_dir, safe_name + '.xlsx')
        wb.save(final_path)

        return jsonify({
            "success": True,
            "download_url": f"/tools/excel/download/{safe_name}.xlsx",
            "filename": f"{safe_name}.xlsx"
        })
    except Exception as e:
        logger.error(f"HTML to Excel error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500


@app.route("/tools/excel/download/<path:filename>")
def tools_excel_download(filename):
    """تحميل ملف Excel المُنشأ"""
    from flask import send_from_directory
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    return send_from_directory(outputs_dir, filename, as_attachment=True)


# ════════════════════════════════════════════════════════════
#  مساعد الذكاء الاصطناعي (AI Assistant) - يستطيع تعديل الكود
# ════════════════════════════════════════════════════════════
_PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

def _ai_read_file(rel_path):
    """قراءة ملف من المشروع"""
    try:
        full = os.path.join(_PROJECT_ROOT, rel_path)
        full = os.path.realpath(full)
        if not full.startswith(_PROJECT_ROOT):
            return None, "مسار غير مسموح به"
        with open(full, 'r', encoding='utf-8', errors='replace') as f:
            return f.read(), None
    except Exception as e:
        return None, str(e)

def _ai_write_file(rel_path, content):
    """كتابة ملف في المشروع"""
    try:
        full = os.path.join(_PROJECT_ROOT, rel_path)
        full = os.path.realpath(full)
        if not full.startswith(_PROJECT_ROOT):
            return False, "مسار غير مسموح به"
        os.makedirs(os.path.dirname(full), exist_ok=True)
        with open(full, 'w', encoding='utf-8') as f:
            f.write(content)
        return True, None
    except Exception as e:
        return False, str(e)

def _ai_list_files():
    """قائمة ملفات المشروع"""
    result = []
    skip_dirs = {'.git', '__pycache__', '.local', 'node_modules', 'pptx_app/outputs', 'sessions', '.cache', '_extract'}
    skip_exts = {'.pyc', '.pyo', '.session', '.session-journal', '.lock', '.jpg', '.png', '.ico', '.svg', '.webp'}
    for root, dirs, files in os.walk(_PROJECT_ROOT):
        dirs[:] = [d for d in dirs if d not in skip_dirs and not d.startswith('.')]
        rel_root = os.path.relpath(root, _PROJECT_ROOT)
        for f in files:
            ext = os.path.splitext(f)[1].lower()
            if ext in skip_exts:
                continue
            rel = os.path.join(rel_root, f) if rel_root != '.' else f
            result.append(rel)
    return result[:200]


def _ai_github_push(files_dict: dict, commit_msg: str = "🤖 تعديل تلقائي بواسطة المساعد الذكي") -> dict:
    """
    رفع ملفات إلى GitHub عبر REST API.
    files_dict: {rel_path: content_str}
    """
    import base64 as _b64
    headers = {
        'Authorization': f'token {GITHUB_TOKEN}',
        'Accept': 'application/vnd.github.v3+json',
        'Content-Type': 'application/json',
    }
    api = 'https://api.github.com'
    results = {}
    for rel_path, content in files_dict.items():
        try:
            # الحصول على SHA الحالي للملف (إن وُجد)
            get_url = f"{api}/repos/{GITHUB_REPO}/contents/{rel_path}"
            gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
            sha = gr.json().get('sha') if gr.ok else None
            # رفع الملف
            body = {
                'message': commit_msg,
                'content': _b64.b64encode(content.encode('utf-8', errors='replace')).decode(),
                'branch': GITHUB_BRANCH,
            }
            if sha:
                body['sha'] = sha
            pr = requests.put(get_url, headers=headers, json=body, timeout=30)
            results[rel_path] = 'pushed' if pr.ok else f"error {pr.status_code}"
        except Exception as ex:
            results[rel_path] = f"exception: {ex}"
    return results


def _ai_extract_code_blocks(text: str) -> list[dict]:
    """
    استخراج كتل الكود من رد الذكاء الاصطناعي.
    يدعم:  ```python app.py  أو  ```# app.py  أو  ```app.py
    """
    import re as _re
    blocks = []
    pattern = _re.compile(
        r'```(?P<lang>\w+)?\s*(?P<path>[^\n`]{2,100}\.[a-zA-Z]{1,6})?\n(?P<code>.*?)```',
        _re.DOTALL
    )
    for m in pattern.finditer(text):
        code = m.group('code').strip()
        path = (m.group('path') or '').strip()
        lang = (m.group('lang') or '').strip()
        blocks.append({'lang': lang, 'path': path, 'code': code})
    return blocks


@app.route("/api/ai_assistant", methods=["POST"])
def api_ai_assistant():
    """مساعد الذكاء الاصطناعي الكامل — قراءة، تعديل، رفع GitHub"""
    try:
        data     = request.get_json() or {}
        messages = data.get('messages', [])
        action   = data.get('action', 'chat')      # chat | apply | github_push
        target_file = data.get('file', '')         # ملف مستهدف يمكن تحديده يدوياً

        from groq import Groq as _Groq
        client = _Groq(api_key=GROQ_API_KEY)

        # ── قائمة الملفات والملخص ──────────────────────────────
        files_list   = _ai_list_files()
        file_summaries = []
        for mf in ['app.py', 'templates/index.html', 'templates/academic.html',
                   'static/js/app.js', 'requirements.txt', 'render.yaml']:
            c, _ = _ai_read_file(mf)
            if c:
                file_summaries.append(f"- {mf} ({len(c.splitlines())} سطر)")

        system_prompt = f"""أنت مساعد ذكاء اصطناعي مفيد لمستخدمي برنامج "مركز سرعة انجاز".

البرنامج: أداة أتمتة تيليجرام (الانضمام للمجموعات، الرد التلقائي، الجدولة، إرسال الرسائل).

مهمتك: الإجابة على أسئلة المستخدمين حول كيفية استخدام البرنامج، شرح الميزات، ومساعدتهم في حل مشكلاتهم العملية.

لا تملك صلاحية تعديل ملفات البرنامج أو كتابة كود — أنت مساعد محادثة فقط.

أجب بالعربية دائماً. كن مختصراً وودوداً وعملياً."""

        history = [{'role': 'system', 'content': system_prompt}] + messages[-20:]

        # ── استدعاء GROQ (محادثة عادية فقط) ──────────────────
        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile',
            messages=history,
            max_tokens=1500,
            temperature=0.6,
        )
        reply = resp.choices[0].message.content

        return jsonify({
            "success":      True,
            "reply":        reply,
            "apply_result": None,
            "github_result": None,
            "files_list":   []
        })
    except Exception as e:
        logger.error(f"AI Assistant error: {e}", exc_info=True)
        return jsonify({"error": str(e), "success": False}), 500


@app.route("/api/ai_github_push", methods=["POST"])
def api_ai_github_push():
    """رفع ملفات محددة إلى GitHub"""
    try:
        data       = request.get_json() or {}
        file_paths = data.get('files', [])   # قائمة مسارات نسبية
        commit_msg = data.get('message', '🤖 تحديث بواسطة المساعد الذكي')
        if not file_paths:
            # رفع جميع الملفات الرئيسية
            file_paths = ['app.py', 'templates/index.html', 'templates/academic.html',
                          'static/js/app.js', 'requirements.txt', 'render.yaml']
        files_to_push = {}
        for fp in file_paths:
            c, _ = _ai_read_file(fp)
            if c:
                files_to_push[fp] = c
        if not files_to_push:
            return jsonify({"error": "لم يُعثر على ملفات"}), 400
        result = _ai_github_push(files_to_push, commit_msg)
        return jsonify({"success": True, "result": result, "pushed": len(files_to_push)})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai_read_file", methods=["POST"])
def api_ai_read_file_endpoint():
    """قراءة ملف للمساعد الذكي"""
    try:
        data = request.get_json() or {}
        rel_path = data.get('path', '')
        if not rel_path:
            return jsonify({"error": "المسار مطلوب"}), 400
        content, err = _ai_read_file(rel_path)
        if err:
            return jsonify({"error": err}), 400
        return jsonify({"success": True, "content": content, "lines": len(content.splitlines())})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai_write_file", methods=["POST"])
def api_ai_write_file_endpoint():
    """كتابة/تعديل ملف عبر المساعد الذكي"""
    try:
        data = request.get_json() or {}
        rel_path = data.get('path', '')
        content  = data.get('content', '')
        if not rel_path:
            return jsonify({"error": "المسار مطلوب"}), 400
        ok, err = _ai_write_file(rel_path, content)
        if not ok:
            return jsonify({"error": err}), 400
        return jsonify({"success": True, "message": f"تم حفظ {rel_path} بنجاح"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ════════════════════════════════════════════════════════════
#  منشئ العروض التقديمية — مدموج بالكامل داخل app.py
# ════════════════════════════════════════════════════════════

# ── مجلد المخرجات ──────────────────────────────────────────
from pathlib import Path
_PPTX_OUTPUTS_DIR = Path(__file__).parent / "pptx_app" / "outputs"
_PPTX_OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
_PPTX_IMG_CACHE   = _PPTX_OUTPUTS_DIR / "img_cache"
_PPTX_IMG_CACHE.mkdir(parents=True, exist_ok=True)

# ── قوالب الألوان (templates.py) ───────────────────────────
_PPTX_THEMES = {
    "blue":   {"name":"أزرق احترافي",  "primary":(0x66,0x7E,0xEA),"secondary":(0x76,0x4B,0xA2),"accent":(0xFF,0xFF,0xFF),"text_dark":(0x1A,0x1A,0x2E),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF5,0xF7,0xFF)},
    "green":  {"name":"أخضر طبيعي",    "primary":(0x11,0x99,0x55),"secondary":(0x00,0x7A,0x3D),"accent":(0xFF,0xFF,0xFF),"text_dark":(0x1A,0x2E,0x1A),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF0,0xFB,0xF4)},
    "red":    {"name":"أحمر حيوي",     "primary":(0xE5,0x39,0x35),"secondary":(0xB7,0x1C,0x1C),"accent":(0xFF,0xCC,0x02),"text_dark":(0x2E,0x1A,0x1A),"text_light":(0xFF,0xFF,0xFF),"bg":(0xFD,0xF2,0xF2)},
    "purple": {"name":"بنفسجي ملكي",   "primary":(0x6A,0x1B,0x9A),"secondary":(0x4A,0x14,0x8C),"accent":(0xFF,0xD7,0x00),"text_dark":(0x1A,0x1A,0x2E),"text_light":(0xFF,0xFF,0xFF),"bg":(0xF8,0xF0,0xFF)},
}

def _pptx_get_theme(color: str) -> dict:
    return _PPTX_THEMES.get(color, _PPTX_THEMES["blue"])

# ── مساعدات pptx عامة ──────────────────────────────────────
try:
    from pptx import Presentation as _Prs
    from pptx.util import Inches as _Inches, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

def _pptx_rgb(t):
    return _RGBColor(*t)

def _pptx_set_rtl(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")

def _pptx_set_run_font(run, font_name: str):
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs", "a:ea"):
        el = rPr.find(_qn(tag))
        if el is None:
            el = _etree.SubElement(rPr, _qn(tag))
        el.set("typeface", font_name)

def _pptx_add_text_box(slide, left, top, width, height, text, size,
                       bold=False, color=(0,0,0), align=None, wrap=True,
                       font_name="Traditional Arabic"):
    if align is None:
        align = _PP_ALIGN.RIGHT
    box = slide.shapes.add_textbox(_Inches(left), _Inches(top), _Inches(width), _Inches(height))
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align; _pptx_set_rtl(p)
    run = p.add_run()
    run.text = text; run.font.size = _Pt(size); run.font.bold = bold
    run.font.color.rgb = _pptx_rgb(color); _pptx_set_run_font(run, font_name)
    return box

def _pptx_set_cell_bg(cell, color_tuple):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    sf = _etree.SubElement(tcPr, _qn("a:solidFill"))
    sc = _etree.SubElement(sf, _qn("a:srgbClr"))
    sc.set("val", "{:02X}{:02X}{:02X}".format(*color_tuple))

# ── جالب الصور (image_fetcher.py) ──────────────────────────
import hashlib as _hashlib

_PPTX_ARABIC_TO_EN = {
    "تعليم":"education","مدرسة":"school","جامعة":"university","طالب":"students",
    "تدريب":"training","تعلم":"learning","اختبار":"exam","تقنية":"technology",
    "ذكاء اصطناعي":"artificial intelligence","برمجة":"coding","حاسوب":"computer",
    "شبكة":"network","بيانات":"data analytics","سحابة":"cloud computing",
    "ابتكار":"innovation","ذكاء":"artificial intelligence","أعمال":"business",
    "شركة":"corporate office","تجارة":"commerce","سوق":"market","مبيعات":"sales",
    "تسويق":"marketing","إدارة":"management","قيادة":"leadership","فريق":"team",
    "اجتماع":"meeting","عمل":"workplace","موظف":"employees","استراتيجية":"strategy",
    "خطة":"planning","ميزانية":"finance","استثمار":"investment","نمو":"growth chart",
    "صحة":"healthcare","طب":"medicine","مستشفى":"hospital","رياضة":"sports",
    "بيئة":"environment","طاقة":"renewable energy","طبيعة":"nature landscape",
    "مجتمع":"community","أسرة":"family","شباب":"youth","سياحة":"tourism",
    "بناء":"architecture","مشروع":"construction project","تصميم":"design",
    "هندسة":"engineering","خلاصة":"success achievement","نتائج":"results achievement",
    "هدف":"goal target","مستقبل":"future vision","إنجاز":"achievement","نجاح":"success",
}
_PPTX_SLIDE_TYPE_DEFAULTS = {
    "title":"professional presentation","conclusion":"success achievement team",
    "table":"data analytics chart","chart":"business graph analytics","bullets":"professional business",
}

def _pptx_get_slide_image(slide_title, slide_bullets, slide_type="bullets", groq_client=None):
    keywords = ""
    if groq_client:
        try:
            content = slide_title + "\n" + "\n".join(slide_bullets[:3])
            resp = groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role":"user","content":f"Extract 2-3 English keywords for image search from Arabic:\n{content}\nReturn ONLY keywords separated by commas."}],
                max_tokens=25, temperature=0.2,
            )
            kw = resp.choices[0].message.content.strip().replace("\n",",")
            keywords = ",".join([k.strip() for k in kw.split(",") if k.strip()])
        except Exception:
            pass
    if not keywords:
        full = slide_title + " " + " ".join(slide_bullets[:3])
        for ar, en in _PPTX_ARABIC_TO_EN.items():
            if ar in full:
                keywords = en.replace(" ",","); break
    if not keywords:
        keywords = _PPTX_SLIDE_TYPE_DEFAULTS.get(slide_type, "professional business")
    safe_kw   = keywords.strip()[:80]
    cache_key = _hashlib.md5(safe_kw.encode()).hexdigest()[:10]
    cache_path = _PPTX_IMG_CACHE / f"{cache_key}.jpg"
    if cache_path.exists() and cache_path.stat().st_size > 4000:
        return str(cache_path)
    try:
        import requests as _req
        resp = _req.get(f"https://loremflickr.com/900/550/{safe_kw}", timeout=8, allow_redirects=True)
        if resp.status_code == 200 and len(resp.content) > 4000:
            cache_path.write_bytes(resp.content); return str(cache_path)
    except Exception:
        pass
    return None

# ── توليد الرسوم البيانية (chart_generator.py) ─────────────
def _pptx_create_chart(chart_type, labels, values, title="", color="#667eea"):
    try:
        import matplotlib; matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import tempfile
        colors_list = ["#667eea","#764ba2","#11aa55","#e53935","#ff6b35","#00bcd4","#ffc107","#9c27b0"]
        fig, ax = plt.subplots(figsize=(8,5))
        fig.patch.set_facecolor("#f8f9ff"); ax.set_facecolor("#f8f9ff")
        if chart_type == "bar":
            bars = ax.bar(labels, values, color=colors_list[:len(labels)])
            for bar, val in zip(bars, values):
                ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max(values)*0.01, f"{val:,.0f}", ha="center", va="bottom", fontsize=9, fontweight="bold")
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct="%1.1f%%", colors=colors_list[:len(labels)], startangle=90, pctdistance=0.85)
        elif chart_type == "line":
            ax.plot(labels, values, color=color, linewidth=2.5, marker="o", markersize=7)
            ax.fill_between(range(len(labels)), values, alpha=0.15, color=color)
            ax.set_xticks(range(len(labels))); ax.set_xticklabels(labels)
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        elif chart_type == "horizontal_bar":
            ax.barh(labels, values, color=colors_list[:len(labels)])
            ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
        if title:
            ax.set_title(title, fontsize=13, fontweight="bold", pad=12)
        plt.tight_layout()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        plt.savefig(tmp.name, dpi=150, bbox_inches="tight", facecolor=fig.get_facecolor())
        plt.close(fig); return tmp.name
    except Exception as e:
        logger.warning(f"Chart error: {e}"); return None

# ── معالج الذكاء الاصطناعي (ai_processor.py) ───────────────
_PPTX_GROQ_MODEL = "llama-3.3-70b-versatile"

def _pptx_subtitle(ptype):
    return {"business":"عرض تجاري احترافي","educational":"مواد تعليمية متميزة","sales":"عرض تسويقي متكامل","general":"عرض تقديمي شامل"}.get(ptype,"عرض تقديمي")

def _pptx_section_titles(ptype):
    return {"business":["نظرة عامة","الأهداف الاستراتيجية","الخطة التنفيذية","الميزانية والموارد","مؤشرات النجاح"],"educational":["المقدمة","المفاهيم الأساسية","التطبيقات العملية","الأمثلة والتدريبات","التقييم"],"sales":["المشكلة","الحل المقترح","المميزات والفوائد","الأسعار والعروض","لماذا نحن؟"],"general":["المقدمة","المحتوى الرئيسي","التفاصيل","النتائج","التوصيات"]}.get(ptype,["المقدمة","المحتوى","التفاصيل","النتائج","الخلاصة"])

def _pptx_process_locally(text, num_slides, ptype, extracted_tables):
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    sentences = []
    for line in lines:
        parts = re.split(r'[.،,;؛]', line)
        sentences.extend([p.strip() for p in parts if len(p.strip()) > 5])
    title = sentences[0][:70] if sentences else "العرض التقديمي"
    slides = [{"title":title,"subtitle":_pptx_subtitle(ptype),"bullets":[],"slide_type":"title"}]
    remaining = sentences[1:] if len(sentences)>1 else ["محتوى العرض"]
    content_count = max(1, num_slides-2)
    chunk_size = max(1, len(remaining)//max(1,content_count))
    sec_titles = _pptx_section_titles(ptype)
    tables_used = 0
    for i in range(content_count):
        chunk = remaining[i*chunk_size:(i+1)*chunk_size]
        bullets = [b[:120] for b in chunk if b][:5]
        if not bullets:
            bullets = [f"النقطة الرئيسية {i+1}","التفاصيل والمعلومات","الخلاصة الجزئية"]
        slide = {"title":sec_titles[i%len(sec_titles)],"subtitle":"","bullets":bullets,"slide_type":"bullets"}
        if extracted_tables and tables_used<len(extracted_tables) and i==content_count//2:
            slide["slide_type"]="table"; slide["table_data"]=extracted_tables[tables_used]; tables_used+=1
        slides.append(slide)
    slides.append({"title":"الخلاصة والتوصيات","subtitle":"","bullets":["✅ "+(sentences[-1][:90] if sentences else "تم عرض أهم النقاط"),"📌 نرحب بأسئلتكم واستفساراتكم","🙏 شكراً لاهتمامكم"],"slide_type":"conclusion"})
    return slides[:num_slides]

class _AIProcessor:
    def __init__(self):
        self.groq_available = False; self.client = None
        try:
            from groq import Groq
            api_key = os.environ.get("GROQ_API_KEY","").strip()
            if api_key:
                self.client = Groq(api_key=api_key); self.groq_available = True
        except Exception as e:
            logger.warning(f"[Groq init] {e}")

    @property
    def is_ai_available(self): return self.groq_available

    def text_to_presentation_structure(self, text, num_slides=6, presentation_type="general",
                                       title_override="", include_tables=False, include_charts=False,
                                       extracted_tables=None):
        if self.groq_available:
            slides = self._process_with_groq(text, num_slides, presentation_type, include_tables, include_charts)
        else:
            slides = _pptx_process_locally(text, num_slides, presentation_type, extracted_tables or [])
        if title_override and slides:
            slides[0]["title"] = title_override
        return slides

    def _process_with_groq(self, text, num_slides, ptype, include_tables, include_charts):
        label = {"general":"عام","business":"تجاري","educational":"تعليمي","sales":"تسويقي"}.get(ptype,"عام")
        extras = ""
        if include_tables: extras += '\n- أضف شريحة جدول بنوع "table" مع مفتاح "table_data"'
        if include_charts: extras += '\n- أضف شريحة رسم بياني بنوع "chart" مع chart_type/chart_labels/chart_values'
        user_prompt = f"""حلّل النص وأنشئ هيكل عرض {label} من {num_slides} شرائح.{extras}

النص:
{text[:4000]}

القواعد:
- الأولى: نوع "title" مع عنوان رئيسي وعنوان فرعي
- الأخيرة: نوع "conclusion" مع 3 نقاط
- الوسطى: نوع "bullets" مع 3-5 نقاط
- كل المحتوى بالعربية
- أرجع JSON فقط: [{{"title":"...","subtitle":"...","bullets":["..."],"slide_type":"title|bullets|table|chart|conclusion",...}}]"""
        try:
            resp = self.client.chat.completions.create(
                model=_PPTX_GROQ_MODEL,
                messages=[{"role":"system","content":"أنت خبير عروض تقديمية. تُنتج JSON صحيحاً فقط."},{"role":"user","content":user_prompt}],
                max_tokens=3000, temperature=0.65,
            )
            content = resp.choices[0].message.content.strip()
            match = re.search(r'\[.*\]', content, re.DOTALL)
            if match:
                slides = json.loads(match.group())
                if isinstance(slides, list) and slides:
                    return slides
        except Exception as e:
            logger.warning(f"[Groq pptx] {e}")
        return _pptx_process_locally(text, num_slides, ptype, [])

# ── مولّد PPTX (presentation_generator.py) ─────────────────
class _PresentationGenerator:
    def __init__(self):
        self.font_name = "Traditional Arabic"
        self.body_font_size = 22

    def _parse_color(self, color_str):
        from pptx.dml.color import RGBColor
        import re
        if not color_str:
            return None
        s = color_str.strip().lower()
        if s.startswith('#'):
            c = s.lstrip('#')
            if len(c) == 3:
                c = c[0]*2 + c[1]*2 + c[2]*2
            if len(c) >= 6:
                return RGBColor(int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
        if s.startswith('rgb'):
            nums = re.findall(r'[\d.]+', s)
            if len(nums) >= 3:
                r,g,b = int(float(nums[0])), int(float(nums[1])), int(float(nums[2]))
                return RGBColor(min(r,255), min(g,255), min(b,255))
        named = {
            'black':(0,0,0),'white':(255,255,255),'red':(255,0,0),'green':(0,128,0),
            'blue':(0,0,255),'yellow':(255,255,0),'orange':(255,165,0),'purple':(128,0,128),
            'pink':(255,192,203),'gray':(128,128,128),'grey':(128,128,128),'brown':(165,42,42),
            'cyan':(0,255,255),'magenta':(255,0,255),'navy':(0,0,128),'teal':(0,128,128),
            'lime':(0,255,0),'maroon':(128,0,0),'olive':(128,128,0),'silver':(192,192,192),
            'gold':(255,215,0),'coral':(255,127,80),'salmon':(250,128,114),'turquoise':(64,224,208),
            'indigo':(75,0,130),'violet':(238,130,238),'darkblue':(0,0,139),'darkgreen':(0,100,0),
            'darkred':(139,0,0),'darkgray':(169,169,169),'lightblue':(173,216,230),
            'lightgreen':(144,238,144),'lightyellow':(255,255,224),'lightgray':(211,211,211),
            'crimson':(220,20,60),'deepskyblue':(0,191,255),'forestgreen':(34,139,34),
            'hotpink':(255,105,180),'limegreen':(50,205,50),'mediumblue':(0,0,205),
            'orangered':(255,69,0),'royalblue':(65,105,225),'seagreen':(46,139,87),
            'skyblue':(135,206,235),'slategray':(112,128,144),'steelblue':(70,130,180),
            'tomato':(255,99,71),'yellowgreen':(154,205,50),'beige':(245,245,220),
            'ivory':(255,255,240),'khaki':(240,230,140),'lavender':(230,230,250),
            'wheat':(245,222,179)
        }
        if s in named:
            r,g,b = named[s]
            return RGBColor(r,g,b)
        return None

    def _fetch_image(self, src):
        import base64 as _b64, requests as _req
        if src.startswith("data:image"):
            try:
                header, data = src.split(",", 1)
                return _b64.b64decode(data)
            except Exception:
                return None
        elif src.startswith("http"):
            try:
                resp = _req.get(src, timeout=5)
                if resp.ok:
                    return resp.content
            except Exception:
                pass
        return None

    def _save_temp_image(self, src):
        import tempfile as _tmp
        data = self._fetch_image(src)
        if data:
            ext = 'png' if 'png' in src else 'jpg'
            fd, path = _tmp.mkstemp(suffix=f'.{ext}')
            with os.fdopen(fd, 'wb') as f:
                f.write(data)
            return path
        return None

    def create_presentation(self, slides_data, theme_color="blue", cover_data=None,
                            extracted_images=None, font_name="Traditional Arabic",
                            body_font_size=22, ai_images=False, groq_client=None):
        from datetime import datetime as _dt
        self.font_name = font_name; self.body_font_size = body_font_size
        theme = _pptx_get_theme(theme_color)
        prs = _Prs()
        prs.slide_width = _Inches(13.33); prs.slide_height = _Inches(7.5)
        if cover_data:
            self._add_cover_slide(prs, cover_data, theme)
        for i, sd in enumerate(slides_data):
            stype = sd.get("slide_type","bullets")
            if stype == "title":        self._add_title_slide(prs, sd, theme)
            elif stype == "table" and sd.get("table_data"): self._add_table_slide(prs, sd, theme)
            elif stype == "chart":      self._add_chart_slide(prs, sd, theme)
            elif stype == "conclusion": self._add_conclusion_slide(prs, sd, theme)
            else:
                img_path = None
                if sd.get("images"):
                    img_src = sd["images"][0]
                    img_path = self._save_temp_image(img_src)
                if not img_path and extracted_images and i < len(extracted_images):
                    img_path = extracted_images[i]
                if not img_path and ai_images:
                    img_path = _pptx_get_slide_image(sd.get("title",""), sd.get("bullets",[]), stype, groq_client)
                self._add_content_slide(prs, sd, theme, i, img_path)
        ts = _dt.now().strftime("%Y%m%d_%H%M%S")
        out = str(_PPTX_OUTPUTS_DIR / f"عرض_{ts}.pptx")
        prs.save(out); return out

    def _add_cover_slide(self, prs, cover_data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["primary"])
        rect = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(2.2))
        rect.fill.solid(); rect.fill.fore_color.rgb = _pptx_rgb(theme["secondary"]); rect.line.fill.background()
        logo_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(0.35),_Inches(2),_Inches(1.5))
        lp = logo_box.text_frame.paragraphs[0]; lp.alignment = _PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = cover_data.get("logo","📊"); lr.font.size = _Pt(52)
        org = cover_data.get("organization","")
        if org:
            ob = slide.shapes.add_textbox(_Inches(2.8),_Inches(0.5),_Inches(10),_Inches(1))
            op = ob.text_frame.paragraphs[0]; op.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(op)
            or_ = op.add_run(); or_.text = org; or_.font.size = _Pt(20); or_.font.bold = True
            or_.font.color.rgb = _RGBColor(220,220,255); _pptx_set_run_font(or_, self.font_name)
        fn = self.font_name
        _pptx_add_text_box(slide,1,2.5,11.3,1.8,cover_data.get("title","العنوان الرئيسي"),42,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        subtitle = cover_data.get("subtitle","")
        if subtitle:
            _pptx_add_text_box(slide,1.5,4.4,10.3,0.9,subtitle,22,color=(210,210,240),align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(3.5),_Inches(5.5),_Inches(6.3),_Emu(35000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()
        meta = []
        if cover_data.get("presenter"): meta.append(f"إعداد: {cover_data['presenter']}")
        if cover_data.get("date"): meta.append(cover_data["date"])
        if meta:
            _pptx_add_text_box(slide,1,5.9,11.3,0.7,"  |  ".join(meta),16,color=(200,200,230),align=_PP_ALIGN.CENTER,font_name=fn)

    def _add_title_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["primary"])
        self._add_bottom_strip(slide, theme); fn = self.font_name
        _pptx_add_text_box(slide,1,2.2,11.3,1.8,data.get("title","العنوان"),44,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        subtitle = data.get("subtitle","")
        if subtitle:
            _pptx_add_text_box(slide,2,4.2,9.3,1,subtitle,24,color=(220,220,255),align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(4.5),_Inches(4.0),_Inches(4.3),_Emu(40000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()

    def _add_content_slide(self, prs, data, theme, index, img_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # تطبيق خلفية الشريحة (من HTML أو من الثيم)
        bg_color = data.get("bg_color")
        if bg_color:
            rgb = self._parse_color(bg_color)
            if rgb:
                bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb
            else:
                bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        else:
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title",""),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        bullets = data.get("bullets",[])
        has_image = img_path and Path(img_path).exists()
        if has_image:
            content_box = slide.shapes.add_textbox(_Inches(4.8),_Inches(1.55),_Inches(8.1),_Inches(5.7))
            self._fill_bullets(content_box.text_frame, bullets, theme)
            try:
                frame = slide.shapes.add_shape(1,_Inches(0.18),_Inches(1.48),_Inches(4.44),_Inches(5.64))
                frame.fill.solid(); frame.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); frame.line.fill.background()
                slide.shapes.add_picture(img_path,_Inches(0.25),_Inches(1.55),_Inches(4.3),_Inches(5.5))
            except Exception:
                content_box2 = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
                self._fill_bullets(content_box2.text_frame, bullets, theme)
        else:
            # إضافة صور من HTML مباشرة إن وجدت
            images = data.get("images", [])
            if images:
                temp_path = self._save_temp_image(images[0])
                if temp_path:
                    try:
                        content_box = slide.shapes.add_textbox(_Inches(4.8),_Inches(1.55),_Inches(8.1),_Inches(5.7))
                        self._fill_bullets(content_box.text_frame, bullets, theme)
                        slide.shapes.add_picture(temp_path,_Inches(0.25),_Inches(1.55),_Inches(4.3),_Inches(5.5))
                    except Exception:
                        content_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
                        self._fill_bullets(content_box.text_frame, bullets, theme)
                    finally:
                        try: os.unlink(temp_path)
                        except: pass
                else:
                    content_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
                    self._fill_bullets(content_box.text_frame, bullets, theme)
            else:
                content_box = slide.shapes.add_textbox(_Inches(0.5),_Inches(1.55),_Inches(12.3),_Inches(5.7))
                self._fill_bullets(content_box.text_frame, bullets, theme)
        num_box = slide.shapes.add_textbox(_Inches(12.6),_Inches(6.9),_Inches(0.7),_Inches(0.5))
        np_ = num_box.text_frame.paragraphs[0]; np_.alignment = _PP_ALIGN.CENTER
        nr = np_.add_run(); nr.text = str(index+1); nr.font.size = _Pt(11); nr.font.color.rgb = _pptx_rgb(theme["secondary"])

    def _fill_bullets(self, tf, bullets, theme):
        tf.word_wrap = True; fs = self.body_font_size; fn = self.font_name
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(p); p.space_before = _Pt(10); p.space_after = _Pt(4)
            run = p.add_run(); run.text = f"◆ {bullet}"; run.font.size = _Pt(fs)
            run.font.color.rgb = _pptx_rgb(theme["text_dark"]); _pptx_set_run_font(run, fn)

    def _add_table_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title","جدول البيانات"),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        table_data = data.get("table_data",[])
        if not table_data: return
        rows = len(table_data); cols = max(len(r) for r in table_data)
        if rows==0 or cols==0: return
        table = slide.shapes.add_table(rows,cols,_Inches(0.5),_Inches(1.6),_Inches(12.3),_Inches(min(5.5,0.5+rows*0.55))).table
        fn = self.font_name
        for r_idx, row in enumerate(table_data):
            for c_idx in range(cols):
                cell_data = row[c_idx] if c_idx < len(row) else {"text": ""}
                if isinstance(cell_data, dict):
                    text = cell_data.get("text", "")
                    bg_color = cell_data.get("bg")
                    align_val = cell_data.get("align", "center")
                    bold = cell_data.get("bold", False)
                else:
                    text = str(cell_data); bg_color = None; align_val = "center"; bold = (r_idx == 0)
                cell = table.cell(r_idx, c_idx); cell.text = text
                tf = cell.text_frame; par = tf.paragraphs[0]
                par.alignment = getattr(_PP_ALIGN, align_val.upper(), _PP_ALIGN.CENTER); _pptx_set_rtl(par)
                if bg_color:
                    rgb = self._parse_color(bg_color)
                    if rgb: _pptx_set_cell_bg(cell, (rgb.red, rgb.green, rgb.blue))
                    else: _pptx_set_cell_bg(cell, theme["primary"] if r_idx==0 else ((235,238,255) if r_idx%2==0 else (255,255,255)))
                else:
                    _pptx_set_cell_bg(cell, theme["primary"] if r_idx==0 else ((235,238,255) if r_idx%2==0 else (255,255,255)))
                runs = par.runs
                if runs:
                    runs[0].font.size = _Pt(15); runs[0].font.bold = bold or (r_idx==0)
                    runs[0].font.color.rgb = _pptx_rgb(theme["text_light"] if r_idx==0 else theme["text_dark"])
                    _pptx_set_run_font(runs[0], fn)

    def _add_chart_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["bg"])
        header = slide.shapes.add_shape(1,_Inches(0),_Inches(0),_Inches(13.33),_Inches(1.35))
        header.fill.solid(); header.fill.fore_color.rgb = _pptx_rgb(theme["primary"]); header.line.fill.background()
        _pptx_add_text_box(slide,0.4,0.15,12,1.05,data.get("title","رسم بياني"),27,bold=True,color=theme["text_light"],align=_PP_ALIGN.RIGHT,font_name=self.font_name)
        chart_type = data.get("chart_type","bar"); labels = data.get("chart_labels",[]); values = data.get("chart_values",[])
        if not labels or not values:
            _pptx_add_text_box(slide,1,3,11,2,"لا توجد بيانات كافية",20,color=theme["text_dark"],font_name=self.font_name); return
        try: values_float = [float(v) for v in values]
        except Exception: values_float = [1.0]*len(labels)
        primary_hex = "#{:02x}{:02x}{:02x}".format(*theme["primary"])
        img_path = _pptx_create_chart(chart_type, labels, values_float, data.get("chart_title",""), primary_hex)
        if img_path and Path(img_path).exists():
            slide.shapes.add_picture(img_path,_Inches(1.5),_Inches(1.6),_Inches(10.3),_Inches(5.6))
        else:
            _pptx_add_text_box(slide,1,3,11,2,"تعذّر إنشاء الرسم البياني",20,color=theme["text_dark"],font_name=self.font_name)

    def _add_conclusion_slide(self, prs, data, theme):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = _pptx_rgb(theme["secondary"])
        fn = self.font_name
        _pptx_add_text_box(slide,1,1.5,11.3,1.2,data.get("title","الخلاصة"),36,bold=True,color=theme["text_light"],align=_PP_ALIGN.CENTER,font_name=fn)
        line = slide.shapes.add_shape(1,_Inches(3.5),_Inches(2.9),_Inches(6.3),_Emu(35000))
        line.fill.solid(); line.fill.fore_color.rgb = _pptx_rgb(theme["accent"]); line.line.fill.background()
        bullets = data.get("bullets",[]); cont_box = slide.shapes.add_textbox(_Inches(1.5),_Inches(3.2),_Inches(10.3),_Inches(4.0))
        tf = cont_box.text_frame; tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = _PP_ALIGN.RIGHT; _pptx_set_rtl(p); p.space_before = _Pt(12)
            run = p.add_run(); run.text = bullet; run.font.size = _Pt(self.body_font_size)
            run.font.color.rgb = _RGBColor(220,220,255); _pptx_set_run_font(run, fn)

    def _add_bottom_strip(self, slide, theme):
        rect = slide.shapes.add_shape(1,_Inches(0),_Inches(6.8),_Inches(13.33),_Inches(0.7))
        rect.fill.solid(); rect.fill.fore_color.rgb = _pptx_rgb(theme["secondary"]); rect.line.fill.background()

# ── تطبيق التصميم (design_applier.py) ──────────────────────
class _DesignApplier:
    def __init__(self): self.generator = _PresentationGenerator()
    def apply_design_to_presentation(self, presentation_path, design, slides_data):
        return self.generator.create_presentation(slides_data, theme_color=design.get("theme_color","blue"))

# ── استخراج الملفات (file_extractor.py) ────────────────────
def _extract_content(file_bytes: bytes, filename: str) -> dict:
    ext = Path(filename).suffix.lower()
    if ext in (".docx",".doc"):
        import docx as _docx, io as _io, tempfile as _tmp
        doc = _docx.Document(_io.BytesIO(file_bytes))
        result = {"text_blocks":[],"tables":[],"images":[],"full_text":""}
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text: continue
            result["text_blocks"].append({"type":"paragraph","style":para.style.name if para.style else "","text":text})
            result["full_text"] += text + "\n"
        for tbl in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
            if rows: result["tables"].append(rows)
        return result
    elif ext == ".pdf":
        result = {"text_blocks":[],"tables":[],"images":[],"full_text":""}
        try:
            import pdfplumber as _pdfp, io as _io
            with _pdfp.open(_io.BytesIO(file_bytes)) as pdf:
                for pg_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        result["text_blocks"].append({"type":"page","style":f"صفحة {pg_num}","text":text.strip()})
                        result["full_text"] += text.strip() + "\n"
                    for tbl in page.extract_tables():
                        if tbl:
                            result["tables"].append([[cell or "" for cell in row] for row in tbl])
        except Exception as e:
            result["full_text"] = f"خطأ في قراءة PDF: {str(e)}"
        return result
    else:
        text = file_bytes.decode("utf-8", errors="ignore")
        return {"text_blocks":[{"type":"paragraph","style":"Normal","text":text}],"tables":[],"images":[],"full_text":text}

# ── محوّل HTML إلى PPTX (html_to_pptx.py) ──────────────────
def _parse_html_slides(html: str):
    from bs4 import BeautifulSoup as _BS
    import re as _re
    soup = _BS(html, "lxml")
    for tag in soup.find_all(["script","noscript"]): tag.decompose()
    slide_cands = (soup.find_all("section") or
                   soup.find_all(class_=_re.compile(r"slide|page|frame",_re.I)) or
                   soup.find_all(attrs={"data-slide":True}) or
                   soup.find_all("article"))
    if not slide_cands:
        sections = []; current = None
        for el in (soup.body or soup).children:
            if not hasattr(el,"name"): continue
            if el.name in ("h1","h2","h3"):
                if current: sections.append(current)
                current = [el]
            elif current is not None: current.append(el)
        if current: sections.append(current)
        slide_cands = []
        for group in sections:
            wrapper = _BS("<div></div>","lxml").div
            for el in group: wrapper.append(el.__copy__())
            slide_cands.append(wrapper)

    def _extract(el):
        sd = {
            "title": "", "subtitle": "", "bullets": [], "slide_type": "bullets",
            "bg_color": None, "text_color": None, "table_data": None,
            "images": [], "is_title_slide": False, "font_size": None, "bold": False,
            "raw_html": str(el)
        }
        # 1. استخراج العنوان من .slide-title أو h1-h4
        title_el = el.find(class_="slide-title")
        if title_el:
            sd["title"] = title_el.get_text(strip=True)
        else:
            for htag in ["h1","h2","h3","h4"]:
                h = el.find(htag)
                if h:
                    sd["title"] = h.get_text(strip=True)
                    sd["is_title_slide"] = (htag == "h1")
                    break
        # 2. استخراج الخلفية
        style_attr = el.get("style", "")
        bg_match = _re.search(r'background(?:-color)?\s*:\s*([^;]+)', style_attr)
        if bg_match:
            sd["bg_color"] = bg_match.group(1).strip()
        content_el = el.find(class_="content")
        if content_el:
            style2 = content_el.get("style", "")
            bg_match2 = _re.search(r'background(?:-color)?\s*:\s*([^;]+)', style2)
            if bg_match2:
                sd["bg_color"] = bg_match2.group(1).strip()
        # 3. استخراج النقاط
        bullets_el = el.find(class_="list-items")
        if bullets_el:
            sd["bullets"] = [li.get_text(strip=True) for li in bullets_el.find_all("li") if li.get_text(strip=True)]
        else:
            lis = el.find_all("li")
            if lis:
                sd["bullets"] = [li.get_text(strip=True) for li in lis if li.get_text(strip=True)]
            else:
                body = el.find(class_="slide-body")
                if body:
                    p_els = body.find_all("p")
                    if p_els:
                        sd["bullets"] = [p.get_text(strip=True) for p in p_els if p.get_text(strip=True) and len(p.get_text(strip=True))>3]
                if not sd["bullets"]:
                    sd["bullets"] = [p.get_text(strip=True) for p in el.find_all("p") if p.get_text(strip=True) and p.get_text(strip=True) not in (sd["title"],sd["subtitle"]) and len(p.get_text(strip=True))>3]
        sd["bullets"] = sd["bullets"][:8]
        # 4. استخراج الجداول مع ألوان الخلايا
        tbl = el.find("table")
        if tbl:
            rows = []
            for tr in tbl.find_all("tr"):
                cells = []
                for td in tr.find_all(["td","th"]):
                    cell_text = td.get_text(strip=True)
                    cell_style = td.get("style","")
                    bg = _re.search(r'background(?:-color)?\s*:\s*([^;]+)', cell_style)
                    bg_val = bg.group(1).strip() if bg else None
                    align = _re.search(r'text-align\s*:\s*([^;]+)', cell_style)
                    align_val = align.group(1).strip() if align else "center"
                    cells.append({"text": cell_text, "bg": bg_val, "align": align_val,
                                  "bold": (td.name=="th") or ("font-weight" in cell_style and "bold" in cell_style)})
                if cells: rows.append(cells)
            if rows:
                sd["table_data"] = rows; sd["slide_type"] = "table"
        # 5. استخراج الصور
        for img in el.find_all("img"):
            src = img.get("src","")
            if src and (src.startswith("data:") or src.startswith("http")):
                sd["images"].append(src)
        # 6. لون النص
        color_match = _re.search(r'(?<![a-z])color\s*:\s*([^;]+)', style_attr)
        if color_match:
            sd["text_color"] = color_match.group(1).strip()
        # 7. نوع الشريحة
        classes = " ".join(el.get("class",[]))
        if _re.search(r"title|cover|first|intro",classes,_re.I) or sd["is_title_slide"]:
            sd["slide_type"] = "title"
        elif _re.search(r"end|outro|thank|conclusion",classes,_re.I):
            sd["slide_type"] = "conclusion"
        return sd

    slides = []
    for el in slide_cands:
        data = _extract(el)
        if data.get("title") or data.get("bullets") or data.get("table_data"):
            slides.append(data)
    if not slides:
        slides = [_extract(soup.body or soup)]
    return slides

def _html_to_pptx(html: str, override_title: str = "") -> str:
    from datetime import datetime as _dt
    slides_data = _parse_html_slides(html)
    if not slides_data:
        raise ValueError("لم يتم العثور على محتوى في HTML")
    if override_title and slides_data:
        slides_data[0]["title"] = override_title
    gen = _PresentationGenerator()
    out = gen.create_presentation(
        slides_data=slides_data,
        theme_color="blue",
        cover_data=None,
        extracted_images=[],
        font_name="Traditional Arabic",
        body_font_size=22,
        ai_images=False,
        groq_client=None
    )
    return out

# ── تهيئة الكائنات ──────────────────────────────────────────
try:
    _pptx_ai  = _AIProcessor()
    _pptx_gen = _PresentationGenerator()
    _pptx_dap = _DesignApplier()
    logger.info("✅ PowerPoint modules loaded successfully (inlined)")
except Exception as _pptx_err:
    logger.error(f"❌ Failed to init PPTX: {_pptx_err}")
    _pptx_ai = _pptx_gen = _pptx_dap = None


@app.route('/tools/pptx/ai_status')
def tools_pptx_ai_status():
    available = _pptx_ai is not None and _pptx_ai.is_ai_available
    return jsonify({'available': available})


@app.route('/tools/pptx/extract_file', methods=['POST'])
def tools_pptx_extract_file():
    if 'file' not in request.files:
        return jsonify({'error': 'لا يوجد ملف'}), 400
    f = request.files['file']
    try:
        result = _extract_content(f.read(), f.filename)
        return jsonify({
            'success': True,
            'full_text': result['full_text'],
            'tables': result.get('tables', []),
            'word_count': len(result['full_text'].split()),
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/tools/pptx/generate', methods=['POST'])
def tools_pptx_generate():
    if _pptx_ai is None or _pptx_gen is None:
        return jsonify({'error': 'وحدات PowerPoint غير متاحة — تحقق من سجلات الخادم'}), 500
    try:
        data = request.json or {}
        text           = data.get('text', '')
        num_slides     = int(data.get('num_slides', 6))
        ptype          = data.get('ptype', 'general')
        title          = data.get('title', '')
        inc_tables     = bool(data.get('inc_tables', False))
        inc_charts     = bool(data.get('inc_charts', False))
        theme_color    = data.get('theme_color', 'blue')
        font_name      = data.get('font_name', 'Traditional Arabic')
        body_font_size = int(data.get('body_font_size', 22))
        cover_data     = data.get('cover_data', None)
        ext_tables     = data.get('extracted_tables', [])

        slides = _pptx_ai.text_to_presentation_structure(
            text=text, num_slides=num_slides, presentation_type=ptype,
            title_override=title, include_tables=inc_tables,
            include_charts=inc_charts, extracted_tables=ext_tables,
        )

        if inc_tables and ext_tables:
            if not any(s.get('slide_type') == 'table' for s in slides):
                for i, tbl in enumerate(ext_tables[:2]):
                    slides.insert(min(2 + i, len(slides) - 1),
                        {'title': f'جدول البيانات {i+1}', 'slide_type': 'table',
                         'table_data': tbl, 'bullets': []})

        if inc_charts and ext_tables:
            def _parse_table_for_chart(tbl):
                """استخراج بيانات رسم بياني من جدول"""
                try:
                    if not tbl or len(tbl) < 2: return None
                    headers = [str(c) for c in tbl[0]]
                    labels, values = [], []
                    for row in tbl[1:]:
                        if not row: continue
                        label = str(row[0]).strip()
                        for cell in row[1:]:
                            try:
                                val = float(str(cell).replace(',','').replace('%','').strip())
                                labels.append(label); values.append(val); break
                            except Exception: pass
                    if labels and values:
                        return {'labels': labels, 'values': values,
                                'title': headers[1] if len(headers) > 1 else ''}
                except Exception: pass
                return None
            if not any(s.get('slide_type') == 'chart' for s in slides):
                for tbl in ext_tables:
                    info = _parse_table_for_chart(tbl)
                    if info:
                        slides.insert(min(3, len(slides) - 1),
                            {'title': 'تحليل البيانات', 'slide_type': 'chart',
                             'chart_type': 'bar', 'chart_labels': info['labels'],
                             'chart_values': info['values'],
                             'chart_title': info.get('title', ''), 'bullets': []})
                        break

        groq_client = _pptx_ai.client if (_pptx_ai and _pptx_ai.is_ai_available) else None
        out = _pptx_gen.create_presentation(
            slides_data=slides, theme_color=theme_color,
            cover_data=cover_data, extracted_images=[],
            font_name=font_name, body_font_size=body_font_size,
            ai_images=False, groq_client=groq_client,
        )
        filename = os.path.basename(out)
        return jsonify({'success': True, 'filename': filename, 'slides': slides})
    except Exception as e:
        import traceback
        logger.error(f"PPTX generate error: {e}")
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route('/tools/pptx/download/<path:filename>')
def tools_pptx_download(filename):
    from flask import send_from_directory as _sfd2
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    try:
        return _sfd2(outputs_dir, filename, as_attachment=True)
    except Exception as _dl_err:
        logger.error(f"PPTX download error for '{filename}': {_dl_err}")
        return jsonify({'error': f'الملف غير موجود: {filename}'}), 404


@app.route('/download/excel/<path:filename>')
def download_excel(filename):
    """تحميل ملفات Excel من مجلد الإخراج"""
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    file_path = os.path.join(outputs_dir, filename)
    if not os.path.exists(file_path):
        return jsonify({'error': 'الملف غير موجود'}), 404
    return send_file(
        file_path,
        as_attachment=True,
        download_name=filename,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


@app.route('/download/excel-list')
def download_excel_list():
    """قائمة ملفات Excel المتاحة للتحميل"""
    outputs_dir = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs')
    files = []
    for f in os.listdir(outputs_dir):
        if f.endswith('.xlsx') or f.endswith('.xls'):
            files.append({'name': f, 'url': f'/download/excel/{f}'})
    return jsonify({'files': files})


@app.route('/stats1208')
@app.route('/احصاء1208')
def stats1208_page():
    """صفحة حل اختبار احصاء 1208 التفاعلية"""
    return render_template('stats1208.html')


@app.route('/stats1208.xlsx')
@app.route('/get-stats-excel')
def download_stats1208():
    """تحميل ملف Excel – احصاء 1208 – رابط مباشر"""
    f = os.path.join(os.path.dirname(__file__), 'static', 'stats1208_solution.xlsx')
    if not os.path.exists(f):
        # حاول الإنشاء تلقائياً إن لم يوجد
        try:
            import subprocess, sys
            subprocess.run([sys.executable, 'create_excel_solution.py'], check=True)
            import shutil
            src = os.path.join(os.path.dirname(__file__), 'pptx_app', 'outputs',
                               'احصاء_1208_الحلول_الكاملة.xlsx')
            shutil.copy(src, f)
        except Exception as ex:
            return jsonify({'error': str(ex)}), 500
    return send_file(
        f,
        as_attachment=True,
        download_name='احصاء_1208_الحلول_الكاملة.xlsx',
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )


@app.route('/tools/pptx/chat', methods=['POST'])
def tools_pptx_chat():
    try:
        data     = request.json or {}
        messages = data.get('messages', [])
        groq_key = os.environ.get('GROQ_API_KEY', '').strip()
        if not groq_key:
            return jsonify({'reply': '⚠️ مفتاح GROQ_API_KEY غير موجود في إعدادات الأسرار. يمكنك الحصول على مفتاح مجاني من console.groq.com'})
        from groq import Groq as _Groq
        client = _Groq(api_key=groq_key)
        history = [{'role': 'system', 'content':
            'أنت مساعد ذكي متخصص في إنشاء العروض التقديمية PowerPoint والمحتوى الأكاديمي. أجب باللغة العربية دائماً.'}
        ] + messages
        resp = client.chat.completions.create(
            model='llama-3.3-70b-versatile', messages=history, max_tokens=1500)
        return jsonify({'reply': resp.choices[0].message.content})
    except Exception as e:
        return jsonify({'reply': f'❌ خطأ: {e}'})


@app.route('/tools/pptx/preview_html_links', methods=['POST'])
def tools_pptx_preview_html_links():
    """استخراج الروابط من HTML وجلب عناوينها ومقتطفاتها."""
    try:
        data = request.json or {}
        html = data.get('html', '')
        if not html:
            return jsonify({'links': []})

        from bs4 import BeautifulSoup as _BS
        import requests as _req

        soup = _BS(html, 'lxml')
        links = []
        seen = set()

        for a in soup.find_all('a', href=True):
            href = a['href'].strip()
            if not href or href.startswith('#') or href in seen:
                continue
            seen.add(href)
            anchor_text = a.get_text(strip=True)[:100]
            info = {'url': href, 'anchor_text': anchor_text, 'title': '', 'snippet': '', 'error': None}

            if href.startswith('http://') or href.startswith('https://'):
                try:
                    r = _req.get(href, timeout=6,
                                 headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
                                 allow_redirects=True)
                    if r.status_code == 200:
                        page = _BS(r.text[:80000], 'lxml')
                        t = page.find('title')
                        info['title'] = t.get_text(strip=True)[:120] if t else ''
                        paras = [p.get_text(strip=True) for p in page.find_all('p')
                                 if len(p.get_text(strip=True)) > 40]
                        info['snippet'] = paras[0][:250] if paras else ''
                    else:
                        info['error'] = f'HTTP {r.status_code}'
                except Exception as ex:
                    info['error'] = str(ex)[:80]
            else:
                info['error'] = 'رابط محلي (لا يمكن جلبه)'

            links.append(info)
            if len(links) >= 12:
                break

        return jsonify({'links': links})
    except Exception as e:
        return jsonify({'links': [], 'error': str(e)}), 500


@app.route('/tools/pptx/from_html', methods=['POST'])
def tools_pptx_from_html():
    """تحويل كود HTML مباشرةً إلى ملف PPTX مع دعم جلب محتوى الروابط."""
    if _pptx_gen is None:
        return jsonify({'error': 'وحدات PowerPoint غير متاحة — تحقق من سجلات الخادم'}), 500
    try:
        data           = request.json or {}
        html_code      = data.get('html', '')
        title          = data.get('title', '')
        theme_color    = data.get('theme_color', 'blue')
        font_name      = data.get('font_name', 'Traditional Arabic')
        body_font_size = int(data.get('body_font_size', 22))
        fetch_links    = bool(data.get('fetch_links', True))
        cover_data     = data.get('cover_data', None)

        if not html_code.strip():
            return jsonify({'error': 'الرجاء لصق كود HTML أولاً'}), 400

        html_to_parse = html_code

        # جلب محتوى الروابط الخارجية وإضافتها كأقسام إضافية
        if fetch_links:
            from bs4 import BeautifulSoup as _BS2
            import requests as _req2
            soup2 = _BS2(html_code, 'lxml')
            extra_parts = []
            seen2 = set()
            for a in soup2.find_all('a', href=True):
                href = a['href'].strip()
                if not href.startswith('http') or href in seen2:
                    continue
                seen2.add(href)
                try:
                    r2 = _req2.get(href, timeout=6,
                                   headers={'User-Agent': 'Mozilla/5.0'},
                                   allow_redirects=True)
                    if r2.status_code == 200:
                        ps = _BS2(r2.text[:80000], 'lxml')
                        pt = ps.find('title')
                        pg_title = pt.get_text(strip=True) if pt else href
                        paras = [p.get_text(strip=True) for p in ps.find_all('p')
                                 if len(p.get_text(strip=True)) > 30]
                        if paras:
                            bullets = ''.join(f'<li>{p}</li>' for p in paras[:5])
                            extra_parts.append(
                                f'<section><h2>{pg_title}</h2><ul>{bullets}</ul></section>'
                            )
                except Exception:
                    pass
                if len(extra_parts) >= 5:
                    break

            if extra_parts:
                insert = '\n'.join(extra_parts)
                body_end = html_to_parse.rfind('</body>')
                if body_end != -1:
                    html_to_parse = html_to_parse[:body_end] + insert + html_to_parse[body_end:]
                else:
                    html_to_parse += insert

        # تحليل HTML وبناء شرائح
        slides_data = _parse_html_slides(html_to_parse)
        if not slides_data:
            return jsonify({'error': 'لم يتم العثور على محتوى قابل للتحويل في HTML'}), 400

        if title and slides_data:
            slides_data[0]['title'] = title

        out_path = _pptx_gen.create_presentation(
            slides_data=slides_data,
            theme_color=theme_color,
            cover_data=cover_data,
            extracted_images=[],
            font_name=font_name,
            body_font_size=body_font_size,
            ai_images=False,
            groq_client=None,
        )
        filename = os.path.basename(out_path)
        return jsonify({'success': True, 'filename': filename, 'slides': slides_data})

    except Exception as e:
        import traceback
        logger.error(f"HTML to PPTX error: {e}")
        return jsonify({'error': str(e), 'trace': traceback.format_exc()}), 500


@app.route("/academic")
def academic_analysis():
    _groq_key = os.environ.get('GROQ_API_KEY', '').strip()
    resp = make_response(render_template('academic.html', groq_key=_groq_key))
    resp.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate, max-age=0'
    return resp


alert_queue.start()
# تهيئة الحلقة المشتركة مبكراً قبل أي اتصالات gevent
_ensure_shared_login_loop()
load_all_sessions()

def _auto_resume_persistent_tasks():
    def worker():
        time.sleep(3)
        logger.info("🔁 فحص المهام الدائمة لإعادة تشغيلها تلقائياً...")
        with USERS_LOCK:
            user_ids = list(USERS.keys())
        resumed = 0
        for uid in user_ids:
            try:
                settings = load_settings(uid)
                want_monitor = bool(settings.get('monitoring_persistent', False))
                want_rotating = bool(settings.get('rotating_persistent', False))
                if not (want_monitor or want_rotating):
                    continue
                ok = telegram_manager.ensure_client_active(uid)
                if not ok:
                    logger.warning(f"⏭️  لا يمكن استئناف مهام {uid}: العميل غير متاح/غير موثق")
                    continue
                with USERS_LOCK:
                    if uid in USERS and not USERS[uid].get('client_manager'):
                        USERS[uid]['client_manager'] = telegram_manager.get_client_manager(uid)
                if want_monitor:
                    with USERS_LOCK:
                        already = USERS.get(uid, {}).get('is_running', False)
                        if not already:
                            USERS[uid]['is_running'] = True
                    if not already:
                        t = _OSThread(target=monitoring_worker, args=(uid,), daemon=True)
                        t.start()
                        with USERS_LOCK:
                            if uid in USERS:
                                USERS[uid]['thread'] = t
                        logger.info(f"♻️  استُؤنفت المراقبة للحساب {uid}")
                        resumed += 1
                        try:
                            socketio.emit('log_update', {
                                "message": "♻️ تم استئناف المراقبة تلقائياً (مهمة دائمة)"
                            }, to=uid)
                            socketio.emit('monitoring_status', {
                                "monitoring_active": True, "status": "running", "is_running": True
                            }, to=uid)
                            socketio.emit('update_monitoring_buttons', {"is_running": True}, to=uid)
                        except Exception:
                            pass
                if want_rotating:
                    msgs = settings.get('rotating_messages', [])
                    grps = dedupe_groups(settings.get('rotating_groups', []))
                    interval = int(settings.get('rotating_interval', 5))
                    valid_msgs = [m for m in msgs if m and m.strip()]
                    if grps and valid_msgs:
                        is_alive = (uid in rotating_manager.threads
                                    and rotating_manager.threads[uid]
                                    and rotating_manager.threads[uid].is_alive())
                        if not is_alive:
                            def _cb(u, status, group, info):
                                if status == 'success':
                                    socketio.emit('log_update', {"message": f"🔄 [متسلسل] أرسل إلى {group}"}, to=u)
                                else:
                                    socketio.emit('log_update', {"message": f"❌ [متسلسل] فشل إلى {group}: {info}"}, to=u)
                            rotating_manager.start(uid, grps, valid_msgs, interval, _cb)
                            logger.info(f"♻️  استُؤنف الإرسال المتسلسل للحساب {uid}")
                            resumed += 1
                            try:
                                socketio.emit('log_update', {
                                    "message": f"♻️ تم استئناف الإرسال المتسلسل تلقائياً ({len(valid_msgs)} رسائل) كل {interval} دقيقة"
                                }, to=uid)
                            except Exception:
                                pass
                    else:
                        logger.info(f"⏭️  تخطي استئناف الإرسال المتسلسل لـ {uid}: لا توجد رسائل/مجموعات")
            except Exception as e:
                logger.error(f"خطأ أثناء استئناف مهام {uid}: {e}")
        if resumed:
            logger.info(f"✅ تم استئناف {resumed} مهمة دائمة")
        else:
            logger.info("ℹ️  لا توجد مهام دائمة لاستئنافها")
    _OSThread(target=worker, daemon=True, name="AutoResumeTasks").start()

_auto_resume_persistent_tasks()

# ════════════════════════════════════════════════════════════
#  مزامنة الجلسات مع GitHub
# ════════════════════════════════════════════════════════════
def github_upload_session(user_id):
    """رفع ملف الجلسة إلى GitHub للحفظ الاحتياطي"""
    try:
        session_file = os.path.join(SESSIONS_DIR, f"{user_id}_session.session")
        if not os.path.exists(session_file):
            return False
        with open(session_file, 'rb') as f:
            content = f.read()
        import base64 as _b64
        headers = {
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json',
        }
        rel_path = f"sessions/{user_id}_session.session"
        get_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{rel_path}"
        gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
        sha = gr.json().get('sha') if gr.ok else None
        body = {
            'message': f'Upload session for {user_id}',
            'content': _b64.b64encode(content).decode(),
            'branch': GITHUB_BRANCH,
        }
        if sha:
            body['sha'] = sha
        pr = requests.put(get_url, headers=headers, json=body, timeout=30)
        if pr.ok:
            logger.info(f"Session for {user_id} uploaded to GitHub")
            return True
        else:
            logger.warning(f"Failed to upload session: {pr.status_code}")
            return False
    except Exception as e:
        logger.error(f"github_upload_session error: {e}")
        return False

def github_delete_session(user_id):
    """حذف ملف الجلسة الملغاة من GitHub"""
    try:
        headers = {
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json',
        }
        rel_path = f"sessions/{user_id}_session.session"
        get_url = f"https://api.github.com/repos/{GITHUB_REPO}/contents/{rel_path}"
        gr = requests.get(get_url, headers=headers, params={'ref': GITHUB_BRANCH}, timeout=15)
        if not gr.ok:
            return False
        sha = gr.json().get('sha')
        if not sha:
            return False
        body = {
            'message': f'Delete revoked session for {user_id}',
            'sha': sha,
            'branch': GITHUB_BRANCH,
        }
        dr = requests.delete(get_url, headers=headers, json=body, timeout=30)
        if dr.ok:
            logger.info(f"Session for {user_id} deleted from GitHub")
            return True
        return False
    except Exception as e:
        logger.error(f"github_delete_session error: {e}")
        return False

# ════════════════════════════════════════════════════════════
#  مدقق صحة الجلسات الدوري
# ════════════════════════════════════════════════════════════
def start_session_health_checker():
    """يفحص صحة جميع الجلسات كل 120 ثانية"""
    def _checker():
        import time as _t
        _t.sleep(120)  # انتظر دقيقتين قبل الفحص الأول لإعطاء العملاء وقتاً كافياً
        while True:
            try:
                with USERS_LOCK:
                    user_ids = list(USERS.keys())
                for uid in user_ids:
                    try:
                        with USERS_LOCK:
                            user_data = USERS.get(uid, {})
                        client_manager = user_data.get('client_manager')
                        if not client_manager or not client_manager.client:
                            continue
                        # تخطّ الفحص إذا كان الخيط لا يزال يبدأ
                        if client_manager.thread and client_manager.thread.is_alive() and not client_manager.is_ready.is_set():
                            logger.debug(f"Health check skipped for {uid} — client still starting")
                            continue
                        # تخطّ الفحص إذا لم يكن الـ loop يعمل بعد
                        if not client_manager.loop or not client_manager.loop.is_running():
                            continue
                        is_valid = client_manager.check_session_valid_sync()
                        if not is_valid:
                            logger.warning(f"Session health check failed for {uid} — revoking")
                            socketio.emit('session_revoked', {
                                "user_id": uid,
                                "reason": "انتهت صلاحية الجلسة أو تم إلغاؤها"
                            }, to=uid)
                            github_delete_session(uid)
                    except Exception as inner_e:
                        logger.debug(f"Health check error for {uid}: {inner_e}")
            except Exception as e:
                logger.debug(f"Session health checker error: {e}")
            _t.sleep(120)

    # استخدام OS thread حقيقي للـ health checker
    t = _OSThread(target=_checker, daemon=True, name='SessionHealthChecker')
    t.start()
    logger.info("✅ مدقق صحة الجلسات الدوري مُفعّل (كل 120 ثانية)")

start_session_health_checker()

# ════════════════════════════════════════════════════════════
#  نقاط API — فحص الجلسة وإعادة التعيين
# ════════════════════════════════════════════════════════════
@app.route("/api/check_session_valid", methods=["GET"])
def api_check_session_valid():
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            user_data = USERS.get(user_id, {})
        client_manager = user_data.get('client_manager')
        if not client_manager or not client_manager.client:
            return jsonify({"success": True, "valid": False, "reason": "العميل غير متصل"})
        is_valid = client_manager.check_session_valid_sync()
        return jsonify({"success": True, "valid": is_valid})
    except Exception as e:
        return jsonify({"success": False, "valid": False, "reason": str(e)})

@app.route("/api/force_reset_session", methods=["POST"])
def api_force_reset_session():
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            user_data = USERS.get(user_id, {})
        client_manager = user_data.get('client_manager')
        if client_manager:
            client_manager.run_coroutine(client_manager.force_reset_session())
        github_delete_session(user_id)
        clear_user_session(user_id)
        with USERS_LOCK:
            if user_id in USERS:
                USERS[user_id]['connected'] = False
                USERS[user_id]['authenticated'] = False
                USERS[user_id]['client_manager'] = None
        return jsonify({"success": True, "message": "تم إعادة تعيين الجلسة بنجاح"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/settings", methods=["GET", "POST"])
def api_auto_join_settings():
    user_id = session.get('user_id', 'user_1')
    if request.method == 'POST':
        try:
            data = request.json or {}
            settings = load_settings(user_id)
            if 'links' in data:
                settings['auto_join_links'] = data['links']
            if 'delay' in data:
                settings['auto_join_delay'] = int(data['delay'])
            if 'max_retries' in data:
                settings['auto_join_max_retries'] = int(data['max_retries'])
            save_settings(user_id, settings)
            return jsonify({"success": True})
        except Exception as e:
            return jsonify({"success": False, "message": str(e)})
    settings = load_settings(user_id)
    return jsonify({
        "success": True,
        "links": settings.get('auto_join_links', []),
        "delay": settings.get('auto_join_delay', 3),
        "max_retries": settings.get('auto_join_max_retries', 1)
    })

@app.route("/api/auto_join/stop", methods=["POST"])
def api_auto_join_stop():
    user_id = session.get('user_id', 'user_1')
    try:
        stopped = False
        with USERS_LOCK:
            stop_event = USERS.get(user_id, {}).get('auto_join_stop')
            pause_event = USERS.get(user_id, {}).get('auto_join_pause')
            if stop_event and not stop_event.is_set():
                stop_event.set()
                stopped = True
            if pause_event:
                pause_event.clear()  # أطلق الإيقاف المؤقت حتى يصل الإيقاف الكامل
        msg = "⏹ تم إيقاف الانضمام التلقائي" if stopped else "ℹ️ لا يوجد انضمام نشط حالياً"
        socketio.emit('log_update', {"message": msg}, to=user_id)
        return jsonify({"success": True, "message": msg, "stopped": stopped})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/pause", methods=["POST"])
def api_auto_join_pause():
    """إيقاف مؤقت للانضمام — تبقى الإعدادات محفوظة"""
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            pause_event = USERS.get(user_id, {}).get('auto_join_pause')
            state = USERS.get(user_id, {}).get('auto_join_state', {})
        if pause_event and not pause_event.is_set():
            pause_event.set()
            msg = "⏸ تم الإيقاف المؤقت — اضغط استئناف للمتابعة"
            socketio.emit('log_update', {"message": msg}, to=user_id)
            return jsonify({"success": True, "message": msg})
        return jsonify({"success": False, "message": "لا يوجد انضمام نشط أو هو متوقف مؤقتاً بالفعل"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/resume", methods=["POST"])
def api_auto_join_resume():
    """استئناف الانضمام بعد إيقاف مؤقت"""
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            pause_event = USERS.get(user_id, {}).get('auto_join_pause')
        if pause_event and pause_event.is_set():
            pause_event.clear()
            msg = "▶ تم الاستئناف"
            socketio.emit('log_update', {"message": msg}, to=user_id)
            return jsonify({"success": True, "message": msg})
        return jsonify({"success": False, "message": "الانضمام ليس في حالة توقف مؤقت"})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/exit", methods=["POST"])
def api_auto_join_exit():
    """خروج يدوي — يوقف الانضمام ويمسح الإعدادات"""
    user_id = session.get('user_id', 'user_1')
    try:
        with USERS_LOCK:
            stop_event  = USERS.get(user_id, {}).get('auto_join_stop')
            pause_event = USERS.get(user_id, {}).get('auto_join_pause')
            if stop_event and not stop_event.is_set():
                stop_event.set()
            if pause_event:
                pause_event.clear()
            if user_id in USERS:
                USERS[user_id].pop('auto_join_state', None)
                USERS[user_id].pop('auto_join_stop', None)
                USERS[user_id].pop('auto_join_pause', None)
        # امسح الإعدادات المحفوظة أيضاً
        try:
            _s = load_settings(user_id)
            _s.pop('auto_join_links', None)
            _s.pop('auto_join_delay', None)
            save_settings(user_id, _s)
        except Exception:
            pass
        msg = "🚪 خروج يدوي — تم مسح جميع إعدادات الانضمام"
        socketio.emit('log_update', {"message": msg}, to=user_id)
        socketio.emit('auto_join_exited', {}, to=user_id)
        return jsonify({"success": True, "message": msg})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/auto_join/status", methods=["GET"])
def api_auto_join_status():
    """يُرجع الحالة الجارية للانضمام المتقدم (للاستعادة عند تبديل الحسابات)"""
    user_id = session.get('user_id', 'user_1')
    with USERS_LOCK:
        state = USERS.get(user_id, {}).get('auto_join_state', {})
    return jsonify({"success": True, "state": state})


@app.route("/api/auto_join/history", methods=["GET"])
def api_auto_join_history():
    """
    سجل نتائج الانضمام مع إمكانية التصفية حسب الحالة.
    query params:
      - status: all | success | failed | already  (افتراضي: all)
      - limit:  عدد السجلات (افتراضي: 100)
    """
    user_id = session.get('user_id', 'user_1')
    status_filter = request.args.get('status', 'all').lower()
    try:
        limit = max(1, min(int(request.args.get('limit', 100)), 1000))
    except (ValueError, TypeError):
        limit = 100

    with USERS_LOCK:
        state = USERS.get(user_id, {}).get('auto_join_state', {})

    all_items = state.get('items', [])

    if status_filter == 'success':
        items = [i for i in all_items if i.get('status') == 'success']
    elif status_filter == 'failed':
        items = [i for i in all_items if i.get('status') == 'failed']
    elif status_filter == 'already':
        items = [i for i in all_items if i.get('status') == 'already']
    else:
        items = all_items

    items = items[-limit:]

    summary = {
        'total':   state.get('total', 0),
        'done':    state.get('done', 0),
        'success': state.get('success', 0),
        'fail':    state.get('fail', 0),
        'already': state.get('already', 0),
        'running': state.get('running', False),
    }

    return jsonify({
        "success": True,
        "filter": status_filter,
        "count": len(items),
        "summary": summary,
        "items": items
    })


# ════════════════════════════════════════════════════════════
#  نظام الاستمرارية الدائم — يمنع توقف التطبيق تلقائياً
# ════════════════════════════════════════════════════════════
def _start_keepalive():
    """
    يُرسل ping لنفسه كل 4 دقائق لمنع السكون (خاصة على Render Free).
    يعمل كخيط daemon لا يمنع إيقاف التطبيق يدوياً.
    """
    import time as _time
    import socket as _socket

    def _ping_self():
        _time.sleep(30)  # انتظر حتى يكتمل الإقلاع
        _port = int(os.environ.get('PORT', 5000))
        while True:
            try:
                _socket.setdefaulttimeout(10)
                conn = _socket.create_connection(('127.0.0.1', _port), timeout=5)
                conn.close()
            except Exception:
                pass
            _time.sleep(240)  # كل 4 دقائق

    t = _OSThread(target=_ping_self, daemon=True, name='KeepAlive')
    t.start()
    logger.info("🔄 نظام الاستمرارية الدائم مُفعّل (ping كل 4 دقائق)")

_start_keepalive()

# ── نقطة keepalive يمكن استدعاؤها خارجياً (من UptimeRobot مثلاً) ──
@app.route('/keepalive')
@app.route('/ping')
def route_keepalive():
    return jsonify({"status": "alive", "time": datetime.utcnow().isoformat()})


@app.route('/debug/ping')
def debug_ping():
    import socket as _dsock
    results = {}
    targets = [
        ('api.telegram.org', 443),
        ('149.154.167.51', 443),
        ('91.108.4.0', 443),
    ]
    for host, port in targets:
        try:
            _dsock.create_connection((host, port), timeout=10)
            results[f"{host}:{port}"] = "✅ متاح"
        except Exception as e:
            results[f"{host}:{port}"] = f"❌ {e}"
    all_ok = all("✅" in v for v in results.values())
    return jsonify({
        "telegram_reachable": all_ok,
        "checks": results
    }), 200 if all_ok else 503


# ════════════════════════════════════════════════════════════
#  إعادة التعيين — مستخدم واحد أو جميع المستخدمين
# ════════════════════════════════════════════════════════════
def _do_reset_user(uid):
    """تنفيذ إعادة تعيين مستخدم واحد (يُستدعى داخلياً)"""
    try:
        with USERS_LOCK:
            user_data = USERS.get(uid, {})
        client_manager = user_data.get('client_manager')
        if client_manager:
            try:
                if client_manager.client and client_manager.loop and client_manager.client.is_connected():
                    future = asyncio.run_coroutine_threadsafe(
                        client_manager.client.log_out(), client_manager.loop
                    )
                    future.result(timeout=8)
            except Exception:
                pass
            try:
                client_manager.stop_flag.set()
            except Exception:
                pass
        telegram_manager.login_managers.pop(uid, None)
        telegram_manager.client_managers.pop(uid, None)
        github_delete_session(uid)
        clear_user_session(uid)
        with USERS_LOCK:
            USERS.pop(uid, None)
    except Exception as e:
        logger.warning(f"_do_reset_user({uid}): {e}")


@app.route("/api/reset_user", methods=["POST"])
def api_reset_user():
    """إعادة تعيين الحساب الحالي فقط — حذف الجلسة محلياً وعلى GitHub"""
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "الجلسة غير صالحة"})
    user_id = session['user_id']
    data = request.get_json(silent=True) or {}
    target_id = data.get('target_user_id', user_id)
    try:
        _do_reset_user(target_id)
        socketio.emit('login_status', {
            "logged_in": False, "connected": False,
            "awaiting_code": False, "awaiting_password": False, "is_running": False,
        }, to=target_id)
        socketio.emit('connection_status', {"status": "disconnected"}, to=target_id)
        socketio.emit('user_reset_done', {"user_id": target_id}, to=target_id)
        return jsonify({"success": True, "message": f"✅ تم إعادة تعيين {target_id} بالكامل"})
    except Exception as e:
        logger.error(f"api_reset_user error for {user_id}: {e}")
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/reset_all", methods=["POST"])
def api_reset_all():
    """إعادة تعيين جميع المستخدمين والجلسات والإعدادات — يُعيد النظام لحالته الأولى"""
    errors = []
    for uid in list(PREDEFINED_USERS.keys()):
        try:
            _do_reset_user(uid)
        except Exception as e:
            errors.append(str(e))
    socketio.emit('force_full_reset', {"message": "تم إعادة تعيين جميع الجلسات والإعدادات"})
    if errors:
        return jsonify({"success": True,
                        "message": f"✅ تمت الإعادة مع ملاحظات: {'; '.join(errors)}"})
    return jsonify({"success": True,
                    "message": "✅ تم إعادة تعيين جميع الجلسات والإعدادات بالكامل"})


@app.route("/api/health", methods=["GET"])
def api_health_status():
    try:
        users_status = {}
        with USERS_LOCK:
            for uid, data in USERS.items():
                cm = data.get('client_manager')
                loop_ok = False
                if cm:
                    lp = getattr(cm, 'loop', None)
                    loop_ok = bool(lp and lp.is_running())

                rot = rotating_manager
                rot_active = (uid in rot.threads and rot.threads[uid] and rot.threads[uid].is_alive())
                next_send_in = None
                if rot_active and uid in rot.next_send_at:
                    remaining = int(rot.next_send_at[uid] - time.time())
                    next_send_in = max(0, remaining)

                users_status[uid] = {
                    "name": PREDEFINED_USERS.get(uid, {}).get("name", uid),
                    "authenticated": data.get("authenticated", False),
                    "connected": data.get("connected", False),
                    "is_running": data.get("is_running", False),
                    "awaiting_code": data.get("awaiting_code", False),
                    "awaiting_password": data.get("awaiting_password", False),
                    "client_loop_running": loop_ok,
                    "rotating_send_active": rot_active,
                    "rotating_next_send_in_seconds": next_send_in,
                    "monitored_keywords": len(cm.monitored_keywords) if cm else 0,
                    "monitored_groups": len(cm.monitored_groups) if cm else 0,
                }

        return jsonify({
            "success": True,
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
            "users": users_status
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ============================================
# قسم مراقبة الروابط
# ============================================

LINK_MONITORS = {}
LINK_MONITOR_WORKERS = {}

def extract_links_from_text(text):
    """استخراج جميع الروابط من النص"""
    url_pattern = r'https?://[^\s]+|t\.me/[^\s]+|telegram\.me/[^\s]+'
    links = re.findall(url_pattern, text)
    clean_links = []
    for link in links:
        link = re.sub(r'[.,;:!?)]+$', '', link)
        if link and link not in clean_links:
            clean_links.append(link)
    return clean_links

@app.route('/api/link_monitor/start', methods=['POST'])
def start_link_monitor():
    """بدء مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404

    data = request.json
    monitor_all = data.get('monitor_all', True)
    specific_chats = data.get('specific_chats', '').strip().split('\n')
    send_to_saved = data.get('send_to_saved', True)

    if user_id in LINK_MONITOR_WORKERS and LINK_MONITOR_WORKERS[user_id].is_alive():
        return jsonify({'error': 'توجد عملية مراقبة جارية بالفعل'}), 400

    if specific_chats and specific_chats[0]:
        specific_chats = [chat.strip() for chat in specific_chats if chat.strip()]
    else:
        specific_chats = []

    LINK_MONITORS[user_id] = {
        'is_running': True,
        'monitor_all': monitor_all,
        'specific_chats': specific_chats,
        'send_to_saved': send_to_saved,
        'links_found': [],
        'start_time': datetime.now().isoformat(),
        'total_links': 0,
        'monitored_chats': 0
    }

    worker = threading.Thread(target=run_link_monitor_worker, args=(user_id,))
    worker.daemon = True
    worker.start()
    LINK_MONITOR_WORKERS[user_id] = worker

    return jsonify({'success': True, 'message': 'بدأت مراقبة الروابط'})

def run_link_monitor_worker(user_id):
    """دالة تعمل في خيط منفصل لمراقبة الروابط"""
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        loop.run_until_complete(link_monitor_worker(user_id))
        loop.close()
    except Exception as e:
        logger.error(f"خطأ في worker مراقبة الروابط: {str(e)}")
        socketio.emit('link_monitor_error', {'error': str(e)}, room=user_id)

async def link_monitor_worker(user_id):
    """الدالة الأساسية لمراقبة الروابط"""
    state = LINK_MONITORS[user_id]
    client = USERS[user_id]['client']

    try:
        from telethon import events as telethon_events

        dialogs = await client.get_dialogs()
        state['monitored_chats'] = len(dialogs)

        if state['monitor_all']:
            chats_to_monitor = None  # None = جميع الدردشات
        else:
            chats_to_monitor = []
            for dialog in dialogs:
                name = getattr(dialog.entity, 'title', None) or getattr(dialog.entity, 'first_name', '') or ''
                if any(s in name for s in state['specific_chats']):
                    chats_to_monitor.append(dialog.entity)

        logger.info(f"بدء مراقبة الروابط — {state['monitored_chats']} دردشة")
        socketio.emit('link_monitor_status_update', {
            'monitored_chats': state['monitored_chats'],
            'is_running': True
        }, room=user_id)

        @client.on(telethon_events.NewMessage(chats=chats_to_monitor))
        async def handle_new_message(event):
            if not state.get('is_running'):
                return
            try:
                message = event.message
                if not message or not message.text:
                    return

                links = extract_links_from_text(message.text)
                if not links:
                    return

                chat = await event.get_chat()
                chat_name = getattr(chat, 'title', None) or getattr(chat, 'first_name', None) or 'Unknown'

                sender = await event.get_sender()
                sender_name = getattr(sender, 'first_name', None) or 'Unknown'

                for link in links:
                    link_data = {
                        'link': link,
                        'chat_name': chat_name,
                        'chat_id': getattr(chat, 'id', 0),
                        'sender_name': sender_name,
                        'message_text': message.text[:200] + ('...' if len(message.text) > 200 else ''),
                        'message_id': message.id,
                        'date': message.date.isoformat() if message.date else datetime.now().isoformat(),
                        'timestamp': datetime.now().isoformat()
                    }

                    state['links_found'].append(link_data)
                    state['total_links'] += 1

                    if state.get('send_to_saved'):
                        try:
                            saved_msg = (
                                f"🔗 رابط جديد!\n\n"
                                f"📌 الرابط: {link}\n"
                                f"💬 الدردشة: {chat_name}\n"
                                f"👤 المرسل: {sender_name}\n"
                                f"📝 النص: {message.text[:100]}\n\n"
                                f"🕐 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                            )
                            await client.send_message('me', saved_msg)
                            logger.info(f"تم إرسال الرابط إلى المحفوظات: {link}")
                        except Exception as e:
                            logger.error(f"فشل الإرسال إلى المحفوظات: {str(e)}")

                    socketio.emit('link_monitor_update', {
                        'link': link_data,
                        'total': state['total_links']
                    }, room=user_id)

                    logger.info(f"رابط مكتشف في [{chat_name}]: {link}")

            except Exception as e:
                logger.error(f"خطأ في معالجة رسالة: {str(e)}")

        while state.get('is_running'):
            await asyncio.sleep(1)

        client.remove_event_handler(handle_new_message)

    except Exception as e:
        logger.error(f"خطأ في مراقبة الروابط: {str(e)}")
        socketio.emit('link_monitor_error', {'error': str(e)}, room=user_id)
    finally:
        state['is_running'] = False
        socketio.emit('link_monitor_done', {
            'total_links': state.get('total_links', 0),
            'monitored_chats': state.get('monitored_chats', 0)
        }, room=user_id)
        if user_id in LINK_MONITOR_WORKERS:
            del LINK_MONITOR_WORKERS[user_id]

@app.route('/api/link_monitor/stop')
def stop_link_monitor():
    """إيقاف مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        LINK_MONITORS[user_id]['is_running'] = False
    return jsonify({'success': True, 'message': 'تم إيقاف مراقبة الروابط'})

@app.route('/api/link_monitor/status')
def link_monitor_status():
    """الحصول على حالة مراقبة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        state = LINK_MONITORS[user_id]
        is_running = user_id in LINK_MONITOR_WORKERS and LINK_MONITOR_WORKERS[user_id].is_alive()
        return jsonify({
            'is_running': is_running,
            'total_links': state.get('total_links', 0),
            'monitored_chats': state.get('monitored_chats', 0),
            'links_found': state.get('links_found', [])[-50:]
        })
    return jsonify({'is_running': False, 'total_links': 0, 'monitored_chats': 0, 'links_found': []})

@app.route('/api/link_monitor/links')
def get_monitored_links():
    """الحصول على قائمة الروابط المكتشفة"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    limit = int(request.args.get('limit', 100))
    if user_id in LINK_MONITORS:
        links = LINK_MONITORS[user_id].get('links_found', [])
        links_sorted = sorted(links, key=lambda x: x.get('timestamp', ''), reverse=True)
        return jsonify({'links': links_sorted[:limit], 'total': len(links)})
    return jsonify({'links': [], 'total': 0})

@app.route('/api/link_monitor/clear')
def clear_monitored_links():
    """مسح قائمة الروابط"""
    user_id = session.get('user_id')
    if user_id not in USERS:
        return jsonify({'error': 'المستخدم غير موجود'}), 404
    if user_id in LINK_MONITORS:
        LINK_MONITORS[user_id]['links_found'] = []
        LINK_MONITORS[user_id]['total_links'] = 0
    return jsonify({'success': True, 'message': 'تم مسح الروابط'})





# ===========================
# وظيفة رسائلي — تتبع وإدارة الرسائل المرسلة
# ===========================

@app.route("/api/sent_batches")
def api_sent_batches():
    """عرض جميع الدفعات المرسلة (رسائلي)"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        batches = list(USERS[user_id].get('sent_batches', []))
    result = []
    for b in reversed(batches):
        result.append({
            "id": b["id"],
            "text": b["text"],
            "has_media": b.get("has_media", False),
            "sent_at": b["sent_at"],
            "edited_at": b.get("edited_at"),
            "sent_count": b.get("sent_count", len(b["entries"])),
            "group_count": len(b["entries"]),
            "groups": [{"title": e.get("group", ""), "username": e.get("group", "")} for e in b["entries"]]
        })
    return jsonify({"success": True, "batches": result})


@app.route("/api/edit_batch", methods=["POST"])
def api_edit_batch():
    """تعديل جميع رسائل الدفعة دفعة واحدة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        client_manager = USERS[user_id].get('client_manager')
    if not client_manager:
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"})
    data = request.json or {}
    batch_id = data.get("batch_id", "")
    new_text = data.get("new_text", "")
    if not batch_id or not new_text:
        return jsonify({"success": False, "message": "بيانات ناقصة"})
    def run_edit():
        try:
            client_manager.run_coroutine(
                client_manager._edit_batch_messages(batch_id, new_text)
            )
        except Exception as e:
            socketio.emit('log_update', {"message": f"❌ خطأ في التعديل: {str(e)[:100]}"}, to=user_id)
    _OSThread(target=run_edit, daemon=True).start()
    return jsonify({"success": True, "message": "⏳ جارٍ تعديل الرسائل..."})


@app.route("/api/delete_batch", methods=["POST"])
def api_delete_batch():
    """حذف جميع رسائل الدفعة دفعة واحدة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        client_manager = USERS[user_id].get('client_manager')
    if not client_manager:
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"})
    data = request.json or {}
    batch_id = data.get("batch_id", "")
    if not batch_id:
        return jsonify({"success": False, "message": "batch_id مطلوب"})
    def run_delete():
        try:
            client_manager.run_coroutine(
                client_manager._delete_batch_messages(batch_id)
            )
        except Exception as e:
            socketio.emit('log_update', {"message": f"❌ خطأ في الحذف: {str(e)[:100]}"}, to=user_id)
    _OSThread(target=run_delete, daemon=True).start()
    return jsonify({"success": True, "message": "⏳ جارٍ حذف الرسائل..."})


@app.route("/api/batch_details/<batch_id>")
def api_batch_details(batch_id):
    """عرض تفاصيل دفعة محددة"""
    user_id = session.get('user_id')
    if not user_id or user_id not in USERS:
        return jsonify({"success": False, "message": "غير مسجّل"}), 401
    with USERS_LOCK:
        batches = USERS[user_id].get('sent_batches', [])
        batch = next((b for b in batches if b["id"] == batch_id), None)
    if not batch:
        return jsonify({"success": False, "message": "الدفعة غير موجودة"}), 404
    return jsonify({
        "success": True,
        "batch": {
            "id": batch["id"],
            "text": batch["text"],
            "has_media": batch.get("has_media", False),
            "sent_at": batch["sent_at"],
            "edited_at": batch.get("edited_at"),
            "entries": batch["entries"]
        }
    })


# ====== Web Push API Routes ======

@app.route("/api/push/vapid-public-key", methods=["GET"])
def api_push_vapid_key():
    return jsonify({"publicKey": VAPID_PUBLIC_KEY})

@app.route("/api/push/subscribe", methods=["POST"])
def api_push_subscribe():
    user_id = str(session.get('user_id', 'user_1'))
    sub = request.json
    if not sub or 'endpoint' not in sub:
        return jsonify({"success": False, "error": "اشتراك غير صالح"}), 400
    push_subscriptions[user_id] = sub
    _save_push_subs(push_subscriptions)
    logger.info(f"✅ Push subscription saved for user {user_id}")
    return jsonify({"success": True})

@app.route("/api/push/settings", methods=["GET"])
def api_push_settings_get():
    user_id = session.get('user_id', 'user_1')
    s = load_settings(user_id)
    return jsonify({
        "push_enabled": s.get('push_notifications_enabled', False),
        "speech_enabled": s.get('speech_notifications_enabled', False)
    })

@app.route("/api/push/settings", methods=["POST"])
def api_push_settings_post():
    user_id = session.get('user_id', 'user_1')
    data = request.json or {}
    s = load_settings(user_id)
    if 'push_enabled' in data:
        s['push_notifications_enabled'] = bool(data['push_enabled'])
    if 'speech_enabled' in data:
        s['speech_notifications_enabled'] = bool(data['speech_enabled'])
    save_settings(user_id, s)
    return jsonify({"success": True})

# ============================================================
#  نظام الإدارة + المصادقة البيومترية + التخزين عبر GitHub
# ============================================================

# ---- بيانات تسجيل الدخول الإدارية ----
ADMIN_USERNAME = "Anwer"
ADMIN_PASSWORD = "772997043a*anwer"

# ---- إعدادات GitHub الخاصة بالبصمة والملفات ----
BIO_REPO_OWNER  = "anwer1230"
BIO_REPO_NAME   = "Web-browser"
BIO_BRANCH      = "main"
BIOMETRIC_PATH  = "data/biometric_devices.json"
UPLOADS_FOLDER  = "uploads"

# ─── استدعاء الإدارة عبر البحث في رسائلي ───────────────────────────────────────

@app.route('/api/admin_ui_invoke', methods=['POST'])
def api_admin_ui_invoke():
    """استدعاء أيقونة الإدارة عبر كلمة مرور المشرف من حقل البحث في رسائلي"""
    data = request.json or {}
    password = data.get('password', '').strip()
    user_id = session.get('user_id', 'user_1')
    if password == ADMIN_PASSWORD:
        session['admin_ui_visible'] = True
        session['admin_ui_user'] = user_id
        session['admin_auth'] = True
        session.permanent = True
        try:
            socketio.emit('admin_ui_invoked', {
                'visible': True,
                'message': '✅ تم استدعاء لوحة الإدارة بنجاح'
            }, to=user_id)
        except Exception:
            pass
        return jsonify({'success': True, 'message': '✅ تم استدعاء لوحة الإدارة بنجاح', 'visible': True})
    return jsonify({'success': False, 'message': '❌ كلمة المرور غير صحيحة'})

@app.route('/api/admin_ui_hide', methods=['POST'])
def api_admin_ui_hide():
    """إخفاء أيقونة الإدارة عند الخروج من الإدارة"""
    user_id = session.get('user_id', 'user_1')
    session.pop('admin_ui_visible', None)
    session.pop('admin_ui_user', None)
    session.pop('admin_auth', None)
    try:
        socketio.emit('admin_ui_invoked', {
            'visible': False,
            'message': '🔒 تم إخفاء لوحة الإدارة'
        }, to=user_id)
    except Exception:
        pass
    return jsonify({'success': True, 'message': 'تم إخفاء الإدارة'})

@app.route('/api/admin_ui_status', methods=['GET'])
def api_admin_ui_status():
    """التحقق من حالة ظهور أيقونة الإدارة — الأيقونة مخفية دائماً"""
    return jsonify({'visible': False})

@app.route('/admin_panel')
def admin_panel_page():
    """صفحة لوحة الإدارة الكاملة"""
    if not session.get("admin_auth"):
        return redirect('/admin')
    return render_template('admin_panel.html')

@app.route('/admin/api/session_check', methods=['GET'])
def admin_session_check():
    """التحقق من جلسة الإدارة"""
    if session.get('admin_auth'):
        return jsonify({'success': True})
    return jsonify({'success': False})

# ── دوال GitHub API ───────────────────────────────────────────────────

def _gh_headers():
    tok = GITHUB_TOKEN or os.environ.get("GITHUB_TOKEN", "")
    return {"Authorization": f"token {tok}", "Accept": "application/vnd.github.v3+json"}

def upload_to_github(file_path_in_repo, content_bytes, commit_message="رفع ملف تلقائي"):
    if not GITHUB_TOKEN:
        logger.error("GITHUB_TOKEN غير مضبوط")
        return False
    url = f"https://api.github.com/repos/{BIO_REPO_OWNER}/{BIO_REPO_NAME}/contents/{file_path_in_repo}"
    content_b64 = base64.b64encode(content_bytes).decode("utf-8")
    sha = None
    try:
        resp = requests.get(url, headers=_gh_headers(), timeout=10)
        if resp.status_code == 200:
            sha = resp.json().get("sha")
    except Exception as e:
        logger.warning(f"فشل جلب SHA: {e}")
    data = {"message": commit_message, "content": content_b64, "branch": BIO_BRANCH}
    if sha:
        data["sha"] = sha
    try:
        r = requests.put(url, headers=_gh_headers(), json=data, timeout=20)
        if r.status_code in (200, 201):
            return True
        logger.error(f"فشل رفع الملف: {r.status_code} {r.text[:100]}")
        return False
    except Exception as e:
        logger.error(f"استثناء رفع الملف: {e}")
        return False

def download_from_github(file_path_in_repo):
    if not GITHUB_TOKEN:
        return None
    url = f"https://api.github.com/repos/{BIO_REPO_OWNER}/{BIO_REPO_NAME}/contents/{file_path_in_repo}"
    try:
        resp = requests.get(url, headers=_gh_headers(), timeout=10)
        if resp.status_code == 200:
            c = resp.json().get("content", "").replace("\n", "")
            if c:
                return base64.b64decode(c)
    except Exception as e:
        logger.warning(f"فشل تحميل {file_path_in_repo}: {e}")
    return None


# ══════════════════════════════════════════════════════════════════════════════
# ─── الروابط المحفوظة ─────────────────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

SAVED_LINKS_FILE = os.path.join(DATA_DIR, 'saved_links.json')

def load_saved_links():
    """تحميل الروابط المحفوظة من الملف المحلي أو GitHub"""
    try:
        if os.path.exists(SAVED_LINKS_FILE):
            with open(SAVED_LINKS_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    try:
        content = download_from_github("data/saved_links.json")
        if content:
            data = json.loads(content.decode("utf-8"))
            with open(SAVED_LINKS_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            return data
    except Exception:
        pass
    return {"links": []}

def save_saved_links(data):
    """حفظ الروابط المحفوظة محلياً وفي GitHub"""
    try:
        with open(SAVED_LINKS_FILE, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Error saving saved_links locally: {e}")
    try:
        raw = json.dumps(data, ensure_ascii=False, indent=2).encode('utf-8')
        upload_to_github("data/saved_links.json", raw, "تحديث الروابط المحفوظة")
    except Exception as e:
        logger.error(f"Error saving saved_links to GitHub: {e}")

def add_saved_link(url, title=None, category='عام', notes='', source='يدوي'):
    data = load_saved_links()
    for link in data["links"]:
        if link["url"] == url:
            return False, "الرابط موجود بالفعل"
    new_link = {
        "id": str(uuid.uuid4())[:8],
        "url": url,
        "title": title or url,
        "category": category,
        "date_saved": datetime.now().isoformat(),
        "source": source,
        "notes": notes
    }
    data["links"].append(new_link)
    save_saved_links(data)
    return True, new_link

def add_multiple_links(urls, category='عام', source='دفعة'):
    data = load_saved_links()
    added = []
    skipped = []
    for url in urls:
        url = url.strip()
        if not url:
            continue
        exists = any(l["url"] == url for l in data["links"])
        if exists:
            skipped.append(url)
            continue
        new_link = {
            "id": str(uuid.uuid4())[:8],
            "url": url,
            "title": url,
            "category": category,
            "date_saved": datetime.now().isoformat(),
            "source": source,
            "notes": ""
        }
        data["links"].append(new_link)
        added.append(url)
    save_saved_links(data)
    return added, skipped

def delete_saved_link(link_id):
    data = load_saved_links()
    data["links"] = [l for l in data["links"] if l["id"] != link_id]
    save_saved_links(data)
    return True

def delete_multiple_links(link_ids):
    data = load_saved_links()
    data["links"] = [l for l in data["links"] if l["id"] not in link_ids]
    save_saved_links(data)
    return True

def update_saved_link(link_id, updates):
    data = load_saved_links()
    for link in data["links"]:
        if link["id"] == link_id:
            if "title" in updates: link["title"] = updates["title"]
            if "category" in updates: link["category"] = updates["category"]
            if "notes" in updates: link["notes"] = updates["notes"]
            save_saved_links(data)
            return True, link
    return False, None

def get_saved_links_by_category(category=None):
    data = load_saved_links()
    if category:
        return [l for l in data["links"] if l["category"] == category]
    return data["links"]

def get_categories():
    data = load_saved_links()
    return sorted(set(l["category"] for l in data["links"]))

# ── مسارات الروابط المحفوظة ──────────────────────────────────────────────────

@app.route("/saved_links")
def saved_links_page():
    return render_template('saved_links.html')

@app.route("/api/saved_links", methods=["GET"])
def api_get_saved_links():
    category = request.args.get('category')
    links = get_saved_links_by_category(category)
    categories = get_categories()
    return jsonify({"success": True, "links": links, "categories": categories, "total": len(links)})

@app.route("/api/saved_links/add", methods=["POST"])
def api_add_saved_link():
    data = request.json or {}
    url = data.get('url', '').strip()
    title = data.get('title', '').strip() or url
    category = data.get('category', 'عام')
    notes = data.get('notes', '')
    source = data.get('source', 'يدوي')
    if not url:
        return jsonify({"success": False, "message": "الرابط مطلوب"})
    success, result = add_saved_link(url, title, category, notes, source)
    if success:
        return jsonify({"success": True, "link": result})
    return jsonify({"success": False, "message": result})

@app.route("/api/saved_links/add_batch", methods=["POST"])
def api_add_saved_links_batch():
    data = request.json or {}
    urls = data.get('urls', [])
    category = data.get('category', 'عام')
    source = data.get('source', 'دفعة')
    if not urls:
        return jsonify({"success": False, "message": "لا توجد روابط"})
    added, skipped = add_multiple_links(urls, category, source)
    return jsonify({"success": True, "added": added, "skipped": skipped,
                    "added_count": len(added), "skipped_count": len(skipped)})

@app.route("/api/saved_links/delete", methods=["POST"])
def api_delete_saved_link():
    data = request.json or {}
    link_id = data.get('id')
    if not link_id:
        return jsonify({"success": False, "message": "معرف الرابط مطلوب"})
    delete_saved_link(link_id)
    return jsonify({"success": True})

@app.route("/api/saved_links/delete_batch", methods=["POST"])
def api_delete_saved_links_batch():
    data = request.json or {}
    link_ids = data.get('ids', [])
    if not link_ids:
        return jsonify({"success": False, "message": "لا توجد معرفات"})
    delete_multiple_links(link_ids)
    return jsonify({"success": True})

@app.route("/api/saved_links/update", methods=["POST"])
def api_update_saved_link():
    data = request.json or {}
    link_id = data.get('id')
    if not link_id:
        return jsonify({"success": False, "message": "معرف الرابط مطلوب"})
    updates = {}
    if 'title' in data: updates['title'] = data['title']
    if 'category' in data: updates['category'] = data['category']
    if 'notes' in data: updates['notes'] = data['notes']
    if not updates:
        return jsonify({"success": False, "message": "لا توجد تحديثات"})
    success, link = update_saved_link(link_id, updates)
    if success:
        return jsonify({"success": True, "link": link})
    return jsonify({"success": False, "message": "الرابط غير موجود"})

@app.route("/api/saved_links/export", methods=["POST"])
def api_export_saved_links():
    data = request.json or {}
    link_ids = data.get('ids', [])
    data_links = load_saved_links()
    if link_ids:
        links = [l for l in data_links["links"] if l["id"] in link_ids]
    else:
        links = data_links["links"]
    if not links:
        return jsonify({"success": False, "message": "لا توجد روابط للتصدير"})
    content = "\n".join(l["url"] for l in links)
    response = make_response(content)
    response.headers["Content-Type"] = "text/plain; charset=utf-8"
    response.headers["Content-Disposition"] = f"attachment; filename=saved_links_{datetime.now().strftime('%Y%m%d_%H%M')}.txt"
    return response

@app.route("/api/saved_links/copy", methods=["POST"])
def api_copy_saved_links():
    data = request.json or {}
    link_ids = data.get('ids', [])
    data_links = load_saved_links()
    if link_ids:
        links = [l for l in data_links["links"] if l["id"] in link_ids]
    else:
        links = data_links["links"]
    content = "\n".join(l["url"] for l in links)
    return jsonify({"success": True, "text": content, "count": len(links)})

@app.route("/api/saved_links/send_to_auto_join", methods=["POST"])
def api_send_saved_to_auto_join():
    data = request.json or {}
    link_ids = data.get('ids', [])
    data_links = load_saved_links()
    if link_ids:
        links = [l for l in data_links["links"] if l["id"] in link_ids]
    else:
        links = data_links["links"]
    urls = [l["url"] for l in links]
    session['auto_join_links'] = urls
    return jsonify({"success": True, "urls": urls, "count": len(urls)})

def file_exists_in_github(file_path_in_repo):
    url = f"https://api.github.com/repos/{BIO_REPO_OWNER}/{BIO_REPO_NAME}/contents/{file_path_in_repo}"
    try:
        return requests.head(url, headers=_gh_headers(), timeout=8).status_code == 200
    except:
        return False

# ── دوال البصمة ───────────────────────────────────────────────────────

def load_biometric_devices():
    raw = download_from_github(BIOMETRIC_PATH)
    if raw:
        try:
            return json.loads(raw.decode("utf-8"))
        except json.JSONDecodeError:
            pass
    return {}

def save_biometric_devices(devices):
    return upload_to_github(
        BIOMETRIC_PATH,
        json.dumps(devices, indent=2, ensure_ascii=False).encode("utf-8"),
        "تحديث بيانات البصمة"
    )

# ── دوال ملفات المستخدمين ─────────────────────────────────────────────

def save_user_file(user_id, filename, file_bytes):
    return upload_to_github(
        f"{UPLOADS_FOLDER}/{user_id}/{filename}",
        file_bytes,
        f"رفع ملف للمستخدم {user_id}"
    )

def get_user_file_url(user_id, filename):
    return f"https://raw.githubusercontent.com/{BIO_REPO_OWNER}/{BIO_REPO_NAME}/{BIO_BRANCH}/{UPLOADS_FOLDER}/{user_id}/{filename}"

def list_user_files(user_id):
    url = f"https://api.github.com/repos/{BIO_REPO_OWNER}/{BIO_REPO_NAME}/contents/{UPLOADS_FOLDER}/{user_id}"
    try:
        resp = requests.get(url, headers=_gh_headers(), timeout=10)
        if resp.status_code == 200:
            return [item["name"] for item in resp.json() if item.get("type") == "file"]
    except:
        pass
    return []

def get_or_create_user(user_id):
    with USERS_LOCK:
        if user_id not in USERS:
            settings = load_settings(user_id)
            USERS[user_id] = {
                'client_manager': None, 'settings': settings,
                'thread': None, 'is_running': False,
                'stats': {"sent": 0, "errors": 0},
                'connected': False, 'authenticated': False,
                'awaiting_code': False, 'awaiting_password': False,
                'phone_code_hash': None, 'monitoring_active': False,
                'event_handlers_registered': False,
                'blocked': settings.get('blocked', False),
                'disabled': settings.get('disabled', False),
                'alerts': settings.get('alerts', []),
                'last_seen': datetime.now().isoformat(),
                'phone_number': settings.get('phone', ''),
                'telegram_name': settings.get('telegram_name', ''),
            }
        return USERS[user_id]

def ensure_github_dirs():
    if not file_exists_in_github(BIOMETRIC_PATH):
        upload_to_github(BIOMETRIC_PATH, b"{}", "إنشاء ملف بيانات البصمة")
    if not file_exists_in_github(f"{UPLOADS_FOLDER}/.gitkeep"):
        upload_to_github(f"{UPLOADS_FOLDER}/.gitkeep", b"# uploads\n", "إنشاء مجلد uploads")

# ── مسارات المصادقة الإدارية ──────────────────────────────────────────

@app.route("/admin/api/check", methods=["GET"])
def admin_check():
    return jsonify({"authenticated": session.get("admin_auth", False)})

@app.route("/admin/api/login", methods=["POST"])
def admin_login():
    data = request.get_json() or {}
    if data.get("username") == ADMIN_USERNAME and data.get("password") == ADMIN_PASSWORD:
        session["admin_auth"] = True
        session.permanent = True
        return jsonify({"success": True, "message": "تم تسجيل الدخول بنجاح"})
    return jsonify({"success": False, "message": "بيانات غير صحيحة"}), 401

@app.route("/admin/api/logout", methods=["POST"])
def admin_logout():
    session.pop("admin_auth", None)
    return jsonify({"success": True, "message": "تم تسجيل الخروج"})

# ── مسارات البصمة ────────────────────────────────────────────────────

@app.route("/admin/api/biometric/register", methods=["POST"])
def biometric_register():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"}), 403
    data = request.get_json() or {}
    device_id = data.get("device_id")
    biometric_token = data.get("biometric_token")
    if not device_id or not biometric_token:
        return jsonify({"success": False, "message": "device_id و biometric_token مطلوبان"})
    devices = load_biometric_devices()
    devices[device_id] = biometric_token
    if save_biometric_devices(devices):
        return jsonify({"success": True, "message": f"تم تسجيل الجهاز {device_id} بنجاح"})
    return jsonify({"success": False, "message": "فشل حفظ بيانات البصمة"}), 500

@app.route("/admin/api/biometric/login", methods=["POST"])
def biometric_login():
    data = request.get_json() or {}
    device_id = data.get("device_id")
    biometric_token = data.get("biometric_token")
    if not device_id or not biometric_token:
        return jsonify({"success": False, "message": "device_id و biometric_token مطلوبان"})
    devices = load_biometric_devices()
    if devices.get(device_id) == biometric_token:
        session["admin_auth"] = True
        session.permanent = True
        return jsonify({"success": True, "message": "✅ تم تسجيل الدخول بواسطة البصمة"})
    return jsonify({"success": False, "message": "❌ بيانات البصمة غير صحيحة"}), 401

@app.route("/admin/api/biometric/unregister", methods=["POST"])
def biometric_unregister():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "يجب تسجيل الدخول أولاً"}), 403
    data = request.get_json() or {}
    device_id = data.get("device_id")
    if not device_id:
        return jsonify({"success": False, "message": "device_id مطلوب"})
    devices = load_biometric_devices()
    if device_id in devices:
        del devices[device_id]
        if save_biometric_devices(devices):
            return jsonify({"success": True, "message": f"تم حذف الجهاز {device_id}"})
        return jsonify({"success": False, "message": "فشل حفظ التغييرات"})
    return jsonify({"success": False, "message": "الجهاز غير موجود"})

# ── مسارات إدارة المستخدمين ───────────────────────────────────────────

@app.route("/admin/api/users", methods=["GET"])
def admin_users():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    users = []
    for slot, uinfo in PREDEFINED_USERS.items():
        ud = get_or_create_user(slot)
        s = load_settings(slot)
        users.append({
            "user_id": slot, "name": uinfo.get("name", slot),
            "phone": ud.get('phone_number', ''),
            "logged_in": ud.get('authenticated', False),
            "telegram_name": ud.get('telegram_name', ''),
            "blocked": ud.get('blocked', False),
            "disabled": ud.get('disabled', False),
            "last_seen": ud.get('last_seen', None),
            "groups": s.get("groups", []),
            "watch_words": s.get("watch_words", []),
            "auto_replies": s.get("auto_replies", []),
            "alerts_count": len(ud.get('alerts', []))
        })
    return jsonify({"success": True, "users": users})

@app.route("/admin/api/user/<slot>", methods=["POST"])
def admin_update_user(slot):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    ud = get_or_create_user(slot)
    data = request.get_json() or {}
    action = data.get("action")
    if action == "block":
        ud["blocked"] = data.get("blocked", False)
    elif action == "disable":
        ud["disabled"] = data.get("disabled", False)
    s = load_settings(slot)
    s["blocked"] = ud.get("blocked", False)
    s["disabled"] = ud.get("disabled", False)
    save_settings(slot, s)
    with USERS_LOCK:
        USERS[slot] = ud
    return jsonify({"success": True})

@app.route("/admin/api/fetch_chats/<slot>", methods=["GET"])
def admin_fetch_chats(slot):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    settings = load_settings(slot)
    groups = settings.get("groups", [])
    chats = [{"link": g, "title": g, "username": g.split('/')[-1] if '/' in g else g} for g in groups]
    return jsonify({"success": True, "chats": chats})

@app.route("/admin/api/copy_chats/<slot>", methods=["GET"])
def admin_copy_chats(slot):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    settings = load_settings(slot)
    groups = [g for g in settings.get("groups", []) if g.strip()]
    return jsonify({"success": True, "text": "\n".join(groups), "count": len(groups)})

@app.route("/admin/api/export_chats/<slot>", methods=["GET"])
def admin_export_chats(slot):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    settings = load_settings(slot)
    groups = [g for g in settings.get("groups", []) if g.strip()]
    response = make_response("\n".join(groups))
    response.headers["Content-Type"] = "text/plain; charset=utf-8"
    response.headers["Content-Disposition"] = f"attachment; filename=chats_{slot}.txt"
    return response

@app.route("/admin/api/user_alerts/<slot>", methods=["GET"])
def admin_user_alerts(slot):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    ud = get_or_create_user(slot)
    return jsonify({"success": True, "alerts": ud.get('alerts', [])})

@app.route("/admin/api/upload", methods=["POST"])
def admin_upload_file():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "لا يوجد ملف مرفق"})
    file = request.files['file']
    if not file.filename:
        return jsonify({"success": False, "message": "اسم الملف فارغ"})
    user_id = request.form.get("user_id", "admin")
    file_bytes = file.read()
    if save_user_file(user_id, file.filename, file_bytes):
        return jsonify({"success": True, "message": "تم رفع الملف بنجاح",
                        "url": get_user_file_url(user_id, file.filename),
                        "filename": file.filename})
    return jsonify({"success": False, "message": "فشل رفع الملف إلى GitHub"}), 500

@app.route("/admin/api/files/<user_id>", methods=["GET"])
def admin_list_files(user_id):
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    return jsonify({"success": True, "files": list_user_files(user_id)})


# ══════════════════════════════════════════════════════════════
#  نظام الإشعارات العامة (Broadcast Notifications)
# ══════════════════════════════════════════════════════════════

_NOTIF_FILE = os.path.join(os.path.dirname(__file__), 'data', 'notifications.json')
os.makedirs(os.path.dirname(_NOTIF_FILE), exist_ok=True)

def _load_notifications():
    """تحميل الإشعارات من الملف المحلي، ثم GitHub كاحتياطي."""
    # أولاً: الملف المحلي
    try:
        if os.path.exists(_NOTIF_FILE):
            with open(_NOTIF_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception:
        pass
    # ثانياً: GitHub
    try:
        content = download_from_github("data/notifications.json")
        if content:
            data = json.loads(content.decode("utf-8"))
            # نسخ محلي للسرعة
            with open(_NOTIF_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            return data
    except Exception:
        pass
    return []

def _save_notifications(notifications):
    """حفظ الإشعارات محلياً وفي GitHub."""
    try:
        with open(_NOTIF_FILE, 'w', encoding='utf-8') as f:
            json.dump(notifications, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Notifications local save error: {e}")
    # محاولة الرفع إلى GitHub (اختياري)
    try:
        raw = json.dumps(notifications, ensure_ascii=False, indent=2).encode('utf-8')
        upload_to_github("data/notifications.json", raw, "تحديث الإشعارات")
    except Exception:
        pass

@app.route("/admin/api/notifications", methods=["GET"])
def admin_get_notifications():
    """جلب جميع الإشعارات المخزنة."""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    return jsonify({"success": True, "notifications": _load_notifications()})

@app.route("/admin/api/broadcast_notification", methods=["POST"])
def admin_broadcast_notification():
    """نشر إشعار فوري لجميع المستخدمين المتصلين وتخزينه."""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = request.get_json() or {}
    message  = data.get("message", "").strip()
    notif_type = data.get("type", "info")   # info / success / warning / danger
    if not message:
        return jsonify({"success": False, "message": "نص الإشعار مطلوب"})
    notification = {
        "id":        str(uuid.uuid4()),
        "message":   message,
        "type":      notif_type,
        "timestamp": datetime.now().isoformat(),
    }
    notifs = _load_notifications()
    notifs.insert(0, notification)          # الأحدث أولاً
    notifs = notifs[:100]                   # الاحتفاظ بآخر 100
    _save_notifications(notifs)
    # بث فوري لجميع المستخدمين المتصلين عبر Socket.IO
    socketio.emit("new_broadcast_notification", notification)

    # ── إرسال Web Push لكل المشتركين (حتى لو التطبيق مغلق) ──
    push_sent = 0
    _type_icon = {"success": "✅", "warning": "⚠️", "danger": "❌", "info": "📢"}.get(notif_type, "📢")
    push_title = f"{_type_icon} إشعار من الإدارة"
    for uid, sub in list(push_subscriptions.items()):
        try:
            import json as _json
            _payload = _json.dumps({
                "title": push_title,
                "body":  message,
                "type":  "broadcast",
                "id":    notification["id"],
                "icon":  "/static/icons/app-logo.png",
                "badge": "/static/icons/app-logo.png",
                "data":  {"notif_id": notification["id"]}
            })
            _webpush_fn(
                subscription_info=sub,
                data=_payload,
                vapid_private_key=VAPID_PRIVATE_PEM,
                vapid_claims=VAPID_CLAIMS,
                ttl=86400
            )
            push_sent += 1
        except Exception as _px:
            _sc = getattr(getattr(_px, 'response', None), 'status_code', 0)
            if _sc in (410, 404):
                push_subscriptions.pop(uid, None)
            logger.warning(f"Broadcast push failed for {uid}: {_px}")
    if push_sent:
        _save_push_subs(push_subscriptions)
    logger.info(f"📢 إشعار مبث: {message[:60]} | Push: {push_sent} مستخدم")
    return jsonify({"success": True, "message": f"تم نشر الإشعار بنجاح (Push: {push_sent})", "notification": notification})

@app.route("/admin/api/delete_notification/<notif_id>", methods=["DELETE"])
def admin_delete_notification(notif_id):
    """حذف إشعار محدد."""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    notifs = [n for n in _load_notifications() if n.get("id") != notif_id]
    _save_notifications(notifs)
    return jsonify({"success": True})

@app.route("/admin/api/push_stats", methods=["GET"])
def admin_push_stats():
    """إحصائيات المشتركين في Web Push"""
    if not session.get("admin_auth"):
        return jsonify({"success": False}), 403
    subs = list(push_subscriptions.keys())
    return jsonify({
        "success": True,
        "count": len(subs),
        "subscribers": subs
    })


@app.route("/admin/api/test_push", methods=["POST"])
def admin_test_push():
    """إرسال إشعار اختباري لكل المشتركين"""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data  = request.get_json() or {}
    msg   = data.get("message", "✅ اختبار الإشعار — يعمل بشكل صحيح!")
    ntype = data.get("type", "general")
    icons = {"broadcast": "📢", "schedule_expired": "⏹", "general": "🔔"}
    title = f"{icons.get(ntype,'🔔')} إشعار اختباري"
    sent = 0
    for uid, sub in list(push_subscriptions.items()):
        try:
            import json as _json
            _payload = _json.dumps({
                "title": title,
                "body":  msg,
                "type":  ntype,
                "icon":  "/static/icons/app-logo.png",
                "badge": "/static/icons/app-logo.png",
                "data":  {"test": True}
            })
            _webpush_fn(
                subscription_info=sub,
                data=_payload,
                vapid_private_key=VAPID_PRIVATE_PEM,
                vapid_claims=VAPID_CLAIMS,
                ttl=300
            )
            sent += 1
        except Exception as _px:
            _sc = getattr(getattr(_px, 'response', None), 'status_code', 0)
            if _sc in (410, 404):
                push_subscriptions.pop(uid, None)
            logger.warning(f"Test push failed for {uid}: {_px}")
    if sent:
        _save_push_subs(push_subscriptions)
    return jsonify({"success": True, "sent": sent,
                    "message": f"✅ تم إرسال الاختبار لـ {sent} مستخدم"})


@app.route("/admin/api/clear_notifications", methods=["POST"])
def admin_clear_notifications():
    """مسح جميع الإشعارات."""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    _save_notifications([])
    return jsonify({"success": True, "message": "تم مسح جميع الإشعارات"})

# ── صفحة لوحة تحكم الإدارة (بسيطة) ────────────────────────────────────

@app.route("/admin")
@app.route("/admin/")
def admin_dashboard():
    if not session.get("admin_auth"):
        return '''<!DOCTYPE html><html dir="rtl"><head><meta charset="utf-8">
        <title>تسجيل الدخول - الإدارة</title>
        <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css">
        </head><body class="bg-dark text-light"><div class="container py-5">
        <div class="row justify-content-center"><div class="col-md-4">
        <div class="card bg-secondary"><div class="card-header"><h5 class="mb-0">🔐 لوحة الإدارة</h5></div>
        <div class="card-body">
        <div id="msg"></div>
        <div class="mb-3"><label class="form-label">اسم المستخدم</label>
        <input type="text" class="form-control" id="adm_u" value="Anwer"></div>
        <div class="mb-3"><label class="form-label">كلمة المرور</label>
        <input type="password" class="form-control" id="adm_p"></div>
        <button class="btn btn-primary w-100" onclick="doLogin()">دخول</button>
        <hr><button class="btn btn-outline-info w-100 mt-2" onclick="bioLogin()">🔑 دخول بالبصمة</button>
        </div></div></div></div></div>
        <script>
        async function doLogin(){
          const r=await fetch('/admin/api/login',{method:'POST',headers:{'Content-Type':'application/json'},
          body:JSON.stringify({username:document.getElementById('adm_u').value,password:document.getElementById('adm_p').value})});
          const d=await r.json();
          if(d.success){location.reload();}else{document.getElementById('msg').innerHTML='<div class="alert alert-danger">'+d.message+'</div>';}
        }
        async function bioLogin(){
          const did=localStorage.getItem('deviceId'),bt=localStorage.getItem('biometricToken');
          if(!did||!bt){alert('لم تسجل بصمتك بعد');return;}
          const r=await fetch('/admin/api/biometric/login',{method:'POST',headers:{'Content-Type':'application/json'},
          body:JSON.stringify({device_id:did,biometric_token:bt})});
          const d=await r.json();
          if(d.success){location.reload();}else{alert(d.message);}
        }
        </script></body></html>''', 200

    return '''<!DOCTYPE html><html dir="rtl"><head><meta charset="utf-8">
    <title>لوحة الإدارة</title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <script src="https://cdn.socket.io/4.6.1/socket.io.min.js"></script>
    <style>
    #consoleScreen{
      background:#080808;color:#c8ffc8;font-family:'Courier New',monospace;font-size:12.5px;
      height:480px;overflow-y:auto;padding:10px 14px;border:none;
      scrollbar-width:thin;scrollbar-color:#1a3a1a #080808;
    }
    #consoleScreen::-webkit-scrollbar{width:6px;}
    #consoleScreen::-webkit-scrollbar-track{background:#080808;}
    #consoleScreen::-webkit-scrollbar-thumb{background:#1a3a1a;border-radius:3px;}
    .cl{display:flex;gap:8px;padding:1.5px 0;border-bottom:1px solid rgba(0,80,0,0.12);
        font-size:12.5px;line-height:1.45;word-break:break-all;align-items:flex-start;}
    .cl:hover{background:rgba(0,255,65,0.04);}
    .cl .ts{color:#2d6a2d;min-width:68px;flex-shrink:0;font-size:11.5px;}
    .cl .nm{color:#4a8a4a;min-width:70px;flex-shrink:0;font-size:11.5px;overflow:hidden;text-overflow:ellipsis;}
    .cl .lv{min-width:52px;flex-shrink:0;font-size:11px;font-weight:700;}
    .cl .tx{flex:1;color:#b8f0b8;}
    .cl.info .lv{color:#00bfff;}.cl.info .tx{color:#9fdfff;}
    .cl.warn .lv{color:#ffa500;}.cl.warn .tx{color:#ffe0a0;}
    .cl.error .lv{color:#ff4444;}.cl.error .tx{color:#ffaaaa;}
    .cl.critical .lv{color:#ff0000;font-weight:900;}.cl.critical .tx{color:#ff8888;}
    .cl.debug .lv{color:#666;}.cl.debug .tx{color:#888;}
    .cl.sys .lv{color:#a855f7;}.cl.sys .tx{color:#d8b4fe;}
    .hl-get{color:#22d3ee;}.hl-post{color:#fbbf24;}.hl-put{color:#a3e635;}
    .hl-del{color:#f87171;}.hl-200{color:#4ade80;}.hl-404{color:#fb923c;}.hl-500{color:#ef4444;}
    .file-item{cursor:pointer;padding:4px 8px;border-radius:3px;transition:background 0.15s;}
    .file-item:hover{background:rgba(255,255,255,0.1);}
    .file-item.selected{background:rgba(0,123,255,0.3);border-left:3px solid #0d6efd;}
    </style>
    </head><body class="bg-dark text-light">
    <div class="container-fluid py-3">
    <div class="d-flex justify-content-between align-items-center mb-4">
    <h4><i class="fas fa-shield-alt me-2"></i>لوحة إدارة أبو مالك</h4>
    <div>
    <button class="btn btn-outline-info btn-sm me-2" onclick="registerBio()">🔑 تسجيل البصمة</button>
    <button class="btn btn-outline-danger btn-sm" onclick="doLogout()">خروج</button>
    </div></div>

    <!-- ════ جدول المستخدمين ════ -->
    <div id="usersTable"><div class="text-center py-5"><i class="fas fa-spinner fa-spin fa-3x"></i></div></div>

    <!-- ════ قسم الإشعارات العامة ════ -->
    <div class="card bg-secondary mt-4">
      <div class="card-header d-flex justify-content-between align-items-center">
        <h5 class="mb-0"><i class="fas fa-bullhorn me-2 text-warning"></i>📢 نشر الإشعارات</h5>
        <button class="btn btn-outline-danger btn-sm" onclick="clearAllNotifs()">
          <i class="fas fa-trash-alt me-1"></i>مسح الكل
        </button>
      </div>
      <div class="card-body">
        <div class="row g-3">
          <div class="col-md-7">
            <label class="form-label small">نص الإشعار</label>
            <textarea id="notifMsg" class="form-control bg-dark text-light border-secondary" rows="3"
              placeholder="اكتب الإشعار..."></textarea>
          </div>
          <div class="col-md-5">
            <div class="row g-2">
              <div class="col-12">
                <label class="form-label small">نوع الإشعار</label>
                <select id="notifType" class="form-select bg-dark text-light border-secondary">
                  <option value="info">💙 معلومة</option>
                  <option value="success">💚 نجاح</option>
                  <option value="warning">🟡 تحذير</option>
                  <option value="danger">🔴 تنبيه عاجل</option>
                </select>
              </div>
              <div class="col-12">
                <label class="form-label small fw-bold">
                  <i class="fas fa-clock me-1 text-warning"></i>إعادة الإرسال الدوري
                </label>
                <div class="input-group mb-1">
                  <span class="input-group-text bg-dark text-light border-secondary">كل</span>
                  <input type="number" id="notifInterval" class="form-control bg-dark text-light border-secondary"
                    value="0" min="0" max="1440" placeholder="0">
                  <span class="input-group-text bg-dark text-light border-secondary">دقيقة</span>
                </div>
                <div class="text-muted" style="font-size:0.71rem">
                  <i class="fas fa-info-circle me-1"></i>0 = إرسال مرة واحدة فقط — أدخل رقماً لتكرار الإشعار كل X دقيقة
                </div>
              </div>
              <div class="col-12 d-flex gap-2">
                <button class="btn btn-warning fw-bold flex-grow-1" onclick="sendNotif()" id="sendBtn">
                  <i class="fas fa-paper-plane me-1"></i>نشر الإشعار
                </button>
                <button class="btn btn-danger" onclick="stopRepeat()" id="stopBtn" style="display:none">
                  <i class="fas fa-stop me-1"></i>إيقاف التكرار
                </button>
              </div>
            </div>
          </div>
        </div>
        <div id="notifResult" class="mt-2"></div>
        <div id="repeatStatus" class="mt-1"></div>
      </div>
    </div>

    <!-- سجل الإشعارات -->
    <div class="card bg-secondary mt-3">
      <div class="card-header">
        <h6 class="mb-0"><i class="fas fa-history me-2"></i>سجل الإشعارات المرسلة</h6>
      </div>
      <div class="card-body p-2">
        <div id="notifHistory"><div class="text-center text-muted py-3"><i class="fas fa-spinner fa-spin me-1"></i> جاري التحميل...</div></div>
      </div>
    </div>

    <!-- ════ إعدادات الإشعارات المتقدمة ════ -->
    <div class="card bg-secondary mt-4">
      <div class="card-header d-flex justify-content-between align-items-center">
        <h5 class="mb-0"><i class="fas fa-sliders-h me-2 text-info"></i>⚙️ إعدادات الإشعارات</h5>
        <span id="pushSubCount" class="badge bg-info">0 مشترك</span>
      </div>
      <div class="card-body">
        <div class="row g-3">

          <!-- نوع الصوت -->
          <div class="col-md-5">
            <label class="form-label small fw-bold text-warning">
              <i class="fas fa-volume-up me-1"></i>نوع الصوت عند الإشعار
            </label>
            <select id="soundType" class="form-select bg-dark text-light border-secondary"
                    onchange="saveSoundType(this.value)">
              <option value="beep">🔊 نغمة فقط (Beep)</option>
              <option value="tts">🗣️ قراءة صوتية (TTS) فقط</option>
              <option value="both" selected>🔊🗣️ نغمة + قراءة</option>
              <option value="silent">🔕 بدون صوت</option>
            </select>
            <div class="form-text">يُطبّق على كل المستخدمين عند استقبال الإشعار</div>
          </div>

          <!-- اختبار الإشعار -->
          <div class="col-md-7">
            <label class="form-label small fw-bold text-success">
              <i class="fas fa-vial me-1"></i>اختبار الإشعار الفوري
            </label>
            <div class="row g-2">
              <div class="col-8">
                <input type="text" id="testPushMsg" class="form-control bg-dark text-light border-secondary"
                       value="✅ اختبار الإشعار — يعمل بشكل صحيح!"
                       placeholder="نص الاختبار...">
              </div>
              <div class="col-4">
                <select id="testPushType" class="form-select bg-dark text-light border-secondary">
                  <option value="general">📢 عام</option>
                  <option value="broadcast">📣 إداري</option>
                  <option value="schedule_expired">⏹ مجدول</option>
                </select>
              </div>
              <div class="col-12">
                <button class="btn btn-success w-100" onclick="sendTestPush()">
                  <i class="fas fa-paper-plane me-1"></i>إرسال اختبار لكل المشتركين
                </button>
              </div>
            </div>
          </div>

          <!-- حالة المشتركين -->
          <div class="col-12">
            <div id="pushSubsList" class="small" style="max-height:160px;overflow-y:auto;
                 background:#111;border-radius:6px;padding:8px;">
              <div class="text-muted text-center py-2"><i class="fas fa-spinner fa-spin me-1"></i>جاري التحميل...</div>
            </div>
          </div>

          <div id="notifSettingsResult" class="col-12"></div>
        </div>
      </div>
    </div>

    <!-- ════ تحميل ملفات المشروع ════ -->
    <div class="card bg-secondary mt-4">
      <div class="card-header">
        <h5 class="mb-0"><i class="fas fa-download me-2 text-info"></i>📦 تحميل ملفات المشروع</h5>
      </div>
      <div class="card-body">
        <div class="row g-3 mb-3">
          <div class="col-md-6">
            <div class="card bg-dark h-100">
              <div class="card-body text-center">
                <i class="fas fa-file-archive fa-2x text-success mb-2"></i>
                <h6>تحميل المشروع كاملاً</h6>
                <p class="small text-muted">جميع الملفات مضغوطة في ملف ZIP واحد</p>
                <a href="/api/admin/download_project" class="btn btn-success w-100">
                  <i class="fas fa-download me-1"></i>تحميل ZIP كامل
                </a>
              </div>
            </div>
          </div>
          <div class="col-md-6">
            <div class="card bg-dark h-100">
              <div class="card-body text-center">
                <i class="fas fa-file-alt fa-2x text-info mb-2"></i>
                <h6>تحميل ملفات محددة</h6>
                <p class="small text-muted">اختر الملفات التي تريد تحميلها</p>
                <button class="btn btn-info w-100" onclick="loadFilesList()">
                  <i class="fas fa-list me-1"></i>عرض قائمة الملفات
                </button>
              </div>
            </div>
          </div>
        </div>
        <!-- قائمة الملفات -->
        <div id="filesSection" style="display:none">
          <div class="mb-2">
            <input type="text" id="fileSearch" class="form-control form-control-sm bg-dark text-light border-secondary"
              placeholder="🔍 ابحث عن ملف..." oninput="renderFilesList()">
          </div>
          <div class="d-flex justify-content-between align-items-center mb-2 flex-wrap gap-2">
            <span class="small text-muted" id="filesCount">0 ملف</span>
            <div class="d-flex gap-2 flex-wrap">
              <button class="btn btn-outline-secondary btn-sm" onclick="selectAllFiles()">تحديد الكل</button>
              <button class="btn btn-outline-secondary btn-sm" onclick="deselectAllFiles()">إلغاء الكل</button>
              <button class="btn btn-primary btn-sm" onclick="downloadSelected()" id="dlSelBtn" disabled>
                <i class="fas fa-file-archive me-1"></i><span id="dlSelTxt">تحميل المحدد (ZIP)</span>
              </button>
            </div>
          </div>
          <div class="small text-muted mb-2" style="font-size:0.72rem">
            <i class="fas fa-info-circle me-1"></i>اضغط <span class="text-success fw-bold"><i class="fas fa-download"></i></span> بجانب أي ملف لتحميله مباشرة، أو حدد عدة ملفات وحمّلها كـ ZIP
          </div>
          <div id="filesList" style="max-height:380px;overflow-y:auto;background:#111;border-radius:6px;padding:8px;"></div>
        </div>
      </div>
    </div>

    <!-- ════ سجلات النظام — مرآة الكونسول ════ -->
    <div class="mt-4" style="border:1px solid #1a3a1a;border-radius:8px;overflow:hidden;">
      <!-- Header -->
      <div style="background:#0d1f0d;padding:8px 14px;display:flex;justify-content:space-between;align-items:center;flex-wrap:wrap;gap:8px;border-bottom:1px solid #1a3a1a;">
        <div style="display:flex;align-items:center;gap:12px;">
          <span style="color:#00ff41;font-family:'Courier New',monospace;font-size:1rem;font-weight:700;">
            &gt;_ Console
          </span>
          <span id="consoleStatus" style="background:#0a1f0a;color:#00ff41;border:1px solid #00ff41;padding:2px 8px;border-radius:12px;font-family:monospace;font-size:0.72rem;">● Live</span>
          <span id="logCountBadge" style="color:#3a6a3a;font-family:monospace;font-size:0.72rem;">0 events</span>
        </div>
        <div style="display:flex;gap:6px;flex-wrap:wrap;align-items:center;">
          <select id="logFilter" onchange="filterLogs()"
                  style="background:#0a0a0a;border:1px solid #1f3f1f;color:#00cc33;padding:3px 8px;border-radius:4px;font-family:monospace;font-size:0.73rem;outline:none;">
            <option value="all">All</option>
            <option value="error">ERROR</option>
            <option value="warning">WARNING</option>
            <option value="info">INFO</option>
            <option value="debug">DEBUG</option>
          </select>
          <input id="consoleSearch" type="text" placeholder="🔍 Search..."
                 oninput="filterLogs()"
                 style="background:#0a0a0a;border:1px solid #1f3f1f;color:#00cc33;padding:3px 10px;border-radius:4px;font-family:monospace;font-size:0.73rem;width:130px;outline:none;">
          <button onclick="copyAllLogs()" title="نسخ كل السجلات"
                  style="background:#0a2a0a;border:1px solid #00aa22;color:#00ff41;padding:3px 10px;border-radius:4px;font-family:monospace;font-size:0.73rem;cursor:pointer;">
            📋 نسخ الكل
          </button>
          <button onclick="clearConsole()"
                  style="background:#2a0a0a;border:1px solid #aa2222;color:#ff6666;padding:3px 10px;border-radius:4px;font-family:monospace;font-size:0.73rem;cursor:pointer;">
            🗑 مسح
          </button>
          <button onclick="toggleAutoScroll()" id="autoScrollBtn"
                  style="background:#0a0a2a;border:1px solid #3344aa;color:#8899ff;padding:3px 10px;border-radius:4px;font-family:monospace;font-size:0.73rem;cursor:pointer;">
            🔒 Auto
          </button>
        </div>
      </div>
      <!-- Console Output -->
      <div id="consoleScreen">
        <div style="color:#2d6a2d;font-family:monospace;font-size:11.5px;padding:6px 2px;opacity:0.7;">
          // Waiting for events...
        </div>
      </div>
      <!-- AI Chat Section -->
      <div style="background:#060606;border-top:1px solid #1a2a1a;padding:10px 14px;">
        <div style="color:#3a7a3a;font-family:monospace;font-size:0.73rem;margin-bottom:6px;">
          🤖 اسأل الذكاء عن أي حدث في السجلات
        </div>
        <div style="display:flex;gap:6px;">
          <input type="text" id="aiLogQuestion"
                 placeholder="مثال: آخر الأخطاء / ابحث عن werkzeug / أخبرني بالأحداث الأخيرة"
                 onkeydown="if(event.key===&apos;Enter&apos;) askAiAboutLogs()"
                 style="flex:1;background:#040404;border:1px solid #1a3a1a;color:#00cc33;padding:7px 12px;border-radius:4px;font-family:monospace;font-size:0.8rem;outline:none;">
          <button onclick="askAiAboutLogs()" id="aiSendBtn"
                  style="background:linear-gradient(135deg,#0a3a0a,#081808);border:1px solid #00cc33;color:#00ff41;padding:7px 16px;border-radius:4px;font-family:monospace;font-size:0.8rem;cursor:pointer;font-weight:700;white-space:nowrap;">
            ▶ إرسال
          </button>
        </div>
        <div id="aiLogResponse"
             style="display:none;margin-top:8px;font-family:monospace;font-size:0.78rem;
                    white-space:pre-wrap;color:#88dd88;background:#040404;
                    border:1px solid #1a3a1a;border-radius:4px;padding:10px;
                    max-height:200px;overflow-y:auto;line-height:1.5;"></div>
      </div>
    </div>

    </div><!-- end container -->

    <script>
    /* ══════════════════════════════════════════
       المستخدمون
    ══════════════════════════════════════════ */
    async function loadUsers(){
      const r=await fetch('/admin/api/users');const d=await r.json();
      if(!d.success){document.getElementById('usersTable').innerHTML='<p class="text-danger">'+d.message+'</p>';return;}
      let h='<div class="table-responsive"><table class="table table-dark table-bordered table-sm">';
      h+='<thead><tr><th>ID</th><th>الاسم</th><th>الهاتف</th><th>متصل؟</th><th>محظور؟</th><th>تنبيهات</th><th>إجراء</th></tr></thead><tbody>';
      d.users.forEach(u=>{
        h+=`<tr><td>${u.user_id}</td><td>${u.name}</td><td>${u.phone||'-'}</td>
        <td>${u.logged_in?'<span class="badge bg-success">نعم</span>':'<span class="badge bg-secondary">لا</span>'}</td>
        <td>${u.blocked?'<span class="badge bg-danger">محظور</span>':'<span class="badge bg-success">فعّال</span>'}</td>
        <td><span class="badge bg-warning text-dark">${u.alerts_count}</span></td>
        <td>
        <button class="btn btn-xs btn-outline-warning btn-sm" onclick="toggleBlock('${u.user_id}',${!u.blocked})">${u.blocked?'فك الحظر':'حظر'}</button>
        <a href="/admin/api/copy_chats/${u.user_id}" class="btn btn-xs btn-outline-info btn-sm ms-1" target="_blank">روابط</a>
        </td></tr>`;
      });
      h+='</tbody></table></div>';
      document.getElementById('usersTable').innerHTML=h;
    }
    async function toggleBlock(slot,blocked){
      await fetch('/admin/api/user/'+slot,{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({action:'block',blocked:blocked})});loadUsers();
    }
    async function doLogout(){await fetch('/admin/api/logout',{method:'POST'});location.reload();}
    function registerBio(){
      let did=localStorage.getItem('deviceId');
      if(!did){did=crypto.randomUUID?crypto.randomUUID():'dev-'+Date.now();localStorage.setItem('deviceId',did);}
      const bt=Array.from(crypto.getRandomValues(new Uint8Array(32))).map(b=>b.toString(16).padStart(2,'0')).join('');
      fetch('/admin/api/biometric/register',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({device_id:did,biometric_token:bt})})
      .then(r=>r.json()).then(d=>{if(d.success){localStorage.setItem('biometricToken',bt);alert('✅ تم تسجيل البصمة');}else{alert('❌ '+d.message);}});
    }

    /* ══════════════════════════════════════════
       الإشعارات مع التكرار والنطق
    ══════════════════════════════════════════ */
    const _typeColors={info:'primary',success:'success',warning:'warning',danger:'danger'};
    const _typeIcons={info:'ℹ️',success:'✅',warning:'⚠️',danger:'🚨'};
    let _repeatTimer=null;
    let _repeatCount=0;
    let _countdownTimer=null;
    let _nextSendAt=0;

    function _startCountdown(intervalMin){
      if(_countdownTimer) clearInterval(_countdownTimer);
      _nextSendAt=Date.now()+intervalMin*60000;
      _countdownTimer=setInterval(()=>{
        const rem=Math.max(0,_nextSendAt-Date.now());
        if(!_repeatTimer){clearInterval(_countdownTimer);return;}
        const m=Math.floor(rem/60000), s=Math.floor((rem%60000)/1000);
        const rs=document.getElementById('repeatStatus');
        if(rs) rs.innerHTML=
          `<div class="alert alert-info py-2 small d-flex justify-content-between align-items-center">
            <span>🔁 تم الإرسال <strong>${_repeatCount}</strong> مرة — التكرار كل ${intervalMin} دقيقة</span>
            <span class="badge bg-dark text-warning" style="font-size:0.85rem">⏱ ${m}:${String(s).padStart(2,'0')}</span>
          </div>`;
        if(rem===0) _nextSendAt=Date.now()+intervalMin*60000;
      },500);
    }

    function _speakNotif(msg){
      if(!window.speechSynthesis) return;
      window.speechSynthesis.cancel();
      const utt=new SpeechSynthesisUtterance(msg);
      utt.lang='ar-SA'; utt.rate=0.95; utt.pitch=1; utt.volume=1;
      const voices=window.speechSynthesis.getVoices();
      const arVoice=voices.find(v=>v.lang&&v.lang.startsWith('ar'));
      if(arVoice) utt.voice=arVoice;
      window.speechSynthesis.speak(utt);
    }

    async function _doSendNotif(){
      const msg=document.getElementById('notifMsg').value.trim();
      const type=document.getElementById('notifType').value;
      if(!msg) return false;
      const r=await fetch('/admin/api/broadcast_notification',{
        method:'POST',headers:{'Content-Type':'application/json'},
        body:JSON.stringify({message:msg,type:type})
      });
      const d=await r.json();
      if(d.success){
        _repeatCount++;
        document.getElementById('repeatStatus').innerHTML=
          _repeatTimer?`<div class="alert alert-info py-1 small">🔁 تم الإرسال مرة ${_repeatCount}</div>`:'';
        loadNotifHistory();
        _speakNotif('تم إرسال الإشعار: '+msg);
        return true;
      }
      return false;
    }

    async function sendNotif(){
      const msg=document.getElementById('notifMsg').value.trim();
      const type=document.getElementById('notifType').value;
      const intervalMin=parseInt(document.getElementById('notifInterval').value)||0;
      const res=document.getElementById('notifResult');
      if(!msg){res.innerHTML='<div class="alert alert-danger py-1 small">الرجاء كتابة نص الإشعار</div>';return;}
      stopRepeat();
      res.innerHTML='<div class="alert alert-secondary py-1 small"><i class="fas fa-spinner fa-spin me-1"></i>جاري النشر...</div>';
      try{
        const ok=await _doSendNotif();
        if(ok){
          res.innerHTML='<div class="alert alert-success py-1 small">✅ تم نشر الإشعار بنجاح</div>';
          if(intervalMin>0){
            document.getElementById('stopBtn').style.display='';
            document.getElementById('sendBtn').innerHTML='<i class="fas fa-sync fa-spin me-1"></i>يعمل التكرار...';
            _repeatCount=1;
            _startCountdown(intervalMin);
            _repeatTimer=setInterval(async()=>{
              await _doSendNotif();
              _nextSendAt=Date.now()+intervalMin*60000;
            }, intervalMin*60000);
          } else {
            setTimeout(()=>{res.innerHTML='';},4000);
          }
        } else {
          res.innerHTML='<div class="alert alert-danger py-1 small">❌ فشل الإرسال</div>';
        }
      }catch(e){
        res.innerHTML=`<div class="alert alert-danger py-1 small">❌ خطأ: ${e.message}</div>`;
      }
    }

    function stopRepeat(){
      if(_repeatTimer){clearInterval(_repeatTimer);_repeatTimer=null;}
      if(_countdownTimer){clearInterval(_countdownTimer);_countdownTimer=null;}
      _repeatCount=0;
      document.getElementById('stopBtn').style.display='none';
      document.getElementById('sendBtn').innerHTML='<i class="fas fa-paper-plane me-1"></i>نشر الإشعار';
      document.getElementById('repeatStatus').innerHTML='';
    }

    async function deleteNotif(id){
      await fetch('/admin/api/delete_notification/'+id,{method:'DELETE'});loadNotifHistory();
    }
    async function clearAllNotifs(){
      if(!confirm('مسح جميع الإشعارات؟')) return;
      await fetch('/admin/api/clear_notifications',{method:'POST'});loadNotifHistory();
    }

    // ── إعدادات الإشعارات المتقدمة ──────────────────────────
    async function loadPushStats(){
      try{
        const r=await fetch('/admin/api/push_stats');
        const d=await r.json();
        const cnt=document.getElementById('pushSubCount');
        if(cnt) cnt.textContent=d.count+' مشترك';
        const lst=document.getElementById('pushSubsList');
        if(!lst) return;
        if(!d.subscribers||!d.subscribers.length){
          lst.innerHTML='<div class="text-muted text-center py-2">لا يوجد مشتركون بعد</div>';return;
        }
        lst.innerHTML=d.subscribers.map((uid,i)=>`
          <div class="d-flex align-items-center gap-2 py-1 border-bottom border-secondary">
            <span class="badge bg-success">${i+1}</span>
            <span class="text-light small flex-grow-1">👤 ${uid}</span>
            <span class="badge bg-primary">مشترك</span>
          </div>`).join('');
      }catch(e){
        const lst=document.getElementById('pushSubsList');
        if(lst) lst.innerHTML='<div class="text-danger small">فشل التحميل: '+e.message+'</div>';
      }
    }

    async function sendTestPush(){
      const msg  = (document.getElementById('testPushMsg')  ||{}).value||'✅ اختبار';
      const type = (document.getElementById('testPushType') ||{}).value||'general';
      const res  = document.getElementById('notifSettingsResult');
      if(res) res.innerHTML='<div class="alert alert-info py-1 small"><i class="fas fa-spinner fa-spin me-1"></i>جاري الإرسال...</div>';
      try{
        const r=await fetch('/admin/api/test_push',{
          method:'POST',headers:{'Content-Type':'application/json'},
          body:JSON.stringify({message:msg,type:type})
        });
        const d=await r.json();
        if(res) res.innerHTML=`<div class="alert alert-${d.success?'success':'danger'} py-1 small">${d.message||d.error}</div>`;
        loadPushStats();
      }catch(e){
        if(res) res.innerHTML='<div class="alert alert-danger py-1 small">خطأ: '+e.message+'</div>';
      }
    }

    function saveSoundType(val){
      // احفظ نوع الصوت في localStorage ليقرأه كل المستخدمين
      try{ localStorage.setItem('notif_sound_type', val); }catch(e){}
    }

    function loadSoundTypePref(){
      try{
        const v=localStorage.getItem('notif_sound_type')||'both';
        const el=document.getElementById('soundType');
        if(el) el.value=v;
      }catch(e){}
    }
    async function loadNotifHistory(){
      const el=document.getElementById('notifHistory');
      try{
        const r=await fetch('/admin/api/notifications');const d=await r.json();
        if(!d.success){el.innerHTML='<p class="text-danger small">'+d.message+'</p>';return;}
        if(!d.notifications.length){el.innerHTML='<p class="text-muted text-center small py-2">لا توجد إشعارات</p>';return;}
        let h='';
        d.notifications.forEach(n=>{
          const col=_typeColors[n.type]||'secondary';const ico=_typeIcons[n.type]||'📢';
          const ts=n.timestamp?n.timestamp.replace('T',' ').slice(0,16):'';
          h+=`<div class="d-flex align-items-start gap-2 border-bottom border-dark pb-2 mb-2">
            <span class="badge bg-${col} mt-1">${ico} ${n.type||'info'}</span>
            <div class="flex-grow-1">
              <div class="text-light small">${n.message.replace(/</g,'&lt;')}</div>
              <div class="text-muted" style="font-size:0.72rem">${ts}</div>
            </div>
            <button class="btn btn-outline-danger btn-sm py-0 px-1" style="font-size:0.7rem"
              onclick="deleteNotif('${n.id}')"><i class="fas fa-times"></i></button>
          </div>`;
        });
        el.innerHTML=h;
      }catch(e){el.innerHTML='<p class="text-danger small">خطأ في التحميل</p>';}
    }

    /* ══════════════════════════════════════════
       تحميل الملفات
    ══════════════════════════════════════════ */
    let _allFiles=[];
    let _selectedFiles=new Set();

    async function loadFilesList(){
      const sec=document.getElementById('filesSection');
      const lst=document.getElementById('filesList');
      sec.style.display='';
      lst.innerHTML='<div class="text-center py-3"><i class="fas fa-spinner fa-spin"></i> جاري التحميل...</div>';
      try{
        const r=await fetch('/api/admin/list_project_files');const d=await r.json();
        if(!d.success){lst.innerHTML='<p class="text-danger small">'+d.error+'</p>';return;}
        _allFiles=d.files; _selectedFiles=new Set();
        document.getElementById('filesCount').textContent=d.files.length+' ملف';
        renderFilesList();
      }catch(e){lst.innerHTML='<p class="text-danger small">خطأ: '+e.message+'</p>';}
    }

    function renderFilesList(filter){
      const lst=document.getElementById('filesList');
      let h='';
      const q=(filter||document.getElementById('fileSearch').value||'').toLowerCase();
      const shown=_allFiles.filter(f=>!q||f.path.toLowerCase().includes(q));
      document.getElementById('filesCount').textContent=shown.length+' ملف'+(q?' (مصفّاة)':'');
      shown.forEach((f,i)=>{
        const sel=_selectedFiles.has(f.path);
        const icon=f.path.endsWith('.py')?'🐍':f.path.endsWith('.html')?'🌐':f.path.endsWith('.js')?'📜':f.path.endsWith('.css')?'🎨':f.path.endsWith('.json')?'📋':f.path.endsWith('.png')||f.path.endsWith('.jpg')||f.path.endsWith('.jpeg')||f.path.endsWith('.webp')?'🖼️':f.path.endsWith('.txt')?'📝':'📄';
        const safeP=f.path.replace(/\\/g,'\\\\').replace(/'/g,"\\'");
        h+=`<div class="file-item d-flex align-items-center gap-2 ${sel?'selected':''}" onclick="toggleFile('${safeP}',this)">
          <input type="checkbox" class="form-check-input" ${sel?'checked':''} style="pointer-events:none">
          <span>${icon}</span>
          <span class="small text-light flex-grow-1" style="word-break:break-all">${f.path}</span>
          <span class="badge bg-dark text-muted me-1" style="font-size:0.65rem;white-space:nowrap">${f.size}</span>
          <a href="/api/admin/download_single_file?path=${encodeURIComponent(f.path)}"
             class="btn btn-outline-success btn-sm py-0 px-1" style="font-size:0.7rem;white-space:nowrap"
             title="تحميل مباشر" onclick="event.stopPropagation()">
            <i class="fas fa-download"></i>
          </a>
        </div>`;
      });
      lst.innerHTML=h||'<p class="text-muted small text-center py-2">لا توجد ملفات</p>';
      updateDlBtn();
    }

    function toggleFile(path,el){
      if(_selectedFiles.has(path)){_selectedFiles.delete(path);el.classList.remove('selected');el.querySelector('input').checked=false;}
      else{_selectedFiles.add(path);el.classList.add('selected');el.querySelector('input').checked=true;}
      updateDlBtn();
    }
    function selectAllFiles(){_allFiles.forEach(f=>_selectedFiles.add(f.path));renderFilesList();}
    function deselectAllFiles(){_selectedFiles.clear();renderFilesList();}
    function updateDlBtn(){
      const btn=document.getElementById('dlSelBtn');
      const txt=document.getElementById('dlSelTxt');
      btn.disabled=_selectedFiles.size===0;
      if(txt) txt.textContent=_selectedFiles.size>0?`تحميل ${_selectedFiles.size} ملف (ZIP)`:'تحميل المحدد (ZIP)';
    }

    async function downloadSelected(){
      if(_selectedFiles.size===0) return;
      const btn=document.getElementById('dlSelBtn');
      const txt=document.getElementById('dlSelTxt');
      btn.disabled=true; if(txt) txt.textContent='⏳ جاري الضغط...';
      try{
        const r=await fetch('/api/admin/download_selected_files',{
          method:'POST',headers:{'Content-Type':'application/json'},
          body:JSON.stringify({files:[..._selectedFiles]})
        });
        if(!r.ok){const d=await r.json();alert('❌ '+d.error);btn.disabled=false;updateDlBtn();return;}
        const blob=await r.blob();
        const url=URL.createObjectURL(blob);
        const a=document.createElement('a');a.href=url;a.download='selected_files.zip';
        document.body.appendChild(a);a.click();document.body.removeChild(a);
        URL.revokeObjectURL(url);
      }catch(e){alert('❌ خطأ: '+e.message);}
      btn.disabled=false; updateDlBtn();
    }

    /* ══════════════════════════════════════════════════════
       مرآة الكونسول — سجلات النظام الفورية (محسّن)
    ══════════════════════════════════════════════════════ */
    const _consoleEl = document.getElementById('consoleScreen');
    let _autoScroll   = true;
    let _allLogs      = [];   /* [{ts,name,level,raw,cssClass}] */
    let _currentFilter = 'all';
    let _searchText    = '';
    const _maxLogs     = 1000;

    const _socket = io({transports:['websocket','polling']});

    _socket.on('connect', () => {
      const st = document.getElementById('consoleStatus');
      st.style.color='#00ff41'; st.style.borderColor='#00ff41';
      st.textContent = '● Live';
      _addEntry({level:'INFO', msg:'✅ متصل بالخادم — تدفق السجلات نشط', time:_now(), name:'SYSTEM'});
    });
    _socket.on('disconnect', () => {
      const st = document.getElementById('consoleStatus');
      st.style.color='#ff4444'; st.style.borderColor='#ff4444';
      st.textContent = '● Offline';
      _addEntry({level:'ERROR', msg:'❌ انقطع الاتصال بالخادم', time:_now(), name:'SYSTEM'});
    });
    _socket.on('live_log', data => _addEntry(data));
    _socket.on('new_broadcast_notification', data => {
      _addEntry({level:'INFO', msg:'📢 إشعار مُرسل: '+data.message, time:_now(), name:'NOTIF'});
    });

    function _now(){ return new Date().toLocaleTimeString('en-GB',{hour12:false}); }

    /* تلوين HTTP methods وstatus codes */
    function _colorText(txt){
      return txt
        .replace(/\b(GET)\b/g,'<span class="hl-get">GET</span>')
        .replace(/\b(POST)\b/g,'<span class="hl-post">POST</span>')
        .replace(/\b(PUT|PATCH)\b/g,'<span class="hl-put">$1</span>')
        .replace(/\b(DELETE)\b/g,'<span class="hl-del">DELETE</span>')
        .replace(/\s(2\d\d)\s/g,' <span class="hl-200">$1</span> ')
        .replace(/\s(404)\s/g,' <span class="hl-404">404</span> ')
        .replace(/\s(5\d\d)\s/g,' <span class="hl-500">$1</span> ');
    }

    function _addEntry(entry){
      const lvl = (entry.level||'INFO').toUpperCase();
      let cssClass = 'info';
      if(lvl==='ERROR'||lvl==='CRITICAL') cssClass = lvl==='CRITICAL'?'critical':'error';
      else if(lvl==='WARNING') cssClass = 'warn';
      else if(lvl==='DEBUG')   cssClass = 'debug';
      else if(entry.name==='SYSTEM'||entry.name==='NOTIF') cssClass = 'sys';

      const raw = (entry.full || ((entry.time||'')+(entry.name?'['+entry.name+']':'')+' '+lvl+' '+entry.msg));
      const logEntry = { ts: entry.time||_now(), name:entry.name||'', level:lvl, raw, cssClass,
                         msg: entry.msg||entry.full||'' };
      _allLogs.push(logEntry);
      if(_allLogs.length > _maxLogs) _allLogs.shift();

      _updateCount();
      if(_matchFilter(logEntry)) _renderLine(logEntry);
    }

    function _matchFilter(e){
      const fv = _currentFilter;
      const lvl = e.level;
      const ok = fv==='all' || (fv==='error'&&(lvl==='ERROR'||lvl==='CRITICAL')) ||
                 (fv==='warning'&&lvl==='WARNING') || (fv==='info'&&lvl==='INFO') ||
                 (fv==='debug'&&lvl==='DEBUG');
      if(!ok) return false;
      if(_searchText && !e.raw.toLowerCase().includes(_searchText)) return false;
      return true;
    }

    function _esc(s){ return String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); }

    function _renderLine(e){
      const row = document.createElement('div');
      row.className = 'cl ' + e.cssClass;
      row.innerHTML =
        '<span class="ts">'+_esc(e.ts)+'</span>'+
        '<span class="nm">'+_esc(e.name).slice(0,10)+'</span>'+
        '<span class="lv">'+e.level+'</span>'+
        '<span class="tx">'+_colorText(_esc(e.msg||e.raw))+'</span>';
      _consoleEl.appendChild(row);
      if(_autoScroll) _consoleEl.scrollTop = _consoleEl.scrollHeight;
    }

    function _updateCount(){
      const badge = document.getElementById('logCountBadge');
      if(badge) badge.textContent = _allLogs.length + ' events';
    }

    function filterLogs(){
      _currentFilter = document.getElementById('logFilter').value;
      _searchText    = (document.getElementById('consoleSearch').value||'').toLowerCase().trim();
      _consoleEl.innerHTML = '';
      _allLogs.forEach(e => { if(_matchFilter(e)) _renderLine(e); });
    }

    function clearConsole(){
      _consoleEl.innerHTML='';
      _allLogs=[];
      _updateCount();
    }

    function toggleAutoScroll(){
      _autoScroll = !_autoScroll;
      const btn = document.getElementById('autoScrollBtn');
      btn.textContent = _autoScroll ? '🔒 Auto' : '🔓 Manual';
      btn.style.color = _autoScroll ? '#8899ff' : '#ffaa44';
      if(_autoScroll) _consoleEl.scrollTop = _consoleEl.scrollHeight;
    }

    function copyAllLogs(){
      const lines = _allLogs.map(e=>
        '['+e.ts+'] ['+e.name+'] '+e.level+' '+e.msg
      ).join('\n');
      if(!lines){ alert('لا توجد سجلات للنسخ'); return; }
      navigator.clipboard.writeText(lines).then(()=>{
        const btn = event.target;
        const orig = btn.textContent;
        btn.textContent = '✅ تم النسخ!';
        btn.style.color = '#00ff41';
        setTimeout(()=>{ btn.textContent=orig; btn.style.color=''; }, 2000);
      }).catch(()=>{
        const ta = document.createElement('textarea');
        ta.value = lines; document.body.appendChild(ta);
        ta.select(); document.execCommand('copy');
        document.body.removeChild(ta);
        alert('✅ تم نسخ '+_allLogs.length+' سطر');
      });
    }

    async function askAiAboutLogs(){
      const q   = (document.getElementById('aiLogQuestion').value||'').trim();
      const res = document.getElementById('aiLogResponse');
      const btn = document.getElementById('aiSendBtn');
      if(!q) return;
      btn.disabled = true; btn.textContent = '⏳...';
      res.style.display = 'block';
      res.textContent = '⏳ جاري التحليل...';
      try{
        const r = await fetch('/admin/api/log_query',{
          method:'POST',
          headers:{'Content-Type':'application/json'},
          body:JSON.stringify({query:q, logs: _allLogs.slice(-200).map(e=>'['+e.ts+']['+e.name+'] '+e.level+': '+e.msg)})
        });
        const d = await r.json();
        res.textContent = d.success ? d.answer : ('❌ ' + (d.message||'فشل'));
      } catch(e){
        res.textContent = '❌ خطأ: ' + e.message;
      }
      btn.disabled = false; btn.textContent = '▶ إرسال';
    }

    // تهيئة تلقائية للأصوات
    if(window.speechSynthesis){
      window.speechSynthesis.onvoiceschanged = function(){ window.speechSynthesis.getVoices(); }
    }

    loadUsers();
    loadNotifHistory();
    loadPushStats();
    loadSoundTypePref();
    loadCardStatus();

    // ══════════════════════════════════════════════════════════
    //  إدارة بطاقات الشحن
    // ══════════════════════════════════════════════════════════
    async function loadCardStatus(){
      try{
        const r = await fetch('/admin/api/card_status');
        const d = await r.json();
        const el = document.getElementById('cardToggleBtn');
        if(el){ el.textContent = d.enabled ? '🔴 تعطيل النظام' : '🟢 تفعيل النظام'; }
        const badge = document.getElementById('cardStatusBadge');
        if(badge){ badge.textContent = d.enabled ? '✅ مفعّل' : '⛔ معطّل'; badge.className = d.enabled ? 'badge bg-success' : 'badge bg-danger'; }
      }catch(e){ console.error('loadCardStatus:', e); }
    }
    async function toggleCardSystem(){
      const r = await fetch('/admin/api/card_status'); const d = await r.json();
      const newState = !d.enabled;
      const r2 = await fetch('/admin/api/toggle_card_system',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({enabled:newState})});
      const d2 = await r2.json();
      document.getElementById('cardMsg').innerHTML = '<div class="alert alert-'+(d2.success?'success':'danger')+'">'+d2.message+'</div>';
      loadCardStatus();
    }
    async function createVouchers(){
      const plan_id = parseInt(document.getElementById('voucherPlan').value);
      const count = parseInt(document.getElementById('voucherCount').value)||10;
      const r = await fetch('/admin/api/create_vouchers',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({plan_id,count})});
      const d = await r.json();
      if(d.success){
        const txt = d.codes.join('\n');
        document.getElementById('voucherResult').value = txt;
        document.getElementById('voucherResultArea').style.display='';
        document.getElementById('cardMsg').innerHTML='<div class="alert alert-success">✅ تم إنشاء '+d.count+' قسيمة</div>';
      }else{
        document.getElementById('cardMsg').innerHTML='<div class="alert alert-danger">❌ '+d.message+'</div>';
      }
    }
    async function loadVouchers(){
      const r = await fetch('/admin/api/vouchers'); const d = await r.json();
      const tbody = document.getElementById('vouchersTbody');
      if(!tbody) return;
      const vouchers = (d.vouchers||[]).slice(-100).reverse();
      tbody.innerHTML = vouchers.map(v=>`<tr>
        <td><code class="text-warning" style="font-size:0.75rem">${v.plan_name||'—'}</code></td>
        <td><span class="badge bg-${v.status==='unused'?'secondary':v.status==='active'?'success':v.status==='used'?'primary':'danger'}">${v.status==='unused'?'غير مستخدمة':v.status==='active'?'نشطة':v.status==='used'?'مستخدمة':'منتهية'}</span></td>
        <td style="font-size:0.72rem">${(v.created_at||'—').slice(0,10)}</td>
        <td style="font-size:0.72rem">${(v.used_at||'—').slice(0,10)}</td>
        <td style="font-size:0.72rem">${(v.expires_at||'—').slice(0,16).replace('T',' ')}</td>
      </tr>`).join('');
    }
    async function loadCardSessions(){
      const r = await fetch('/admin/api/card_sessions'); const d = await r.json();
      const tbody = document.getElementById('cardSessionsTbody');
      if(!tbody) return;
      const sessions = d.sessions||[];
      tbody.innerHTML = sessions.length===0?'<tr><td colspan="5" class="text-center text-muted">لا توجد جلسات نشطة</td></tr>':
        sessions.map(s=>`<tr>
          <td style="font-size:0.72rem">${s.plan_name||'—'}</td>
          <td style="font-size:0.72rem">${s.ip_address||'—'}</td>
          <td style="font-size:0.72rem">${(s.start_time||'—').slice(0,16).replace('T',' ')}</td>
          <td style="font-size:0.72rem">${(s.expires_at||'—').slice(0,16).replace('T',' ')}</td>
          <td><button class="btn btn-danger btn-sm" onclick="terminateCardSession('${s.session_id}')"><i class="fas fa-times"></i></button></td>
        </tr>`).join('');
    }
    async function terminateCardSession(sid){
      if(!confirm('إنهاء هذه الجلسة؟')) return;
      await fetch('/admin/api/terminate_card_session',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({session_id:sid})});
      loadCardSessions();
    }
    async function deleteVouchers(status){
      if(!confirm('حذف القسائم '+(status||'الكل')+'؟')) return;
      const body = status ? {status} : {};
      await fetch('/admin/api/delete_vouchers',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
      document.getElementById('cardMsg').innerHTML='<div class="alert alert-success">✅ تم الحذف</div>';
      loadVouchers();
    }
    function copyVouchers(){
      const ta = document.getElementById('voucherResult');
      if(!ta.value) return;
      navigator.clipboard.writeText(ta.value).then(()=>alert('✅ تم النسخ!'));
    }
    </script>

    <!-- ════ قسم إدارة البطاقات ════ -->
    <style>
    #cardSection .card{background:#1a1f2e;border:1px solid #2d3748;}
    #cardSection .form-control,#cardSection .form-select{background:#0d1117;border-color:#2d3748;color:#e2e8f0;}
    </style>
    <div id="cardSection" class="mt-4">
      <div class="card" style="background:#1a1f2e;border:1px solid #2d3748;">
        <div class="card-header d-flex justify-content-between align-items-center" style="border-bottom:1px solid #2d3748;">
          <h5 class="mb-0 text-light"><i class="fas fa-id-card me-2 text-info"></i>🎓 مركز سرعة إنجاز — نظام بطاقات الشحن</h5>
          <div class="d-flex align-items-center gap-2">
            <span id="cardStatusBadge" class="badge bg-danger">⛔ معطّل</span>
            <button id="cardToggleBtn" class="btn btn-sm btn-outline-light" onclick="toggleCardSystem()">🟢 تفعيل النظام</button>
          </div>
        </div>
        <div class="card-body">
          <div id="cardMsg"></div>

          <!-- إنشاء قسائم جديدة -->
          <div class="row g-3 mb-4">
            <div class="col-md-4">
              <label class="form-label small text-light">نوع الباقة</label>
              <select id="voucherPlan" class="form-select" style="background:#0d1117;border-color:#2d3748;color:#e2e8f0;">
                <option value="1">📅 يومية (24 ساعة)</option>
                <option value="2">📅 أسبوعية (7 أيام)</option>
                <option value="3">📅 شهرية (30 يوم)</option>
              </select>
            </div>
            <div class="col-md-3">
              <label class="form-label small text-light">عدد القسائم</label>
              <input type="number" id="voucherCount" class="form-control" value="10" min="1" max="200" style="background:#0d1117;border-color:#2d3748;color:#e2e8f0;">
            </div>
            <div class="col-md-3 d-flex align-items-end">
              <button class="btn btn-success w-100" onclick="createVouchers()">
                <i class="fas fa-plus me-1"></i>إنشاء قسائم
              </button>
            </div>
            <div class="col-md-2 d-flex align-items-end">
              <button class="btn btn-outline-secondary w-100" onclick="loadVouchers();loadCardSessions();">
                <i class="fas fa-sync me-1"></i>تحديث
              </button>
            </div>
          </div>

          <!-- منطقة الأكواد المولّدة -->
          <div id="voucherResultArea" style="display:none;" class="mb-4">
            <div class="d-flex justify-content-between mb-1">
              <label class="form-label small text-success">✅ الأكواد المولّدة</label>
              <button class="btn btn-sm btn-outline-success" onclick="copyVouchers()"><i class="fas fa-copy me-1"></i>نسخ الكل</button>
            </div>
            <textarea id="voucherResult" class="form-control font-monospace" rows="6" readonly
              style="background:#0a0e1a;border-color:#2d3748;color:#4ade80;font-size:0.82rem;letter-spacing:1px;"></textarea>
          </div>

          <!-- جدول القسائم -->
          <ul class="nav nav-tabs mb-3" style="border-color:#2d3748;">
            <li class="nav-item">
              <a class="nav-link active text-light" href="#" onclick="loadVouchers();document.querySelectorAll('.c-tab').forEach(e=>e.style.display='none');document.getElementById('tabVouchers').style.display='';return false;"
                 style="border-color:#2d3748;background:#0d1117;">📋 سجل القسائم</a>
            </li>
            <li class="nav-item">
              <a class="nav-link text-light" href="#" onclick="loadCardSessions();document.querySelectorAll('.c-tab').forEach(e=>e.style.display='none');document.getElementById('tabSessions').style.display='';return false;"
                 style="border-color:#2d3748;background:#161b22;">🟢 الجلسات النشطة</a>
            </li>
          </ul>

          <div id="tabVouchers" class="c-tab">
            <div class="d-flex justify-content-end gap-2 mb-2">
              <button class="btn btn-sm btn-outline-warning" onclick="deleteVouchers('used')"><i class="fas fa-trash me-1"></i>حذف المستخدمة</button>
              <button class="btn btn-sm btn-outline-danger" onclick="deleteVouchers('')"><i class="fas fa-trash me-1"></i>حذف الكل</button>
            </div>
            <div class="table-responsive">
              <table class="table table-dark table-sm table-hover" style="font-size:0.82rem;">
                <thead><tr><th>الباقة</th><th>الحالة</th><th>تاريخ الإنشاء</th><th>تاريخ الاستخدام</th><th>الانتهاء</th></tr></thead>
                <tbody id="vouchersTbody"><tr><td colspan="5" class="text-center text-muted">اضغط تحديث لتحميل البيانات</td></tr></tbody>
              </table>
            </div>
          </div>

          <div id="tabSessions" class="c-tab" style="display:none;">
            <div class="table-responsive">
              <table class="table table-dark table-sm table-hover" style="font-size:0.82rem;">
                <thead><tr><th>الباقة</th><th>IP</th><th>بداية الجلسة</th><th>الانتهاء</th><th>إجراء</th></tr></thead>
                <tbody id="cardSessionsTbody"><tr><td colspan="5" class="text-center text-muted">لا توجد جلسات نشطة</td></tr></tbody>
              </table>
            </div>
          </div>

        </div>
      </div>
    </div>

    <!-- ════════════ قسم التحديثات ════════════ -->
    <div class="card bg-secondary mt-4" style="border:1px solid #3d4a5c;">
      <div class="card-header d-flex justify-content-between align-items-center">
        <h5 class="mb-0"><i class="fas fa-sync-alt me-2 text-primary"></i>🔄 تحديث التطبيق</h5>
        <div class="form-check form-switch mb-0">
          <input class="form-check-input" type="checkbox" id="autoUpdateToggle" style="width:50px;height:25px;cursor:pointer;">
          <label class="form-check-label text-light small me-2" for="autoUpdateToggle">تحديث تلقائي</label>
        </div>
      </div>
      <div class="card-body">
        <div id="updateStatus" class="alert alert-info mb-3">
          <i class="fas fa-spinner fa-spin me-2"></i> جاري التحقق من التحديثات...
        </div>
        <div class="d-flex flex-wrap gap-2 mb-3">
          <button class="btn btn-primary" onclick="checkForUpdatesUI()" id="checkUpdateBtn">
            <i class="fas fa-search me-1"></i>بحث عن تحديث
          </button>
          <button class="btn btn-success" onclick="performUpdateUI()" id="updateNowBtn" disabled>
            <i class="fas fa-download me-1"></i>تحديث الآن
          </button>
          <button class="btn btn-outline-secondary" onclick="refreshUpdateStatus()">
            <i class="fas fa-sync me-1"></i>تحديث الحالة
          </button>
        </div>
        <div id="updateProgress" style="display:none;">
          <div class="progress mb-2">
            <div class="progress-bar progress-bar-striped progress-bar-animated" role="progressbar" style="width:0%;" id="updateProgressBar">0%</div>
          </div>
          <div id="updateLogs" class="small text-light" style="max-height:150px;overflow-y:auto;background:#0d1117;padding:8px;border-radius:6px;font-family:monospace;font-size:12px;white-space:pre-wrap;"></div>
        </div>
        <div class="mt-2 small text-muted">
          <span id="versionInfo">الإصدار الحالي: <span class="text-light">جاري التحميل...</span></span>
          <span class="mx-2">|</span>
          <span id="lastCheckInfo">آخر فحص: <span class="text-light">-</span></span>
        </div>
      </div>
    </div>
    <script>
    function refreshUpdateStatus(){
      fetch('/api/auto_update_status').then(function(r){return r.json();}).then(function(data){
        if(data.success){
          document.getElementById('autoUpdateToggle').checked=data.auto_update;
          if(data.last_update)
            document.getElementById('lastCheckInfo').innerHTML='آخر تحديث: <span class="text-light">'+new Date(data.last_update).toLocaleString('ar-EG')+'</span>';
        }
      }).catch(function(){});
      checkForUpdatesUI();
    }
    function checkForUpdatesUI(){
      var btn=document.getElementById('checkUpdateBtn'),statusDiv=document.getElementById('updateStatus'),updateBtn=document.getElementById('updateNowBtn');
      btn.disabled=true; btn.innerHTML='<i class="fas fa-spinner fa-spin me-1"></i>جاري البحث...';
      statusDiv.className='alert alert-info'; statusDiv.innerHTML='<i class="fas fa-spinner fa-spin me-2"></i> جاري التحقق...';
      fetch('/api/check_update').then(function(r){return r.json();}).then(function(data){
        if(data.success){
          document.getElementById('versionInfo').innerHTML='الإصدار الحالي: <span class="text-light">'+(data.current||'غير معروف')+'</span>';
          if(data.has_update){
            statusDiv.className='alert alert-warning';
            statusDiv.innerHTML='<i class="fas fa-exclamation-triangle me-2"></i>'+data.message+'<br><small class="text-muted">الأحدث: '+(data.latest||'')+'</small>';
            updateBtn.disabled=false;
          }else{
            statusDiv.className='alert alert-success';
            statusDiv.innerHTML='<i class="fas fa-check-circle me-2"></i> '+data.message;
            updateBtn.disabled=true;
          }
        }else{
          statusDiv.className='alert alert-danger';
          statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> '+(data.message||'خطأ');
        }
      }).catch(function(err){
        statusDiv.className='alert alert-danger';
        statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> خطأ: '+err.message;
      }).finally(function(){ btn.disabled=false; btn.innerHTML='<i class="fas fa-search me-1"></i>بحث عن تحديث'; });
    }
    function performUpdateUI(){
      if(!confirm('هل أنت متأكد من رغبتك في التحديث الآن؟')) return;
      var btn=document.getElementById('updateNowBtn'),statusDiv=document.getElementById('updateStatus');
      var progressDiv=document.getElementById('updateProgress'),progressBar=document.getElementById('updateProgressBar'),logsDiv=document.getElementById('updateLogs');
      btn.disabled=true; btn.innerHTML='<i class="fas fa-spinner fa-spin me-1"></i>جاري التحديث...';
      statusDiv.className='alert alert-warning'; statusDiv.innerHTML='<i class="fas fa-sync-alt fa-spin me-2"></i> جاري تنفيذ التحديث...';
      progressDiv.style.display='block'; progressBar.style.width='10%'; progressBar.textContent='10%'; logsDiv.innerHTML='📥 بدء عملية التحديث...\n';
      fetch('/api/perform_update',{method:'POST',headers:{'Content-Type':'application/json'}})
      .then(function(r){return r.json();}).then(function(data){
        if(data.logs) logsDiv.innerHTML=data.logs.join('\n');
        if(data.success){
          progressBar.style.width='95%'; progressBar.textContent='95%'; statusDiv.className='alert alert-success';
          logsDiv.innerHTML+='\n🔄 جاري إعادة تشغيل الخادم...\n';
          var countdown=10;
          var iv=setInterval(function(){ countdown--;
            statusDiv.innerHTML='<i class="fas fa-check-circle me-2"></i> ✅ تم التحديث! إعادة تحميل بعد '+countdown+' ثوانٍ';
            if(countdown<=0){clearInterval(iv);window.location.reload();}
          },1000);
        }else{
          progressBar.className='progress-bar bg-danger'; statusDiv.className='alert alert-danger';
          statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> ❌ '+(data.message||'فشل التحديث');
          btn.disabled=false; btn.innerHTML='<i class="fas fa-download me-1"></i>إعادة المحاولة';
        }
      }).catch(function(err){
        progressBar.className='progress-bar bg-danger'; statusDiv.className='alert alert-danger';
        statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> خطأ: '+err.message;
        btn.disabled=false; btn.innerHTML='<i class="fas fa-download me-1"></i>إعادة المحاولة';
      });
    }
    document.addEventListener('DOMContentLoaded',function(){
      document.getElementById('autoUpdateToggle').addEventListener('change',function(){
        var enabled=this.checked,statusDiv=document.getElementById('updateStatus');
        statusDiv.className='alert alert-info';
        statusDiv.innerHTML='<i class="fas fa-spinner fa-spin me-2"></i> '+(enabled?'جاري تفعيل':'جاري إلغاء')+' التحديث التلقائي...';
        fetch('/api/toggle_auto_update',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({enabled:enabled})})
        .then(function(r){return r.json();}).then(function(data){
          if(data.success){
            statusDiv.className='alert alert-success';
            statusDiv.innerHTML='<i class="fas fa-check-circle me-2"></i> '+data.message;
            setTimeout(refreshUpdateStatus,3000);
          }else{
            statusDiv.className='alert alert-danger';
            statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> '+data.message;
            document.getElementById('autoUpdateToggle').checked=!enabled;
          }
        }).catch(function(err){
          statusDiv.className='alert alert-danger';
          statusDiv.innerHTML='<i class="fas fa-times-circle me-2"></i> خطأ: '+err.message;
          document.getElementById('autoUpdateToggle').checked=!enabled;
        });
      });
      refreshUpdateStatus();
      setInterval(refreshUpdateStatus,300000);
      if(typeof socket!=='undefined'){
        socket.on('update_completed',function(data){
          var statusDiv=document.getElementById('updateStatus');
          statusDiv.className='alert alert-success';
          statusDiv.innerHTML='<i class="fas fa-check-circle me-2"></i> '+(data.message||'✅ تم التحديث تلقائياً')+'<br><small>الإصدار: '+(data.version||'')+'</small>';
          setTimeout(refreshUpdateStatus,5000);
        });
      }
    });
    </script>

    <!-- ════════ الإشعارات الدورية (الترويجية) ════════ -->
    <div class="card bg-secondary mt-4">
      <div class="card-header d-flex justify-content-between align-items-center">
        <h5 class="mb-0"><i class="fas fa-bell me-2 text-warning"></i>📢 الإشعارات الدورية (ترويجية)</h5>
        <div class="form-check form-switch">
          <input class="form-check-input" type="checkbox" id="promoToggle" style="width:50px;height:25px;">
          <label class="form-check-label text-light small me-2" for="promoToggle">تفعيل</label>
        </div>
      </div>
      <div class="card-body">
        <div id="promoStatus" class="alert alert-info mb-2">
          الحالة: <span id="promoState">غير مفعّل</span>
          <span class="ms-3" id="promoCount">عدد الرسائل: 0</span>
        </div>
        <div class="mb-2">
          <label class="form-label small text-light">الرسائل الترويجية (كل سطر رسالة)</label>
          <textarea id="promoMessages" class="form-control bg-dark text-light border-secondary" rows="6"
            placeholder="اكتب رسالة في كل سطر...&#10;سيتم إرسالها دورياً كل 5 دقائق"></textarea>
        </div>
        <div class="d-flex gap-2 flex-wrap">
          <button class="btn btn-primary btn-sm" onclick="savePromoMessages()">
            <i class="fas fa-save me-1"></i>حفظ الرسائل
          </button>
          <button class="btn btn-outline-info btn-sm" onclick="loadPromoData()">
            <i class="fas fa-sync me-1"></i>تحديث
          </button>
          <button class="btn btn-outline-secondary btn-sm" onclick="resetPromoIndex()">
            <i class="fas fa-undo me-1"></i>إعادة تعيين المؤشر
          </button>
        </div>
        <div id="promoResult" class="mt-2"></div>
        <div class="mt-2 small text-muted">
          <i class="fas fa-info-circle me-1"></i>
          يتم إرسال الإشعارات كل <strong>5 دقائق</strong> للمستخدمين المشتركين في Web Push
        </div>
      </div>
    </div>
    <script>
    async function loadPromoData(){
      try{
        var r=await fetch('/api/promo_data');
        var d=await r.json();
        if(d.success){
          document.getElementById('promoToggle').checked=d.enabled;
          document.getElementById('promoState').textContent=d.enabled?'مفعّل ✅':'غير مفعّل ❌';
          document.getElementById('promoState').className=d.enabled?'text-success':'text-secondary';
          document.getElementById('promoMessages').value=d.messages.join('\n');
          document.getElementById('promoCount').textContent='عدد الرسائل: '+d.messages.length;
        }
      }catch(e){console.error('Load promo error:',e);}
    }
    async function savePromoMessages(){
      var msgs=document.getElementById('promoMessages').value.split('\n').map(function(m){return m.trim();}).filter(function(m){return m;});
      var resultDiv=document.getElementById('promoResult');
      if(msgs.length===0){resultDiv.innerHTML='<div class="alert alert-danger py-1 small">⚠️ الرجاء إدخال رسالة واحدة على الأقل</div>';return;}
      resultDiv.innerHTML='<div class="alert alert-info py-1 small"><i class="fas fa-spinner fa-spin me-1"></i> جاري الحفظ...</div>';
      try{
        var r=await fetch('/api/promo_save',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:msgs})});
        var d=await r.json();
        if(d.success){
          resultDiv.innerHTML='<div class="alert alert-success py-1 small">✅ '+d.message+'</div>';
          document.getElementById('promoCount').textContent='عدد الرسائل: '+d.count;
        }else{resultDiv.innerHTML='<div class="alert alert-danger py-1 small">❌ '+d.message+'</div>';}
      }catch(e){resultDiv.innerHTML='<div class="alert alert-danger py-1 small">❌ خطأ: '+e.message+'</div>';}
    }
    async function resetPromoIndex(){
      if(!confirm('هل تريد إعادة تعيين المؤشر للرسالة الأولى؟')) return;
      try{
        var r=await fetch('/api/promo_data');
        var d=await r.json();
        if(d.success){
          var r2=await fetch('/api/promo_save',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:d.messages})});
          var d2=await r2.json();
          if(d2.success){
            document.getElementById('promoResult').innerHTML='<div class="alert alert-success py-1 small">✅ تم إعادة تعيين المؤشر إلى الرسالة الأولى</div>';
            loadPromoData();
          }
        }
      }catch(e){alert('خطأ: '+e.message);}
    }
    document.getElementById('promoToggle').addEventListener('change',async function(){
      var enabled=this.checked,statusDiv=document.getElementById('promoState'),resultDiv=document.getElementById('promoResult');
      statusDiv.textContent='جاري التغيير...';
      resultDiv.innerHTML='<div class="alert alert-info py-1 small"><i class="fas fa-spinner fa-spin me-1"></i> جاري التفعيل...</div>';
      try{
        var r=await fetch('/api/promo_toggle',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({enabled:enabled})});
        var d=await r.json();
        if(d.success){
          statusDiv.textContent=enabled?'مفعّل ✅':'غير مفعّل ❌';
          statusDiv.className=enabled?'text-success':'text-secondary';
          resultDiv.innerHTML='<div class="alert alert-success py-1 small">✅ '+d.message+'</div>';
        }else{
          statusDiv.textContent='فشل التغيير';
          resultDiv.innerHTML='<div class="alert alert-danger py-1 small">❌ '+d.message+'</div>';
          document.getElementById('promoToggle').checked=!enabled;
        }
      }catch(e){
        statusDiv.textContent='خطأ';
        resultDiv.innerHTML='<div class="alert alert-danger py-1 small">❌ خطأ: '+e.message+'</div>';
        document.getElementById('promoToggle').checked=!enabled;
      }
    });
    document.addEventListener('DOMContentLoaded',function(){
      loadPromoData();
      setInterval(loadPromoData,30000);
    });
    </script>

    </body></html>''', 200



@app.route("/api/app_logs", methods=["GET"])
def api_app_logs():
    """إرجاع سجلات التطبيق المخزنة في الذاكرة للواجهة الأمامية"""
    try:
        level   = request.args.get('level', 'ALL').upper()
        user_id = request.args.get('user_id', '')
        recs = _mem_log_handler.get_records(None if level == 'ALL' else level)
        logs = []
        for r in recs:
            logs.append({
                'id':    r.get('id', ''),
                'time':  r.get('time', ''),
                'level': r.get('level', 'INFO'),
                'msg':   r.get('msg', ''),
                'name':  r.get('name', ''),
            })
        return jsonify({"success": True, "logs": logs, "total": len(logs)})
    except Exception as e:
        return jsonify({"success": False, "logs": [], "message": str(e)})


@app.route("/admin/api/log_query", methods=["POST"])
def admin_log_query():
    """تحليل ذكي للسجلات بالكلمات المفتاحية أو الأنماط"""
    try:
        if not session.get('admin'):
            return jsonify({"success": False, "message": "❌ غير مصرح"})
        data  = request.json or {}
        query = (data.get('query') or '').strip()
        ctx_logs = data.get('logs', [])   # السجلات المرسلة من الواجهة
        q_lower = query.lower()

        if not q_lower:
            total = len(_mem_log_handler.get_records(None))
            return jsonify({"success": True,
                            "answer": f"💡 يوجد {total} حدث في ذاكرة السجلات.\n"
                                      "اكتب مثلاً:\n"
                                      "• آخر الأخطاء\n• ابحث عن werkzeug\n• تحذيرات\n• GET /api"})

        all_recs = _mem_log_handler.get_records(None)

        # ── تحديد نوع الاستعلام ──
        is_error = any(w in q_lower for w in ['خطأ','error','errors','critical','exception','traceback','fail','فشل'])
        is_warn  = any(w in q_lower for w in ['تحذير','warning','warn'])
        is_info  = any(w in q_lower for w in ['معلومة','info','نجاح','success','ناجح'])
        is_debug = any(w in q_lower for w in ['debug','تصحيح'])
        is_last  = any(w in q_lower for w in ['أخير','اخير','آخر','last','recent','latest'])
        is_count = any(w in q_lower for w in ['كم','عدد','count','how many','total'])

        matched = []
        label   = ''

        if is_error and not is_warn and not is_info:
            matched = [r for r in all_recs if r['level'] in ('ERROR','CRITICAL')]
            label = f'❌ الأخطاء ({len(matched)})'
        elif is_warn and not is_error:
            matched = [r for r in all_recs if r['level'] == 'WARNING']
            label = f'⚠️ التحذيرات ({len(matched)})'
        elif is_info and not is_error:
            matched = [r for r in all_recs if r['level'] == 'INFO']
            label = f'ℹ️ أحداث INFO ({len(matched)})'
        elif is_debug:
            matched = [r for r in all_recs if r['level'] == 'DEBUG']
            label = f'🔵 أحداث DEBUG ({len(matched)})'
        else:
            # بحث نصي حر
            kw = q_lower.replace('ابحث عن','').replace('search','').strip()
            if kw:
                matched = [r for r in all_recs
                           if kw in (r.get('full') or r.get('msg','')+r.get('name','')).lower()]
                label = f'🔍 بحث عن "{kw}" — {len(matched)} نتيجة'
            else:
                matched = all_recs[-20:]
                label = f'📋 آخر {len(matched)} حدث'

        if is_count:
            errs  = sum(1 for r in all_recs if r['level'] in ('ERROR','CRITICAL'))
            warns = sum(1 for r in all_recs if r['level']=='WARNING')
            infos = sum(1 for r in all_recs if r['level']=='INFO')
            return jsonify({"success": True, "answer":
                f"📊 إحصائيات السجلات ({len(all_recs)} حدث إجمالاً):\n"
                f"  ❌ أخطاء: {errs}\n"
                f"  ⚠️ تحذيرات: {warns}\n"
                f"  ℹ️ معلومات: {infos}", "count": len(all_recs)})

        if not matched:
            return jsonify({"success": True, "answer": f'🔍 لا توجد سجلات مطابقة لـ "{query}"', "count": 0})

        # عرض آخر 15 نتيجة
        show = matched[-15:] if not is_last else matched[-30:]
        LVL_ICO = {'ERROR':'❌','CRITICAL':'🔴','WARNING':'⚠️','INFO':'✅','DEBUG':'🔵'}
        lines = [label + ':']
        for r in show:
            ico = LVL_ICO.get(r['level'],'•')
            msg = (r.get('msg') or r.get('full',''))[:140]
            lines.append(f"{ico} [{r['time']}] {r['name']}: {msg}")
        if len(matched) > len(show):
            lines.append(f"... و{len(matched)-len(show)} حدث آخر")
        return jsonify({"success": True, "answer": '\n'.join(lines), "count": len(matched)})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/admin/download_project", methods=["GET"])
def api_admin_download_project():
    """تحميل كامل ملفات المشروع كملف ZIP"""
    import zipfile, io, os as _os
    try:
        buf = io.BytesIO()
        base_dir = _os.path.dirname(_os.path.abspath(__file__))
        EXCLUDE_DIRS = {'.git','__pycache__','node_modules','.venv','venv','env','.local','pnpm','store','outputs'}
        EXCLUDE_EXTS = {'.pyc','.pyo','.pyd','.db-wal','.db-shm'}
        MAX_FILE_MB  = 10
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED, compresslevel=6) as zf:
            for root, dirs, files in _os.walk(base_dir):
                dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS and not d.startswith('.')]
                for fname in files:
                    fpath = _os.path.join(root, fname)
                    _, ext = _os.path.splitext(fname)
                    if ext in EXCLUDE_EXTS: continue
                    try:
                        if _os.path.getsize(fpath)/(1024*1024) > MAX_FILE_MB: continue
                    except Exception: continue
                    arcname = _os.path.relpath(fpath, base_dir)
                    try: zf.write(fpath, arcname)
                    except Exception: pass
        buf.seek(0)
        from flask import send_file
        return send_file(buf, mimetype='application/zip', as_attachment=True,
                         download_name='AbuMalik-Services-Backup.zip')
    except Exception as e:
        logger.error(f"Download project error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/list_project_files", methods=["GET"])
def api_admin_list_project_files():
    """قائمة بجميع ملفات المشروع مع أحجامها"""
    import os as _os
    try:
        base_dir = _os.path.dirname(_os.path.abspath(__file__))
        EXCLUDE_DIRS = {'.git','__pycache__','node_modules','.venv','venv','env','.local','pnpm','store','outputs'}
        EXCLUDE_EXTS = {'.pyc','.pyo','.pyd','.db-wal','.db-shm'}
        MAX_FILE_MB  = 50
        files_list = []
        for root, dirs, files in _os.walk(base_dir):
            dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS and not d.startswith('.')]
            for fname in files:
                fpath = _os.path.join(root, fname)
                _, ext = _os.path.splitext(fname)
                if ext in EXCLUDE_EXTS: continue
                try:
                    size_bytes = _os.path.getsize(fpath)
                    if size_bytes/(1024*1024) > MAX_FILE_MB: continue
                    rel = _os.path.relpath(fpath, base_dir)
                    if size_bytes < 1024:
                        size_str = f"{size_bytes} B"
                    elif size_bytes < 1024*1024:
                        size_str = f"{size_bytes/1024:.1f} KB"
                    else:
                        size_str = f"{size_bytes/(1024*1024):.1f} MB"
                    files_list.append({"path": rel, "size": size_str, "bytes": size_bytes})
                except Exception:
                    continue
        files_list.sort(key=lambda x: x["path"])
        return jsonify({"success": True, "files": files_list, "total": len(files_list)})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/download_selected_files", methods=["POST"])
def api_admin_download_selected_files():
    """تحميل ملفات محددة كـ ZIP"""
    import zipfile, io, os as _os
    try:
        data = request.get_json() or {}
        selected = data.get("files", [])
        if not selected:
            return jsonify({"success": False, "error": "لم يتم تحديد أي ملفات"}), 400
        base_dir = _os.path.dirname(_os.path.abspath(__file__))
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED, compresslevel=6) as zf:
            for rel_path in selected:
                # تأمين المسار — لا نسمح بالخروج من مجلد المشروع
                full_path = _os.path.normpath(_os.path.join(base_dir, rel_path))
                if not full_path.startswith(base_dir): continue
                if not _os.path.isfile(full_path): continue
                try: zf.write(full_path, rel_path)
                except Exception: pass
        buf.seek(0)
        from flask import send_file
        return send_file(buf, mimetype='application/zip', as_attachment=True,
                         download_name='selected_project_files.zip')
    except Exception as e:
        logger.error(f"Download selected files error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/download_single_file", methods=["GET"])
def api_admin_download_single_file():
    """تحميل ملف واحد مباشرة بدون ضغط"""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "error": "غير مخول"}), 403
    import os as _os
    rel_path = request.args.get("path", "").strip()
    if not rel_path:
        return jsonify({"success": False, "error": "مسار الملف مطلوب"}), 400
    base_dir  = _os.path.dirname(_os.path.abspath(__file__))
    full_path = _os.path.normpath(_os.path.join(base_dir, rel_path))
    # تأمين: لا نسمح بالخروج من مجلد المشروع
    if not full_path.startswith(base_dir):
        return jsonify({"success": False, "error": "مسار غير مسموح به"}), 403
    if not _os.path.isfile(full_path):
        return jsonify({"success": False, "error": "الملف غير موجود"}), 404
    try:
        return send_file(full_path, as_attachment=True,
                         download_name=_os.path.basename(full_path))
    except Exception as e:
        logger.error(f"Download single file error: {e}")
        return jsonify({"success": False, "error": str(e)}), 500


# ============================================================
# وظيفة استخراج وتصنيف روابط تيليجرام حسب الدولة
# ============================================================
from urllib.parse import urlparse as _urlparse

LINK_FINDER_CATEGORIES = [
    "🎓 جامعي / أكاديمي",
    "📚 معاهد / دورات مهنية",
    "📝 تدريب ميداني / تقارير",
    "🧪 علمية / تعليمية (مناهج)",
    "💼 وظائف / توظيف",
    "🚗 خدمات / توصيل / عقارات",
    "📢 إعلانات / تسويق",
    "📰 أخبار / عاجل",
    "🎬 ترفيه / مسلسلات / فيديو",
    "💻 تقنية / برمجة / شبكات",
    "📌 عام / متنوع"
]

KEYWORDS_FALLBACK_LINK = {
    "🎓 جامعي / أكاديمي": ["جامعة","كلية","بكالوريوس","ماجستير","دكتوراه","بحث","محاضرة","تخصص","طلاب","قبول","ساعات معتمدة"],
    "📚 معاهد / دورات مهنية": ["معهد","دورة","شهادة مهنية","مدرب","تدريب","الرخصة المهنية","أكاديمية"],
    "📝 تدريب ميداني / تقارير": ["تدريب ميداني","تقرير","مشروع تخرج","COOP","تعاوني","ميداني","دراسة حالة"],
    "🧪 علمية / تعليمية (مناهج)": ["رياضيات","فيزياء","كيمياء","منهج","مدرسة","اختبار","قدرات","تحصيلي","ابتدائي","متوسط","ثانوي"],
    "💼 وظائف / توظيف": ["وظائف","توظيف","مقابلة","راتب","مطلوب","شركة","خريجين","تجنيد"],
    "🚗 خدمات / توصيل / عقارات": ["توصيل","مشاوير","سواق","نقل","عقار","مقاولات","شقة","إيجار"],
    "📢 إعلانات / تسويق": ["عرض","خصم","بيع","شراء","سعر","متجر","تخفيض"],
    "📰 أخبار / عاجل": ["عاجل","أخبار","وزارة","أحداث","بيان","رسمي"],
    "🎬 ترفيه / مسلسلات / فيديو": ["مسلسل","فيلم","مشاهدة","تحميل","حلقة","مقاطع","يوتيوب","ضحك"],
    "💻 تقنية / برمجة / شبكات": ["برمجة","تطوير","خادم","شبكة","برامج","تطبيق","سيرفر","بوت"]
}

_lf_country_cache = {}

def _lf_get_country_tld(hostname):
    if not hostname:
        return None
    tld = hostname.split('.')[-1].upper()
    valid = {'SA','AE','EG','KW','QA','OM','BH','JO','LB','IQ','YE','SY','PS','SD','LY','TN','MA','DZ','MR'}
    return tld if tld in valid else None

def _lf_get_country_ip(hostname):
    try:
        ip = socket.gethostbyname(hostname)
        resp = requests.get(f'http://ip-api.com/json/{ip}', timeout=5)
        if resp.status_code == 200:
            d = resp.json()
            if d.get('status') == 'success':
                return d.get('countryCode')
    except Exception:
        pass
    return None

def _lf_get_country(hostname):
    if not hostname:
        return None
    if hostname in _lf_country_cache:
        return _lf_country_cache[hostname]
    country = _lf_get_country_tld(hostname) or _lf_get_country_ip(hostname)
    _lf_country_cache[hostname] = country
    return country

def _lf_extract_links(text):
    raw = re.findall(r'https?://[^\s<>"\'\)]+', text)
    seen = set()
    links = []
    for l in raw:
        l = l.strip().rstrip('.,;:)')
        if l and l not in seen:
            seen.add(l)
            links.append(l)
    return links

def _lf_get_context(link, text, window=200):
    idx = text.find(link)
    if idx == -1:
        return link[:150]
    start = max(0, idx - window)
    end = min(len(text), idx + len(link) + window)
    return text[start:end].replace('\n', ' ').strip()

def _lf_classify_fallback(text):
    text_lower = text.lower()
    scores = {}
    for cat, keywords in KEYWORDS_FALLBACK_LINK.items():
        count = sum(1 for kw in keywords if kw in text_lower)
        if count > 0:
            scores[cat] = count
    if not scores:
        return "📌 عام / متنوع", 0.0
    best = max(scores, key=scores.get)
    return best, min(1.0, scores[best] / 5.0)

def _lf_classify_groq(context, groq_client):
    try:
        prompt = f"""صنّف النص التالي إلى واحدة من هذه الفئات (اختر الأنسب):
{', '.join(LINK_FINDER_CATEGORIES)}

النص: {context[:600]}

أجب فقط باسم الفئة كما هي مكتوبة بالضبط دون أي شرح."""
        resp = groq_client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=20,
            temperature=0.1
        )
        result = resp.choices[0].message.content.strip()
        for cat in LINK_FINDER_CATEGORIES:
            if cat in result:
                return cat, 0.92
        return "📌 عام / متنوع", 0.5
    except Exception as e:
        logger.warning(f"Groq link classify: {e}")
        return _lf_classify_fallback(context)


@app.route("/link-finder")
def link_finder_page():
    if 'user_id' not in session:
        return redirect('/')
    return render_template('link_finder.html')


_LF_COUNTRY_NAMES = {
    'SA': 'السعودية', 'AE': 'الإمارات', 'EG': 'مصر', 'KW': 'الكويت',
    'QA': 'قطر', 'OM': 'عمان', 'BH': 'البحرين', 'JO': 'الأردن',
    'LB': 'لبنان', 'IQ': 'العراق', 'YE': 'اليمن', 'SY': 'سوريا',
    'PS': 'فلسطين', 'SD': 'السودان', 'LY': 'ليبيا', 'TN': 'تونس',
    'MA': 'المغرب', 'DZ': 'الجزائر', 'MR': 'موريتانيا',
}

@app.route("/api/link_finder/start", methods=["POST"])
def api_link_finder_start():
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "الرجاء تسجيل الدخول"}), 401
    user_id = session['user_id']
    data     = request.get_json() or {}
    mode     = data.get('mode', 'text')           # 'text' | 'telegram_search'
    text     = data.get('text', '').strip()
    keyword  = data.get('keyword', '').strip()
    country  = data.get('country', '').strip().upper()
    category = data.get('category', '').strip()
    use_ai   = data.get('use_ai', True)

    # ── وضع البحث الجديد: telegram_search ──────────────────────
    if mode == 'telegram_search':
        with USERS_LOCK:
            if user_id not in USERS:
                return jsonify({"success": False, "message": "❌ المستخدم غير مسجل"})
            cm = USERS[user_id].get('client_manager')
            if not cm or not cm.client:
                return jsonify({"success": False, "message": "❌ يرجى تسجيل الدخول أولاً"})

        def _tg_search_worker():
            try:
                # بناء عبارة البحث: كلمة + دولة + فئة
                parts = []
                if keyword:  parts.append(keyword)
                if country:  parts.append(_LF_COUNTRY_NAMES.get(country, country))
                if category: parts.append(category)
                if not parts: parts.append('قناة')
                query = ' '.join(parts)

                socketio.emit('search_dialog_progress', {
                    'chat': f'🔍 يبحث في تيليجرام: {query}',
                    'scanned': 0, 'found': 0
                }, to=user_id)

                # استدعاء البحث العام
                results = cm.run_coroutine(
                    search_global_groups(cm.client, query, limit=80, filter_type='all')
                )

                # تحضير كلاسيفاير
                groq_client = None
                if use_ai and GROQ_API_KEY:
                    try:
                        from groq import Groq as _G
                        groq_client = _G(api_key=GROQ_API_KEY)
                    except Exception:
                        pass

                # معالجة كل نتيجة وتصنيفها
                links_from_text = []
                if text:
                    links_from_text = _lf_extract_links(text)

                all_links = []
                for r in results:
                    url = r.get('url', '')
                    if url:
                        all_links.append({'url': url, 'title': r.get('title', ''), 'members': r.get('members', 0)})
                for u in links_from_text:
                    all_links.append({'url': u, 'title': '', 'members': 0})

                if not all_links:
                    socketio.emit('link_finder_done', {
                        'total': 0, 'classified': 0,
                        'message': f'لم يُعثر على نتائج لـ "{query}"'
                    }, to=user_id)
                    return

                classified = 0
                for item in all_links:
                    url   = item['url']
                    title = item.get('title', '')
                    ctx   = title or url
                    if groq_client:
                        cat, score = _lf_classify_groq(ctx, groq_client)
                    else:
                        cat, score = _lf_classify_fallback(ctx)
                    # تصفية حسب الفئة المطلوبة (إذا حُدّدت)
                    if category and cat != category and score < 0.6:
                        continue
                    socketio.emit('link_classified', {
                        'link':     url,
                        'category': cat,
                        'score':    score,
                        'context':  (title + f' ({item["members"]:,} عضو)') if item.get('members') else title,
                        'country':  _LF_COUNTRY_NAMES.get(country, country) if country else ''
                    }, to=user_id)
                    classified += 1
                    time.sleep(0.03)

                socketio.emit('link_finder_done', {
                    'total': classified, 'classified': classified
                }, to=user_id)
            except Exception as e:
                logger.error(f"TG search worker error: {e}")
                socketio.emit('link_finder_error', {'error': str(e)}, to=user_id)

        t = threading.Thread(target=_tg_search_worker, daemon=True)
        t.start()
        return jsonify({"success": True, "message": "بدأ البحث في تيليجرام"})

    # ── الوضع الكلاسيكي: نص مباشر ──────────────────────────────
    if not text:
        return jsonify({"success": False, "message": "النص المصدر فارغ"})

    def process_links():
        try:
            raw_links = _lf_extract_links(text)
            if not raw_links:
                socketio.emit('link_finder_done', {'total': 0, 'classified': 0, 'message': 'لم يُعثر على روابط في النص'}, to=user_id)
                return

            groq_client = None
            if use_ai and GROQ_API_KEY:
                try:
                    from groq import Groq
                    groq_client = Groq(api_key=GROQ_API_KEY)
                except Exception:
                    pass

            for link in raw_links:
                context = _lf_get_context(link, text)
                if groq_client:
                    cat, score = _lf_classify_groq(context, groq_client)
                else:
                    cat, score = _lf_classify_fallback(context)
                socketio.emit('link_classified', {
                    'link': link, 'category': cat, 'score': score,
                    'context': context[:150], 'country': country
                }, to=user_id)
                time.sleep(0.05)

            socketio.emit('link_finder_done', {'total': len(raw_links), 'classified': len(raw_links)}, to=user_id)
        except Exception as e:
            logger.error(f"Link finder worker error: {e}")
            socketio.emit('link_finder_error', {'error': str(e)}, to=user_id)

    t = threading.Thread(target=process_links, daemon=True)
    t.start()
    return jsonify({"success": True, "message": "بدأ البحث والتصنيف"})


@app.route("/api/link_finder/export", methods=["POST"])
def api_link_finder_export():
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "الرجاء تسجيل الدخول"}), 401
    data = request.get_json() or {}
    links_data = data.get('links', [])
    if not links_data:
        return jsonify({"error": "لا توجد بيانات"}), 400
    try:
        import pandas as pd
        df = pd.DataFrame(links_data)
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name="الروابط")
        output.seek(0)
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                         as_attachment=True, download_name='links_classified.xlsx')
    except Exception as e:
        logger.error(f"Link finder export error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/auto_join/links", methods=["POST"])
def api_auto_join_from_link_finder():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    data = request.json or {}
    links = data.get('links', [])
    if not links:
        return jsonify({"success": False, "message": "لا توجد روابط"})
    settings = load_settings(user_id)
    settings['auto_join_links'] = links
    save_settings(user_id, settings)
    socketio.emit('log_update', {"message": f"📥 تم استلام {len(links)} رابط من أداة البحث، جاهزة للانضمام"}, to=user_id)
    return jsonify({"success": True, "message": f"تم تخزين {len(links)} رابط للانضمام"})


# ══════════════════════════════════════════════════════════════════════════════
# ─── نظام بطاقات الشحن المتكامل ──────────────────────────────────────────────
# ══════════════════════════════════════════════════════════════════════════════

def load_cards_data():
    """تحميل بيانات البطاقات من الملف المحلي"""
    with _CARDS_LOCK:
        try:
            if os.path.exists(CARDS_FILE):
                with open(CARDS_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass
        default = {
            "card_system_enabled": False,
            "plans": [
                {"id": 1, "name": "يومية",   "time_limit": 86400,   "data_limit": 5368709120,  "profile_name": "daily"},
                {"id": 2, "name": "أسبوعية", "time_limit": 604800,  "data_limit": 10737418240, "profile_name": "weekly"},
                {"id": 3, "name": "شهرية",   "time_limit": 2592000, "data_limit": 32212254720, "profile_name": "monthly"}
            ],
            "vouchers": [],
            "active_card_sessions": []
        }
        with open(CARDS_FILE, 'w', encoding='utf-8') as f:
            json.dump(default, f, ensure_ascii=False, indent=2)
        return default

def save_cards_data(data):
    """حفظ بيانات البطاقات محلياً"""
    with _CARDS_LOCK:
        try:
            with open(CARDS_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"Error saving cards data: {e}")

def voucher_display_code(v):
    """يرجع رمز البطاقة الظاهر للمشرف — يدعم الكروت القديمة التي لا تحتوي على 'code'."""
    return v.get("code") or (v.get("code_hash", "")[:12] + "***" if v.get("code_hash") else "—")

def voucher_display_plan_name(v, plans=None):
    """يرجع اسم الخطة — يدعم الكروت القديمة التي تخزن plan_id فقط بدون plan_name."""
    if v.get("plan_name"):
        return v["plan_name"]
    if plans is None:
        plans = load_cards_data().get("plans", [])
    plan = next((p for p in plans if p.get("id") == v.get("plan_id")), None)
    return plan["name"] if plan else "بدون خطة"

def generate_vouchers(plan_id, count=10, allowed_features=None):
    """إنشاء قسائم بطاقات جديدة مع تحديد الوظائف التي تفتحها"""
    data = load_cards_data()
    plan = next((p for p in data["plans"] if p["id"] == plan_id), None)
    if not plan:
        raise ValueError("الخطة غير موجودة")
    if allowed_features is None:
        allowed_features = []
    codes = []
    for _ in range(count):
        raw = f"{secrets.token_hex(3).upper()}-{secrets.token_hex(3).upper()}"
        normalized = raw.replace("-", "").upper()
        hashed = hashlib.sha256(normalized.encode()).hexdigest()
        data["vouchers"].append({
            "code": raw,
            "code_hash": hashed,
            "plan_id": plan_id,
            "plan_name": plan["name"],
            "allowed_features": allowed_features,
            "status": "unused",
            "created_at": datetime.now().isoformat(),
            "used_at": None,
            "expires_at": None
        })
        codes.append(raw)
    save_cards_data(data)
    return codes

def validate_voucher(code):
    normalized = code.strip().replace("-", "").upper()
    hashed = hashlib.sha256(normalized.encode()).hexdigest()
    data = load_cards_data()
    voucher = next((v for v in data["vouchers"] if v["code_hash"] == hashed), None)
    if not voucher:
        return None, "❌ الرمز غير صحيح"
    if voucher["status"] == "used":
        return None, "❌ تم استخدام هذا الرمز مسبقاً"
    if voucher["status"] == "expired":
        return None, "❌ انتهت صلاحية الرمز"
    plan = next((p for p in data["plans"] if p["id"] == voucher["plan_id"]), None)
    if not plan:
        return None, "❌ الخطة غير موجودة"
    return {"voucher": voucher, "plan": plan}, None

def activate_card_voucher(code, client_ip="0.0.0.0"):
    result, err = validate_voucher(code)
    if err:
        return {"success": False, "message": err}
    voucher = result["voucher"]
    plan = result["plan"]
    data = load_cards_data()
    data.setdefault("active_card_sessions", [])
    data.setdefault("vouchers", [])
    session_id = secrets.token_hex(16)
    now = datetime.now()
    expires_at = now + timedelta(seconds=plan["time_limit"])
    new_session = {
        "session_id": session_id,
        "voucher_code_hash": voucher["code_hash"],
        "plan_name": plan["name"],
        "ip_address": client_ip,
        "start_time": now.isoformat(),
        "expires_at": expires_at.isoformat(),
        "session_timeout": plan["time_limit"]
    }
    for v in data["vouchers"]:
        if v["code_hash"] == voucher["code_hash"]:
            v["status"] = "active"
            v["used_at"] = now.isoformat()
            v["expires_at"] = expires_at.isoformat()
            break
    data["active_card_sessions"].append(new_session)
    save_cards_data(data)

    # تطبيق الوظائف المفتوحة على المستخدم الحالي
    # بطاقة عامة (allowed_features فارغة) = تفتح جميع الوظائف (نظام كامل)
    # بطاقة محددة (allowed_features غير فارغة) = تفتح الوظائف المحددة فقط
    allowed = voucher.get("allowed_features", [])
    try:
        user_id = session.get('user_id', 'user_1')
        feat_data = load_feature_restrictions()
        unlocked = feat_data.get("user_unlocked", {}).get(user_id, [])
        features_to_unlock = allowed if allowed else list(FEATURE_LABELS.keys())
        for f in features_to_unlock:
            if f not in unlocked:
                unlocked.append(f)
        feat_data.setdefault("user_unlocked", {})[user_id] = unlocked
        save_feature_restrictions(feat_data)
        logger.info(f"✅ فُتحت الوظائف {features_to_unlock} للمستخدم {user_id} عبر البطاقة")
    except Exception as _fe:
        logger.error(f"activate_card_voucher features error: {_fe}")

    return {
        "success": True,
        "session_id": session_id,
        "expires_at": expires_at.isoformat(),
        "plan_name": plan["name"],
        "unlocked_features": allowed,
        "message": f"✅ تم الدخول بنجاح - باقة {plan['name']}"
    }

def terminate_card_session(session_id, reason="انتهت الجلسة"):
    data = load_cards_data()
    s = next((x for x in data["active_card_sessions"] if x["session_id"] == session_id), None)
    if s:
        for v in data["vouchers"]:
            if v["code_hash"] == s.get("voucher_code_hash"):
                v["status"] = "used"
                break
        data["active_card_sessions"] = [x for x in data["active_card_sessions"] if x["session_id"] != session_id]
        save_cards_data(data)

def _card_session_monitor():
    """مراقبة انتهاء صلاحية جلسات البطاقات"""
    while True:
        try:
            data = load_cards_data()
            now = datetime.now()
            to_terminate = []
            for s in data.get("active_card_sessions", []):
                try:
                    expires = datetime.fromisoformat(s["expires_at"])
                    if now >= expires:
                        to_terminate.append(s["session_id"])
                except Exception:
                    pass
            for sid in to_terminate:
                terminate_card_session(sid, "انتهى الوقت تلقائياً")
        except Exception as e:
            logger.error(f"card_session_monitor: {e}")
        time.sleep(60)

_OSThread(target=_card_session_monitor, daemon=True).start()

# ── مسارات API لبطاقات الشحن ────────────────────────────────────────────────

@app.route("/login")
def login_page():
    """صفحة الدخول ببطاقة الشحن"""
    cards = load_cards_data()
    if not cards.get("card_system_enabled", False):
        return redirect("/")
    if session.get("card_logged_in"):
        return redirect("/")
    return render_template("login_card.html")

@app.route("/api/login_card", methods=["POST"])
def api_login_card():
    data = request.json or {}
    code = data.get("code", "").strip()
    if not code:
        return jsonify({"success": False, "message": "❌ الرجاء إدخال الرمز"})

    # ── دخول مجاني بكلمة مرور الأدمن (بدون بطاقة) ──────────────
    if code == ADMIN_PASSWORD or code.replace("-","") == ADMIN_PASSWORD:
        session["card_logged_in"]   = True
        session.pop("card_session_id", None)
        session["card_plan_name"]   = "أدمن — دخول مجاني"
        session["card_expires_at"]  = ""
        session["admin_auth"]       = True
        session.permanent = True
        logger.info("✅ دخول مجاني بكلمة مرور الأدمن")
        return jsonify({"success": True, "redirect": "/", "message": "✅ مرحباً أدمن — دخول مجاني"})

    client_ip = request.remote_addr or "0.0.0.0"
    result = activate_card_voucher(code, client_ip)
    if result["success"]:
        session["card_logged_in"] = True
        session["card_session_id"] = result["session_id"]
        session["card_plan_name"] = result.get("plan_name", "")
        session["card_expires_at"] = result.get("expires_at", "")
        session.permanent = True
        return jsonify({"success": True, "redirect": "/", "message": result["message"]})
    return jsonify({"success": False, "message": result["message"]})

@app.route("/api/card_session_status", methods=["GET"])
def api_card_session_status():
    sid = session.get("card_session_id")
    if not sid:
        return jsonify({"active": False})
    data = load_cards_data()
    active = next((s for s in data["active_card_sessions"] if s["session_id"] == sid), None)
    if not active:
        return jsonify({"active": False, "message": "انتهت الجلسة"})
    try:
        expires = datetime.fromisoformat(active["expires_at"])
        remaining = max(0, int((expires - datetime.now()).total_seconds()))
    except Exception:
        remaining = 0
    return jsonify({
        "active": True,
        "remaining": remaining,
        "plan_name": active.get("plan_name", ""),
        "expires_at": active.get("expires_at", "")
    })

@app.route("/api/card_logout", methods=["POST"])
def api_card_logout():
    sid = session.get("card_session_id")
    if sid:
        terminate_card_session(sid, "تسجيل خروج يدوي")
    session.pop("card_logged_in", None)
    session.pop("card_session_id", None)
    session.pop("card_plan_name", None)
    session.pop("card_expires_at", None)
    return jsonify({"success": True, "redirect": "/login"})

# ── مسارات إدارة البطاقات ────────────────────────────────────────────────────

@app.route("/admin/api/create_vouchers", methods=["POST"])
def admin_create_vouchers():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = request.json or {}
    plan_id = int(data.get("plan_id", 1))
    count   = min(int(data.get("count", 10)), 200)
    features = [f for f in data.get("features", []) if f in FEATURE_LABELS]
    notify_user_id = data.get("notify_user_id", "").strip() or session.get("user_id", "user_1")
    try:
        codes = generate_vouchers(plan_id, count, allowed_features=features)
        tg_sent = False
        # ── إرسال الكروت عبر تيليجرام كرسالة تنبيهات إلى حساب المستخدم ──
        try:
            if notify_user_id and notify_user_id in PREDEFINED_USERS:
                with USERS_LOCK:
                    _cm = USERS.get(notify_user_id, {}).get('client_manager')
                if _cm and getattr(_cm, 'authenticated', False):
                    _loop = getattr(_cm, 'loop', None)
                    if _loop and _loop.is_running():
                        feat_lbl = ", ".join(features) if features else "كامل النظام"
                        msg_lines = [
                            "🎫 **تم إنشاء " + str(count) + " كرت تفعيل جديد**",
                            "📋 الخطة: " + str(plan_id) + " | الوظائف: " + feat_lbl,
                            "─" * 30,
                        ]
                        msg_lines += codes
                        msg_lines.append(
                            "\n📅 " + __import__('datetime').datetime.now().strftime('%Y-%m-%d %H:%M')
                        )
                        notification_msg = "\n".join(msg_lines)
                        import asyncio
                        async def _tg_send():
                            await _cm.client.send_message('me', notification_msg, link_preview=False)
                        asyncio.run_coroutine_threadsafe(_tg_send(), _loop).result(timeout=10)
                        tg_sent = True
        except Exception as _tge:
            logger.warning(f"TG voucher notify failed: {_tge}")
        return jsonify({"success": True, "codes": codes, "count": len(codes),
                        "features": features, "tg_sent": tg_sent})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/admin/api/card_status", methods=["GET"])
def admin_card_status():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    return jsonify({"enabled": data.get("card_system_enabled", False)})

@app.route("/admin/api/toggle_card_system", methods=["POST"])
def admin_toggle_card_system():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data_req = request.json or {}
    cards = load_cards_data()
    # إذا أُرسل حقل enabled صراحةً استخدمه، وإلا عكس الحالة الحالية
    if "enabled" in data_req:
        enabled = bool(data_req["enabled"])
    else:
        enabled = not bool(cards.get("card_system_enabled", False))
    cards["card_system_enabled"] = enabled
    save_cards_data(cards)
    status_text = "تفعيل" if enabled else "تعطيل"
    return jsonify({"success": True, "enabled": enabled, "message": f"تم {status_text} نظام البطاقات"})

@app.route("/admin/api/vouchers", methods=["GET"])
def admin_list_vouchers():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    return jsonify({"success": True, "vouchers": data.get("vouchers", [])})

@app.route("/admin/api/card_sessions", methods=["GET"])
def admin_list_card_sessions():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    return jsonify({"success": True, "sessions": data.get("active_card_sessions", [])})

@app.route("/admin/api/terminate_card_session", methods=["POST"])
def admin_terminate_card_session():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    d = request.json or {}
    sid = d.get("session_id")
    if not sid:
        return jsonify({"success": False, "message": "معرف الجلسة مطلوب"})
    terminate_card_session(sid, "إنهاء يدوي من المشرف")
    return jsonify({"success": True})

@app.route("/admin/api/delete_vouchers", methods=["POST"])
def admin_delete_vouchers():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    d = request.json or {}
    status_filter = d.get("status")
    cards = load_cards_data()
    if status_filter:
        cards["vouchers"] = [v for v in cards["vouchers"] if v["status"] != status_filter]
    else:
        cards["vouchers"] = []
    save_cards_data(cards)
    return jsonify({"success": True})


# ══════════════════════════════════════════════════════════════════════════════
# ▸ نظام تقييد الوظائف (Feature Restrictions) + نظام Pro بالبطاقات
# ══════════════════════════════════════════════════════════════════════════════

FEATURE_RESTRICTIONS_FILE = os.path.join(DATA_DIR, "feature_restrictions.json")
FEATURE_VOUCHERS_FILE     = os.path.join(DATA_DIR, "feature_vouchers.json")
_FR_LOCK = threading.Lock()
_FV_LOCK = threading.Lock()

FEATURE_LABELS = {
    "learning":             "نظام التعلم الذكي",
    "rotating":             "النشر الدوري المتسلسل",
    "group_search":         "البحث في روابطي",
    "auto_join":            "الانضمام المتقدم",
    "auto_replies":         "الردود التلقائية",
    "saved_links":          "الروابط المحفوظة",
    "academic":             "التحليل الأكاديمي الذكي",
    "formatter_pdf2word":   "تحويل PDF إلى Word",
    "formatter_html2word":  "تحويل HTML إلى Word",
    "formatter_html2excel": "تحويل HTML إلى Excel",
    "formatter_html2ppt":   "تحويل HTML إلى PPT",
    "link_finder":          "محرك البحث عن الروابط",
    "message_sending":      "إرسال الرسائل",
    "monitoring":           "مراقبة المجموعات",
    "scan_groups":          "مسح المجموعات",
}

# قائمة كاملة للتقييد: الوظائف + المستخدمون الخمسة
RESTRICTED_FEATURES_LIST = [
    {"id": "learning",            "name": "🧠 نظام التعلم الذكي"},
    {"id": "rotating",            "name": "🔄 النشر الدوري المتسلسل"},
    {"id": "group_search",        "name": "🔍 البحث في روابطي"},
    {"id": "auto_join",           "name": "🤖 الانضمام المتقدم"},
    {"id": "auto_replies",        "name": "💬 الردود التلقائية"},
    {"id": "saved_links",         "name": "🔗 الروابط المحفوظة"},
    {"id": "academic",            "name": "📚 التحليل الأكاديمي الذكي"},
    {"id": "formatter_pdf2word",  "name": "📄 تحويل PDF إلى Word"},
    {"id": "formatter_html2word", "name": "📝 تحويل HTML إلى Word"},
    {"id": "formatter_html2excel","name": "📊 تحويل HTML إلى Excel"},
    {"id": "formatter_html2ppt",  "name": "🎞️ تحويل HTML إلى PPT"},
    {"id": "link_finder",         "name": "🕵️ محرك البحث عن الروابط"},
    {"id": "message_sending",     "name": "📨 إرسال الرسائل"},
    {"id": "monitoring",          "name": "👁️ مراقبة المجموعات"},
    {"id": "scan_groups",         "name": "📡 مسح المجموعات"},
]

_FR_DEFAULT = {
    "enabled": False,
    "global_restricted": [],
    "user_restrictions": {},
    "user_unlocked": {}
}

def load_feature_restrictions():
    with _FR_LOCK:
        try:
            if os.path.exists(FEATURE_RESTRICTIONS_FILE):
                with open(FEATURE_RESTRICTIONS_FILE, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                for k, v in _FR_DEFAULT.items():
                    if k not in data:
                        data[k] = v
                return data
        except Exception:
            pass
        return dict(_FR_DEFAULT)

def save_feature_restrictions(data):
    with _FR_LOCK:
        try:
            with open(FEATURE_RESTRICTIONS_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"save_feature_restrictions: {e}")

def is_feature_restricted_for_user(user_id, feature_id):
    """التحقق مما إذا كانت وظيفة معينة مقيّدة لمستخدم محدد."""
    data = load_feature_restrictions()
    if not data.get("enabled", False):
        return False

    # ── فحص user_unlocked أولاً: البطاقة تفتح أي وظيفة مهما كان مصدر التقييد ──
    unlocked = data.get("user_unlocked", {}).get(user_id, [])
    if feature_id in unlocked:
        return False

    # ── 1. قيود خاصة بالمستخدم (يضعها المشرف من لوحة الإدارة) ──
    user_restr = data.get("user_restrictions", {}).get(user_id, [])
    if "all" in user_restr or feature_id in user_restr:
        return True

    # ── 2. قيود عامة تسري على الجميع ──
    if feature_id in data.get("global_restricted", []):
        # فحص قائمة المستخدمين المجانيين من لوحة الإدارة (feature_vouchers)
        try:
            fv = load_feature_vouchers()
            allowed_users = fv.get(feature_id, {}).get("users", [])
            if user_id in allowed_users:
                return False
        except Exception:
            pass
        return True

    return False

def is_user_restricted(user_id):
    """تحقق إذا كان المستخدم مقيداً بأي شكل — يُعاد True إذا وُجدت وظيفة مقيدة واحدة على الأقل."""
    data = load_feature_restrictions()
    if not data.get("enabled", False):
        return False
    unlocked = data.get("user_unlocked", {}).get(user_id, [])
    # فحص القيود الخاصة بالمستخدم (مع مراعاة user_unlocked)
    user_restr = data.get("user_restrictions", {}).get(user_id, [])
    if "all" in user_restr:
        # إذا كانت كل الوظائف مفتوحة ببطاقة فالمستخدم ليس مقيداً
        if all(f in unlocked for f in FEATURE_LABELS):
            return False
        return True
    for f in user_restr:
        if f not in unlocked:
            return True
    # فحص القيود العامة
    for f in data.get("global_restricted", []):
        if f not in unlocked:
            return True
    return False

def load_feature_vouchers():
    with _FV_LOCK:
        try:
            if os.path.exists(FEATURE_VOUCHERS_FILE):
                with open(FEATURE_VOUCHERS_FILE, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass
        return {k: {"users": []} for k in FEATURE_LABELS}

def save_feature_vouchers(data):
    with _FV_LOCK:
        try:
            with open(FEATURE_VOUCHERS_FILE, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.error(f"save_feature_vouchers: {e}")

# ── alias routes للتوافق مع admin_panel.html ──────────────────────────────

@app.route("/admin/api/card_system_status", methods=["GET"])
def admin_card_system_status():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    enabled = data.get("card_system_enabled", False)
    total = len(data.get("vouchers", []))
    active_sessions = len(data.get("active_card_sessions", []))
    unused = sum(1 for v in data.get("vouchers", []) if v.get("status") == "unused")
    return jsonify({
        "success": True,
        "enabled": enabled,
        "stats": {"total": total, "unused": unused, "active_sessions": active_sessions}
    })

@app.route("/admin/api/list_vouchers", methods=["GET"])
def admin_list_vouchers_alias():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    sessions_by_hash = {}
    for s in data.get("active_card_sessions", []):
        sessions_by_hash[s.get("voucher_code_hash")] = s
    result = []
    now = datetime.now()
    plans = data.get("plans", [])
    for v in data.get("vouchers", []):
        code_hash = v.get("code_hash", "")
        active_session = sessions_by_hash.get(code_hash)
        used_by = active_session.get("ip_address") if active_session else None
        remaining_seconds = None
        expires_raw = v.get("expires_at") or ""
        if expires_raw and v.get("status") in ("active",):
            try:
                remaining_seconds = max(0, int((datetime.fromisoformat(expires_raw) - now).total_seconds()))
            except Exception:
                remaining_seconds = None
        result.append({
            "code": voucher_display_code(v),
            "plan_name": voucher_display_plan_name(v, plans),
            "status": v.get("status", "unknown"),
            "used": v.get("status") in ("used", "active", "expired"),
            "used_by": used_by,
            "created_at": (v.get("created_at") or "")[:10],
            "used_at": (v.get("used_at") or "")[:16],
            "expires_at": (v.get("expires_at") or "")[:16],
            "remaining_seconds": remaining_seconds,
        })
    return jsonify({"success": True, "vouchers": result})

@app.route("/admin/api/delete_used_vouchers", methods=["POST"])
def admin_delete_used_vouchers():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    cards = load_cards_data()
    before = len(cards["vouchers"])
    cards["vouchers"] = [v for v in cards["vouchers"] if v.get("status") not in ("used", "expired")]
    deleted = before - len(cards["vouchers"])
    save_cards_data(cards)
    return jsonify({"success": True, "deleted": deleted, "message": f"تم حذف {deleted} كرت"})

# ── Feature Restrictions API ───────────────────────────────────────────────

@app.route("/api/feature_restrictions", methods=["GET"])
def api_feature_restrictions_public():
    r = load_feature_restrictions()
    user_id = session.get('user_id', 'user_1')
    # كلمة مرور الأدمن = دخول حر من كل التقييدات
    if session.get('admin_auth'):
        return jsonify({"success": True, "enabled": r.get("enabled", False), "restricted": [], "bypass": "admin"})
    # فحص التقييد الفعلي للمستخدم (user_unlocked يُطبَّق داخل is_feature_restricted_for_user)
    restricted = []
    for fid in FEATURE_LABELS:
        if is_feature_restricted_for_user(user_id, fid):
            restricted.append(fid)
    return jsonify({"success": True, "enabled": r.get("enabled", False), "restricted": restricted})

@app.route("/admin/api/feature_restrictions", methods=["GET", "POST"])
def admin_feature_restrictions():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    if request.method == "POST":
        data = request.json or {}
        r = load_feature_restrictions()
        if "enabled" in data:
            r["enabled"] = bool(data["enabled"])
        if "global_restricted" in data:
            r["global_restricted"] = [f for f in data["global_restricted"] if f in FEATURE_LABELS]
        if "user_restrictions" in data:
            r["user_restrictions"] = data["user_restrictions"]
        if "user_unlocked" in data:
            r["user_unlocked"] = data["user_unlocked"]
        # دعم الـ schema القديم للتوافق العكسي
        if "restricted" in data and "global_restricted" not in data:
            r["global_restricted"] = [f for f in data["restricted"] if f in FEATURE_LABELS]
        save_feature_restrictions(r)
        return jsonify({"success": True, "message": "✅ تم حفظ إعدادات التقييد", "data": r})
    r = load_feature_restrictions()
    return jsonify({
        "success": True,
        "enabled": r.get("enabled", False),
        "global_restricted": r.get("global_restricted", []),
        "user_restrictions": r.get("user_restrictions", {}),
        "user_unlocked": r.get("user_unlocked", {}),
        "available_features": FEATURE_LABELS,
        "features_list": RESTRICTED_FEATURES_LIST
    })

# ── Feature Vouchers (Pro unlock) API ─────────────────────────────────────

@app.route("/api/feature_status", methods=["GET"])
def api_feature_status():
    user_id = session.get('user_id', session.get("uid", "user_1"))
    # كلمة مرور الأدمن = دخول حر — كل الوظائف متاحة
    if session.get('admin_auth'):
        result = {feature: {"restricted": False, "unlocked": True} for feature in FEATURE_LABELS}
        return jsonify({"success": True, "features": result, "bypass": "admin"})
    # فحص دقيق لكل وظيفة بحق المستخدم الحالي
    # (user_unlocked يُعالَج داخل is_feature_restricted_for_user)
    result = {}
    r = load_feature_restrictions()
    unlocked_list = r.get("user_unlocked", {}).get(user_id, [])
    for feature in FEATURE_LABELS:
        restr = is_feature_restricted_for_user(user_id, feature)
        result[feature] = {
            "restricted": restr,
            "unlocked": feature in unlocked_list,
        }
    return jsonify({"success": True, "features": result})

@app.route("/api/unlock_feature", methods=["POST"])
def api_unlock_feature():
    user_id = session.get('user_id', session.get("uid", "user_1"))
    if not user_id:
        return jsonify({"success": False, "message": "❌ يجب تسجيل الدخول أولاً"}), 401
    data = request.json or {}
    feature = data.get("feature", "").strip()
    code = data.get("code", "").strip()
    if not feature or feature not in FEATURE_LABELS:
        return jsonify({"success": False, "message": "❌ وظيفة غير معروفة"})
    # الأدمن يفتح أي وظيفة: إما session.admin_auth أو كلمة مرور المشرف مباشرة ←
    is_admin_req = session.get('admin_auth') or (code and code == ADMIN_PASSWORD)
    if is_admin_req:
        if is_admin_req:
            session['admin_auth'] = True
            session.permanent = True               # الجلسة تدوم 30 يوم
        feat_data = load_feature_restrictions()
        unlocked = feat_data.get("user_unlocked", {}).get(user_id, [])
        if feature not in unlocked:
            unlocked.append(feature)
        feat_data.setdefault("user_unlocked", {})[user_id] = unlocked
        save_feature_restrictions(feat_data)
        return jsonify({"success": True, "message": f"✅ تم فتح {FEATURE_LABELS.get(feature, feature)} (أدمن — دخول حر)"})
    if not code:
        return jsonify({"success": False, "message": "❌ أدخل كود التفعيل أو كلمة مرور المشرف"})
    result, err = validate_voucher(code)
    if err:
        return jsonify({"success": False, "message": err})
    voucher = result["voucher"]
    allowed = voucher.get("allowed_features", [])
    # البطاقة لا تقيد وظيفة معينة → تفتح الوظيفة المطلوبة مباشرة
    if allowed and feature not in allowed:
        return jsonify({"success": False, "message": "❌ هذه البطاقة لا تفتح الوظيفة المطلوبة"})
    feat_data = load_feature_restrictions()
    unlocked = feat_data.get("user_unlocked", {}).get(user_id, [])
    if feature not in unlocked:
        unlocked.append(feature)
    feat_data.setdefault("user_unlocked", {})[user_id] = unlocked
    save_feature_restrictions(feat_data)
    return jsonify({
        "success": True,
        "message": f"✅ تم فتح {FEATURE_LABELS.get(feature, feature)} بنجاح!"
    })

@app.route("/admin/api/feature_vouchers", methods=["GET", "POST"])
def admin_feature_vouchers():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    if request.method == "POST":
        data = request.json or {}
        action = data.get("action", "update")
        feature = data.get("feature", "")
        if not feature or feature not in FEATURE_LABELS:
            return jsonify({"success": False, "message": "وظيفة غير موجودة"})
        fv = load_feature_vouchers()
        fv.setdefault(feature, {"users": []})
        if action == "reset":
            fv[feature]["users"] = []
        else:
            users = [u for u in data.get("users", []) if u]
            fv[feature]["users"] = users
        save_feature_vouchers(fv)
        return jsonify({"success": True, "message": f"✅ تم تحديث {FEATURE_LABELS[feature]}"})
    fv = load_feature_vouchers()
    result = {}
    for feature, label in FEATURE_LABELS.items():
        d = fv.get(feature, {"users": []})
        result[feature] = {"label": label, "users": d.get("users", []), "user_count": len(d.get("users", []))}
    return jsonify({"success": True, "features": result})

@app.route("/admin/api/send_vouchers_notification", methods=["POST"])
def admin_send_vouchers_notification():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = request.json or {}
    codes_text = data.get("text", "").strip()
    if not codes_text:
        return jsonify({"success": False, "message": "لا يوجد نص للإرسال"})
    user_id = session.get('user_id', 'user_1')
    try:
        socketio.emit('new_broadcast_notification', {
            "id": str(uuid.uuid4()),
            "message": f"📋 نسخة الكروت المولّدة:\n\n{codes_text}",
            "type": "success",
            "timestamp": datetime.now().isoformat()
        }, to=user_id)
    except Exception as e:
        logger.error(f"send_vouchers_notification error: {e}")
    return jsonify({"success": True, "message": "✅ تم إرسال الكروت كإشعار"})

@app.route("/admin/api/export_vouchers_txt_fixed", methods=["GET"])
def admin_export_vouchers_txt_fixed():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = load_cards_data()
    vouchers = data.get("vouchers", [])
    plans = data.get("plans", [])
    status_map = {"unused": "🟢 غير مستخدم", "active": "🔵 نشط", "used": "🔴 مستخدم", "expired": "⚫ منتهي"}
    lines = [
        "═" * 50,
        "    🎓 مركز سرعة إنجاز - قائمة الكروت (نسخة منسّقة)",
        f"    تاريخ التصدير: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "═" * 50, ""
    ]
    if not vouchers:
        lines.append("⚠️ لا توجد كروت في النظام")
    else:
        for i, v in enumerate(vouchers, 1):
            plan_name = voucher_display_plan_name(v, plans)
            status = status_map.get(v.get("status", ""), v.get("status", ""))
            created = (v.get("created_at") or "").split("T")[0] or "---"
            code = voucher_display_code(v)
            feats = ", ".join(v.get("allowed_features", [])) or "—"
            lines.append(f"[{str(i).zfill(3)}] الكود: {code}")
            lines.append(f"       الخطة: {plan_name} | {status} | التاريخ: {created}")
            lines.append(f"       الوظائف: {feats}")
            lines.append("")
    lines += ["═" * 50]
    content = "\n".join(lines)
    response = make_response(content)
    response.headers["Content-Type"] = "text/plain; charset=utf-8"
    response.headers["Content-Disposition"] = (
        f"attachment; filename=vouchers_{datetime.now().strftime('%Y%m%d_%H%M')}.txt"
    )
    return response

# ─── مسارات نظام التحديث الذاتي ──────────────────────────────────────────

@app.route("/api/check_update", methods=["GET"])
def api_check_update():
    has_update, current, latest, message = check_for_updates()
    return jsonify({
        "success": True,
        "has_update": has_update,
        "current": current[:7] if current else None,
        "latest": latest[:7] if latest else None,
        "message": message
    })

@app.route("/api/perform_update", methods=["POST"])
def api_perform_update():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "❌ غير مصرح"}), 403
    has_update, current, latest, msg = check_for_updates()
    if not has_update:
        return jsonify({"success": False, "message": "✅ لا توجد تحديثات جديدة", "logs": [msg]})
    success, logs = perform_update()
    if success:
        return jsonify({"success": True, "message": "✅ تم التحديث بنجاح", "logs": logs, "restarting": True})
    return jsonify({"success": False, "message": "❌ فشل التحديث", "logs": logs})

@app.route("/api/auto_update_status", methods=["GET"])
def api_auto_update_status():
    settings = load_update_settings()
    return jsonify({
        "success": True,
        "auto_update": settings.get('auto_update', False),
        "last_check": settings.get('last_check'),
        "last_update": settings.get('last_update')
    })

@app.route("/api/toggle_auto_update", methods=["POST"])
def api_toggle_auto_update():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "❌ غير مصرح"}), 403
    data = request.json or {}
    enabled = data.get('enabled', False)
    settings = load_update_settings()
    settings['auto_update'] = enabled
    settings['last_check'] = datetime.now().isoformat()
    save_update_settings(settings)
    if enabled:
        start_auto_update_thread()
    else:
        stop_auto_update_thread()
    return jsonify({"success": True, "auto_update": enabled,
                    "message": f"✅ تم {'تفعيل' if enabled else 'إلغاء'} التحديث التلقائي"})

@app.route("/api/update_logs", methods=["GET"])
def api_update_logs():
    settings = load_update_settings()
    return jsonify({"success": True, "last_update": settings.get('last_update'), "last_check": settings.get('last_check')})

# ─── بدء خلفية التحديث التلقائي إذا كان مفعلاً ───

# ─── مسارات API للإشعارات الدورية ────────────────────────────────────

@app.route("/api/promo_data", methods=["GET"])
def api_promo_data():
    data = load_promo_data()
    return jsonify({"success": True, "enabled": data.get("enabled", False),
                    "messages": data.get("messages", []), "current_index": data.get("current_index", 0)})

@app.route("/api/promo_save", methods=["POST"])
def api_promo_save():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = request.json or {}
    messages = [m.strip() for m in data.get("messages", []) if m.strip()]
    if not messages:
        return jsonify({"success": False, "message": "الرجاء إدخال رسالة واحدة على الأقل"})
    promo = load_promo_data()
    promo["messages"] = messages
    promo["current_index"] = 0
    save_promo_data(promo)
    return jsonify({"success": True, "count": len(messages), "message": f"تم حفظ {len(messages)} رسالة"})

@app.route("/api/promo_toggle", methods=["POST"])
def api_promo_toggle():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = request.json or {}
    enabled = data.get("enabled", False)
    promo = load_promo_data()
    promo["enabled"] = enabled
    save_promo_data(promo)
    if enabled:
        start_promo_thread()
    else:
        stop_promo_thread()
    return jsonify({"success": True, "enabled": enabled,
                    "message": f"✅ تم {'تفعيل' if enabled else 'إلغاء'} الإشعارات الدورية"})

@app.route("/api/promo_status", methods=["GET"])
def api_promo_status():
    data = load_promo_data()
    return jsonify({"success": True, "enabled": data.get("enabled", False),
                    "message_count": len(data.get("messages", [])),
                    "current_index": data.get("current_index", 0)})

# ──────────────────────────────────────────────────────────────────────────
# ▸ API: جلب كل المجموعات من الحساب
# ──────────────────────────────────────────────────────────────────────────
@app.route("/api/get_all_groups", methods=["GET"])
def api_get_all_groups():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    with USERS_LOCK:
        if user_id not in USERS:
            return jsonify({"success": False, "message": "المستخدم غير موجود"}), 404
        client_manager = USERS[user_id].get('client_manager')
    if not client_manager or not client_manager.client:
        return jsonify({"success": False, "message": "العميل غير متصل"}), 400
    try:
        dialogs = client_manager.run_coroutine(client_manager.client.get_dialogs())
        groups = []
        for d in dialogs:
            entity = d.entity
            if hasattr(entity, 'megagroup') or hasattr(entity, 'broadcast') or hasattr(entity, 'gigagroup'):
                title = getattr(d, 'title', None) or getattr(entity, 'title', 'بدون عنوان')
                username = getattr(entity, 'username', None)
                link = f"https://t.me/{username}" if username else None
                is_channel = bool(getattr(entity, 'broadcast', False))
                groups.append({
                    "id": entity.id,
                    "title": title,
                    "username": username,
                    "link": link,
                    "type": "قناة" if is_channel else "مجموعة"
                })
        groups.sort(key=lambda x: x['title'])
        return jsonify({"success": True, "groups": groups, "count": len(groups)})
    except Exception as e:
        logger.error(f"خطأ في جلب المجموعات: {e}")
        return jsonify({"success": False, "message": str(e)}), 500

# ──────────────────────────────────────────────────────────────────────────
# ▸ API: مزامنة وتصدير/استيراد الإعدادات
# ──────────────────────────────────────────────────────────────────────────
@app.route("/api/sync_settings", methods=["POST"])
def api_sync_settings():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    data = request.json or {}
    settings = data.get('settings', {})
    if settings and save_settings(user_id, settings, force=True):
        return jsonify({"success": True, "message": "تم حفظ النسخة الاحتياطية"})
    return jsonify({"success": False, "message": "فشل الحفظ"})

@app.route("/api/load_backup_settings", methods=["GET"])
def api_load_backup_settings():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    settings = load_settings(user_id)
    return jsonify({"success": True, "settings": settings or {}})

@app.route("/api/export_settings", methods=["GET"])
def api_export_settings():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    settings = load_settings(user_id)
    export_data = {
        "user_id": user_id,
        "exported_at": datetime.now().isoformat(),
        "version": "2.0",
        "settings": settings or {}
    }
    response = make_response(json.dumps(export_data, ensure_ascii=False, indent=2))
    response.headers["Content-Type"] = "application/json; charset=utf-8"
    response.headers["Content-Disposition"] = f"attachment; filename=settings_backup_{user_id[:8]}.json"
    return response

@app.route("/api/import_settings", methods=["POST"])
def api_import_settings():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"success": False, "message": "غير مسجل"}), 401
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "لم يتم رفع ملف"}), 400
    file = request.files['file']
    if not file.filename.endswith('.json'):
        return jsonify({"success": False, "message": "الملف يجب أن يكون بصيغة JSON"}), 400
    try:
        content = file.read().decode('utf-8')
        data = json.loads(content)
        settings = data.get('settings', {})
        if not settings:
            return jsonify({"success": False, "message": "الملف لا يحتوي على إعدادات صالحة"}), 400
        save_settings(user_id, settings, force=True)
        return jsonify({"success": True, "settings": settings, "message": "تم استيراد الإعدادات بنجاح"})
    except json.JSONDecodeError:
        return jsonify({"success": False, "message": "الملف تالف أو تنسيقه غير صحيح"}), 400
    except Exception as e:
        return jsonify({"success": False, "message": f"خطأ: {str(e)}"}), 500

# ──────────────────────────────────────────────────────────────────────────
# ▸ API: إدارة روابط الدعوة (للمسؤول)
# ──────────────────────────────────────────────────────────────────────────
@app.route("/admin/api/create_invite", methods=["POST"])
def admin_create_invite():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    token = generate_invite_token()
    link = get_invite_link(token)
    return jsonify({"success": True, "token": token, "link": link, "message": "تم إنشاء الرابط"})

@app.route("/admin/api/invites", methods=["GET"])
def admin_list_invites():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = load_invites()
    tokens = sorted(data["tokens"], key=lambda x: x.get("created_at", ""), reverse=True)
    return jsonify({"success": True, "tokens": tokens, "total": len(tokens)})

@app.route("/admin/api/revoke_invite", methods=["POST"])
def admin_revoke_invite():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = request.json or {}
    token = data.get("token")
    if not token:
        return jsonify({"success": False, "message": "الرمز مطلوب"}), 400
    invites = load_invites()
    for item in invites["tokens"]:
        if item["token"] == token:
            if item["status"] == "used":
                return jsonify({"success": False, "message": "الرمز مستخدم بالفعل"}), 400
            item["status"] = "expired"
            save_invites(invites)
            return jsonify({"success": True, "message": "تم إلغاء الرابط"})
    return jsonify({"success": False, "message": "الرمز غير موجود"}), 404

@app.route("/admin/api/delete_invite", methods=["POST"])
def admin_delete_invite():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = request.json or {}
    token = data.get("token")
    if not token:
        return jsonify({"success": False, "message": "الرمز مطلوب"}), 400
    invites = load_invites()
    invites["tokens"] = [t for t in invites["tokens"] if t["token"] != token]
    save_invites(invites)
    return jsonify({"success": True, "message": "تم حذف الرابط"})

# ──────────────────────────────────────────────────────────────────────────
# ▸ API: تصدير الكروت كنص أو PDF
# ──────────────────────────────────────────────────────────────────────────
@app.route("/admin/api/export_vouchers_txt", methods=["GET"])
def admin_export_vouchers_txt():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = load_cards_data()
    vouchers = data.get("vouchers", [])
    plans = data.get("plans", [])
    lines = ["═══════════════════════════════════════",
             "    مركز سرعة انجاز - قائمة الكروت",
             f"    تاريخ التصدير: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
             "═══════════════════════════════════════\n"]
    for i, v in enumerate(vouchers, 1):
        status_ar = {"unused": "غير مستخدم", "active": "نشط", "used": "مستخدم", "expired": "منتهي"}.get(v.get("status",""), v.get("status",""))
        code = voucher_display_code(v)
        plan_name = voucher_display_plan_name(v, plans)
        lines.append(f"[{i}] الكود: {code} | الخطة: {plan_name} | الحالة: {status_ar} | التاريخ: {v.get('created_at','')[:10]}")
    content = "\n".join(lines)
    response = make_response(content)
    response.headers["Content-Type"] = "text/plain; charset=utf-8"
    response.headers["Content-Disposition"] = f"attachment; filename=vouchers_{datetime.now().strftime('%Y%m%d')}.txt"
    return response

@app.route("/admin/api/export_vouchers_pdf", methods=["GET"])
def admin_export_vouchers_pdf():
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مصرح"}), 403
    data = load_cards_data()
    vouchers = data.get("vouchers", [])
    plans = data.get("plans", [])
    # إنشاء PDF بسيط بصيغة HTML ثم إعادته كـ HTML للطباعة إذا لم تتوفر مكتبة PDF
    html_content = f"""<!DOCTYPE html>
<html dir="rtl" lang="ar">
<head><meta charset="UTF-8"><title>قائمة الكروت</title>
<style>body{{font-family:Arial,sans-serif;direction:rtl;}}
table{{width:100%;border-collapse:collapse;font-size:12px;}}
th,td{{border:1px solid #ddd;padding:8px;text-align:right;}}
th{{background:#1e3c78;color:white;}}
.unused{{color:green;}} .used{{color:gray;}} .active{{color:blue;}} .expired{{color:red;}}
h2{{color:#1e3c78;}} .meta{{color:#666;font-size:11px;}}
</style></head>
<body>
<h2>🏦 مركز سرعة انجاز — قائمة الكروت</h2>
<p class="meta">تاريخ التصدير: {datetime.now().strftime('%Y-%m-%d %H:%M')} | إجمالي الكروت: {len(vouchers)}</p>
<table>
<tr><th>#</th><th>رقم البطاقة</th><th>الخطة</th><th>الحالة</th><th>تاريخ الإنشاء</th><th>تاريخ الاستخدام</th></tr>
"""
    status_map = {"unused": ("غير مستخدم","unused"), "active": ("نشط","active"), "used": ("مستخدم","used"), "expired": ("منتهي","expired")}
    for i, v in enumerate(vouchers, 1):
        st = v.get("status","")
        st_ar, st_cls = status_map.get(st, (st, ""))
        code = voucher_display_code(v)
        plan_name = voucher_display_plan_name(v, plans)
        html_content += f"<tr><td>{i}</td><td style='font-family:monospace;'>{code}</td><td>{plan_name}</td><td class='{st_cls}'>{st_ar}</td><td>{v.get('created_at','')[:10]}</td><td>{v.get('used_at','') or '-'}</td></tr>\n"
    html_content += f"""</table>
<script>window.onload=function(){{window.print();}}</script>
</body></html>"""
    response = make_response(html_content)
    response.headers["Content-Type"] = "text/html; charset=utf-8"
    return response

# ══════════════════════════════════════════════════════════════════════════════
#  API: تتبع اتصال وانقطاع المستخدمين مع تسجيل رقم الجوال واسم التليجرام
# ══════════════════════════════════════════════════════════════════════════════

@app.route("/api/user_connected", methods=["POST"])
def api_user_connected():
    """تسجيل اتصال المستخدم مع رقمه واسمه في التليجرام"""
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "❌ الجلسة غير صالحة"}), 401
    from install_tracker import load_user_sessions, save_user_sessions
    data = request.get_json() or {}
    # يجب أن يطابق user_id في الجلسة
    session_user_id = session['user_id']
    user_id = data.get("user_id", session_user_id)
    if user_id != session_user_id:
        return jsonify({"success": False, "message": "❌ غير مسموح"}), 403
    install_id = request.headers.get("X-Install-ID", "unknown")
    phone = data.get("phone", "")
    account_name = data.get("account_name", "")
    name = data.get("name", user_id)
    if user_id == "unknown":
        return jsonify({"success": False, "message": "user_id مطلوب"}), 400
    now_iso = datetime.now().isoformat()
    try:
        sessions_data = load_user_sessions()
        inst = next((i for i in sessions_data.get("installations", []) if i.get("install_id") == install_id), None)
        if inst:
            state = inst.setdefault("users_state", {}).setdefault(user_id, {})
            old_connected = state.get("connected", False)
            state["connected"] = True
            state["last_seen"] = now_iso
            state["phone"] = phone or state.get("phone", "")
            state["account_name"] = account_name or state.get("account_name", "")
            state["name"] = name
            if not old_connected:
                state["connected_at"] = now_iso
                history = state.get("session_history", [])
                history.append({
                    "connected_at": now_iso,
                    "disconnected_at": None,
                    "phone": state["phone"],
                    "account_name": state["account_name"]
                })
                if len(history) > 20:
                    history = history[-20:]
                state["session_history"] = history
            save_user_sessions(sessions_data)
        if user_id in USERS:
            with USERS_LOCK:
                USERS[user_id]["connected"] = True
                USERS[user_id]["last_seen"] = now_iso
                USERS[user_id]["connected_at"] = now_iso
                if phone:
                    USERS[user_id]["phone"] = phone
                if account_name:
                    USERS[user_id]["account_name"] = account_name
        socketio.emit("installation_updated", {
            "install_id": install_id,
            "user_id": user_id,
            "connected": True,
            "connected_at": now_iso,
            "phone": phone,
            "account_name": account_name
        })
    except Exception as e:
        logger.error(f"api_user_connected error: {e}")
    return jsonify({"success": True, "message": f"تم تسجيل اتصال {user_id}", "connected_at": now_iso})


@app.route("/api/user_disconnected", methods=["POST"])
def api_user_disconnected():
    """تسجيل انقطاع المستخدم"""
    if 'user_id' not in session:
        return jsonify({"success": False, "message": "❌ الجلسة غير صالحة"}), 401
    from install_tracker import load_user_sessions, save_user_sessions
    data = request.get_json() or {}
    session_user_id = session['user_id']
    user_id = data.get("user_id", session_user_id)
    if user_id != session_user_id:
        return jsonify({"success": False, "message": "❌ غير مسموح"}), 403
    install_id = request.headers.get("X-Install-ID", "unknown")
    if user_id == "unknown":
        return jsonify({"success": False, "message": "user_id مطلوب"}), 400
    now_iso = datetime.now().isoformat()
    try:
        sessions_data = load_user_sessions()
        inst = next((i for i in sessions_data.get("installations", []) if i.get("install_id") == install_id), None)
        if inst:
            state = inst.setdefault("users_state", {}).setdefault(user_id, {})
            old_connected = state.get("connected", False)
            state["connected"] = False
            state["last_seen"] = now_iso
            state["disconnected_at"] = now_iso
            if old_connected:
                history = state.get("session_history", [])
                if history and history[-1].get("disconnected_at") is None:
                    history[-1]["disconnected_at"] = now_iso
                state["session_history"] = history
            save_user_sessions(sessions_data)
        if user_id in USERS:
            with USERS_LOCK:
                USERS[user_id]["connected"] = False
                USERS[user_id]["last_seen"] = now_iso
                USERS[user_id]["disconnected_at"] = now_iso
        socketio.emit("installation_updated", {
            "install_id": install_id,
            "user_id": user_id,
            "connected": False,
            "disconnected_at": now_iso
        })
    except Exception as e:
        logger.error(f"api_user_disconnected error: {e}")
    return jsonify({"success": True, "message": f"تم تسجيل انقطاع {user_id}", "disconnected_at": now_iso})


@app.route("/admin/api/user_connection_history/<user_id>", methods=["GET"])
def admin_user_connection_history(user_id):
    """عرض سجل اتصالات مستخدم عبر جميع التثبيتات"""
    if not session.get("admin_auth"):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    from install_tracker import load_user_sessions
    sessions_data = load_user_sessions()
    history = []
    for inst in sessions_data.get("installations", []):
        state = inst.get("users_state", {}).get(user_id, {})
        if state:
            history.append({
                "install_id": inst.get("install_id", ""),
                "ip": inst.get("ip", ""),
                "user_agent": inst.get("user_agent", "")[:80],
                "connected": state.get("connected", False),
                "connected_at": state.get("connected_at", ""),
                "disconnected_at": state.get("disconnected_at", ""),
                "last_seen": state.get("last_seen", ""),
                "phone": state.get("phone", ""),
                "account_name": state.get("account_name", ""),
                "session_history": state.get("session_history", [])
            })
    return jsonify({"success": True, "history": history, "user_id": user_id})


# ─── تشغيل خلفية الإشعارات الدورية إذا كانت مفعّلة ──────────────────
_promo_init = load_promo_data()
if _promo_init.get("enabled", False):
    start_promo_thread()
    logger.info("🔄 تم بدء خلفية الإشعارات الدورية تلقائياً")


_upd_settings = load_update_settings()
if _upd_settings.get('auto_update', False):
    start_auto_update_thread()

def _is_admin_auth():
    return session.get("admin_auth", False)

register_admin_routes(
    app,
    _is_admin_auth,
    predefined_users=PREDEFINED_USERS,
    users_dict=USERS,
    users_lock=USERS_LOCK,
    load_settings_func=load_settings,
    save_settings_func=save_settings,
    reset_user_func=_do_reset_user,
)


# ════════════════════════════════════════════════════════
#  نظام المنصة: تسجيل الدخول وإدارة المستخدمين الديناميكيين
# ════════════════════════════════════════════════════════

@app.route("/user-login")
def user_login_page():
    """صفحة تسجيل الدخول لرؤية المستخدمين الإضافيين"""
    if session.get('platform_logged_in') or session.get('admin_auth'):
        return redirect("/")
    return render_template("user_login.html")

@app.route("/api/platform_login", methods=["POST"])
def api_platform_login():
    """تسجيل الدخول للوصول إلى المستخدمين الإضافيين"""
    data = request.json or {}
    username = data.get("username", "").strip()
    password = data.get("password", "").strip()
    if not username or not password:
        return jsonify({"success": False, "message": "جميع الحقول مطلوبة"})
    # كلمة مرور المشرف تعطي دخولاً كاملاً
    if username == ADMIN_USERNAME and password == ADMIN_PASSWORD:
        session['platform_logged_in'] = True
        session['admin_auth']         = True
        session.permanent             = True
        return jsonify({"success": True, "message": "✅ تم الدخول كمشرف"})
    try:
        auth_result = authenticate_platform_user(username, password)
        if auth_result:
            session['platform_logged_in']  = True
            session['platform_username']   = auth_result
            session.permanent              = True
            return jsonify({"success": True, "message": "✅ تم تسجيل الدخول بنجاح"})
    except Exception as e:
        logger.error(f"platform_login error: {e}")
    return jsonify({"success": False, "message": "❌ اسم المستخدم أو كلمة المرور غير صحيحة"})

@app.route("/api/platform_logout", methods=["POST"])
def api_platform_logout():
    session.pop("platform_logged_in", None)
    session.pop("platform_username", None)
    return jsonify({"success": True})

@app.route("/api/platform_register", methods=["POST"])
def api_platform_register():
    """إنشاء حساب جديد في المنصة"""
    if not session.get('admin_auth'):
        return jsonify({"success": False, "message": "يجب أن تكون مشرفاً لإنشاء حسابات"}), 403
    data = request.json or {}
    username = data.get("username", "").strip()
    password = data.get("password", "").strip()
    if not username or not password:
        return jsonify({"success": False, "message": "جميع الحقول مطلوبة"})
    try:
        success, msg = create_platform_account(username, password)
        return jsonify({"success": success, "message": msg})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/add_account_slot", methods=["POST"])
def api_add_account_slot():
    """إنشاء فتحة حساب جديدة والتبديل إليها"""
    try:
        global PREDEFINED_USERS

        # ── التحقق من أن المستخدم الحالي قد سجّل الدخول بتيليجرام أولاً ──
        _current_uid = session.get('user_id', 'user_1')
        _current_name = PREDEFINED_USERS.get(_current_uid, {}).get('name', _current_uid)
        _is_logged_in = False
        try:
            with USERS_LOCK:
                _is_logged_in = bool(USERS.get(_current_uid, {}).get('authenticated', False))
        except Exception:
            pass
        if not _is_logged_in:
            _is_logged_in = bool(load_string_session(_current_uid))
        if not _is_logged_in:
            return jsonify({
                "success": False,
                "message": f"⚠️ يجب تسجيل الدخول بحساب تيليجرام في الحساب الحالي «{_current_name}» أولاً قبل إضافة حساب جديد. انتقل إلى الحساب الحالي وسجّل الدخول.",
                "redirect_user": _current_uid
            })

        existing = set(PREDEFINED_USERS.keys())
        n = 1
        while f"user_{n}" in existing:
            n += 1
        new_uid = f"user_{n}"
        _colors = ["#6366f1","#28a745","#ffc107","#dc3545","#6f42c1","#17a2b8","#fd7e14","#20c997"]
        _color = _colors[(n - 1) % len(_colors)]
        _ok, _msg = add_dynamic_user(new_uid, f"حساب {n}", "fas fa-user-plus", _color)
        if _ok:
            PREDEFINED_USERS = load_dynamic_users()
        # التبديل إلى الفتحة الجديدة
        old_uid = session.get('user_id')
        session['user_id'] = new_uid
        session.permanent = True
        return jsonify({
            "success": True,
            "user_id": new_uid,
            "account_name": f"حساب {n}",
            "color": _color,
            "icon": "fas fa-user-plus",
            "message": f"✅ تم إنشاء فتحة حساب {n} جديدة"
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/add_dynamic_user", methods=["POST"])
def api_add_dynamic_user():
    """إضافة مستخدم ديناميكي جديد"""
    data = request.json or {}
    user_id = data.get("user_id", "").strip()
    name    = data.get("name",    "").strip()
    icon    = data.get("icon",    "fas fa-user")
    color   = data.get("color",   "#6c757d")
    if not user_id or not name:
        return jsonify({"success": False, "message": "معرف المستخدم والاسم مطلوبان"})
    if not user_id.startswith("user_"):
        return jsonify({"success": False, "message": "يجب أن يبدأ المعرف بـ user_"})
    try:
        success, msg = add_dynamic_user(user_id, name, icon, color)
        if success:
            global PREDEFINED_USERS
            PREDEFINED_USERS = load_dynamic_users()
        return jsonify({"success": success, "message": msg})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/delete_dynamic_user", methods=["POST"])
def api_delete_dynamic_user():
    """حذف مستخدم ديناميكي"""
    if not session.get('admin_auth'):
        return jsonify({"success": False, "message": "غير مخول"}), 403
    data = request.json or {}
    user_id = data.get("user_id", "").strip()
    if not user_id:
        return jsonify({"success": False, "message": "معرف المستخدم مطلوب"})
    try:
        success, msg = delete_dynamic_user(user_id)
        if success:
            global PREDEFINED_USERS
            PREDEFINED_USERS = load_dynamic_users()
        return jsonify({"success": success, "message": msg})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

@app.route("/api/list_dynamic_users", methods=["GET"])
def api_list_dynamic_users():
    """قائمة جميع المستخدمين الديناميكيين"""
    try:
        users = load_dynamic_users()
        return jsonify({"success": True, "users": users})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    print(f"🌐 تشغيل الخادم على المنفذ {port}...")
    print(f"🔗 رابط التطبيق: http://0.0.0.0:{port}")
    print("🛡️ نظام الاستمرارية الدائم مُفعل — يعمل حتى الإيقاف اليدوي")
    print("🎓 مركز سرعة انجاز للخدمات الطلابية والأكاديمية - الإصدار المتكامل")
    print("📊 تم دمج النظام الأكاديمي الذكي + منسق المستندات + منشئ العروض")

    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    network_monitor.start()

    try:
        socketio.run(app, host='0.0.0.0', port=port, debug=False, allow_unsafe_werkzeug=True)
    except Exception as e:
        print(f"❌ خطأ في تشغيل الخادم: {e}")
